#!/usr/bin/env python3
"""
Generate diverse OOXML test files using python-docx, openpyxl, and python-pptx.
Creates files with various features: tables, images, styles, charts, etc.
"""
import hashlib, json, os, re, sys
from pathlib import Path

def generate_docx_files(output_dir):
    """Generate diverse docx test files."""
    try:
        from docx import Document
        from docx.shared import Inches, Pt, Cm, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.table import WD_TABLE_ALIGNMENT
    except ImportError:
        print("  python-docx not installed, skipping docx generation")
        return []

    files = []
    docx_dir = output_dir / "docx"
    docx_dir.mkdir(parents=True, exist_ok=True)

    # 1. Simple document with various heading levels
    doc = Document()
    doc.add_heading("Test Document - Heading Levels", 0)
    for i in range(1, 5):
        doc.add_heading(f"Heading Level {i}", i)
        doc.add_paragraph(f"This is paragraph text under heading level {i}. " * 3)
    p = docx_dir / "gen_headings.docx"
    doc.save(str(p))
    files.append(p)

    # 2. Japanese text document
    doc = Document()
    doc.add_heading("日本語テスト文書", 0)
    doc.add_paragraph("これは日本語のテスト文書です。漢字、ひらがな、カタカナを含みます。")
    doc.add_paragraph("吾輩は猫である。名前はまだ無い。どこで生れたかとんと見当がつかぬ。")
    doc.add_paragraph("何でも薄暗いじめじめした所でニャーニャー泣いていた事だけは記憶している。")
    doc.add_paragraph("東京都千代田区霞が関1-2-3 電話番号：03-1234-5678")
    for font_name in ["MS Gothic", "MS Mincho", "Yu Gothic", "Meiryo"]:
        p = doc.add_paragraph()
        run = p.add_run(f"フォント: {font_name} - あいうえおアイウエオ漢字")
        run.font.name = font_name
        run.font.size = Pt(12)
    p = docx_dir / "gen_japanese.docx"
    doc.save(str(p))
    files.append(p)

    # 3. Table document
    doc = Document()
    doc.add_heading("Table Test", 0)
    table = doc.add_table(rows=10, cols=5)
    table.style = "Table Grid"
    for i, row in enumerate(table.rows):
        for j, cell in enumerate(row.cells):
            if i == 0:
                cell.text = f"Header {j+1}"
            else:
                cell.text = f"Row {i} Col {j+1}"
    # Merged cells table
    doc.add_heading("Merged Cells", 1)
    table2 = doc.add_table(rows=4, cols=4)
    table2.style = "Table Grid"
    table2.cell(0, 0).merge(table2.cell(0, 1)).text = "Merged"
    table2.cell(1, 0).text = "A"
    table2.cell(1, 1).text = "B"
    p = docx_dir / "gen_tables.docx"
    doc.save(str(p))
    files.append(p)

    # 4. Styled text document
    doc = Document()
    doc.add_heading("Text Styling Test", 0)
    p = doc.add_paragraph()
    run = p.add_run("Bold text. ")
    run.bold = True
    run = p.add_run("Italic text. ")
    run.italic = True
    run = p.add_run("Underlined text. ")
    run.underline = True
    run = p.add_run("Red text. ")
    run.font.color.rgb = RGBColor(0xFF, 0, 0)
    run = p.add_run("Large text. ")
    run.font.size = Pt(24)
    run = p.add_run("Small text. ")
    run.font.size = Pt(8)
    p = docx_dir / "gen_styles.docx"
    doc.save(str(p))
    files.append(p)

    # 5. Multi-paragraph document
    doc = Document()
    doc.add_heading("Lorem Ipsum", 0)
    for i in range(20):
        doc.add_paragraph(
            f"Paragraph {i+1}: Lorem ipsum dolor sit amet, consectetur adipiscing elit. "
            "Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua. "
            "Ut enim ad minim veniam, quis nostrud exercitation ullamco laboris. "
        )
    p = docx_dir / "gen_long.docx"
    doc.save(str(p))
    files.append(p)

    # 6. Lists document
    doc = Document()
    doc.add_heading("Lists Test", 0)
    doc.add_heading("Bullet List", 1)
    for item in ["First item", "Second item", "Third item", "Fourth item"]:
        doc.add_paragraph(item, style="List Bullet")
    doc.add_heading("Numbered List", 1)
    for item in ["Step one", "Step two", "Step three"]:
        doc.add_paragraph(item, style="List Number")
    p = docx_dir / "gen_lists.docx"
    doc.save(str(p))
    files.append(p)

    # 7-16. More variant documents
    templates = [
        ("gen_report.docx", "Annual Report 2024", [
            ("Executive Summary", "The fiscal year showed strong performance across all metrics."),
            ("Financial Overview", "Revenue increased by 15% year-over-year."),
            ("Outlook", "We expect continued growth in the coming fiscal year."),
        ]),
        ("gen_memo.docx", "Internal Memo", [
            ("Subject", "Updated meeting schedule for Q2"),
            ("Details", "Please note the following changes to our regular meetings."),
        ]),
        ("gen_proposal.docx", "Project Proposal", [
            ("Background", "This proposal outlines the development plan."),
            ("Objectives", "1. Improve performance\n2. Reduce costs\n3. Enhance UX"),
            ("Timeline", "Phase 1: April-June\nPhase 2: July-September"),
            ("Budget", "Estimated total: $50,000"),
        ]),
        ("gen_minutes.docx", "Meeting Minutes", [
            ("Date", "March 13, 2026"),
            ("Attendees", "Tanaka, Yamada, Suzuki, Sato"),
            ("Agenda", "1. Project updates\n2. Budget review\n3. Action items"),
            ("Decisions", "Approved the Q2 budget allocation."),
        ]),
        ("gen_contract.docx", "Service Agreement", [
            ("Parties", "Party A: Oxi Corp\nParty B: Client Inc"),
            ("Terms", "This agreement is effective from April 1, 2026."),
            ("Scope", "Development and maintenance of web application."),
            ("Payment", "Monthly fee: 500,000 JPY"),
        ]),
    ]
    for fname, title, sections in templates:
        doc = Document()
        doc.add_heading(title, 0)
        for heading, content in sections:
            doc.add_heading(heading, 1)
            doc.add_paragraph(content)
        fp = docx_dir / fname
        doc.save(str(fp))
        files.append(fp)

    # Japanese business documents
    jp_templates = [
        ("gen_jp_report.docx", "令和6年度 事業報告書", [
            ("事業概要", "本年度は以下の事業を実施しました。"),
            ("実績", "売上高：10億円\n営業利益：2億円\n経常利益：1.8億円"),
            ("課題と展望", "来年度に向けて、DX推進を最重点課題とします。"),
        ]),
        ("gen_jp_notice.docx", "お知らせ", [
            ("件名", "システムメンテナンスのお知らせ"),
            ("日時", "2026年4月1日（水）午前2時〜午前6時"),
            ("影響範囲", "全サービスが一時的にご利用いただけなくなります。"),
        ]),
        ("gen_jp_manual.docx", "操作マニュアル", [
            ("はじめに", "本マニュアルはシステムの基本操作について説明します。"),
            ("ログイン方法", "1. ブラウザでURLにアクセス\n2. ID・パスワードを入力\n3. ログインボタンをクリック"),
            ("基本操作", "メニューバーから各機能にアクセスできます。"),
            ("よくある質問", "Q: パスワードを忘れた場合は？\nA: 管理者にお問い合わせください。"),
        ]),
        ("gen_jp_application.docx", "申請書", [
            ("申請者", "氏名：山田太郎\n所属：開発部"),
            ("申請内容", "出張申請\n期間：2026年4月15日〜4月17日\n行先：大阪"),
            ("理由", "顧客との打ち合わせのため"),
        ]),
        ("gen_jp_invoice.docx", "請求書", [
            ("請求先", "株式会社テスト\n東京都港区芝公園1-1-1"),
            ("請求内容", "システム開発費：3,000,000円\nサーバー費用：100,000円\n合計：3,100,000円（税込）"),
            ("支払期限", "2026年4月末日"),
        ]),
    ]
    for fname, title, sections in jp_templates:
        doc = Document()
        doc.add_heading(title, 0)
        for heading, content in sections:
            doc.add_heading(heading, 1)
            doc.add_paragraph(content)
        fp = docx_dir / fname
        doc.save(str(fp))
        files.append(fp)

    print(f"  Generated {len(files)} docx files")
    return files

def generate_xlsx_files(output_dir):
    """Generate diverse xlsx test files."""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter
    except ImportError:
        print("  openpyxl not installed, skipping xlsx generation")
        return []

    files = []
    xlsx_dir = output_dir / "xlsx"
    xlsx_dir.mkdir(parents=True, exist_ok=True)

    # 1. Simple data table
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    headers = ["ID", "Name", "Category", "Value", "Date"]
    for col, h in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col, value=h)
        cell.font = Font(bold=True)
        cell.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        cell.font = Font(bold=True, color="FFFFFF")
    import random
    random.seed(42)
    categories = ["A", "B", "C", "D"]
    for i in range(2, 102):
        ws.cell(row=i, column=1, value=i-1)
        ws.cell(row=i, column=2, value=f"Item_{i-1:04d}")
        ws.cell(row=i, column=3, value=random.choice(categories))
        ws.cell(row=i, column=4, value=round(random.uniform(10, 10000), 2))
        ws.cell(row=i, column=5, value=f"2026-{random.randint(1,12):02d}-{random.randint(1,28):02d}")
    p = xlsx_dir / "gen_data_table.xlsx"
    wb.save(str(p))
    files.append(p)

    # 2. Multi-sheet workbook
    wb = Workbook()
    for sheet_name in ["Sales", "Costs", "Profit", "Summary"]:
        ws = wb.create_sheet(sheet_name)
        for r in range(1, 21):
            for c in range(1, 8):
                ws.cell(row=r, column=c, value=random.randint(100, 99999))
    if "Sheet" in wb.sheetnames:
        del wb["Sheet"]
    p = xlsx_dir / "gen_multi_sheet.xlsx"
    wb.save(str(p))
    files.append(p)

    # 3. Formulas
    wb = Workbook()
    ws = wb.active
    ws.title = "Formulas"
    ws["A1"] = "Value A"
    ws["B1"] = "Value B"
    ws["C1"] = "Sum"
    ws["D1"] = "Product"
    for i in range(2, 22):
        ws[f"A{i}"] = random.randint(1, 100)
        ws[f"B{i}"] = random.randint(1, 100)
        ws[f"C{i}"] = f"=A{i}+B{i}"
        ws[f"D{i}"] = f"=A{i}*B{i}"
    ws["A23"] = "Total"
    ws["C23"] = "=SUM(C2:C21)"
    ws["D23"] = "=SUM(D2:D21)"
    p = xlsx_dir / "gen_formulas.xlsx"
    wb.save(str(p))
    files.append(p)

    # 4. Styled cells
    wb = Workbook()
    ws = wb.active
    ws.title = "Styled"
    colors = ["FF0000", "00FF00", "0000FF", "FFFF00", "FF00FF", "00FFFF"]
    for i, color in enumerate(colors, 1):
        for j in range(1, 6):
            cell = ws.cell(row=i, column=j, value=f"Color {color}")
            cell.fill = PatternFill(start_color=color, end_color=color, fill_type="solid")
    thin_border = Border(
        left=Side(style='thin'), right=Side(style='thin'),
        top=Side(style='thin'), bottom=Side(style='thin')
    )
    for r in range(8, 15):
        for c in range(1, 6):
            cell = ws.cell(row=r, column=c, value=f"Border {r},{c}")
            cell.border = thin_border
    p = xlsx_dir / "gen_styled.xlsx"
    wb.save(str(p))
    files.append(p)

    # 5. Wide spreadsheet
    wb = Workbook()
    ws = wb.active
    ws.title = "Wide"
    for c in range(1, 51):
        ws.cell(row=1, column=c, value=f"Col_{get_column_letter(c)}")
    for r in range(2, 52):
        for c in range(1, 51):
            ws.cell(row=r, column=c, value=random.randint(0, 999))
    p = xlsx_dir / "gen_wide.xlsx"
    wb.save(str(p))
    files.append(p)

    # 6. Japanese content xlsx
    wb = Workbook()
    ws = wb.active
    ws.title = "売上データ"
    headers_jp = ["年月", "部門", "商品名", "数量", "単価", "売上金額"]
    for c, h in enumerate(headers_jp, 1):
        ws.cell(row=1, column=c, value=h)
    departments = ["営業部", "開発部", "総務部", "経理部"]
    products = ["製品A", "製品B", "サービスC", "サポートD"]
    for r in range(2, 52):
        ws.cell(row=r, column=1, value=f"2026年{random.randint(1,12)}月")
        ws.cell(row=r, column=2, value=random.choice(departments))
        ws.cell(row=r, column=3, value=random.choice(products))
        qty = random.randint(1, 100)
        price = random.randint(1000, 50000)
        ws.cell(row=r, column=4, value=qty)
        ws.cell(row=r, column=5, value=price)
        ws.cell(row=r, column=6, value=f"=D{r}*E{r}")
    p = xlsx_dir / "gen_japanese_sales.xlsx"
    wb.save(str(p))
    files.append(p)

    # 7-16: More variants
    for variant in range(10):
        wb = Workbook()
        ws = wb.active
        ws.title = f"Dataset_{variant}"
        cols = random.randint(3, 10)
        rows = random.randint(20, 100)
        for c in range(1, cols+1):
            ws.cell(row=1, column=c, value=f"Field_{c}")
        for r in range(2, rows+1):
            for c in range(1, cols+1):
                if random.random() < 0.3:
                    ws.cell(row=r, column=c, value=f"text_{random.randint(1,999)}")
                else:
                    ws.cell(row=r, column=c, value=round(random.uniform(-1000, 1000), 2))
        p = xlsx_dir / f"gen_dataset_{variant:02d}.xlsx"
        wb.save(str(p))
        files.append(p)

    print(f"  Generated {len(files)} xlsx files")
    return files

def generate_pptx_files(output_dir):
    """Generate diverse pptx test files."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("  python-pptx not installed, skipping pptx generation")
        return []

    files = []
    pptx_dir = output_dir / "pptx"
    pptx_dir.mkdir(parents=True, exist_ok=True)

    # 1. Title slide
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Oxi Project Overview"
    slide.placeholders[1].text = "Document Processing Suite"
    p = pptx_dir / "gen_title.pptx"
    prs.save(str(p))
    files.append(p)

    # 2. Multi-slide presentation
    prs = Presentation()
    titles = [
        "Introduction", "Background", "Problem Statement",
        "Proposed Solution", "Architecture", "Implementation",
        "Results", "Conclusion", "Q&A"
    ]
    for title in titles:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = f"Content for {title}"
        for i in range(3):
            p = tf.add_paragraph()
            p.text = f"Bullet point {i+1} for {title}"
            p.level = 0
    p = pptx_dir / "gen_presentation.pptx"
    prs.save(str(p))
    files.append(p)

    # 3. Japanese presentation
    prs = Presentation()
    jp_slides = [
        ("プロジェクト概要", ["目的と背景", "開発体制", "スケジュール"]),
        ("技術仕様", ["Rust + WebAssembly", "OOXML パーサー", "Canvas レンダリング"]),
        ("進捗報告", ["フェーズ1完了", "テスト実施中", "リリース予定：4月"]),
        ("課題と対策", ["パフォーマンス最適化", "日本語フォント対応", "禁則処理"]),
        ("まとめ", ["100%パース成功率", "LibreOffice比較", "今後の展開"]),
    ]
    for title, bullets in jp_slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = bullets[0]
        for b in bullets[1:]:
            p = tf.add_paragraph()
            p.text = b
    p = pptx_dir / "gen_japanese_pres.pptx"
    prs.save(str(p))
    files.append(p)

    # 4. Table slide
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Data Table"
    rows, cols = 5, 4
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(8), Inches(3)).table
    for c in range(cols):
        table.cell(0, c).text = f"Header {c+1}"
    for r in range(1, rows):
        for c in range(cols):
            table.cell(r, c).text = f"Data {r},{c}"
    p = pptx_dir / "gen_table_slide.pptx"
    prs.save(str(p))
    files.append(p)

    # 5-14: More variants
    topics = [
        ("Sales Report Q1", ["Revenue: $1.2M", "Growth: 15%", "New clients: 25"]),
        ("Product Roadmap", ["v1.0 Released", "v2.0 In Progress", "v3.0 Planned"]),
        ("Team Update", ["Engineering: 12", "Design: 4", "PM: 3"]),
        ("Market Analysis", ["TAM: $50B", "SAM: $5B", "SOM: $500M"]),
        ("Risk Assessment", ["Technical debt", "Resource constraints", "Market changes"]),
        ("Budget Review", ["Planned: $2M", "Actual: $1.8M", "Variance: -10%"]),
        ("Customer Feedback", ["NPS: 72", "CSAT: 4.5/5", "Churn: 2%"]),
        ("Competitive Analysis", ["Feature comparison", "Pricing strategy", "Market position"]),
        ("Training Materials", ["Onboarding", "Best practices", "FAQ"]),
        ("Quarterly Review", ["Achievements", "Challenges", "Next steps"]),
    ]
    for title, bullets in topics:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = "Prepared by Oxi Team"
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"{title} - Part {i+1}"
            body = slide.placeholders[1]
            tf = body.text_frame
            for b in bullets:
                p = tf.add_paragraph()
                p.text = b
        safe = re.sub(r'[^\w]', '_', title.lower())
        fp = pptx_dir / f"gen_{safe}.pptx"
        prs.save(str(fp))
        files.append(fp)

    print(f"  Generated {len(files)} pptx files")
    return files

def main():
    output_dir = Path("./documents")
    output_dir.mkdir(parents=True, exist_ok=True)

    print("=== Generating Test OOXML Files ===")
    all_files = []
    all_files.extend(generate_docx_files(output_dir))
    all_files.extend(generate_xlsx_files(output_dir))
    all_files.extend(generate_pptx_files(output_dir))

    # Update manifest
    manifest_path = output_dir / "manifest.json"
    existing = []
    existing_hashes = set()
    if manifest_path.exists():
        data = json.loads(manifest_path.read_text())
        existing = data.get("documents", [])
        existing_hashes = {d["hash"] for d in existing}

    collected = list(existing)
    added = 0
    for fp in all_files:
        content = fp.read_bytes()
        file_hash = hashlib.md5(content).hexdigest()[:12]
        if file_hash in existing_hashes:
            continue
        existing_hashes.add(file_hash)
        ext = fp.suffix.lower().lstrip('.')
        collected.append({
            "filename": fp.name,
            "source_url": "generated",
            "format": ext,
            "size_bytes": len(content),
            "hash": file_hash,
        })
        added += 1

    counts = {}
    for d in collected:
        counts[d["format"]] = counts.get(d["format"], 0) + 1

    manifest = {"total": len(collected), "counts": counts, "documents": collected}
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2))

    print(f"\nAdded: {added} generated files")
    print(f"Total: {len(collected)} (docx:{counts.get('docx',0)} xlsx:{counts.get('xlsx',0)} pptx:{counts.get('pptx',0)})")

if __name__ == "__main__":
    main()
