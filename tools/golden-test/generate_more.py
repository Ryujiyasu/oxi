#!/usr/bin/env python3
"""Generate more diverse OOXML test files to reach 500+ total."""
import hashlib, json, os, re, sys, random
from pathlib import Path

random.seed(12345)

def gen_docx_batch(output_dir, count=80):
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        print("  python-docx not installed")
        return []

    files = []
    docx_dir = output_dir / "docx"
    docx_dir.mkdir(parents=True, exist_ok=True)

    # Business document templates
    jp_titles = [
        "取締役会議事録", "株主総会招集通知", "事業計画書", "予算申請書",
        "人事評価シート", "研修報告書", "出張報告書", "業務改善提案書",
        "品質管理報告書", "安全衛生委員会議事録", "防災計画書", "情報セキュリティポリシー",
        "個人情報保護方針", "コンプライアンス研修資料", "内部監査報告書", "リスク管理報告書",
        "新商品企画書", "マーケティング戦略書", "顧客満足度調査報告", "社内規程集",
        "就業規則", "給与規程", "退職金規程", "育児休業規程",
        "テレワーク規程", "ハラスメント防止規程", "副業・兼業規程", "懲戒規程",
        "旅費規程", "慶弔見舞金規程", "車両管理規程", "文書管理規程",
        "経費精算マニュアル", "受注管理マニュアル", "クレーム対応マニュアル", "電話対応マニュアル",
        "新入社員ガイド", "管理職ハンドブック", "営業マニュアル", "品質基準書",
    ]
    en_titles = [
        "Board Meeting Minutes", "Annual Report", "Strategic Plan",
        "Budget Proposal", "Performance Review", "Training Report",
        "Travel Report", "Process Improvement", "Quality Report",
        "Safety Committee Minutes", "Disaster Plan", "Security Policy",
        "Privacy Policy", "Compliance Training", "Audit Report",
        "Risk Management Report", "Product Plan", "Marketing Strategy",
        "Customer Survey Results", "Company Handbook",
        "Employee Agreement", "NDA Template", "Service Contract",
        "License Agreement", "Terms of Service", "Privacy Notice",
        "Data Processing Agreement", "SLA Template", "SOW Template",
        "Change Request Form", "Incident Report", "Root Cause Analysis",
        "Project Charter", "Requirements Document", "Test Plan",
        "Release Notes", "User Guide", "API Documentation",
        "Architecture Decision Record", "Technical Specification",
    ]

    all_titles = jp_titles + en_titles
    paragraphs_jp = [
        "本文書は社内利用を目的として作成されたものです。",
        "関係各位におかれましては、ご確認の上、ご対応をお願いいたします。",
        "以下の通り報告いたします。詳細については添付資料をご参照ください。",
        "本件について、下記の通り決議されましたのでお知らせいたします。",
        "上記の件について、慎重に審議した結果、原案通り承認されました。",
        "今後とも引き続きご協力を賜りますようお願い申し上げます。",
        "不明な点がございましたら、担当者までお問い合わせください。",
        "本規程は令和6年4月1日より施行するものとします。",
        "全従業員はこの方針を遵守し、適切な行動を取るものとします。",
        "定期的な見直しを行い、必要に応じて改定するものとします。",
    ]
    paragraphs_en = [
        "This document has been prepared for internal use only.",
        "Please review and provide your feedback by the end of the week.",
        "The following report summarizes our findings and recommendations.",
        "We are pleased to present the results of our quarterly analysis.",
        "Based on our assessment, we recommend the following course of action.",
        "Please do not hesitate to contact us if you have any questions.",
        "This policy is effective immediately and supersedes all previous versions.",
        "All employees are expected to comply with these guidelines.",
        "Regular reviews will be conducted to ensure continued effectiveness.",
        "For additional information, please refer to the appendix.",
    ]

    for i in range(min(count, len(all_titles))):
        doc = Document()
        title = all_titles[i]
        doc.add_heading(title, 0)

        is_jp = i < len(jp_titles)
        paras = paragraphs_jp if is_jp else paragraphs_en

        # Add 3-8 sections
        num_sections = random.randint(3, 8)
        section_names_jp = ["概要", "目的", "背景", "詳細", "対象範囲", "手順", "注意事項", "結論", "今後の予定", "添付資料"]
        section_names_en = ["Overview", "Purpose", "Background", "Details", "Scope", "Procedure", "Notes", "Conclusion", "Next Steps", "Appendix"]
        sections = section_names_jp if is_jp else section_names_en

        for s in range(num_sections):
            doc.add_heading(sections[s % len(sections)], random.choice([1, 2]))
            for _ in range(random.randint(1, 4)):
                doc.add_paragraph(random.choice(paras))

            # Sometimes add a table
            if random.random() < 0.3:
                rows = random.randint(3, 6)
                cols = random.randint(2, 4)
                table = doc.add_table(rows=rows, cols=cols)
                table.style = "Table Grid"
                for r in range(rows):
                    for c in range(cols):
                        table.cell(r, c).text = f"Cell {r+1}-{c+1}"

            # Sometimes add a list
            if random.random() < 0.3:
                for _ in range(random.randint(2, 5)):
                    doc.add_paragraph(
                        random.choice(paras)[:40],
                        style="List Bullet" if random.random() < 0.5 else "List Number"
                    )

        safe = re.sub(r'[^\w]', '_', title)[:30]
        fp = docx_dir / f"gen2_{i:03d}_{safe}.docx"
        doc.save(str(fp))
        files.append(fp)

    print(f"  Generated {len(files)} docx files")
    return files

def gen_xlsx_batch(output_dir, count=80):
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    except ImportError:
        print("  openpyxl not installed")
        return []

    files = []
    xlsx_dir = output_dir / "xlsx"
    xlsx_dir.mkdir(parents=True, exist_ok=True)

    sheet_themes = [
        ("売上月次推移", ["年月", "売上高", "原価", "粗利", "粗利率"]),
        ("従業員名簿", ["社員番号", "氏名", "部門", "役職", "入社日"]),
        ("在庫管理表", ["商品コード", "商品名", "在庫数", "単価", "金額"]),
        ("顧客一覧", ["顧客ID", "会社名", "担当者", "電話番号", "メール"]),
        ("プロジェクト管理", ["タスク", "担当", "開始日", "期限", "進捗率"]),
        ("経費精算", ["日付", "項目", "金額", "備考", "承認"]),
        ("勤怠管理", ["日付", "出勤", "退勤", "休憩", "実働時間"]),
        ("アンケート集計", ["質問", "回答A", "回答B", "回答C", "回答D"]),
        ("Revenue Report", ["Month", "Revenue", "Cost", "Profit", "Margin"]),
        ("Employee Directory", ["ID", "Name", "Department", "Position", "Start Date"]),
        ("Inventory", ["SKU", "Product", "Quantity", "Unit Price", "Total"]),
        ("Customer List", ["Customer ID", "Company", "Contact", "Phone", "Email"]),
        ("Project Tracker", ["Task", "Owner", "Start", "Due", "Status"]),
        ("Expense Report", ["Date", "Item", "Amount", "Category", "Approved"]),
        ("Time Sheet", ["Date", "Clock In", "Clock Out", "Break", "Hours"]),
        ("Survey Results", ["Question", "Strongly Agree", "Agree", "Disagree", "Strongly Disagree"]),
        ("Budget Plan", ["Category", "Q1", "Q2", "Q3", "Q4"]),
        ("KPI Dashboard", ["Metric", "Target", "Actual", "Variance", "Status"]),
        ("Risk Register", ["Risk ID", "Description", "Probability", "Impact", "Mitigation"]),
        ("Vendor List", ["Vendor", "Contact", "Category", "Rating", "Contract End"]),
    ]

    for i in range(count):
        wb = Workbook()
        theme_idx = i % len(sheet_themes)
        title, headers = sheet_themes[theme_idx]

        ws = wb.active
        ws.title = title[:31]  # Excel sheet name limit

        # Header row
        header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        for c, h in enumerate(headers, 1):
            cell = ws.cell(row=1, column=c, value=h)
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = header_fill

        # Data rows
        rows = random.randint(10, 80)
        for r in range(2, rows + 2):
            for c in range(1, len(headers) + 1):
                if random.random() < 0.4:
                    ws.cell(row=r, column=c, value=random.randint(1, 99999))
                elif random.random() < 0.3:
                    ws.cell(row=r, column=c, value=round(random.uniform(0.01, 999.99), 2))
                else:
                    ws.cell(row=r, column=c, value=f"Data_{r}_{c}")

        # Maybe add a second sheet
        if random.random() < 0.4:
            ws2 = wb.create_sheet("Summary")
            ws2["A1"] = "Total"
            ws2["B1"] = f"=SUM('{title[:31]}'!B2:B{rows+1})"

        fp = xlsx_dir / f"gen2_{i:03d}_{re.sub(r'[^\\w]', '_', title)[:20]}.xlsx"
        wb.save(str(fp))
        files.append(fp)

    print(f"  Generated {len(files)} xlsx files")
    return files

def gen_pptx_batch(output_dir, count=50):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("  python-pptx not installed")
        return []

    files = []
    pptx_dir = output_dir / "pptx"
    pptx_dir.mkdir(parents=True, exist_ok=True)

    topics = [
        ("DX推進計画", ["現状分析", "課題整理", "ソリューション提案", "ロードマップ", "予算計画"]),
        ("新規事業提案", ["市場分析", "ビジネスモデル", "競合分析", "収益計画", "リスク分析"]),
        ("AI活用戦略", ["AIトレンド", "活用領域", "導入計画", "人材育成", "KPI設定"]),
        ("サステナビリティ報告", ["環境目標", "社会貢献", "ガバナンス", "ESG指標", "今後の取組"]),
        ("働き方改革", ["現状の課題", "改善施策", "テレワーク推進", "生産性指標", "アクションプラン"]),
        ("サイバーセキュリティ", ["脅威動向", "対策状況", "インシデント対応", "教育訓練", "投資計画"]),
        ("クラウド移行計画", ["現行システム", "移行対象", "アーキテクチャ", "スケジュール", "コスト試算"]),
        ("組織改革提案", ["組織課題", "新体制案", "移行計画", "人材配置", "期待効果"]),
        ("品質改善活動", ["品質データ", "原因分析", "改善施策", "効果測定", "標準化"]),
        ("グローバル展開", ["市場調査", "進出計画", "パートナー戦略", "リスク管理", "投資計画"]),
        ("Digital Transformation", ["Current State", "Challenges", "Solutions", "Roadmap", "Budget"]),
        ("Product Launch", ["Market Opportunity", "Product Features", "Go-to-Market", "Pricing", "Timeline"]),
        ("Q1 Business Review", ["Financial Summary", "Key Achievements", "Challenges", "Customer Wins", "Outlook"]),
        ("Engineering Update", ["Sprint Review", "Technical Debt", "Architecture", "Performance", "Security"]),
        ("Investor Pitch", ["Problem", "Solution", "Market Size", "Business Model", "Ask"]),
        ("Onboarding Program", ["Company Overview", "Culture", "Tools", "Processes", "Resources"]),
        ("Data Analytics Report", ["Data Sources", "Key Metrics", "Trends", "Insights", "Recommendations"]),
        ("Partnership Proposal", ["About Us", "Synergies", "Proposal", "Benefits", "Next Steps"]),
        ("Technology Roadmap", ["Vision", "Current Stack", "Target Architecture", "Migration Plan", "Milestones"]),
        ("Annual Planning", ["Year in Review", "Lessons Learned", "Goals", "Strategy", "Resources"]),
        ("UX Research Findings", ["Methodology", "User Personas", "Pain Points", "Opportunities", "Recommendations"]),
        ("Supply Chain Review", ["Current State", "Bottlenecks", "Optimization", "Risk Mitigation", "Action Plan"]),
        ("Brand Strategy", ["Brand Analysis", "Target Audience", "Positioning", "Messaging", "Channels"]),
        ("Compliance Update", ["Regulatory Changes", "Impact Assessment", "Action Items", "Training", "Timeline"]),
        ("Innovation Workshop", ["Design Thinking", "Ideation", "Prototyping", "Testing", "Scaling"]),
    ]

    for i in range(min(count, len(topics) * 2)):
        prs = Presentation()
        topic_idx = i % len(topics)
        title, sections = topics[topic_idx]

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = f"Version {i+1}.0"

        # Content slides
        for section in sections:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = section
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = f"{section}に関する詳細情報" if topic_idx < 10 else f"Details about {section}"
            for j in range(random.randint(2, 5)):
                p = tf.add_paragraph()
                p.text = f"Point {j+1}: {'具体的な内容をここに記載' if topic_idx < 10 else 'Specific details here'}"

        # Summary slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "まとめ" if topic_idx < 10 else "Summary"
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = "Key takeaways"
        for s in sections[:3]:
            p = tf.add_paragraph()
            p.text = f"- {s}"

        safe = re.sub(r'[^\w]', '_', title)[:25]
        fp = pptx_dir / f"gen2_{i:03d}_{safe}.pptx"
        prs.save(str(fp))
        files.append(fp)

    print(f"  Generated {len(files)} pptx files")
    return files

def main():
    output_dir = Path("./documents")
    output_dir.mkdir(parents=True, exist_ok=True)

    print("=== Generating More Test Files ===")
    all_files = []
    all_files.extend(gen_docx_batch(output_dir, 80))
    all_files.extend(gen_xlsx_batch(output_dir, 80))
    all_files.extend(gen_pptx_batch(output_dir, 50))

    # Update manifest
    manifest_path = output_dir / "manifest.json"
    existing = []
    existing_hashes = set()
    if manifest_path.exists():
        data = json.loads(manifest_path.read_text())
        existing = data.get("documents", [])
        existing_hashes = {d["hash"] for d in existing}

    collected = list(existing)
    added = 0
    for fp in all_files:
        content = fp.read_bytes()
        file_hash = hashlib.md5(content).hexdigest()[:12]
        if file_hash in existing_hashes:
            continue
        existing_hashes.add(file_hash)
        ext = fp.suffix.lower().lstrip('.')
        collected.append({
            "filename": fp.name,
            "source_url": "generated",
            "format": ext,
            "size_bytes": len(content),
            "hash": file_hash,
        })
        added += 1

    counts = {}
    for d in collected:
        counts[d["format"]] = counts.get(d["format"], 0) + 1

    manifest = {"total": len(collected), "counts": counts, "documents": collected}
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2))

    print(f"\nAdded: {added} generated files")
    print(f"Total: {len(collected)} (docx:{counts.get('docx',0)} xlsx:{counts.get('xlsx',0)} pptx:{counts.get('pptx',0)})")

if __name__ == "__main__":
    main()
