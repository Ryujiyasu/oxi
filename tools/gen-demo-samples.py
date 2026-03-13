#!/usr/bin/env python3
"""Generate demo sample files for Oxi Twitter demo video."""

from docx import Document
from docx.shared import Pt, Inches, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
import os

DOCS_DIR = os.path.join(os.path.dirname(__file__), '..', 'docs')


def make_demo_docx():
    """Professional-looking Japanese business document."""
    doc = Document()

    # -- Page setup (A4) --
    section = doc.sections[0]
    section.page_width = Cm(21.0)
    section.page_height = Cm(29.7)
    section.top_margin = Cm(2.5)
    section.bottom_margin = Cm(2.5)
    section.left_margin = Cm(3.0)
    section.right_margin = Cm(3.0)

    # -- Title --
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run('Oxi プロジェクト提案書')
    run.font.size = Pt(22)
    run.font.color.rgb = RGBColor(91, 79, 199)  # Oxi violet
    run.bold = True

    # -- Subtitle --
    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run('Rust + WebAssembly によるドキュメント処理スイート')
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(100, 100, 100)

    # -- Date & Author --
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = meta.add_run('2026年3月13日\n作成者: 安龍二')
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(120, 120, 120)

    doc.add_paragraph()  # spacer

    # -- Section 1 --
    h1 = doc.add_heading('1. エグゼクティブサマリー', level=1)

    doc.add_paragraph(
        'Oxi は、Microsoft Office ファイル (.docx, .xlsx, .pptx) および PDF を '
        'ブラウザ上でネイティブに処理するオープンソースのドキュメントスイートです。'
        'Rust で実装されたコアエンジンを WebAssembly にコンパイルし、'
        'サーバーへのデータ送信なしにすべての処理をクライアントサイドで完結します。'
    )

    # Key points as bullet list
    for point in [
        '100% クライアントサイド処理 — データはブラウザの外に出ない',
        '官公庁ファイル 90 件のパース成功率 100%',
        'デジタル判子 (印鑑) 生成と PAdES 電子署名',
        '日本語組版: 禁則処理 (JIS X 4051) 対応',
        'WASM バイナリサイズ: 約 1.4 MB',
    ]:
        p = doc.add_paragraph(point, style='List Bullet')

    doc.add_paragraph()  # spacer

    # -- Section 2 --
    doc.add_heading('2. 技術アーキテクチャ', level=1)

    doc.add_paragraph(
        'Oxi は以下のクレート構成で設計されています。各クレートは独立してテスト可能で、'
        'WASM ターゲットとネイティブターゲットの両方で動作します。'
    )

    # Architecture table
    table = doc.add_table(rows=8, cols=2, style='Light Grid Accent 1')
    table.alignment = WD_TABLE_ALIGNMENT.CENTER

    headers = ['クレート', '役割']
    for i, h in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = h
        for paragraph in cell.paragraphs:
            for run in paragraph.runs:
                run.bold = True

    data = [
        ('oxi-common', '共通 OOXML ユーティリティ (ZIP, XML)'),
        ('oxidocs-core', '.docx エンジン — パーサー, IR, レイアウト'),
        ('oxicells-core', '.xlsx エンジン — パーサー, IR, エディタ'),
        ('oxislides-core', '.pptx エンジン — パーサー, IR, エディタ'),
        ('oxipdf-core', 'PDF 1.7 エンジン — パーサー, 署名, 生成'),
        ('oxihanko', '判子 (SVG) 生成 + PAdES 署名'),
        ('oxi-wasm', 'WebAssembly バインディング'),
    ]
    for i, (crate, role) in enumerate(data):
        table.rows[i + 1].cells[0].text = crate
        table.rows[i + 1].cells[1].text = role

    doc.add_paragraph()

    # -- Section 3 --
    doc.add_heading('3. 判子 (Hanko) 機能', level=1)

    doc.add_paragraph(
        'oxihanko クレートは、日本のビジネスに不可欠な印鑑をデジタル化します。'
        '丸印・角印・小判型の 3 スタイルに対応し、SVG として高品質に出力。'
        'PDF 電子署名と連携することで、可視スタンプ付きの PAdES 署名を実現します。'
    )

    # Hanko styles table
    table2 = doc.add_table(rows=5, cols=3, style='Light Grid Accent 1')
    table2.alignment = WD_TABLE_ALIGNMENT.CENTER

    for i, h in enumerate(['スタイル', '用途', '特徴']):
        cell = table2.rows[0].cells[i]
        cell.text = h
        for paragraph in cell.paragraphs:
            for run in paragraph.runs:
                run.bold = True

    hanko_data = [
        ('丸印 (Round)', '個人印・認印', '2文字→横書き, 3文字以上→縦書き'),
        ('角印 (Square)', '会社印・法人印', '2×2 グリッド, 右→左配置'),
        ('小判型 (Oval)', '銀行届出印', '楕円スタイル'),
        ('承認印', '日付入り承認', '名前 + 日付 + 区切り線'),
    ]
    for i, (style, use, feat) in enumerate(hanko_data):
        table2.rows[i + 1].cells[0].text = style
        table2.rows[i + 1].cells[1].text = use
        table2.rows[i + 1].cells[2].text = feat

    doc.add_paragraph()

    # -- Section 4 --
    doc.add_heading('4. ロードマップ', level=1)

    doc.add_paragraph(
        '現在 v1 (基盤) フェーズが完了し、v2 (コラボレーション) に向けて開発中です。'
    )

    for item in [
        'v1 (現在): OOXML パース/レンダリング/編集, PDF, 判子',
        'v2: CRDT リアルタイム共同編集, AI アシスト, E2E 暗号化',
        'v3: プラグインシステム, デスクトップ/モバイルアプリ (Tauri)',
        'v4: エンタープライズ — コンプライアンス, 業界特化',
    ]:
        doc.add_paragraph(item, style='List Bullet')

    doc.add_paragraph()

    # -- Section 5 --
    doc.add_heading('5. ライセンス', level=1)

    doc.add_paragraph(
        'MIT License — すべてのソースコードはオープンソースで提供されます。'
        'サードパーティクレートはすべて MIT 互換ライセンスです。'
    )

    # -- Footer-like --
    doc.add_paragraph()
    footer = doc.add_paragraph()
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = footer.add_run('— Oxi: Open-source document suite powered by Rust + WebAssembly —')
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(150, 150, 150)
    run.italic = True

    path = os.path.join(DOCS_DIR, 'sample.docx')
    doc.save(path)
    print(f'Created: {path}')


def make_demo_xlsx():
    """Professional spreadsheet with data."""
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = '売上レポート'

    # Column widths
    ws.column_dimensions['A'].width = 18
    ws.column_dimensions['B'].width = 16
    ws.column_dimensions['C'].width = 16
    ws.column_dimensions['D'].width = 16
    ws.column_dimensions['E'].width = 18

    # Title
    ws.merge_cells('A1:E1')
    ws['A1'] = '2026年度 四半期売上レポート'
    ws['A1'].font = Font(size=16, bold=True, color='5B4FC7')
    ws['A1'].alignment = Alignment(horizontal='center')

    ws.merge_cells('A2:E2')
    ws['A2'] = 'Oxi Document Suite — サンプルデータ'
    ws['A2'].font = Font(size=10, color='888888')
    ws['A2'].alignment = Alignment(horizontal='center')

    # Headers
    headers = ['製品', 'Q1 (万円)', 'Q2 (万円)', 'Q3 (万円)', 'Q4 (万円)']
    header_fill = PatternFill(start_color='5B4FC7', end_color='5B4FC7', fill_type='solid')
    header_font = Font(bold=True, color='FFFFFF', size=11)
    thin_border = Border(
        left=Side(style='thin'), right=Side(style='thin'),
        top=Side(style='thin'), bottom=Side(style='thin')
    )

    for col, h in enumerate(headers, 1):
        cell = ws.cell(row=4, column=col, value=h)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal='center')
        cell.border = thin_border

    # Data
    data = [
        ('OxiDocs (docx)', 1250, 1480, 1720, 1950),
        ('OxiCells (xlsx)', 980, 1100, 1350, 1580),
        ('OxiSlides (pptx)', 450, 620, 780, 920),
        ('OxiPDF', 320, 480, 650, 810),
        ('OxiHanko (判子)', 180, 350, 520, 750),
    ]

    alt_fill = PatternFill(start_color='F0EEFF', end_color='F0EEFF', fill_type='solid')
    num_fmt = '#,##0'

    for i, (product, q1, q2, q3, q4) in enumerate(data):
        row = 5 + i
        values = [product, q1, q2, q3, q4]
        for col, val in enumerate(values, 1):
            cell = ws.cell(row=row, column=col, value=val)
            cell.border = thin_border
            if col >= 2:
                cell.number_format = num_fmt
                cell.alignment = Alignment(horizontal='right')
            if i % 2 == 1:
                cell.fill = alt_fill

    # Total row
    total_row = 5 + len(data)
    ws.cell(row=total_row, column=1, value='合計').font = Font(bold=True)
    ws.cell(row=total_row, column=1).border = thin_border
    for col in range(2, 6):
        cell = ws.cell(row=total_row, column=col)
        cell.value = sum(data[i][col - 1] for i in range(len(data)))
        cell.font = Font(bold=True)
        cell.number_format = num_fmt
        cell.alignment = Alignment(horizontal='right')
        cell.border = thin_border
        total_fill = PatternFill(start_color='E8E5F5', end_color='E8E5F5', fill_type='solid')
        cell.fill = total_fill

    # Note
    ws.cell(row=total_row + 2, column=1, value='※ 上記はデモ用のサンプルデータです').font = Font(size=9, color='999999')

    path = os.path.join(DOCS_DIR, 'sample.xlsx')
    wb.save(path)
    print(f'Created: {path}')


def make_demo_pptx():
    """Simple demo presentation."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor as PptxColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = 'Oxi'
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = PptxColor(91, 79, 199)
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = 'Open-source document suite — Rust + WebAssembly'
    p2.font.size = Pt(24)
    p2.font.color.rgb = PptxColor(100, 100, 100)
    p2.alignment = PP_ALIGN.CENTER

    # Slide 2: Features
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox2 = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf2 = txBox2.text_frame
    p = tf2.paragraphs[0]
    p.text = '主な機能'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PptxColor(91, 79, 199)

    features = [
        '.docx / .xlsx / .pptx / PDF をブラウザで処理',
        '判子 (Hanko) — デジタル印鑑 SVG 生成',
        'PAdES 電子署名',
        '禁則処理 (JIS X 4051)',
        'サーバー不要・100% クライアントサイド',
    ]
    txBox3 = slide2.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10), Inches(4))
    tf3 = txBox3.text_frame
    for i, feat in enumerate(features):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p.text = f'  {feat}'
        p.font.size = Pt(22)
        p.space_after = Pt(12)

    # Slide 3: Roadmap
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox4 = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf4 = txBox4.text_frame
    p = tf4.paragraphs[0]
    p.text = 'ロードマップ'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PptxColor(91, 79, 199)

    roadmap = [
        ('v1 (現在)', '基盤 — OOXML パース, 編集, PDF, 判子'),
        ('v2', 'コラボレーション — CRDT, AI, E2E暗号化'),
        ('v3', 'プラットフォーム — プラグイン, Tauri'),
        ('v4', 'エンタープライズ — コンプライアンス'),
    ]
    txBox5 = slide3.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10), Inches(4))
    tf5 = txBox5.text_frame
    for i, (ver, desc) in enumerate(roadmap):
        p = tf5.paragraphs[0] if i == 0 else tf5.add_paragraph()
        run1 = p.add_run()
        run1.text = f'{ver}: '
        run1.font.size = Pt(22)
        run1.font.bold = True
        run1.font.color.rgb = PptxColor(91, 79, 199)
        run2 = p.add_run()
        run2.text = desc
        run2.font.size = Pt(22)
        p.space_after = Pt(16)

    path = os.path.join(DOCS_DIR, 'sample.pptx')
    prs.save(path)
    print(f'Created: {path}')


if __name__ == '__main__':
    make_demo_docx()
    try:
        make_demo_xlsx()
    except ImportError:
        print('Skipping xlsx (openpyxl not installed)')
    try:
        make_demo_pptx()
    except ImportError:
        print('Skipping pptx (python-pptx not installed)')

    print('\nDone! Files are in docs/')
