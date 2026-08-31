# -*- coding: utf-8 -*-
"""Does `spcFirstLastPara` make PowerPoint honour the FIRST paragraph's spcBef?

`a:bodyPr/@spcFirstLastPara` is not read anywhere in this codebase, and the
layout drops the first paragraph's `spcBef` unconditionally. Three corpus decks
say that is wrong, and say it in a way that tracks rather than being constant:

    d06 s2   declares 10.0pt   PowerPoint draws 9.815pt lower
    d35 s2   declares 10.0pt                    9.834pt
    d16 s26  declares  6.0pt                    6.065pt

This isolates it. Each slide holds ONE text box with two paragraphs, identical
but for the flag and the amount, so the first baseline answers the question on
its own:

    arm    flag  spcBef     prediction if the flag means what it says
    off-0    0     0pt      first baseline at the box top + the face's ascent
    off-6    0     6pt      the same -- the flag is off, the space is dropped
    on-0     1     0pt      the same -- there is no space to honour
    on-6     1     6pt      6pt lower
    ...

The second paragraph's baseline is measured too: `spcBef` between paragraphs is
already implemented, so if the arms disagree THERE the probe is testing two
things at once and the reading is not clean.

    python tools/metrics/gen_pptx_spcfirst_probe.py
    python tools/metrics/export_pptx_spcfirst_probe.py
"""
import os
import sys

from pptx import Presentation
from pptx.util import Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

BASE = os.path.join("pipeline_data", "pptx_probes", "spcfirst_probe")
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# (flag, spcBef in points)
ARMS = [
    (0, 0.0), (0, 6.0), (0, 10.0), (0, 18.0),
    (1, 0.0), (1, 6.0), (1, 10.0), (1, 18.0),
]

SIZE_PT = 24.0
BOX = (72.0, 72.0, 400.0, 300.0)


def sub(parent, tag):
    from lxml import etree
    el = etree.SubElement(parent, f"{{{NS_A}}}{tag}")
    return el


def main() -> None:
    os.makedirs(BASE, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(540)
    blank = prs.slide_layouts[6]

    for flag, spc in ARMS:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Pt(BOX[0]), Pt(BOX[1]), Pt(BOX[2]), Pt(BOX[3]))
        tf = box.text_frame
        tf.word_wrap = True

        body = tf._txBody.find(f"{{{NS_A}}}bodyPr")
        # No insets at all, top anchored, so the first baseline is the only
        # thing the arms can move.
        for attr, value in (("lIns", "0"), ("rIns", "0"), ("tIns", "0"),
                            ("bIns", "0"), ("anchor", "t"), ("wrap", "square"),
                            ("spcFirstLastPara", str(flag))):
            body.set(attr, value)

        for i, text in enumerate(("ALPHA", "BETA")):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = para.add_run()
            run.text = text
            run.font.size = Pt(SIZE_PT)
            run.font.name = "Arial"
            run.font.bold = False
            # spcBef on BOTH paragraphs: the first tests the flag, the second
            # is the control that says the ordinary path still works.
            ppr = para._p.get_or_add_pPr()
            bef = sub(ppr, "spcBef")
            pts = sub(bef, "spcPts")
            pts.set("val", str(int(round(spc * 100))))

    out = os.path.join(BASE, "spcfirst_probe.pptx")
    prs.save(out)
    print(f"wrote {out} with {len(ARMS)} arms")
    for i, (flag, spc) in enumerate(ARMS, 1):
        print(f"   slide {i}: spcFirstLastPara={flag} spcBef={spc:g}pt")


if __name__ == "__main__":
    main()
