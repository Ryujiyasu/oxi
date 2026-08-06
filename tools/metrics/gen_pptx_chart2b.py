# -*- coding: utf-8 -*-
"""chart2b: ONE series named 'Revenue' (to verify: 1-series auto title uses the
series NAME, not a special 'Series 1' literal; and bar width = pitch/(1+1.5))."""
import sys
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

chart_data = CategoryChartData()
chart_data.categories = ["Q1", "Q2", "Q3"]
chart_data.add_series("Revenue", (19.2, 21.4, 16.7))

slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1.0), Inches(1.0), Inches(5.5), Inches(4.0),
    chart_data,
)

out = r"pipeline_data\pptx_probes\chart2b\chart2b.pptx"
import os
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print("saved", out)
