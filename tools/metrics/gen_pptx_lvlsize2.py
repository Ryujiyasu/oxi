# -*- coding: utf-8 -*-
"""Separate the candidates for what a silent run's size falls back to.

d19 s4 (endParaRPr EMPTY, level 46, sibling 24) resolves silent runs to 46 --
the LEVEL. d04 s18 (endParaRPr sz=2400, level 14, sibling 24, and a placeholder
whose idx is 4294967295) resolves its silent run to 24 -- NOT the level. Two
things differ at once, so each arm here changes exactly one:

    endpara_differs   valid title ph, sibling 24, endParaRPr 32
                      -> 44 = level wins; 32 = endParaRPr wins; 24 = sibling
    endpara_body      valid body ph, sibling 12, endParaRPr 36
    bogus_idx         body ph idx=4294967295, sibling 24, NO endParaRPr
                      -> does the broken match alone change the fallback?
    bogus_idx_endpara body ph idx=4294967295, sibling 24, endParaRPr 24
                      -> d04's exact shape
    control           valid title ph, sibling 24, no endParaRPr (expect 44)

    python tools/metrics/gen_pptx_lvlsize2.py
    python tools/metrics/read_pptx_lvlsize_com.py --plan lvlsize2
"""
from __future__ import annotations

import json
import os
import re
import shutil
import zipfile

from pptx import Presentation
from pptx.util import Pt

OUT = os.path.join("tools", "metrics", "lvlsize2.pptx")
PLAN = os.path.join("tools", "metrics", "lvlsize2.json")

# label, layout, ph index, run sizes, endParaRPr size (pt), break the ph idx?
CASES = [
    ("endpara_differs", 0, 0, [None, 24, None], 32, False),
    ("endpara_body", 1, 1, [None, 12, None], 36, False),
    ("bogus_idx", 1, 1, [None, 24, None], None, True),
    ("bogus_idx_endpara", 1, 1, [None, 24, None], 24, True),
    ("control", 0, 0, [None, 24, None], None, False),
]
WORDS = ["alpha", "beta", "gamma"]


def main() -> None:
    pres = Presentation()
    plan = []
    for label, layout, ph_idx, sizes, endpara, bogus in CASES:
        slide = pres.slides.add_slide(pres.slide_layouts[layout])
        ph = slide.placeholders[ph_idx]
        para = ph.text_frame.paragraphs[0]
        for word, size in zip(WORDS, sizes):
            run = para.add_run()
            run.text = word + " "
            if size is not None:
                run.font.size = Pt(size)
        plan.append({"slide": len(pres.slides), "label": label, "sizes": sizes,
                     "endpara": endpara, "bogus": bogus})
    pres.save(OUT)

    # Surgery: endParaRPr and the broken idx are not python-pptx concepts.
    tmp = OUT + ".tmp"
    with zipfile.ZipFile(OUT) as zin, zipfile.ZipFile(tmp, "w",
                                                      zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            m = re.match(r"ppt/slides/slide(\d+)\.xml$", item.filename)
            if m:
                num = int(m.group(1))
                arm = next(a for a in plan if a["slide"] == num)
                xml = data.decode("utf-8")
                if arm["endpara"]:
                    # into the paragraph that holds the arm's runs
                    xml = xml.replace(
                        "</a:r></a:p>",
                        '</a:r><a:endParaRPr lang="en" sz="%d"/></a:p>'
                        % int(arm["endpara"] * 100), 1)
                if arm["bogus"]:
                    xml = re.sub(r'(<p:ph type="body")[^/]*(/>)',
                                 r'\g<1> idx="4294967295"\g<2>', xml, count=1)
                data = xml.encode("utf-8")
            zout.writestr(item, data)
    shutil.move(tmp, OUT)
    with open(PLAN, "w", encoding="utf-8") as fh:
        json.dump(plan, fh, indent=1)
    print("wrote %s (%d slides) and %s" % (OUT, len(plan), PLAN))


if __name__ == "__main__":
    main()
