"""Generate bullet-format measurement repro (spec: bullet / numbered lists).

Slide 1: plain textbox with 3 levels of explicit bullet formats.
Slide 2: plain textbox with buAutoNum (arabicPeriod / romanLcPeriod) at lvl 0/1.
Slide 3: body placeholder (inherits slide-master default bullet) as a contrast.
"""
import sys

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from lxml import etree

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}

A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def set_bullet(paragraph, kind, val=None, level=0):
    """Inject a bullet property set into a paragraph's pPr.

    kind: "none" | "char" | "autoNum"
    val: bullet char (for char) or autoNum type (for autoNum)
    """
    p = paragraph._p
    pPr = p.get_or_add_pPr()
    pPr.set("lvl", str(level))
    # clear any existing bu* elements
    for tag in ("buNone", "buChar", "buAutoNum", "buFont", "buClr", "buSzPct"):
        for el in pPr.findall(f"{{{A}}}{tag}"):
            pPr.remove(el)
    if kind == "none":
        etree.SubElement(pPr, f"{{{A}}}buNone")
    elif kind == "char":
        buFont = etree.SubElement(pPr, f"{{{A}}}buFont")
        buFont.set("typeface", "Arial")
        buChar = etree.SubElement(pPr, f"{{{A}}}buChar")
        buChar.set("char", val)
    elif kind == "autoNum":
        buFont = etree.SubElement(pPr, f"{{{A}}}buFont")
        buFont.set("typeface", "Arial")
        buAutoNum = etree.SubElement(pPr, f"{{{A}}}buAutoNum")
        buAutoNum.set("type", val)


def add_textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(Inches(left / 72.0), Inches(top / 72.0),
                                  Inches(width / 72.0), Inches(height / 72.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(914400)   # 0.1 in
    tf.margin_right = Emu(914400)
    tf.margin_top = Emu(457200)    # 0.05 in
    tf.margin_bottom = Emu(457200)
    return tb, tf


prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ---- Slide 1: explicit buChar at 3 levels ----
s1 = prs.slides.add_slide(blank)
tb, tf = add_textbox(s1, 72, 72, 400, 360)
paras = [
    ("L0 bullet char •", "char", "\u2022", 0),
    ("L0 second bullet", "char", "\u2022", 0),
    ("L1 dash bullet", "char", "\u2013", 1),
    ("L1 second dash", "char", "\u2013", 1),
    ("L2 hyphen bullet", "char", "\u2014", 2),
    ("L2 second hyphen", "char", "\u2014", 2),
]
first = True
for text, kind, val, lvl in paras:
    if first:
        p = tf.paragraphs[0]
        first = False
    else:
        p = tf.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(18)
    set_bullet(p, kind, val, lvl)

# ---- Slide 2: buAutoNum (numbered) at 2 levels ----
s2 = prs.slides.add_slide(blank)
tb, tf = add_textbox(s2, 72, 72, 400, 360)
paras2 = [
    ("1st numbered item", "autoNum", "arabicPeriod", 0),
    ("2nd numbered item", "autoNum", "arabicPeriod", 0),
    ("3rd numbered item", "autoNum", "arabicPeriod", 0),
    ("nested a item", "autoNum", "alphaLcParenR", 1),
    ("nested b item", "autoNum", "alphaLcParenR", 1),
]
first = True
for text, kind, val, lvl in paras2:
    if first:
        p = tf.paragraphs[0]
        first = False
    else:
        p = tf.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(18)
    set_bullet(p, kind, val, lvl)

# ---- Slide 3: body placeholder (inherits master default bullets) ----
# Use the "Title and Content" layout (idx 1) which carries a content placeholder
# whose bullet formatting inherits from the slide master.
s3 = prs.slides.add_slide(prs.slide_layouts[1])
# set the title so the auto placeholders remain
s3.shapes.title.text = "Body placeholder bullets"
body = None
for ph in s3.shapes.placeholders:
    if ph.placeholder_format.idx == 1:
        body = ph
        break
if body is None:
    body = s3.shapes.placeholders[0]
tf = body.text_frame
tf.text = ""
p0 = tf.paragraphs[0]
run = p0.add_run()
run.text = "Body placeholder L0"
p1 = tf.add_paragraph()
run = p1.add_run()
run.text = "Body placeholder L1"
p1.level = 1
p2 = tf.add_paragraph()
run = p2.add_run()
run.text = "Body placeholder L2"
p2.level = 2

prs.save(sys.argv[1] if len(sys.argv) > 1 else "bullet.pptx")
print("saved")
