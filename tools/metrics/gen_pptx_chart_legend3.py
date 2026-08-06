# -*- coding: utf-8 -*-
"""Chart legend spec wave-3 second specimen: THREE-series clustered-column
chart with <c:legend> enabled (python-pptx has_legend), default theme.
Verifies the legend placement rule derived from chart_legend (2 series):
  - vertical: legend block centered on the chart shape (sy + sh/2)
  - horizontal: right-aligned block, right edge = frame_right - 10pt,
    block width = max_label_width + gap(4.62) + swatch(9.89)
  - swatch row pitch = 27.75pt
With 3 series the block is taller, so legend_y0 should shift UP if the
"centered on chart" rule holds (2-series: 197.18; 3-series predicted ~183.3)."""
import sys, os
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

base = r"pipeline_data\pptx_probes\chart_legend3"
os.makedirs(base, exist_ok=True)

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

chart_data = CategoryChartData()
chart_data.categories = ["Q1", "Q2", "Q3"]
chart_data.add_series("Revenue", (19.2, 21.4, 16.7))
chart_data.add_series("Cost", (10.5, 15.0, 12.3))
chart_data.add_series("Profit", (8.7, 6.4, 4.4))

x, y, cx, cy = Inches(1.0), Inches(1.0), Inches(5.5), Inches(4.0)
gframe = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
gframe.chart.has_legend = True

out = os.path.join(base, "chart_legend3.pptx")
prs.save(out)
print("saved:", out, os.path.getsize(out))
