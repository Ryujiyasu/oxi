# -*- coding: utf-8 -*-
"""chart3: THREE series (Revenue/Cost/Profit) to verify:
  - no auto title (series>1)
  - bar width = plot_w / (n_cat x (n_ser + 1.5))  [gapWidth default 150% of a bar]
"""
import sys
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

chart_data = CategoryChartData()
chart_data.categories = ["Q1", "Q2", "Q3"]
chart_data.add_series("Revenue", (19.2, 21.4, 16.7))
chart_data.add_series("Cost", (12.0, 15.0, 10.0))
chart_data.add_series("Profit", (7.2, 6.4, 6.7))

slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1.0), Inches(1.0), Inches(5.5), Inches(4.0),
    chart_data,
)

out = r"pipeline_data\pptx_probes\chart3\chart3.pptx"
import os
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print("saved", out)
