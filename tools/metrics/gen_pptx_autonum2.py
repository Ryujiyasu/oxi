"""buAutoNum second measurement round — resolve open questions.

Slide A (base): same textbox, [no number][autonum marL=0 indent=0]
                [autonum marL=0 indent=-36] -> pins text_area_left + indent geometry.
Slide B (font): text explicitly Arial + buFont=Times New Roman -> does buFont drive
                the number font, or does it follow the text font?
Slide C (startAt): [startAt=5][startAt=5][no startAt] -> the 5,1,2 reading.
Slide D (2-digit): arabicPeriod x11 marL=72 indent=-36 -> number width vs text start.
"""
import sys

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from lxml import etree

A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def add_textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(Inches(left / 72.0), Inches(top / 72.0),
                                  Inches(width / 72.0), Inches(height / 72.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(914400)
    tf.margin_right = Emu(914400)
    tf.margin_top = Emu(457200)
    tf.margin_bottom = Emu(457200)
    return tb, tf


def set_autonum(paragraph, kind, mar_l_pt, indent_pt, lvl=0, start_at=None,
                bufont="Arial"):
    p = paragraph._p
    pPr = p.get_or_add_pPr()
    pPr.set("lvl", str(lvl))
    pPr.set("marL", str(int(mar_l_pt * 12700)))
    pPr.set("indent", str(int(indent_pt * 12700)))
    for tag in ("buNone", "buChar", "buAutoNum", "buFont", "buClr", "buSzPct"):
        for el in pPr.findall(f"{{{A}}}{tag}"):
            pPr.remove(el)
    bf = etree.SubElement(pPr, f"{{{A}}}buFont")
    bf.set("typeface", bufont)
    bn = etree.SubElement(pPr, f"{{{A}}}buAutoNum")
    bn.set("type", kind)
    if start_at is not None:
        bn.set("startAt", str(start_at))


def add_para(tf, text, first, font=None, size=18):
    if first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    if font is not None:
        r.font.name = font
    return p


def fill(prs, paras):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb, tf = add_textbox(s, 72, 72, 400, 360)
    first = True
    for spec in paras:
        text = spec["text"]
        p = add_para(tf, text, first, spec.get("font"))
        first = False
        if "kind" in spec:
            set_autonum(p, spec["kind"], spec.get("mar", 0),
                        spec.get("ind", 0), spec.get("lvl", 0),
                        spec.get("start_at"), spec.get("bufont", "Arial"))
    return s


prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide A: base
fill(prs, [
    {"text": "Plain"},
    {"text": "Auto0", "kind": "arabicPeriod", "mar": 0, "ind": 0},
    {"text": "AutoNeg", "kind": "arabicPeriod", "mar": 0, "ind": -36},
])

# Slide B: number font vs text font
fill(prs, [
    {"text": "ArialB", "font": "Arial", "kind": "arabicPeriod", "mar": 72,
     "ind": -36, "bufont": "Times New Roman"},
    {"text": "ArialC", "font": "Arial", "kind": "arabicPeriod", "mar": 72,
     "ind": -36, "bufont": "Arial"},
])

# Slide C: startAt patterns
fill(prs, [
    {"text": "SA1", "kind": "arabicPeriod", "mar": 72, "ind": -36,
     "start_at": 5},
    {"text": "SA2", "kind": "arabicPeriod", "mar": 72, "ind": -36,
     "start_at": 5},
    {"text": "SA3", "kind": "arabicPeriod", "mar": 72, "ind": -36},
])

# Slide D: 2-digit numbers
fill(prs, [
    {"text": f"N{n}", "kind": "arabicPeriod", "mar": 72, "ind": -36}
    for n in range(1, 12)
])

prs.save(sys.argv[1] if len(sys.argv) > 1 else "autonum2.pptx")
print("saved")
