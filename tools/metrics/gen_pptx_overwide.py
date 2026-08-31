# -*- coding: utf-8 -*-
"""Where does PowerPoint put a line that is WIDER than the box it sits in?

`align_offset` clamps the centring and right-alignment offsets at zero, so a
line too wide for its box starts at the box's left edge. d32 says that is
wrong: seven slides carry a 167.65pt bullet in a 33.24pt box, `algn="ctr"`,
and PowerPoint draws it 12.7pt LEFT of the box -- exactly `(box - line) / 2`
with the negative kept.

    slide   box              line     truth        unclamped     clamped
    s2      1318.264 w33.236 58.694   1305.620     1305.535      1318.264
    s18      534.301 w33.236 58.694    521.590      521.572       534.301
    s21      533.038 w33.236 58.694    520.320      520.309       533.038

Seven arms of one deck, and all of them centred. This settles the same
question for RIGHT alignment, which the corpus does not answer at all, and
gives the rule a repro that does not depend on d32's template.

    arm  alignment  line vs box   prediction if the offset is not clamped
     1   ctr        wider         starts (line-box)/2 LEFT of the box
     2   r          wider         starts (line-box)   LEFT of the box
     3   l          wider         starts at the box (no offset either way)
     4   ctr        narrower      the ordinary case, as a control
     5   r          narrower      likewise

    python tools/metrics/gen_pptx_overwide.py
    python tools/metrics/export_pptx_overwide.py
    python tools/metrics/read_pptx_overwide.py
"""
import os
import sys

from pptx import Presentation
from pptx.util import Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

BASE = os.path.join("pipeline_data", "pptx_probes", "overwide")
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# A 100pt 'W' in Arial advances 94.4pt, so a 40pt box cannot hold it and a
# 300pt box holds it easily. `wrap="none"` keeps the line one line either way.
SIZE_PT = 100.0
BOX_LEFT = 300.0
ARMS = [
    ("ctr", 40.0),
    ("r", 40.0),
    ("l", 40.0),
    ("ctr", 300.0),
    ("r", 300.0),
]


def main() -> None:
    os.makedirs(BASE, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(540)
    blank = prs.slide_layouts[6]

    for algn, box_w in ARMS:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Pt(BOX_LEFT), Pt(200), Pt(box_w), Pt(160))
        tf = box.text_frame
        body = tf._txBody.find(f"{{{NS_A}}}bodyPr")
        for attr, value in (("lIns", "0"), ("rIns", "0"), ("tIns", "0"),
                            ("bIns", "0"), ("anchor", "t"), ("wrap", "none")):
            body.set(attr, value)
        para = tf.paragraphs[0]
        para._p.get_or_add_pPr().set("algn", algn)
        run = para.add_run()
        run.text = "W"
        run.font.size = Pt(SIZE_PT)
        run.font.name = "Arial"

    out = os.path.join(BASE, "overwide.pptx")
    prs.save(out)
    print(f"wrote {out} with {len(ARMS)} arms; box left {BOX_LEFT}pt, 'W' at {SIZE_PT:g}pt")
    for i, (algn, w) in enumerate(ARMS, 1):
        print(f"   slide {i}: algn={algn:3} box {w:g}pt")


if __name__ == "__main__":
    main()
