# -*- coding: utf-8 -*-
"""Which face does PowerPoint use when it cannot serve the one a run names?

Two findings put this question on the table and neither could answer it:

* `pptx_cff_part_census.py` -- PowerPoint REFUSES an embedded part whose
  outlines are CFF, and blind 31's body text (24pt `Open Sauce`) comes out of
  its own PDF as **Calibri**.
* `gen_pptx_cloudfont.py` -- d19 names `Nunito`, which IS in the Office cloud
  cache, and PowerPoint's PDF is Calibri throughout.

Both decks are Canva exports whose theme font is ALSO Calibri, so "the theme's
minor latin" and "Calibri, the Office default" fit the observations equally.
They are different rules for every deck that themes anything else, and Oxi has
to pick one -- dropping the refused part without knowing which cost blind 31
-0.0199 SSIM (`pipeline_data/pptx_benchmark/_cffskip_ab.log`), because GDI's own
substitute is neither.

So this authors a deck whose theme font is NOT Calibri and names families that
cannot be served, and lets PowerPoint's PDF say which face it drew.

    georgia-theme   theme major/minor latin = Georgia (installed, distinctive)
    s1  missing     'Zzyzx Nonesuch QZX'  -- installed nowhere, embedded nowhere
    s2  control     'Georgia'             -- the theme font itself
    s3  control     'Arial'               -- installed, not the theme font
    s4  missing2    'Blibbet Sans'        -- a second miss, to show s1 is a rule

    python tools/metrics/gen_pptx_missfont.py
    python tools/metrics/export_pptx_missfont.py    # PowerPoint COM -> PDF
    python tools/metrics/read_pptx_missfont.py      # which face per arm
"""
from __future__ import annotations

import re
import shutil
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

REPO = Path(__file__).resolve().parents[2]
OUT = REPO / "pipeline_data" / "pptx_probes" / "missfont"
THEME_FACE = "Georgia"

# (label, typeface) -- the label is drawn in the theme face so a reader can pair
# each span with its arm even when the arm's own face is substituted.
ARMS = [
    ("missing", "Zzyzx Nonesuch QZX"),
    ("theme", THEME_FACE),
    ("installed", "Arial"),
    ("missing2", "Blibbet Sans"),
]
TEXT = "Handgloves 0123"


def build(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(6858000)
    blank = prs.slide_layouts[6]
    for label, face in ARMS:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Emu(457200), Emu(457200), Emu(8229600), Emu(1200000))
        tf = box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{label}: {TEXT}"
        run.font.size = Pt(40)
        run.font.name = face
    prs.save(str(path))


def retheme(path: Path) -> None:
    """Point the theme's latin faces at `THEME_FACE`.

    python-pptx's default template themes Calibri, which is exactly the face
    this probe has to tell apart from the fallback, so the deck is rewritten
    part by part rather than edited in place.
    """
    tmp = path.with_suffix(".retheme.pptx")
    with zipfile.ZipFile(path) as zin, zipfile.ZipFile(
        tmp, "w", zipfile.ZIP_DEFLATED
    ) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith("ppt/theme/"):
                text = data.decode("utf-8")
                scheme = re.search(r"<a:fontScheme.*?</a:fontScheme>", text, re.S)
                if scheme:
                    fixed = re.sub(
                        r'(<a:latin typeface=")[^"]*(")',
                        rf"\1{THEME_FACE}\2",
                        scheme.group(0),
                    )
                    text = text.replace(scheme.group(0), fixed)
                data = text.encode("utf-8")
            zout.writestr(item, data)
    shutil.move(str(tmp), str(path))


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    deck = OUT / "missfont.pptx"
    build(deck)
    retheme(deck)
    with zipfile.ZipFile(deck) as z:
        theme = z.read("ppt/theme/theme1.xml").decode("utf-8")
        faces = re.findall(r"<a:latin typeface=\"([^\"]*)\"", theme)
    print(f"wrote {deck}")
    print(f"theme latin faces now: {sorted(set(faces))}")
    for label, face in ARMS:
        print(f"  {label:<10} {face}")


if __name__ == "__main__":
    main()
