# -*- coding: utf-8 -*-
"""Chart spec: the three DOUGHNUT residuals left by the first probe.

  S1: no title at all (autoTitleDeleted=1), no legend
        -> does the ring top follow the pie's sy+11?
  S2: TWO series (a second <c:ser> injected into the XML, as python-pptx
      silently drops it for pie/doughnut)
        -> concentric rings: how are the radii split?
  S3-S8: one series + legend, category labels of increasing width
        -> where does Word start wrapping a legend label to two lines,
           and what is the wrap width?

Frame 72,72,396,288, values 19.2/21.4/16.7, default Office theme."""
import sys, os, re, shutil, zipfile

sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

base = r"pipeline_data\pptx_probes\chart_doughnut_resid"
os.makedirs(base, exist_ok=True)

prs = Presentation()

# The legend label-width ramp: two-word labels so Word can wrap at the space.
RAMP = [
    ["Ab Cd", "Ef", "Gh"],
    ["Abcd Efgh", "Ef", "Gh"],
    ["Abcdefg Hijklmn", "Ef", "Gh"],
    ["Abcdefghij Klmnopqrst", "Ef", "Gh"],
    ["Abcdefghijklm Nopqrstuvwxy", "Ef", "Gh"],
    ["Abcdefghijklmnop Qrstuvwxyzabcd", "Ef", "Gh"],
]


def add(cats=None, title=True, legend=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = cats or ["East", "West", "Midwest"]
    cd.add_series("Series 1", (19.2, 21.4, 16.7))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, Pt(72), Pt(72), Pt(396), Pt(288), cd
    ).chart
    if not title:
        chart.has_title = False
    if legend:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False


add(title=False)  # S1: no title, no legend
add(title=False)  # S2: 2nd series injected below
for cats in RAMP:  # S3..S8
    add(cats=cats, legend=True)

out = os.path.join(base, "chart_doughnut_resid.pptx")
prs.save(out)

# S1/S2 get a bare <c:legend/> from python-pptx?  Strip any legend from S1/S2
# so the ring keeps the frame centre, and give S2 a second series.
SER2 = (
    '<c:ser><c:idx val="1"/><c:order val="1"/>'
    '<c:tx><c:strRef><c:f>Sheet1!$C$1</c:f><c:strCache>'
    '<c:ptCount val="1"/><c:pt idx="0"><c:v>Series 2</c:v></c:pt>'
    "</c:strCache></c:strRef></c:tx>"
    '<c:cat><c:strRef><c:f>Sheet1!$A$2:$A$4</c:f><c:strCache>'
    '<c:ptCount val="3"/>'
    '<c:pt idx="0"><c:v>East</c:v></c:pt>'
    '<c:pt idx="1"><c:v>West</c:v></c:pt>'
    '<c:pt idx="2"><c:v>Midwest</c:v></c:pt>'
    "</c:strCache></c:strRef></c:cat>"
    '<c:val><c:numRef><c:f>Sheet1!$C$2:$C$4</c:f><c:numCache>'
    "<c:formatCode>General</c:formatCode>"
    '<c:ptCount val="3"/>'
    '<c:pt idx="0"><c:v>10.5</c:v></c:pt>'
    '<c:pt idx="1"><c:v>11.2</c:v></c:pt>'
    '<c:pt idx="2"><c:v>8.5</c:v></c:pt>'
    "</c:numCache></c:numRef></c:val></c:ser>"
)

tmp = out + ".tmp"
shutil.copy(out, tmp)
zin = zipfile.ZipFile(tmp, "r")
names = zin.namelist()
charts = sorted(
    (n for n in names if re.match(r"ppt/charts/chart(\d+)\.xml$", n)),
    key=lambda n: int(re.search(r"(\d+)\.xml$", n).group(1)),
)
with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
    for n in names:
        data = zin.read(n)
        if n in (charts[0], charts[1]):
            xml = data.decode("utf-8")
            xml = re.sub(r"<c:legend.*?</c:legend>|<c:legend/>", "", xml, flags=re.S)
            if n == charts[1]:
                xml = xml.replace("</c:ser>", "</c:ser>" + SER2, 1)
                print("  slide2 series injected:", xml.count("<c:ser>"))
            data = xml.encode("utf-8")
        zout.writestr(n, data)
zin.close()
os.remove(tmp)

with zipfile.ZipFile(out) as z:
    for i, n in enumerate(charts, 1):
        x = z.read(n).decode("utf-8")
        print(
            f"  chart{i}: ser={x.count('<c:ser>')} "
            f"legend={'yes' if '<c:legend' in x else 'no'} "
            f"autoTitleDeleted="
            f"{(re.search(r'<c:autoTitleDeleted[^>]*/>', x) or ['(absent)'])[0]}"
        )

print("\nsaved:", out, os.path.getsize(out))
