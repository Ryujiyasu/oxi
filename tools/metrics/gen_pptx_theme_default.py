# gen_pptx_theme_default.py - probe: what font does PowerPoint use for a run with NO font.name?
# V1: run with no rPr at all (no font, no size) / V2: run with size only (no latin typeface)
# V3: run with explicit Arial (control) / V4: title placeholder (majorFont) / V5: body placeholder (minorFont)
# Readout: PDF span font via fitz + theme1.xml minorFont/majorFont latin typeface.
import os
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.enum.text import PP_ALIGN

OUT_DIR = os.path.abspath(r"pipeline_data\pptx_probes\theme_default")
os.makedirs(OUT_DIR, exist_ok=True)

prs = Presentation()
prs.slide_width = Emu(9144000)
prs.slide_height = Emu(6858000)

# V1: textbox run, NO font.name and NO size (bare run -> no rPr)
slide = prs.slides.add_slide(prs.slide_layouts[6])
box = slide.shapes.add_textbox(Emu(914400), Emu(914400), Inches(5), Inches(2))
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Emu(914400); tf.margin_right = Emu(914400)
tf.margin_top = Emu(457200); tf.margin_bottom = Emu(457200)
tf.vertical_anchor = 1
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = "V1 bare run no rPr"

# V2: textbox run with SIZE only (rPr has sz but no latin)
slide = prs.slides.add_slide(prs.slide_layouts[6])
box = slide.shapes.add_textbox(Emu(914400), Emu(914400), Inches(5), Inches(2))
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Emu(914400); tf.margin_right = Emu(914400)
tf.margin_top = Emu(457200); tf.margin_bottom = Emu(457200)
tf.vertical_anchor = 1
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = "V2 size only 18pt"
r.font.size = Pt(18)

# V3: textbox run with explicit Arial (control)
slide = prs.slides.add_slide(prs.slide_layouts[6])
box = slide.shapes.add_textbox(Emu(914400), Emu(914400), Inches(5), Inches(2))
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Emu(914400); tf.margin_right = Emu(914400)
tf.margin_top = Emu(457200); tf.margin_bottom = Emu(457200)
tf.vertical_anchor = 1
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = "V3 explicit Arial 18pt"
r.font.size = Pt(18); r.font.name = "Arial"

# V4: title placeholder (majorFont theme)
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "V4 title placeholder"

# V5: body placeholder (minorFont theme)
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.placeholders[1].text = "V5 body placeholder"

out = os.path.join(OUT_DIR, "theme_default.pptx")
prs.save(out)
print("saved", out)
