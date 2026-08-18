# -*- coding: utf-8 -*-
"""Probe: does PowerPoint draw a LAYOUT shape whose fill is a picture?

`parse_inherited_shapes` refused any layout/master shape carrying an
`a:blipFill` inside its `p:spPr` -- 269 of them in the dev corpus. Lifting the
refusal turned out to be corpus-inert: on d03 the newly emitted shape lands at
draw index 1 and the slide carries its own copy of the same image at index 75,
which covers it exactly. Every deck checked is like that, so the corpus cannot
say whether such a shape is meant to be drawn at all.

This deck can: its layout holds a picture-filled rectangle that NO slide
repeats, so whatever PowerPoint puts on the page is the answer.
"""
from __future__ import annotations

import re
import shutil
import sys
import zipfile
from pathlib import Path

from pptx import Presentation

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\lmfillimg").resolve()
R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def make_png() -> bytes:
    """A 4x4 checker, big enough to see and small enough to inline."""
    import struct
    import zlib

    w = h = 4
    rows = b""
    for y in range(h):
        rows += b"\x00" + b"".join(
            (b"\xff\x40\x40" if (x + y) % 2 else b"\x40\x40\xff") for x in range(w)
        )

    def chunk(tag: bytes, data: bytes) -> bytes:
        return (struct.pack(">I", len(data)) + tag + data
                + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF))

    return (b"\x89PNG\r\n\x1a\n"
            + chunk(b"IHDR", struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0))
            + chunk(b"IDAT", zlib.compress(rows))
            + chunk(b"IEND", b""))


SHAPE = (
    '<p:sp><p:nvSpPr><p:cNvPr id="900" name="fillimg"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
    '<p:spPr><a:xfrm><a:off x="914400" y="914400"/>'
    '<a:ext cx="3657600" cy="2743200"/></a:xfrm>'
    '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
    '<a:blipFill><a:blip r:embed="rIdFillImg"/><a:stretch><a:fillRect/></a:stretch>'
    "</a:blipFill></p:spPr>"
    "<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>"
)


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    base = OUT / "_base.pptx"
    prs = Presentation()
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    box = s.shapes.add_textbox(228600, 114300, 6400800, 300000)
    box.text_frame.text = "the layout below carries a picture-filled rectangle"
    prs.save(base)

    zin = zipfile.ZipFile(base)
    names = zin.namelist()
    parts = {n: zin.read(n) for n in names}
    zin.close()
    # python-pptx's blank layout is layout7 in the default template
    layout_part = next(n for n in names if n.endswith("slideLayout7.xml"))
    rels_part = layout_part.replace("slideLayouts/", "slideLayouts/_rels/") + ".rels"
    with zipfile.ZipFile(OUT / "lmfillimg.pptx", "w", zipfile.ZIP_DEFLATED) as zout:
        for n in names:
            data = parts[n]
            if n == "[Content_Types].xml":
                t = data.decode("utf-8")
                if 'Extension="png"' not in t:
                    t = t.replace("</Types>",
                                  '<Default Extension="png" ContentType="image/png"/></Types>', 1)
                data = t.encode("utf-8")
            elif n == layout_part:
                t = data.decode("utf-8")
                t = re.sub(r"(</p:spTree>)", SHAPE + r"\1", t, count=1)
                data = t.encode("utf-8")
            elif n == rels_part:
                t = data.decode("utf-8")
                t = t.replace("</Relationships>",
                              f'<Relationship Id="rIdFillImg" Type="{R}/image" '
                              'Target="../media/fillimg.png"/></Relationships>', 1)
                data = t.encode("utf-8")
            zout.writestr(n, data)
        zout.writestr("ppt/media/fillimg.png", make_png())
    base.unlink()
    print(f"wrote {OUT / 'lmfillimg.pptx'}")


if __name__ == "__main__":
    main()
