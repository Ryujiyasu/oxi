"""Generate a control repro for Spec #3: placeholder position/size resolution.

Question being measured (Ra-loop, measure-before-implement):
  When a slide placeholder has NO explicit xfrm (empty <p:spPr/>), Word uses the
  referenced slideLayout's matching placeholder xfrm. Does a slide placeholder
  WITH an explicit xfrm keep its OWN geometry (slide wins over layout)?

Design: two slides, both on slide_layouts[0] (Title Slide = layout with ctrTitle
+ subTitle placeholders).
  - slide A: default placeholders (python-pptx writes empty <p:spPr/>)  -> layout xfrm expected
  - slide B: placeholders with EXPLICIT positions (ph.left/top/width/height)
            -> slide xfrm expected (overrides layout)
Plus a text box on slide B for the fallback sanity check.
"""
import os
import sys
sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "..", "pipeline_data", "pptx_probes")
os.makedirs(OUT, exist_ok=True)
path = os.path.join(OUT, "ph_resolution.pptx")

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

layout = prs.slide_layouts[0]  # Title Slide layout (ctrTitle + subTitle ph)

# --- Slide A: default placeholders (no explicit xfrm on the slide) ---
sa = prs.slides.add_slide(layout)
sa.shapes.title.text = "Slide A Title"          # ctrTitle ph
sa.placeholders[1].text = "Slide A body line1\nSlide A body line2"  # subTitle ph

# --- Slide B: placeholders with EXPLICIT geometry (slide should win) ---
sb = prs.slides.add_slide(layout)
t = sb.shapes.title                            # ctrTitle ph
t.left = Inches(0.5)
t.top = Inches(0.75)
t.width = Inches(4.0)
t.height = Inches(1.0)
t.text = "Slide B Title"

b = sb.placeholders[1]                         # subTitle ph
b.left = Inches(1.0)
b.top = Inches(2.0)
b.width = Inches(5.0)
b.height = Inches(1.5)
b.text = "Slide B body"

tb = sb.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(3.0), Inches(1.0))
tb.text = "Slide B textbox"

prs.save(path)
print("saved:", path)
