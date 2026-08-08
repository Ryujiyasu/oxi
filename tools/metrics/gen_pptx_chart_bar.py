# -*- coding: utf-8 -*-
"""Chart spec (item 8): HORIZONTAL bar chart probe (barDir="bar").

Oxi parses <c:barDir val="bar"/> into Chart.bar_dir but the renderer has
always drawn vertical columns, so a horizontal bar chart renders wrong.
This probe exposes Word's horizontal geometry for measurement:

  S1: BAR_CLUSTERED, 1 series   (auto-title condition, mirror of chart1)
  S2: BAR_CLUSTERED, 2 series   (no auto title, mirror of chart2)
  S3: BAR_STACKED,   2 series   (stacking direction, mirror of chart_stacked)
  S4: BAR_CLUSTERED, 2 series + legend      (legend geometry)
  S5: BAR_CLUSTERED, 1 series + data labels (label placement)

All on blank slides (layout 6), default Office theme, frame 72,72,396,288,
categories [East, West, Midwest]."""
import sys, os

sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

base = r"pipeline_data\pptx_probes\chart_bar"
os.makedirs(base, exist_ok=True)

prs = Presentation()


def add(ctype, series, legend=False, dlbls=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = ["East", "West", "Midwest"]
    for name, vals in series:
        cd.add_series(name, vals)
    x, y, cx, cy = Inches(1.0), Inches(1.0), Inches(5.5), Inches(4.0)
    gframe = slide.shapes.add_chart(ctype, x, y, cx, cy, cd)
    chart = gframe.chart
    if legend:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
    if dlbls:
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.show_value = True


ONE = [("Series 1", (19.2, 21.4, 16.7))]
TWO = [("Revenue", (19.2, 21.4, 16.7)), ("Cost", (10.5, 15.0, 12.3))]

add(XL_CHART_TYPE.BAR_CLUSTERED, ONE)
add(XL_CHART_TYPE.BAR_CLUSTERED, TWO)
add(XL_CHART_TYPE.BAR_STACKED, TWO)
add(XL_CHART_TYPE.BAR_CLUSTERED, TWO, legend=True)
add(XL_CHART_TYPE.BAR_CLUSTERED, ONE, dlbls=True)

out = os.path.join(base, "chart_bar.pptx")
prs.save(out)
print("saved:", out, os.path.getsize(out))
