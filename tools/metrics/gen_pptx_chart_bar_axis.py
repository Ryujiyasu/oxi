# -*- coding: utf-8 -*-
"""HORIZONTAL bar value-axis tick sweep.

chart_bar page 3 showed the horizontal value axis uses a COARSER major unit
than the vertical one at the same range (range 40: vertical = 5-unit steps /
8 divisions, horizontal = 10-unit steps / 4 divisions).  This sweep varies the
data range so Word's horizontal division count can be read directly, plus two
frame-width arms to test whether the rule is axis-length dependent.

Each slide: BAR_CLUSTERED, 1 series, 3 categories, frame 72,72,W,288."""
import sys, os

sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

base = r"pipeline_data\pptx_probes\chart_bar_axis"
os.makedirs(base, exist_ok=True)

# (max value, frame width in pt)
ARMS = [
    (2.2, 396.0),
    (4.5, 396.0),
    (8.5, 396.0),
    (12.5, 396.0),
    (18.0, 396.0),
    (22.0, 396.0),
    (34.0, 396.0),
    (45.0, 396.0),
    (78.0, 396.0),
    (130.0, 396.0),
    (240.0, 396.0),
    (480.0, 396.0),
    # length-dependence arms at the same range
    (34.0, 250.0),
    (34.0, 600.0),
]

prs = Presentation()
for mx, fw in ARMS:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = ["A", "B", "C"]
    cd.add_series("S", (mx * 0.5, mx, mx * 0.7))
    slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Pt(72), Pt(72), Pt(fw), Pt(288), cd
    )

out = os.path.join(base, "chart_bar_axis.pptx")
prs.save(out)
print("saved:", out, "arms:", len(ARMS))
for i, (mx, fw) in enumerate(ARMS, 1):
    print(f"  slide {i:2d}: max={mx} frame_w={fw}")
