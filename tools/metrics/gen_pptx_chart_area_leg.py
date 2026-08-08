# -*- coding: utf-8 -*-
"""Chart spec: AREA chart legend band + title levers.

The first area probe (chart_area) had a legend on every arm because
python-pptx writes one by default, so it could not separate:
  * how much the legend band eats  (does plot_right track the label width?)
  * what plot_right is with NO legend
  * where the plot/legend go with an EXPLICIT title

All arms: AREA, 2 series, frame 72,72,396,288."""
import sys, os

sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

base = r"pipeline_data\pptx_probes\chart_area_leg"
os.makedirs(base, exist_ok=True)

CATS = ["Q1", "Q2", "Q3", "Q4"]
V1 = (19.2, 21.4, 16.7, 22.0)
V2 = (10.5, 11.2, 8.5, 12.3)

# (label, series names, legend?, title text or None)
ARMS = [
    ("G1 short names",   ["Ax", "Bx"],                         True,  None),
    ("G2 Ser names",     ["Ser1", "Ser2"],                     True,  None),
    ("G3 long names",    ["Alphabet", "Bracketing"],           True,  None),
    ("G4 NO legend",     ["Ser1", "Ser2"],                     False, None),
    ("G5 explicit title",["Ser1", "Ser2"],                     True,  "Quarterly"),
    ("G6 huge names",    ["Alphabetical Bracket", "Br"],       True,  None),
]

prs = Presentation()
prs.slide_width = Pt(720)
prs.slide_height = Pt(540)
for label, names, legend, title in ARMS:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = CATS
    cd.add_series(names[0], V1)
    cd.add_series(names[1], V2)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.AREA, Pt(72), Pt(72), Pt(396), Pt(288), cd
    ).chart
    if legend:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    else:
        chart.has_title = False

out = os.path.join(base, "chart_area_leg.pptx")
prs.save(out)
for i, (label, *_rest) in enumerate(ARMS, 1):
    print(f"  S{i}: {label}")
print("\nsaved:", out, os.path.getsize(out))
