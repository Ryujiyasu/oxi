# -*- coding: utf-8 -*-
"""Where the FIRST BASELINE sits when `lnSpc` is not 100%.

`Lines(1).BoundTop` answers where the line BOX starts, and the engine matches it
to the hundredth across the corpus. Where the baseline sits INSIDE that box is
not a COM property at all -- but PowerPoint's PDF export states it exactly, as
the origin of each text span, so the export is the oracle here.

The question came from deck 40 slide 1: a 216.15pt Grand Hotel title at
`lnSpc 60%`, whose box top the engine gets right and whose glyphs it draws
**5.2pt high** (both of its lines, by the same amount, so it is the baseline in
the box and not the advance -- the advance agrees with PowerPoint to 0.01pt).
`first_baseline_off` reserves a descent of `natural_descent - quarter*(1 - n)`;
that slide reads closer to `natural_descent * n`. **The two agree at n = 1 and
diverge as n falls**, which is why 31 arms at single spacing could not separate
them.

One shape per arm: one word, one face, one size, one `lnSpc`, two paragraphs so
the step is visible as well as the offset.

    python tools/metrics/gen_pptx_lnspc_baseline.py
    python tools/metrics/read_pptx_lnspc_baseline.py     (exports and reads)
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

OUT = os.path.join("tools", "metrics", "lnspc_baseline.pptx")
# ★The first four have a natural descent SHALLOWER than a quarter of the box
# and the last four DEEPER, which is the branch deck 40's Grand Hotel (0.364)
# takes and the branch the original 31 arms never reached. Every one of them is
# installed on this machine, so PowerPoint reads the same file the engine does.
#
#     Arial 0.228   Georgia 0.234   Verdana 0.202   Calibri 0.250
#     Segoe Script 0.375   Papyrus 0.398   Viner Hand ITC 0.457
#     Javanese Text 0.539
FACES = ["Arial", "Georgia", "Verdana", "Calibri",
         "Segoe Script", "Papyrus", "Viner Hand ITC", "Javanese Text"]
SIZES = [24, 60]
SPACINGS = [0.4, 0.6, 0.8, 1.0, 1.2, 1.5]
WORD = "Hxpg"


def main() -> None:
    pres = Presentation()
    pres.slide_width = Emu(9144000)
    pres.slide_height = Emu(5143500)
    blank = pres.slide_layouts[6]
    arms = []
    for face in FACES:
        for size in SIZES:
            slide = pres.slides.add_slide(blank)
            for i, mult in enumerate(SPACINGS):
                left = Emu(200000 + (i % 3) * 2900000)
                top = Emu(200000 + (i // 3) * 2300000)
                box = slide.shapes.add_textbox(left, top, Emu(2700000), Emu(2200000))
                tf = box.text_frame
                tf.word_wrap = False
                tf.auto_size = MSO_AUTO_SIZE.NONE
                tf.vertical_anchor = MSO_ANCHOR.TOP
                tf.margin_left = tf.margin_right = 0
                tf.margin_top = tf.margin_bottom = 0
                for k in range(2):
                    para = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
                    para.alignment = PP_ALIGN.LEFT
                    para.line_spacing = mult
                    para.space_before = Pt(0)
                    para.space_after = Pt(0)
                    run = para.add_run()
                    run.text = WORD
                    run.font.name = face
                    run.font.size = Pt(size)
                    run.font.bold = False
                    run.font.italic = False
                arms.append((len(pres.slides), face, size, mult))
    pres.save(OUT)
    print("wrote %s: %d slides, %d arms" % (OUT, len(pres.slides), len(arms)))
    print("faces %s x sizes %s x lnSpc %s" % (FACES, SIZES, SPACINGS))


if __name__ == "__main__":
    main()
