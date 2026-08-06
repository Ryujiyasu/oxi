#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generate a minimal single-series pie chart pptx for Ra-loop measurement.

Mirrors gen_pptx_chart1.py (clustered column) but with a PIE chart.
Series 1 = (19.2, 21.4, 16.7), categories East/West/Midwest.
Word PDF -> pipeline_data/pptx_probes/chart_pie/chart_pie.pdf (via COM export).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

OUT = r"pipeline_data\pptx_probes\chart_pie\chart_pie.pptx"


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    chart_data = CategoryChartData()
    chart_data.categories = ["East", "West", "Midwest"]
    chart_data.add_series("Series 1", (19.2, 21.4, 16.7))

    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(1), Inches(1), Inches(5.5), Inches(4.0), chart_data
    )
    chart = gframe.chart
    chart.has_legend = False  # control: measure default legend off

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print(f"Saved {OUT}")


if __name__ == "__main__":
    main()
