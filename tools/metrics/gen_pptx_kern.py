# -*- coding: utf-8 -*-
"""Does PowerPoint kern, and when?

The width gate's last twelve offenders are all lines whose face the engine
resolved correctly, and the worst of them says why. d10 s8 sets 'To-do’s' in Jua
at 34.99pt, and PowerPoint's own character steps read

    'T' 22.5000  'o' 19.3751  '-' 16.5000  'd' 19.3751  'o' **11.6249**

-- the step from that last 'o' to the apostrophe is 7.75pt SHORT of the same
'o' elsewhere in the line. Something pulls the pair together, the engine has no
kerning at all (nothing in `oxislides-core` or the renderer mentions it), and
the deck does not state `a:rPr/@kern`. So PowerPoint kerns by default.

This asks the three questions that decide how to implement it:

    does it kern with no attribute at all          `pair` vs `flat`
    is `@kern` a MINIMUM SIZE, as the schema says   `kern_small` vs `kern_big`
    does the amount scale with the size             `pair` at 20 and 40pt

Arial is the face because it is installed, so the design advances can be read
out of the file and the kern is whatever PowerPoint places on top of them.

    python tools/metrics/gen_pptx_kern.py
    python tools/metrics/read_pptx_kern_com.py
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Emu, Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

REPO = Path(__file__).resolve().parents[2]
OUT = REPO / "pipeline_data" / "pptx_probes" / "kern"
FACE = "Arial"

# (label, text, size, kern attribute in hundredths of a point or None)
ARMS = [
    ("pair40", "AVAV", 40, None),
    ("flat40", "nnnn", 40, None),
    ("pair20", "AVAV", 20, None),
    ("kern1200", "AVAV", 40, 1200),
    ("kern9600", "AVAV", 40, 9600),
    ("kern0", "AVAV", 40, 0),
    ("To", "ToToTo", 40, None),
    ("curly", "o’o’o’", 40, None),
]


def build(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i, (label, text, size, kern) in enumerate(ARMS):
        col, row = i % 2, i // 2
        left = Emu(457200 + col * 4300000)
        top = Emu(685800 + row * 1400000)
        tag = slide.shapes.add_textbox(left, Emu(top.emu - 300000), Emu(3800000), Emu(250000))
        tag.text_frame.paragraphs[0].add_run().text = label
        tag.text_frame.paragraphs[0].runs[0].font.size = Pt(9)

        box = slide.shapes.add_textbox(left, top, Emu(4000000), Emu(1000000))
        tf = box.text_frame
        tf.word_wrap = False
        tf.auto_size = MSO_AUTO_SIZE.NONE
        run = tf.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.name = FACE
        if kern is not None:
            # python-pptx has no setter for it, and the attribute is the whole
            # question, so it goes on the element directly.
            run.font._rPr.set("kern", str(kern))
    prs.save(str(path))


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    deck = OUT / "kern.pptx"
    build(deck)
    print(f"wrote {deck}  (face {FACE})")
    for label, text, size, kern in ARMS:
        print(f"  {label:<10} {text!r:<12} {size}pt  kern={kern}")


if __name__ == "__main__":
    main()
