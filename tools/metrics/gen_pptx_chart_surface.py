# -*- coding: utf-8 -*-
"""Surface-chart probe deck (chart_surface).

python-pptx cannot author any of the four surface types ("XML writer for chart
type SURFACE (83) not yet implemented"), so the deck is built as a LINE chart
and each chart part is rewritten in the zip:

    <c:lineChart>  ->  <c:surfaceChart>   (top view, 2-D contour bands)
                   ->  <c:surface3DChart> (3-D mesh)
    drop chart-level <c:grouping>/<c:varyColors>/<c:marker>
    drop per-series <c:marker>/<c:smooth>   (CT_SurfaceSer has neither)
    insert <c:wireframe val="0|1"/> before the first <c:ser>
    add a third <c:axId> + a <c:serAx> (surface charts need cat/val/ser axes)

CT_SurfaceChart / CT_Surface3DChart child order is
    wireframe?, ser*, bandFmts?, axId, axId, axId

The arms separate the levers:

  P1  surfaceChart     wireframe=0   top view, filled contour bands
  P2  surfaceChart     wireframe=1   top view, wireframe only
  P3  surface3DChart   wireframe=0   3-D filled mesh
  P4  surface3DChart   wireframe=1   3-D wireframe mesh
  P5  surfaceChart     wireframe=0   legend      (band legend geometry)
  P6  surfaceChart     wireframe=0   3 series    (series-axis pitch lever)

Run:  python tools/metrics/gen_pptx_chart_surface.py
"""
import os
import re
import shutil
import sys
import zipfile

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Emu

sys.stdout.reconfigure(encoding="utf-8")

OUT = r"pipeline_data\pptx_probes\chart_surface\chart_surface.pptx"

# 1-based chart part index -> (element name, wireframe value)
ARMS = {
    1: ("surfaceChart", "0"),
    2: ("surfaceChart", "1"),
    3: ("surface3DChart", "0"),
    4: ("surface3DChart", "1"),
    5: ("surfaceChart", "0"),
    6: ("surfaceChart", "0"),
}


def pt(v):
    return Emu(int(round(v * 12700)))


def add(prs, cats, series, legend=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals)
    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, pt(72), pt(72), pt(396), pt(288), cd
    )
    ch = gf.chart
    ch.has_title = False
    ch.has_legend = legend
    return ch


def build():
    prs = Presentation()
    prs.slide_width = pt(720)
    prs.slide_height = pt(540)

    cats = ["Q1", "Q2", "Q3", "Q4"]
    s2 = [("Ser1", (19.2, 21.4, 16.7, 22.0)),
          ("Ser2", (10.5, 11.2, 8.5, 12.3))]
    s3 = s2 + [("Ser3", (5.1, 6.4, 7.2, 4.8))]

    add(prs, cats, s2)                 # P1 top view filled
    add(prs, cats, s2)                 # P2 top view wireframe
    add(prs, cats, s2)                 # P3 3-D filled
    add(prs, cats, s2)                 # P4 3-D wireframe
    add(prs, cats, s2, legend=True)    # P5 legend
    add(prs, cats, s3)                 # P6 three series

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)


def ser_ax(val_ax_id):
    """A minimal CT_SerAx.

    EG_AxShared order: axId, scaling, delete, axPos, ..., crossAx.  The
    crossAx MUST name a real axis -- pointing it at an id that is not in the
    part makes PowerPoint refuse to open the deck outright.
    """
    return (
        '<c:serAx><c:axId val="900000003"/><c:scaling>'
        '<c:orientation val="minMax"/></c:scaling><c:delete val="0"/>'
        '<c:axPos val="b"/><c:majorTickMark val="out"/>'
        '<c:minorTickMark val="none"/><c:tickLblPos val="nextTo"/>'
        '<c:crossAx val="%s"/></c:serAx>' % val_ax_id
    )


def surgery():
    tmp = OUT + ".tmp"
    zin = zipfile.ZipFile(OUT)
    zout = zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED)
    for it in zin.infolist():
        data = zin.read(it.filename)
        m = re.match(r"ppt/charts/chart(\d+)\.xml$", it.filename)
        if m:
            idx = int(m.group(1))
            elem, wire = ARMS[idx]
            x = data.decode("utf-8")

            # chart-level bits CT_Surface*Chart does not allow
            x = re.sub(r"<c:grouping[^>]*/>", "", x)
            x = re.sub(r"<c:varyColors[^>]*/>", "", x)
            # per-series marker + smooth (CT_SurfaceSer has neither)
            x = re.sub(r"<c:marker>.*?</c:marker>", "", x, flags=re.S)
            x = re.sub(r"<c:marker[^>]*/>", "", x)
            x = re.sub(r"<c:smooth[^>]*/>", "", x)

            # lineChart -> surface*, with <c:wireframe> as the first child
            x = x.replace(
                "<c:lineChart>",
                "<c:%s><c:wireframe val=\"%s\"/>" % (elem, wire),
            )
            x = x.replace("</c:lineChart>", "</c:%s>" % elem)

            # third axId inside the plot element (after the existing pair)
            x = re.sub(
                r'(<c:axId val="\d+"/><c:axId val="\d+"/>)(</c:%s>)' % elem,
                r'\1<c:axId val="900000003"/>\2',
                x,
            )
            # a series axis after the value axis (crossing the real valAx)
            val_ax_id = re.search(
                r'<c:valAx><c:axId val="(\d+)"/>', x).group(1)
            x = x.replace("</c:valAx>", "</c:valAx>" + ser_ax(val_ax_id), 1)

            data = x.encode("utf-8")
        zout.writestr(it, data)
    zin.close()
    zout.close()
    shutil.move(tmp, OUT)


def main():
    build()
    surgery()
    z = zipfile.ZipFile(OUT)
    for i in sorted(ARMS):
        x = z.read("ppt/charts/chart%d.xml" % i).decode("utf-8")
        elem, wire = ARMS[i]
        print("P%d %-16s wireframe=%s  serAx=%d  axId=%d  marker=%d smooth=%d"
              % (i, elem, wire,
                 x.count("<c:serAx>"), x.count("<c:axId"),
                 x.count("<c:marker"), x.count("<c:smooth")))
    print("wrote", OUT)


if __name__ == "__main__":
    main()
