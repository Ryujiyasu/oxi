#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Frame-size sweep probe for the pie circle geometry rule.

Six slides, single-series pie (19.2,21.4,16.7), varying frame size x title:
  A: 396x288 no title     B: 396x288 title
  C: 300x200 no title     D: 300x200 title
  E: 500x400 no title     F: 500x400 title
Legend is off everywhere (it is a right-side overlay = no effect, confirmed).
This derives the circle top/bottom/center/r as a function of frame (sx,sy,sw,shh).
"""
import os

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

OUT = r"pipeline_data\pptx_probes\chart_pie3\chart_pie3.pptx"


def add_pie(prs, w_pt, h_pt, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    cd = CategoryChartData()
    cd.categories = ["East", "West", "Midwest"]
    cd.add_series("Series 1", (19.2, 21.4, 16.7))
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(1), Inches(1), Inches(w_pt / 72.0), Inches(h_pt / 72.0), cd
    )
    gframe.chart.has_title = title
    gframe.chart.has_legend = False
    return gframe.chart


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    for w, h, t in ((396, 288, False), (396, 288, True),
                    (300, 200, False), (300, 200, True),
                    (500, 400, False), (500, 400, True)):
        add_pie(prs, w, h, t)
    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print(f"Saved {OUT}")


if __name__ == "__main__":
    main()
