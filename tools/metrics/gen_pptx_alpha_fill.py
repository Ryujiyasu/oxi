# -*- coding: utf-8 -*-
"""Probe: how does PowerPoint composite a shape fill that carries <a:alpha>?

dev corpus has 47 slide-level and 2 inherited shapes whose solidFill carries an
`a:alpha`; Oxi has no alpha on `Shape.fill_color` and paints them opaque.  On
d23 that turns nine `alpha val="0"` full-page rects into solid black slabs.

Each slide is a known opaque backdrop plus one (or two) translucent rects, so
the composited RGB read back from the PDF pins the blend.  python-pptx has no
alpha API, so the `<a:alpha>` is injected into the saved XML -- the shapes are
tagged by a unique fill colour so the injection can find them.
"""
from __future__ import annotations

import re
import shutil
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu

OUT = Path(r"pipeline_data\pptx_probes\alpha_fill").resolve()
SW, SH = 9144000, 5143500

# (label, backdrop hex or None, [(fill hex, alpha permille-of-100000)])
ARMS = [
    ("A1_black_a0",      None,     [("000000", 0)]),
    ("A2_black_a25",     None,     [("000000", 25000)]),
    ("A3_black_a50",     None,     [("000000", 50000)]),
    ("A4_black_a62010",  None,     [("000000", 62010)]),   # d08 layout3
    ("A5_black_a75",     None,     [("000000", 75000)]),
    ("A6_black_a100",    None,     [("000000", 100000)]),
    ("A7_red_black_a50", "FF0000", [("000000", 50000)]),
    ("A8_red_cfd8dc_a49230", "FF0000", [("CFD8DC", 49230)]),  # d16 layout10
    ("A9_stacked_a50_a50", None,   [("000000", 50000), ("0000FF", 50000)]),
    ("A10_green_white_a25", "008000", [("FFFFFF", 25000)]),
]


def rect(slide, x, y, cx, cy, hexval):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cx, cy)
    r.fill.solid()
    r.fill.fore_color.rgb = RGBColor.from_string(hexval)
    r.line.fill.background()
    # The default Office theme puts an outer shadow on autoshapes; PowerPoint
    # rasterises it and draws it UNDER the rect, which would darken the sample.
    r.shadow.inherit = False
    return r


def build() -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(SW), Emu(SH)
    blank = prs.slide_layouts[6]
    for _label, backdrop, fills in ARMS:
        s = prs.slides.add_slide(blank)
        if backdrop:
            rect(s, 0, 0, Emu(SW), Emu(SH), backdrop)
        for hexval, _a in fills:
            # inset a little so the backdrop stays visible at the edge
            rect(s, Emu(SW // 8), Emu(SH // 8), Emu(SW * 3 // 4),
                 Emu(SH * 3 // 4), hexval)
    OUT.mkdir(parents=True, exist_ok=True)
    raw = OUT / "_raw.pptx"
    prs.save(raw)
    return raw


def inject(raw: Path) -> Path:
    """python-pptx cannot write <a:alpha>; splice it into each solidFill."""
    final = OUT / "alpha_fill.pptx"
    with zipfile.ZipFile(raw) as zin, zipfile.ZipFile(
            final, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            m = re.fullmatch(r"ppt/slides/slide(\d+)\.xml", item.filename)
            if m:
                idx = int(m.group(1)) - 1
                if idx < len(ARMS):
                    xml = data.decode("utf-8")
                    for hexval, alpha in ARMS[idx][2]:
                        # only the FILL solidFill -- the line is <a:noFill/>
                        xml = xml.replace(
                            f'<a:solidFill><a:srgbClr val="{hexval}"/>'
                            f'</a:solidFill>',
                            f'<a:solidFill><a:srgbClr val="{hexval}">'
                            f'<a:alpha val="{alpha}"/></a:srgbClr>'
                            f'</a:solidFill>', 1)
                    data = xml.encode("utf-8")
            zout.writestr(item, data)
    raw.unlink()
    return final


def main() -> None:
    if OUT.exists():
        shutil.rmtree(OUT)
    final = inject(build())
    with zipfile.ZipFile(final) as z:
        n = sum(1 for x in z.namelist()
                if re.fullmatch(r"ppt/slides/slide\d+\.xml", x))
        got = sum(z.read(f"ppt/slides/slide{i+1}.xml").decode().count("<a:alpha")
                  for i in range(n))
    want = sum(len(a[2]) for a in ARMS)
    print(f"{final}  slides={n}  a:alpha injected={got}/{want}")
    if got != want:
        raise SystemExit("alpha injection missed a shape")


if __name__ == "__main__":
    main()
