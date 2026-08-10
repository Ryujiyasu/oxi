# -*- coding: utf-8 -*-
"""Legend-label wrap threshold probe: labels of increasing width, one slide each.

Word wraps a doughnut legend category name onto TWO lines at a single-word
boundary when the label is too wide.  This probe sweeps label width so the
keep/wrap flip and the wrapped-band geometry can be measured:

  S1: "A"                          (1 short char)
  S2: "Region"                     (6 chars, short word)
  S3: "Northeastern"               (12 chars)
  S4: "Northeastern Region"        (19 chars, 2 words)
  S5: "Northwestern California"    (24 chars, 2 words)
  S6: "Northeastern Region Coast"  (28 chars, 3 words)

Frame 72,72,396,288, values 19.2/21.4/16.7, DOUGHNUT 1 series, legend RIGHT
(overlay=0 -> the band displaces the ring, which is what makes wrap visible)."""
import sys, os, re, shutil, zipfile

sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

base = r"pipeline_data\pptx_probes\chart_doughnut_wrap"
os.makedirs(base, exist_ok=True)

prs = Presentation()

labels = [
    ["A"],
    ["Region"],
    ["Northeastern"],
    ["Northeastern Region"],
    ["Northwestern California"],
    ["Northeastern Region Coast"],
    # threshold arm: single-word-boundary labels stepping the 2nd word width
    ["Northeastern R"],       # 15ch
    ["Northeastern Re"],      # 16ch
    ["Northeastern Reg"],     # 17ch
    ["Northeastern Regi"],    # 18ch
    ["Northeastern Regio"],   # 19ch
    ["Northeastern Regionn"],  # 20ch, 2 words, longer than S4
    # one-word long label: does Word split a single unbreakable word?
    ["NortheasternRegionX"],  # 20ch, 1 word
    # single-word width sweep (no word boundary => char split if it wraps)
    ["M" * 6],       # 92.34pt
    ["N" * 8],       # 92.95pt
    ["N" * 9],       # 104.57pt  -> 1 line (1-line anchor, measured)
    # intermediate arms to pin the keep/wrap threshold inside (104.57, 107.73]
    ["MMMMMMO"],     # 104.26pt  (Calibri 18: M6 + 'O')
    ["MMMMMMw"],     # 105.21pt  (M6 + 'w')
    ["MMMMMmm"],     # 105.71pt  (M5 + 'm' + 'm')
    ["MMMMMMm"],     # 106.72pt  (M6 + 'm')
    ["M" * 7],       # 107.73pt  -> 2 lines (2-line anchor, measured)
    ["N" * 10],      # 116.19pt
    ["N" * 11],      # 127.81pt
]

for cats in labels:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = cats
    # 3 values so the ring has 3 slices; only the label width matters here
    cd.add_series("Series 1", (19.2, 21.4, 16.7))
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, Pt(72), Pt(72), Pt(396), Pt(288), cd
    )
    chart = gframe.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False

out = os.path.join(base, "chart_doughnut_wrap.pptx")
prs.save(out)
print("saved:", out, os.path.getsize(out))
