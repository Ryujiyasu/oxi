# -*- coding: utf-8 -*-
"""Author the PATH (radial) gradient probe.

`paint_bg_gradient` models `a:path path="circle"` as
`t = |p - focus| / r_max`, a CIRCULAR distance normalised by the farthest
corner. d15 -- the corpus floor once d47 was ruled an environment fault -- says
that is wrong: its background ramp runs up to 30/255 too dark in the middle
while matching exactly at the focus, and linear fits against circular AND
elliptical distance both fail (residual RMS 0.121 / 0.137).

Fitting curves to the deck could not separate the two unknowns -- the ISO-LINE
SHAPE and the RAMP CURVE -- because an error in either looks like an error in
the other. So each is asked separately here:

  geometry   the SAME focus on a square, a wide and a tall shape (arms 5/7/8).
             If the iso-lines are circles the ramp is identical in x and y on
             all three; if they are ellipses scaled to the shape, the wide and
             tall arms stretch with the aspect.
  focus      the focus point itself, at the centre, at two opposite corners,
             and absent (arms 1/2/3/6).
  tileRect   d15's exact `l="-100%" t="-100%"` against the same arm without it
             (arm 4 vs arm 2), which is the one attribute that could rescale
             the whole ramp.

Every arm is ONE rectangle with a two-stop BLACK -> WHITE ramp and no outline,
so the grey level read back from the PDF IS the ramp parameter -- no colour
algebra, and any gamma in the interpolation shows up directly as a bent curve.

Usage:
    python tools/metrics/gen_pptx_gradpath.py
    python tools/metrics/export_pptx_gradpath.py   # PowerPoint COM -> PDF
    python tools/metrics/read_pptx_gradpath.py     # read the PDF back
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

REPO = Path(__file__).resolve().parents[2]
OUT = REPO / "pipeline_data" / "pptx_probes" / "gradpath"

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
EMU_IN = 914400

# Default deck is 10 x 7.5in, so a full-slide shape is 4:3 -- aspect enough to
# tell a circle from an ellipse without being extreme.
FULL = (0, 0, 10 * EMU_IN, 7.5 * EMU_IN)
SQUARE = (2 * EMU_IN, 1 * EMU_IN, 5 * EMU_IN, 5 * EMU_IN)
WIDE = (0, 3 * EMU_IN, 10 * EMU_IN, 2 * EMU_IN)
TALL = (3 * EMU_IN, 0, 2.5 * EMU_IN, 7.5 * EMU_IN)

CENTRE = {"l": "50000", "t": "50000", "r": "50000", "b": "50000"}
TOP_LEFT = {"r": "100000", "b": "100000"}          # d15's own form
BOTTOM_RIGHT = {"l": "100000", "t": "100000"}

# (name, geometry, fillToRect attrs or None, tileRect attrs or None)
ARMS = [
    ("centre_full", FULL, CENTRE, None),
    ("topleft_full", FULL, TOP_LEFT, None),
    ("bottomright_full", FULL, BOTTOM_RIGHT, None),
    ("topleft_tile", FULL, TOP_LEFT, {"l": "-100000", "t": "-100000"}),
    ("centre_square", SQUARE, CENTRE, None),
    ("nofill_full", FULL, None, None),
    ("centre_wide", WIDE, CENTRE, None),
    ("centre_tall", TALL, CENTRE, None),
]


def grad_path_fill(fill_to_rect, tile_rect) -> etree._Element:
    """A black->white two-stop `path="circle"` ramp."""
    gf = etree.Element(f"{{{A}}}gradFill")
    gs_lst = etree.SubElement(gf, f"{{{A}}}gsLst")
    for pos, val in (("0", "000000"), ("100000", "FFFFFF")):
        gs = etree.SubElement(gs_lst, f"{{{A}}}gs")
        gs.set("pos", pos)
        etree.SubElement(gs, f"{{{A}}}srgbClr").set("val", val)
    path = etree.SubElement(gf, f"{{{A}}}path")
    path.set("path", "circle")
    if fill_to_rect is not None:
        ftr = etree.SubElement(path, f"{{{A}}}fillToRect")
        for k, v in fill_to_rect.items():
            ftr.set(k, v)
    if tile_rect is not None:
        tr = etree.SubElement(gf, f"{{{A}}}tileRect")
        for k, v in tile_rect.items():
            tr.set(k, v)
    return gf


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    manifest = []
    for i, (name, (x, y, w, h), ftr, tile) in enumerate(ARMS):
        slide = prs.slides.add_slide(blank)
        shp = slide.shapes.add_shape(1, Emu(x), Emu(y), Emu(w), Emu(h))  # 1 = rect
        shp.line.fill.background()
        sp_pr = shp._element.spPr
        # Drop whatever fill python-pptx put there and install the ramp.
        for tag in ("solidFill", "gradFill", "noFill", "blipFill", "pattFill", "grpFill"):
            for el in sp_pr.findall(f"{{{A}}}{tag}"):
                sp_pr.remove(el)
        geom = sp_pr.find(f"{{{A}}}prstGeom")
        idx = list(sp_pr).index(geom) + 1 if geom is not None else len(sp_pr)
        sp_pr.insert(idx, grad_path_fill(ftr, tile))
        manifest.append(
            {
                "slide": i + 1,
                "name": name,
                "x_emu": x, "y_emu": y, "w_emu": w, "h_emu": h,
                "fillToRect": ftr,
                "tileRect": tile,
            }
        )
    path = OUT / "probe_gradpath.pptx"
    prs.save(str(path))
    (OUT / "manifest.json").write_text(json.dumps(manifest, indent=1), encoding="utf-8")
    print(f"wrote {path} ({len(ARMS)} arms)")
    for m in manifest:
        print(f"  s{m['slide']} {m['name']:<18} fillToRect={m['fillToRect']} tileRect={m['tileRect']}")


if __name__ == "__main__":
    main()
