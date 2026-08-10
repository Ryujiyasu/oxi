# -*- coding: utf-8 -*-
"""Radar-chart probe deck (chart_radar).

Eight arms that separate the levers of a radar plot:

  R1  1 series, 3 categories                -> vertex geometry + radius law
  R2  1 series, 5 categories                -> does the ring follow n_cat?
  R3  2 series                              -> per-series colour / auto title
  R4  1 series + legend                     -> does the legend band shrink the ring?
  R5  1 series RADAR_MARKERS                -> marker shape + size
  R6  1 series RADAR_FILLED                 -> filled polygon
  R7  1 series, wide frame 500              -> radius follows width or height?
  R8  2 series + legend + explicit title    -> title geometry / plot_top

Run:  python tools/metrics/gen_pptx_chart_radar.py
"""
import os
import sys

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Emu

sys.stdout.reconfigure(encoding="utf-8")

OUT = r"pipeline_data\pptx_probes\chart_radar\chart_radar.pptx"


def pt(v):
    return Emu(int(round(v * 12700)))


def add(prs, x, y, w, h, kind, cats, series, legend=False, title=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals)
    gf = slide.shapes.add_chart(kind, pt(x), pt(y), pt(w), pt(h), cd)
    ch = gf.chart
    ch.has_legend = legend
    if legend:
        ch.legend.include_in_layout = False
    if title is not None:
        ch.has_title = True
        ch.chart_title.text_frame.text = title
    return ch


def main():
    prs = Presentation()
    prs.slide_width = pt(720)
    prs.slide_height = pt(540)

    C3 = ["East", "West", "Midwest"]
    C5 = ["East", "West", "Midwest", "North", "South"]
    S3 = (19.2, 21.4, 16.7)
    S5 = (19.2, 21.4, 16.7, 22.0, 18.5)
    T3 = (10.5, 11.2, 8.5)

    R = XL_CHART_TYPE.RADAR
    add(prs, 72, 72, 396, 288, R, C3, [("Ser1", S3)])
    add(prs, 72, 72, 396, 288, R, C5, [("Ser1", S5)])
    add(prs, 72, 72, 396, 288, R, C3, [("Ser1", S3), ("Ser2", T3)])
    add(prs, 72, 72, 396, 288, R, C3, [("Ser1", S3)], legend=True)
    add(prs, 72, 72, 396, 288, XL_CHART_TYPE.RADAR_MARKERS, C5,
        [("Ser1", S5)])
    add(prs, 72, 72, 396, 288, XL_CHART_TYPE.RADAR_FILLED, C5,
        [("Ser1", S5)])
    add(prs, 72, 72, 500, 288, R, C5, [("Ser1", S5)])
    add(prs, 72, 72, 396, 288, R, C5, [("Ser1", S5), ("Ser2", (12.0, 9.0, 14.0, 7.5, 11.0))],
        legend=True, title="Quarterly Mix")

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print("wrote", OUT)


if __name__ == "__main__":
    main()
