# -*- coding: utf-8 -*-
"""Chart spec (item 8) wave-1 repro: minimal clustered-column chart on a blank
slide, default theme. One series, three categories, with data labels + legend
visible (python-pptx defaults), so Word's PDF exposes the full chart chrome
(plot area, axes, labels, legend) for measurement."""
import sys, os
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

base = r"pipeline_data\pptx_probes\chart1"
os.makedirs(base, exist_ok=True)

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

chart_data = CategoryChartData()
chart_data.categories = ["East", "West", "Midwest"]
chart_data.add_series("Series 1", (19.2, 21.4, 16.7))

x, y, cx, cy = Inches(1.0), Inches(1.0), Inches(5.5), Inches(4.0)
gframe = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

out = os.path.join(base, "chart1.pptx")
prs.save(out)
print("saved:", out, os.path.getsize(out))
