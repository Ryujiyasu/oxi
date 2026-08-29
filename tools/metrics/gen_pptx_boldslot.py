# -*- coding: utf-8 -*-
"""Probe: does a DECLARED bold slot stop PowerPoint synthesising bold?

Three decks measured 2026-08-29 disagree about what a bold run advances at:

    blind 04  Inria Sans Light, bold slot present but holding a NON-bold part
              -> PowerPoint = the REGULAR face's design advances (184.16 vs 183.97)
    blind 29  Sniglet, no bold slot at all
              -> PowerPoint = GDI's synthesised bold (168.60 vs 167.91), NOT design (171.08)
    blind 35  Gruppo / Bebas Neue, no bold slot
              -> PowerPoint = GDI's synthesised bold to 0.1% (769.75 vs 768.85)

So the discriminator may be the SLOT's existence rather than the face's weight:
declare a bold and PowerPoint stops widening, even when the part behind it is
not bold. This isolates exactly that, with two decks identical but for one
`<p:bold>` element.

    python tools/metrics/gen_pptx_boldslot.py
    python tools/metrics/export_pptx_boldslot.py      # PowerPoint COM -> PDF
    python tools/metrics/read_pptx_boldslot.py

★`<p:presentation>` needs `embedTrueTypeFonts="true"` or the whole
`embeddedFontLst` is ignored -- a probe whose arms agree is broken until a
positive control says otherwise (S-CFFPART, 2026-08-28).
"""
from __future__ import annotations

import json
import re
import shutil
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

REPO = Path(__file__).resolve().parents[2]
ROOT = REPO / "pipeline_data" / "pptx_benchmark"
OUT = Path(r"pipeline_data\pptx_probes\boldslot").resolve()
FAMILY = "Probe Sans"
TEXT = "Handgloves the quick brown fox"
SIZE = 24


def source_parts() -> tuple[bytes, bytes]:
    """blind 04's Inria Sans Light: its regular part, and the non-bold part its
    bold slot holds. Reusing a real pair keeps the arms honest -- the bold slot
    of the ARM that has one is exactly the kind PowerPoint meets in the wild."""
    manifest = json.loads((ROOT / "manifest.json").read_text(encoding="utf-8"))
    src = ROOT / "pptx" / next(i["local"] for i in manifest if i["idx"] == 4)
    with zipfile.ZipFile(src) as z:
        pres = z.read("ppt/presentation.xml").decode("utf-8", "replace")
        rels = z.read("ppt/_rels/presentation.xml.rels").decode("utf-8", "replace")
        rid = dict(re.findall(r'Id="([^"]+)"[^>]*Target="([^"]+)"', rels))
        blk = re.search(r'<p:embeddedFont><p:font typeface="Inria Sans Light"/>(.*?)</p:embeddedFont>',
                        pres, re.S).group(1)
        reg = re.search(r'<p:regular r:id="([^"]+)"', blk).group(1)
        bold = re.search(r'<p:bold r:id="([^"]+)"', blk).group(1)
        return (z.read("ppt/" + rid[reg].replace("../", "")),
                z.read("ppt/" + rid[bold].replace("../", "")))


def build(dst: Path, with_bold_slot: bool, regular: bytes, bold: bytes) -> None:
    stage = dst.with_suffix(".stage.pptx")
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    box = s.shapes.add_textbox(Emu(457200), Emu(914400), Emu(7772400), Emu(1500000))
    tf = box.text_frame
    tf.word_wrap = False
    r = tf.paragraphs[0].add_run()
    r.text = TEXT
    r.font.size = Pt(SIZE)
    r.font.bold = True
    r.font.name = FAMILY
    prs.save(stage)

    with zipfile.ZipFile(stage) as zin:
        names = zin.namelist()
        data = {n: zin.read(n) for n in names}
    stage.unlink()

    data["ppt/fonts/font1.fntdata"] = regular
    slots = '<p:regular r:id="rIdFont1"/>'
    rel_extra = ('<Relationship Id="rIdFont1" '
                 'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/font" '
                 'Target="fonts/font1.fntdata"/>')
    if with_bold_slot:
        data["ppt/fonts/font2.fntdata"] = bold
        slots += '<p:bold r:id="rIdFont2"/>'
        rel_extra += ('<Relationship Id="rIdFont2" '
                      'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/font" '
                      'Target="fonts/font2.fntdata"/>')

    rels = data["ppt/_rels/presentation.xml.rels"].decode("utf-8")
    data["ppt/_rels/presentation.xml.rels"] = rels.replace(
        "</Relationships>", rel_extra + "</Relationships>").encode("utf-8")

    pres = data["ppt/presentation.xml"].decode("utf-8")
    pres = re.sub(r"<p:presentation ", '<p:presentation embedTrueTypeFonts="true" ', pres, count=1)
    lst = (f'<p:embeddedFontLst><p:embeddedFont><p:font typeface="{FAMILY}"/>'
           f"{slots}</p:embeddedFont></p:embeddedFontLst>")
    # Schema order: embeddedFontLst sits after notesSz, before defaultTextStyle.
    if "<p:defaultTextStyle>" in pres:
        pres = pres.replace("<p:defaultTextStyle>", lst + "<p:defaultTextStyle>", 1)
    else:
        pres = pres.replace("</p:presentation>", lst + "</p:presentation>", 1)
    data["ppt/presentation.xml"] = pres.encode("utf-8")

    ct = data["[Content_Types].xml"].decode("utf-8")
    if 'Extension="fntdata"' not in ct:
        ct = ct.replace("<Default", '<Default Extension="fntdata" '
                        'ContentType="application/x-fontdata"/><Default', 1)
    data["[Content_Types].xml"] = ct.encode("utf-8")

    with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:
        for n in list(data):
            zout.writestr(n, data[n])


def main() -> None:
    shutil.rmtree(OUT, ignore_errors=True)
    OUT.mkdir(parents=True, exist_ok=True)
    regular, bold = source_parts()
    build(OUT / "slot.pptx", True, regular, bold)
    build(OUT / "noslot.pptx", False, regular, bold)
    print(f"wrote {OUT} -- slot.pptx and noslot.pptx  ({FAMILY} {SIZE}pt, bold run)")


if __name__ == "__main__":
    main()
