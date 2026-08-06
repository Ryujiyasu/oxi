# -*- coding: utf-8 -*-
"""Multi-series line chart legend repro: LINE_MARKERS charts with 2 and 3
series on blank slides, default theme, legend enabled, so the multi-series
legend geometry (swatches, labels, row pitch, vertical anchor) is measurable
against the single-series law (legend_y0 = sy + shh/2 + 17.68,
legend_left = plot_right + 15.65) already derived from chart_line."""
import sys, os
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

base = r"pipeline_data\pptx_probes\chart_line2"
os.makedirs(base, exist_ok=True)

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

chart_data = CategoryChartData()
chart_data.categories = ["East", "West", "Midwest", "North", "South"]
chart_data.add_series("Series 1", (19.2, 21.4, 16.7, 22.0, 18.5))
chart_data.add_series("Series 2", (15.3, 17.5, 20.1, 13.9, 16.2))

x, y, cx, cy = Inches(1.0), Inches(1.0), Inches(5.5), Inches(4.0)
gframe = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data)

out = os.path.join(base, "chart_line2.pptx")
prs.save(out)
print("saved:", out, os.path.getsize(out))
