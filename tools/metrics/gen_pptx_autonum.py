"""Generate buAutoNum (auto-numbered list) measurement repro.

Slide 1: hierarchy at lvl 0/1/2 + back-to-lvl0 (counter continue vs reset).
Slide 2: startAt=5 (auto-number start offset).
Slide 3: more number schemes (alphaUcPeriod / arabicParenBoth / romanLcParenR).

Each paragraph gets explicit marL/indent (EMU) + buFont Arial + buAutoNum so
the geometry is independent of the textbox (otherStyle) master defaults.
"""
import sys

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from lxml import etree

A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def add_textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(Inches(left / 72.0), Inches(top / 72.0),
                                  Inches(width / 72.0), Inches(height / 72.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(914400)   # 0.1 in
    tf.margin_right = Emu(914400)
    tf.margin_top = Emu(457200)    # 0.05 in
    tf.margin_bottom = Emu(457200)
    return tb, tf


def set_autonum(paragraph, kind, mar_l_pt, indent_pt, lvl=0, start_at=None,
                bufont="Arial"):
    p = paragraph._p
    pPr = p.get_or_add_pPr()
    pPr.set("lvl", str(lvl))
    pPr.set("marL", str(int(mar_l_pt * 12700)))
    pPr.set("indent", str(int(indent_pt * 12700)))
    for tag in ("buNone", "buChar", "buAutoNum", "buFont", "buClr", "buSzPct"):
        for el in pPr.findall(f"{{{A}}}{tag}"):
            pPr.remove(el)
    bf = etree.SubElement(pPr, f"{{{A}}}buFont")
    bf.set("typeface", bufont)
    bn = etree.SubElement(pPr, f"{{{A}}}buAutoNum")
    bn.set("type", kind)
    if start_at is not None:
        bn.set("startAt", str(start_at))


def fill(prs, paras):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb, tf = add_textbox(s, 72, 72, 400, 360)
    first = True
    for text, kind, lvl, mar, ind, start in paras:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        r = p.add_run()
        r.text = text
        r.font.size = Pt(18)
        set_autonum(p, kind, mar, ind, lvl, start)
    return s


prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: hierarchy + counter continuation/reset
fill(prs, [
    ("Item 1", "arabicPeriod", 0, 72, -36, None),
    ("Item 2", "arabicPeriod", 0, 72, -36, None),
    ("Item 3", "arabicPeriod", 0, 72, -36, None),
    ("Sub a", "alphaLcPeriod", 1, 108, -36, None),
    ("Sub b", "alphaLcPeriod", 1, 108, -36, None),
    ("SubSub I", "romanUcPeriod", 2, 144, -36, None),
    ("Back L0", "arabicPeriod", 0, 72, -36, None),
])

# Slide 2: startAt=5
fill(prs, [
    ("Five", "arabicPeriod", 0, 72, -36, 5),
    ("Six", "arabicPeriod", 0, 72, -36, None),
    ("Seven", "arabicPeriod", 0, 72, -36, None),
])

# Slide 3: more schemes
fill(prs, [
    ("A-one", "alphaUcPeriod", 0, 72, -36, None),
    ("A-two", "alphaUcPeriod", 0, 72, -36, None),
    ("P-one", "arabicParenBoth", 0, 72, -36, None),
    ("P-two", "arabicParenBoth", 0, 72, -36, None),
    ("r-one", "romanLcParenR", 0, 72, -36, None),
    ("r-two", "romanLcParenR", 0, 72, -36, None),
])

prs.save(sys.argv[1] if len(sys.argv) > 1 else "autonum.pptx")
print("saved")
