# -*- coding: utf-8 -*-
"""Does a soft break at the END of a paragraph open a line?

The corpus says no, four times. `pptx_line_audit_com.py` found the only four
paragraphs in dev + blind whose last child is `<a:br/>` -- blind 44 s25 and s36,
d06 s36, d19 s37 -- and PowerPoint sets every one of them with exactly one line
fewer than the engine does. Nothing else in 8170 audited paragraphs disagrees.

A correlation that clean still is not the rule, because those four are all the
same template and all end in a RUN of breaks. This asks PowerPoint the question
one variable at a time, and asks it twice:

    Lines.Count   does the trailing break open a line
    BoundHeight   does it take the vertical space of one anyway

The second question is the one that decides how to implement it. A line that is
not counted but is still reserved moves every centred or bottom-anchored shape
that carries one.

    A  'abc'                     control
    B  'abc' + br                the case
    C  'abc' + br + 'def'        a break that is not trailing -- 2 lines
    D  'abc' + br + br           two trailing breaks
    E  br                        a paragraph that is only a break
    F  'abc' + br + 'def' + br   trailing after a real second line
    G  'abc' + br + br + 'def'   two breaks in the middle

    python tools/metrics/gen_pptx_trailbr.py
    python tools/metrics/read_pptx_trailbr_com.py
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

REPO = Path(__file__).resolve().parents[2]
OUT = REPO / "pipeline_data" / "pptx_probes" / "trailbr"

# Each arm is a list: a string is a run, None is an `<a:br/>`.
ARMS = [
    ("A_plain", ["abc"]),
    ("B_trailing", ["abc", None]),
    ("C_middle", ["abc", None, "def"]),
    ("D_two_trailing", ["abc", None, None]),
    ("E_only_break", [None]),
    ("F_after_second", ["abc", None, "def", None]),
    ("G_two_middle", ["abc", None, None, "def"]),
]


def build(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(6858000)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    for i, (label, parts) in enumerate(ARMS):
        # A column per arm, each box wide enough that nothing wraps -- the only
        # thing that may end a line here is the break itself.
        #
        # ★And a SECOND row with `wrap` left on, because the first version set
        # `word_wrap = False` on every arm and the engine answered "1 line" to
        # all seven inputs. A probe whose arms cannot differ is measuring its
        # own scaffolding: `bodyPr/@wrap="none"` was swallowing the explicit
        # breaks, which is a finding, but only against a row that wraps.
        box = slide.shapes.add_textbox(
            Emu(228600 + i * 1250000), Emu(457200), Emu(1150000), Emu(2400000)
        )
        tf = box.text_frame
        tf.word_wrap = False
        # The label lives in its own shape so the measured shape holds nothing
        # but the arm.
        tag = slide.shapes.add_textbox(
            Emu(228600 + i * 1250000), Emu(228600), Emu(1150000), Emu(228600)
        )
        tag.text_frame.paragraphs[0].add_run().text = label
        tag.text_frame.paragraphs[0].runs[0].font.size = Pt(9)

        wrapped = slide.shapes.add_textbox(
            Emu(228600 + i * 1250000), Emu(3200000), Emu(1150000), Emu(2400000)
        )
        wrapped.text_frame.word_wrap = True
        tag2 = slide.shapes.add_textbox(
            Emu(228600 + i * 1250000), Emu(2971400), Emu(1150000), Emu(228600)
        )
        tag2.text_frame.paragraphs[0].add_run().text = label.replace("_", "W_", 1)
        tag2.text_frame.paragraphs[0].runs[0].font.size = Pt(9)

        for frame in (tf, wrapped.text_frame):
            para = frame.paragraphs[0]
            for part in parts:
                if part is None:
                    para.add_line_break()
                else:
                    run = para.add_run()
                    run.text = part
                    run.font.size = Pt(18)
    prs.save(str(path))


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    deck = OUT / "trailbr.pptx"
    build(deck)
    print(f"wrote {deck}")
    for label, parts in ARMS:
        shape = " + ".join("br" if p is None else repr(p) for p in parts)
        print(f"  {label:<16} {shape}")


if __name__ == "__main__":
    main()
