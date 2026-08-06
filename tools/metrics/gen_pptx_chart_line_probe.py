# -*- coding: utf-8 -*-
"""Line-chart plot-area geometry probe: 5 variants on 5 slides in ONE pptx.
Varies category count / frame size / legend presence / frame height so the
Word-PDF plot-area rule (plot_left/right/top/bot) can be derived from
controlled variation (no-guess rule)."""
import sys, os
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

base = r"pipeline_data\pptx_probes\chart_line_probe"
os.makedirs(base, exist_ok=True)

prs = Presentation()
layouts = prs.slide_layouts[6]

def add(cats, vals, x_in, y_in, w_in, h_in, legend):
    slide = prs.slides.add_slide(layouts)
    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("Series 1", vals)
    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in), cd,
    )
    if not legend:
        gf.chart.has_legend = False
    return gf

# P0 control: 5 cats, 396x288, legend (matches chart_line)
add(["East","West","Midwest","North","South"], (19.2,21.4,16.7,22.0,18.5),
    1.0, 1.0, 5.5, 4.0, True)
# P1: 3 cats, same frame, legend
add(["East","West","Midwest"], (19.2,21.4,16.7), 1.0, 1.0, 5.5, 4.0, True)
# P2: 5 cats, WIDER frame (500pt), legend
add(["East","West","Midwest","North","South"], (19.2,21.4,16.7,22.0,18.5),
    1.0, 1.0, 6.944, 4.0, True)
# P3: 5 cats, same frame, NO legend
add(["East","West","Midwest","North","South"], (19.2,21.4,16.7,22.0,18.5),
    1.0, 1.0, 5.5, 4.0, False)
# P4: 5 cats, TALLER frame (360pt), legend
add(["East","West","Midwest","North","South"], (19.2,21.4,16.7,22.0,18.5),
    1.0, 1.0, 5.5, 5.0, True)
# P5: 5 cats SHORT labels (A..E), 396x288, legend — tests crowding/label-width
#     discriminator for the 78.62 bottom band (P0 vs P1/P2/P3)
add(["A","B","C","D","E"], (19.2,21.4,16.7,22.0,18.5), 1.0, 1.0, 5.5, 4.0, True)
# P6: 4 cats (labels W..Z, widths like 2-5 chars), 396x288, legend — tests
#     whether 4 categories also trigger the large bottom band
add(["W","XY","ZWQ","UVWXY"], (19.2,21.4,16.7,22.0), 1.0, 1.0, 5.5, 4.0, True)

out = os.path.join(base, "chart_line_probe.pptx")
prs.save(out)
print("saved:", out, os.path.getsize(out))
