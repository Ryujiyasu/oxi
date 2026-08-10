# -*- coding: utf-8 -*-
"""Radar geometry sweep (chart_radar_geo).

Every arm has TWO series so Word never draws an auto title; that isolates the
frame/legend levers from the title lever measured in chart_radar.

  G1  396x288 no legend            baseline
  G2  396x200 no legend            height lever
  G3  396x400 no legend            height lever
  G4  396x288 legend               legend lever
  G5  396x288 legend, long names   legend LABEL WIDTH lever
  G6  250x288 no legend            width lever
  G7  180x288 no legend            width lever
  G8  396x288 no legend, 4 cats    bottom-vertex label rule
  G9  396x288 no legend, 6 cats    n_cat lever
  G10 250x288 legend               narrow + legend
"""
import os
import sys

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Emu

sys.stdout.reconfigure(encoding="utf-8")

OUT = r"pipeline_data\pptx_probes\chart_radar_geo\chart_radar_geo.pptx"


def pt(v):
    return Emu(int(round(v * 12700)))


def add(prs, w, h, cats, series, legend=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals)
    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.RADAR, pt(72), pt(72), pt(w), pt(h), cd)
    ch = gf.chart
    ch.has_legend = legend
    if legend:
        ch.legend.include_in_layout = False
    return ch


def main():
    prs = Presentation()
    prs.slide_width = pt(720)
    prs.slide_height = pt(540)

    C4 = ["East", "West", "Midwest", "North"]
    C5 = ["East", "West", "Midwest", "North", "South"]
    C6 = ["East", "West", "Midwest", "North", "South", "Central"]
    A5 = (19.2, 21.4, 16.7, 22.0, 18.5)
    B5 = (10.5, 11.2, 8.5, 12.3, 9.4)
    A4, B4 = A5[:4], B5[:4]
    A6 = A5 + (20.1,)
    B6 = B5 + (10.0,)
    P = [("Ser1", A5), ("Ser2", B5)]

    add(prs, 396, 288, C5, P)
    add(prs, 396, 200, C5, P)
    add(prs, 396, 400, C5, P)
    add(prs, 396, 288, C5, P, legend=True)
    add(prs, 396, 288, C5,
        [("Northeastern Region", A5), ("Southwestern Region", B5)], legend=True)
    add(prs, 250, 288, C5, P)
    add(prs, 180, 288, C5, P)
    add(prs, 396, 288, C4, [("Ser1", A4), ("Ser2", B4)])
    add(prs, 396, 288, C6, [("Ser1", A6), ("Ser2", B6)])
    add(prs, 250, 288, C5, P, legend=True)

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print("wrote", OUT)


if __name__ == "__main__":
    main()
