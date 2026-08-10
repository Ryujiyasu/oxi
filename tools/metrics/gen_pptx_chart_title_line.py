# -*- coding: utf-8 -*-
"""Explicit chart-title probe for LINE charts: LINE_MARKERS with an EXPLICIT
<c:title> (chart.has_title=True + chart_title.text). chart_title_line = 1
series ('Series 1'), chart_title_line2 = 2 series (Revenue/Cost). Both set the
explicit title text to "Quarterly Revenue" to measure how Word draws an
explicit title on a line chart (vs the automatic series-name title at
baseline sy+28.03 / plot_top sy+51.4 that the renderer already handles, and
vs the bar explicit title at baseline sy+24.43 / plot_top sy+45.69)."""
import sys, os, zipfile, re
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE


def make(base_name, n_series):
    base = r"pipeline_data\pptx_probes" + "\\" + base_name
    os.makedirs(base, exist_ok=True)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    chart_data = CategoryChartData()
    chart_data.categories = ["East", "West", "Midwest"]
    if n_series == 1:
        chart_data.add_series("Series 1", (19.2, 21.4, 16.7))
    else:
        chart_data.add_series("Revenue", (19.2, 22.0, 18.0))
        chart_data.add_series("Cost", (10.5, 11.2, 8.5))

    x, y, cx, cy = Inches(1.0), Inches(1.0), Inches(5.5), Inches(4.0)
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data
    )

    chart = gframe.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Quarterly Revenue"

    out = os.path.join(base, base_name + ".pptx")
    prs.save(out)
    print("saved:", out, os.path.getsize(out))

    z = zipfile.ZipFile(out)
    xml = z.read("ppt/charts/chart1.xml").decode("utf-8")
    m = re.search(r"<c:title>.*?</c:title>", xml, re.S)
    print("  TITLE:", (m.group(0)[:300] if m else "NONE"))
    print("  autoTitleDeleted:", "<c:autoTitleDeleted" in xml)


make("chart_title_line", 1)
make("chart_title_line2", 2)
