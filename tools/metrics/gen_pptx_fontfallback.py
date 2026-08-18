# -*- coding: utf-8 -*-
"""Probe: what face does PowerPoint use when the requested font is absent?

d19's master title asks for "Mali", the deck embeds nothing, and PowerPoint's
own PDF contains only Calibri -- while that deck's theme major/minor font is
Arial, so the substitute is NOT the theme font. d02 shows the same for
"Nunito". Oxi instead lets GDI's font mapper choose, and picks a face whose
caps are 15% taller (33.1pt vs 37.9pt at 52pt on d19 slide 1).

Each arm asks for a face that is not installed, with the theme left at its
default, so the reader can see what PowerPoint puts in its place.
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\fontfallback").resolve()
# Names that do not exist on a stock Windows install. The last two are real
# Google fonts the corpus asks for; the others are deliberately invented.
FONTS = ["Mali", "Nunito", "Jua", "Zzyzx Nonexistent", "Fira Sans", "Lobster"]


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for fam in FONTS:
        s = prs.slides.add_slide(blank)
        cap = s.shapes.add_textbox(Emu(228600), Emu(114300), Emu(6400800), Emu(300000))
        cap.text_frame.text = f"requested: {fam}"
        box = s.shapes.add_textbox(Emu(457200), Emu(914400), Emu(7772400), Emu(2000000))
        tf = box.text_frame
        tf.word_wrap = False
        r = tf.paragraphs[0].add_run()
        r.text = "Handgloves ABC xyz 123"
        r.font.size = Pt(40)
        r.font.name = fam
    prs.save(OUT / "fontfallback.pptx")
    print(f"wrote {OUT / 'fontfallback.pptx'}  ({len(FONTS)} arms)")


if __name__ == "__main__":
    main()
