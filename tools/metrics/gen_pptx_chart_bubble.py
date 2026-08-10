# -*- coding: utf-8 -*-
"""Bubble-chart probe deck (chart_bubble).

Eight arms that separate the levers of the bubble renderer:

  U1  1 series, sizes 1/2/4      -> area-vs-radius law (radii ratio 1:1.41:2 if area)
  U2  1 series, sizes 1/4/9      -> confirm (ratio 1:2:3 if area, 1:4:9 if radius)
  U3  2 series                   -> per-series colour + auto-title behaviour
  U4  1 series + legend          -> legend band + swatch shape
  U5  1 series, wide frame 500   -> does the bubble size follow the plot extent?
  U6  1 series + data labels     -> label anchor
  U7  1 series, negative X/Y     -> zero-line anchoring (S:2026-08-10 negative model)
  U8  2 series + legend + title  -> explicit <c:title> geometry

Run:  python tools/metrics/gen_pptx_chart_bubble.py
"""
import os
import sys

from pptx import Presentation
from pptx.chart.data import BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Emu, Pt

sys.stdout.reconfigure(encoding="utf-8")

OUT = r"pipeline_data\pptx_probes\chart_bubble\chart_bubble.pptx"


def pt(v):
    return Emu(int(round(v * 12700)))


def add(prs, x, y, w, h, kind, series, legend=False, labels=False,
        title=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = BubbleChartData()
    for name, pts in series:
        s = cd.add_series(name)
        for (xv, yv, sz) in pts:
            s.add_data_point(xv, yv, sz)
    gf = slide.shapes.add_chart(kind, pt(x), pt(y), pt(w), pt(h), cd)
    ch = gf.chart
    ch.has_legend = legend
    if legend:
        ch.legend.include_in_layout = False
    if labels:
        plot = ch.plots[0]
        plot.has_data_labels = True
        plot.data_labels.show_value = True
    if title is not None:
        ch.has_title = True
        ch.chart_title.text_frame.text = title
    return ch


def main():
    prs = Presentation()
    prs.slide_width = pt(720)
    prs.slide_height = pt(540)

    S1 = [(1.0, 19.2, 1.0), (2.0, 21.4, 2.0), (3.0, 16.7, 4.0)]
    S2 = [(1.0, 19.2, 1.0), (2.0, 21.4, 4.0), (3.0, 16.7, 9.0)]
    T2 = [(1.0, 10.5, 2.0), (2.0, 11.2, 3.0), (3.0, 8.5, 1.0)]
    NEG = [(-2.0, -8.0, 1.0), (1.0, 12.0, 2.0), (3.0, -4.0, 4.0)]

    add(prs, 72, 72, 396, 288, XL_CHART_TYPE.BUBBLE, [("Ser1", S1)])
    add(prs, 72, 72, 396, 288, XL_CHART_TYPE.BUBBLE, [("Ser1", S2)])
    add(prs, 72, 72, 396, 288, XL_CHART_TYPE.BUBBLE,
        [("Ser1", S1), ("Ser2", T2)])
    add(prs, 72, 72, 396, 288, XL_CHART_TYPE.BUBBLE, [("Ser1", S1)],
        legend=True)
    add(prs, 72, 72, 500, 288, XL_CHART_TYPE.BUBBLE, [("Ser1", S1)])
    add(prs, 72, 72, 396, 288, XL_CHART_TYPE.BUBBLE, [("Ser1", S1)],
        labels=True)
    add(prs, 72, 72, 396, 288, XL_CHART_TYPE.BUBBLE, [("Ser1", NEG)])
    add(prs, 72, 72, 396, 288, XL_CHART_TYPE.BUBBLE,
        [("Ser1", S1), ("Ser2", T2)], legend=True, title="Quarterly Mix")

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print("wrote", OUT)


if __name__ == "__main__":
    main()
