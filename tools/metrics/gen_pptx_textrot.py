# -*- coding: utf-8 -*-
"""Author the TURNED TEXT probe.

`create_font_for_wiu` passes 0 for `CreateFontW`'s escapement and orientation,
so a shape's text is always drawn upright no matter what `a:xfrm@rot` says.
d35 slide 34's competitor matrix puts "LOW VALUE 1" / "HIGH VALUE 1" down each
side at -90 degrees and Oxi lays them across the page; both reference renderers
beat Oxi on that slide. **174 text shapes over 17 of the 40 dev decks** are
turned, 133 of them exactly +/-90.

Three things have to be measured before any of it can be drawn:

  1. the direction the baseline runs at each `rot` (is it the shape's angle, and
     in which sense),
  2. what the text is turned ABOUT -- the box centre is the obvious guess and
     the obvious guess has been wrong before,
  3. whether the text is laid out in the shape's OWN box and then turned, or
     laid out in the turned bounding box (they differ the moment the box is not
     square: a 4x1 label at 90 degrees wraps completely differently).

So each arm is one text box on its own slide at a known offset, `<a:noAutofit/>`
and explicit insets so PowerPoint cannot resize or re-inset it, holding ONE
short line in a face with no ambiguity. Block W repeats the geometry with a line
long enough to wrap, which is what separates question 3.

Usage:
    python tools/metrics/gen_pptx_textrot.py
    python tools/metrics/export_pptx_textrot.py    # PowerPoint COM -> PDF
    python tools/metrics/read_pptx_textrot.py      # read the PDF back
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

REPO = Path(__file__).resolve().parents[2]
OUT = REPO / "pipeline_data" / "pptx_probes" / "textrot"

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
P = "http://schemas.openxmlformats.org/presentationml/2006/main"
EMU_IN = 914400
DEG = 60000

# A deliberately NON-square box, centred on the slide so no arm can fall off it
# at any angle: 4.00 x 1.00in about (5.00in, 3.75in) on the default 10 x 7.5in.
BOX_W, BOX_H = 4 * EMU_IN, EMU_IN
BOX_X, BOX_Y = 5 * EMU_IN - BOX_W // 2, 3 * EMU_IN + EMU_IN // 2
INS = 91440  # 0.10in = 7.2pt on every side
SIZE = 24
FACE = "Arial"

ANGLES = [0, 30, 45, 90, 135, 180, 270, -45, -90]
# One line that fits, and one that cannot fit the 4in box's 3.8in text width.
SHORT = "Mx"
LONG = "Mxxxxxxxxxx xxxxxxxxxx xxxxxxxxxx Nx"


def add_arm(prs, name, rot, text, anchor="t", align="l"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Emu(BOX_X), Emu(BOX_Y), Emu(BOX_W), Emu(BOX_H))
    tf = box.text_frame
    body = tf._txBody.find(f"{{{A}}}bodyPr")
    for child in list(body):
        body.remove(child)
    body.set("anchor", anchor)
    body.set("anchorCtr", "0")
    body.set("wrap", "square")
    for k in ("lIns", "rIns", "tIns", "bIns"):
        body.set(k, str(INS))
    etree.SubElement(body, f"{{{A}}}noAutofit")
    p = tf.paragraphs[0]
    p._p.get_or_add_pPr().set("algn", align)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(SIZE)
    r.font.name = FACE
    r.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    xfrm = box._element.spPr.find(f"{{{A}}}xfrm")
    if rot:
        xfrm.set("rot", str(int(round(rot * DEG))))
    # A hairline frame on the SAME xfrm marks where PowerPoint puts the box, so
    # the reader can answer question 2 without trusting any assumption about it.
    frame = slide.shapes.add_shape(1, Emu(BOX_X), Emu(BOX_Y), Emu(BOX_W), Emu(BOX_H))
    frame.fill.background()
    frame.line.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    frame.line.width = Pt(0.75)
    fx = frame._element.spPr.find(f"{{{A}}}xfrm")
    if rot:
        fx.set("rot", str(int(round(rot * DEG))))
    st = frame._element.find(f"{{{P}}}style")
    if st is not None:
        frame._element.remove(st)
    return {"arm": name, "rot": rot, "text": text, "anchor": anchor, "align": align,
            "box": [BOX_X, BOX_Y, BOX_W, BOX_H], "size": SIZE, "face": FACE,
            "ins": INS}


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    arms = []
    n = 0

    def emit(name, **kw):
        nonlocal n
        n += 1
        rec = add_arm(prs, name, **kw)
        rec["slide"] = n
        arms.append(rec)

    # R: the baseline direction and the turning centre.
    for a in ANGLES:
        emit(f"R_rot{a}", rot=a, text=SHORT)
    # W: does the line wrap against the shape's OWN width, or the turned box's?
    for a in (0, 90, -90, 180):
        emit(f"W_rot{a}", rot=a, text=LONG)
    # A: anchor and alignment inside a turned box.
    for a in (0, 90):
        emit(f"A_rot{a}_ctr", rot=a, text=SHORT, anchor="ctr", align="ctr")
        emit(f"A_rot{a}_b_r", rot=a, text=SHORT, anchor="b", align="r")

    path = OUT / "probe_textrot.pptx"
    prs.save(str(path))
    (OUT / "arms.json").write_text(json.dumps(arms, indent=1), encoding="utf-8")
    print(f"wrote {path}  {len(arms)} slides")
    print(f"wrote {OUT / 'arms.json'}")


if __name__ == "__main__":
    main()
