# -*- coding: utf-8 -*-
"""Author a minimal repro for `a:outerShdw` — a shape's drop shadow.

Oxi parses the element only so its colour cannot repaint the shape, and draws
nothing; 130 shadows sit unrendered across 20 dev slides, which average 0.9458
against a 0.9574 corpus. Before drawing one, the four knobs have to be pinned:

  blurRad  how far the penumbra reaches
  dist/dir where the shadow sits relative to the shape
  alpha    how dark it lands on the backdrop
  algn     which corner the (unscaled) shadow hangs from

Each slide sweeps ONE of them with the others held, on plain black squares over
a white background, spaced far enough apart that no two penumbras meet. That
way a scanline through a square's middle reads one shadow at a time.

    python tools/metrics/gen_pptx_shadow.py [out.pptx]
"""
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

OUT = Path(sys.argv[1] if len(sys.argv) > 1
           else r'pipeline_data\pptx_probes\shadow\probe_shadow.pptx')

SIDE_PT = 36.0          # square side
PITCH_PT = 108.0        # centre-to-centre, 3x the side
MARGIN_PT = 54.0

# (slide label, [(blurRad pt, dist pt, dir deg, alpha %), ...])
SWEEPS = [
    ('blur', [(b, 3.0, 0.0, 100.0) for b in (0.0, 2.25, 4.5, 9.0, 18.0, 36.0)]),
    ('dist', [(4.5, d, 0.0, 100.0) for d in (0.0, 0.75, 1.5, 3.0, 6.0, 12.0)]),
    ('dir', [(4.5, 9.0, a, 100.0) for a in (0.0, 45.0, 90.0, 135.0, 180.0, 270.0)]),
    ('alpha', [(4.5, 9.0, 0.0, a) for a in (100.0, 80.0, 54.9, 43.0, 20.0, 10.0)]),
]


def shadow_xml(blur_pt, dist_pt, dir_deg, alpha_pc):
    a = ('<a:alpha val="%d"/>' % round(alpha_pc * 1000)) if alpha_pc < 100 else ''
    return (
        '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:outerShdw blurRad="{round(blur_pt * 12700)}" '
        f'dist="{round(dist_pt * 12700)}" dir="{round(dir_deg * 60000)}" '
        'rotWithShape="0">'
        f'<a:srgbClr val="000000">{a}</a:srgbClr>'
        '</a:outerShdw></a:effectLst>'
    )


def main() -> None:
    from lxml import etree
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    manifest = []

    for si, (label, arms) in enumerate(SWEEPS, start=1):
        slide = prs.slides.add_slide(blank)
        for i, (blur, dist, dirn, alpha) in enumerate(arms):
            x = MARGIN_PT + (i % 3) * PITCH_PT * 1.6
            y = MARGIN_PT + (i // 3) * PITCH_PT * 1.6
            sq = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Pt(x), Pt(y), Pt(SIDE_PT), Pt(SIDE_PT))
            sq.fill.solid()
            sq.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
            sq.line.fill.background()
            sq.shadow.inherit = False
            spPr = sq._element.spPr
            for old in spPr.findall(qn('a:effectLst')):
                spPr.remove(old)
            spPr.append(etree.fromstring(shadow_xml(blur, dist, dirn, alpha)))
            manifest.append({
                'slide': si, 'sweep': label, 'blur_pt': blur, 'dist_pt': dist,
                'dir_deg': dirn, 'alpha_pc': alpha,
                'x_pt': x, 'y_pt': y, 'side_pt': SIDE_PT,
            })

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    OUT.with_suffix('.json').write_text(json.dumps(manifest, indent=1),
                                        encoding='utf-8')
    print('wrote', OUT, 'and', OUT.with_suffix('.json'))
    print(f'{len(SWEEPS)} slides, {len(manifest)} shadowed squares')


if __name__ == '__main__':
    main()
