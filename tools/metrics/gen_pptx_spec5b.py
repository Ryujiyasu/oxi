# -*- coding: utf-8 -*-
"""Spec #5 (paragraph alignment) wave-2 repro: one TextBox with 4 multi-line
paragraphs (Left/Center/Right/Justify), same font/size. Each wraps to ~3 lines."""
import sys, os
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

base = r"pipeline_data\pptx_probes\spec5b_multiline"
os.makedirs(base, exist_ok=True)

TEXT = ("The quick brown fox jumps over the lazy dog. "
        "The quick brown fox jumps over the lazy dog. "
        "The quick brown fox jumps over the lazy dog.")

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(4.0), Inches(1.5))
tf = tb.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.1)
tf.margin_right = Inches(0.1)
tf.margin_top = Inches(0.05)
tf.margin_bottom = Inches(0.05)

items = [
    ("Left", PP_ALIGN.LEFT),
    ("Center", PP_ALIGN.CENTER),
    ("Right", PP_ALIGN.RIGHT),
    ("Justify", PP_ALIGN.JUSTIFY),
]
first = True
for name, align in items:
    if first:
        p = tf.paragraphs[0]
        first = False
    else:
        p = tf.add_paragraph()
    p.alignment = align
    r = p.add_run()
    r.text = TEXT
    r.font.size = Pt(18)
    r.font.name = "Arial"

out = os.path.join(base, "spec5b_multiline.pptx")
prs.save(out)
print("saved:", out, os.path.getsize(out))
