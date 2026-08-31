# -*- coding: utf-8 -*-
"""Where inside a preset shape does PowerPoint put the text?

`layout_text_shape` lays text out in the shape's BOUNDING BOX. A preset
geometry defines its own text rectangle, and for most presets that rectangle is
smaller -- an ellipse holds text only in the box inscribed at 45 degrees, a
homePlate stops short of the point.

Two corpus slides say so, and say it to a tenth of a point:

    d35 s17  homePlate adj=30129  198.64pt box  centred 'first'
             bounding box would put it 7.859pt right of where PowerPoint did
             (observed 7.984)
    d15 s17  homePlate adj=50000  279.28pt box  centred 'LOREM 1'
             bounding box would put it 6.585pt right (observed 6.542)

The census says this is not two slides but a class: text-carrying shapes with a
non-rect preset number 758 ellipse, 165 homePlate, 72 teardrop, 50 pie, 42
roundRect across 14 decks of the two corpora.

Rather than trust the published formulas, this MEASURES the rectangle. Each
preset gets three arms with the same word and no insets:

    algn="l"    the pen lands on the text rect's LEFT
    algn="r"    the pen lands on its RIGHT minus the line's width
    algn="ctr"  a third reading that must agree with the other two

so left and right both come out, and the centre is a check rather than an
assumption. `wrap="none"` keeps it one line in every arm.

    python tools/metrics/gen_pptx_textrect.py
    python tools/metrics/export_pptx_textrect.py
    python tools/metrics/read_pptx_textrect.py
"""
import os
import sys

from lxml import etree
from pptx import Presentation
from pptx.util import Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

BASE = os.path.join("pipeline_data", "pptx_probes", "textrect")
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

TEXT = "Wq"          # short, so even a narrow text rect holds it
SIZE_PT = 18.0
BOX = (200.0, 150.0, 300.0, 200.0)   # left, top, w, h -- w != h so the
                                     # smaller side (`ss`) is identifiable

# (preset, [(adjust name, value)]) -- the counts are text-carrying shapes in
# the dev + wider corpora.
PRESETS = [
    ("rect", []),                      # control: text rect IS the box
    ("ellipse", []),                   # 758
    ("homePlate", []),                 # 165, default adj
    ("homePlate", [("adj", 30129)]),   # d35 s17's value
    ("homePlate", [("adj", 50000)]),   # d15 s17's value
    ("teardrop", []),                  # 72
    ("pie", []),                       # 50
    ("roundRect", []),                 # 42
    ("chevron", []),                   # 8
    ("wedgeRectCallout", []),          # 13
]
ALIGNS = ["l", "r", "ctr"]


def main() -> None:
    os.makedirs(BASE, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(540)
    blank = prs.slide_layouts[6]

    arms = [(p, adj, a) for p, adj in PRESETS for a in ALIGNS]
    for preset, adjs, algn in arms:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Pt(BOX[0]), Pt(BOX[1]), Pt(BOX[2]), Pt(BOX[3]))
        sp = box._element
        # Replace the textbox's `rect` with the preset under test.
        geom = sp.spPr.find(f"{{{NS_A}}}prstGeom")
        geom.set("prst", preset)
        av = geom.find(f"{{{NS_A}}}avLst")
        for name, val in adjs:
            gd = etree.SubElement(av, f"{{{NS_A}}}gd")
            gd.set("name", name)
            gd.set("fmla", f"val {val}")

        tf = box.text_frame
        body = tf._txBody.find(f"{{{NS_A}}}bodyPr")
        # No insets, so the pen lands on the text rect's own edge.
        for attr, value in (("lIns", "0"), ("rIns", "0"), ("tIns", "0"),
                            ("bIns", "0"), ("anchor", "t"), ("wrap", "none")):
            body.set(attr, value)
        para = tf.paragraphs[0]
        para._p.get_or_add_pPr().set("algn", algn)
        run = para.add_run()
        run.text = TEXT
        run.font.size = Pt(SIZE_PT)
        run.font.name = "Arial"

    out = os.path.join(BASE, "textrect.pptx")
    prs.save(out)
    print(f"wrote {out}: {len(arms)} arms "
          f"({len(PRESETS)} presets x {len(ALIGNS)} alignments)")
    print(f"box left {BOX[0]:g}pt width {BOX[2]:g}pt height {BOX[3]:g}pt, "
          f"{TEXT!r} at {SIZE_PT:g}pt")


if __name__ == "__main__":
    main()
