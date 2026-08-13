# gen_pptx_theme_default2.py - probe with a MODIFIED theme: minorFont latin = Times New Roman, majorFont = Georgia.
# Confirms a font-less run resolves to minorFont (body) and a title placeholder to majorFont.
# Also varies the textbox font-less run size to confirm minorFont at any size.
import os
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.enum.text import PP_ALIGN

OUT_DIR = os.path.abspath(r"pipeline_data\pptx_probes\theme_default2")
os.makedirs(OUT_DIR, exist_ok=True)

prs = Presentation()
prs.slide_width = Emu(9144000)
prs.slide_height = Emu(6858000)

# V1: bare run (no rPr) textbox -> minorFont expected
slide = prs.slides.add_slide(prs.slide_layouts[6])
box = slide.shapes.add_textbox(Emu(914400), Emu(914400), Inches(5), Inches(2))
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Emu(914400); tf.margin_right = Emu(914400)
tf.margin_top = Emu(457200); tf.margin_bottom = Emu(457200)
tf.vertical_anchor = 1
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = "V1 bare run no rPr"

# V2: size-only 18pt textbox -> minorFont expected
slide = prs.slides.add_slide(prs.slide_layouts[6])
box = slide.shapes.add_textbox(Emu(914400), Emu(914400), Inches(5), Inches(2))
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Emu(914400); tf.margin_right = Emu(914400)
tf.margin_top = Emu(457200); tf.margin_bottom = Emu(457200)
tf.vertical_anchor = 1
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = "V2 size only 18pt"
r.font.size = Pt(18)

# V3: explicit Arial 18pt control -> ArialMT expected
slide = prs.slides.add_slide(prs.slide_layouts[6])
box = slide.shapes.add_textbox(Emu(914400), Emu(914400), Inches(5), Inches(2))
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Emu(914400); tf.margin_right = Emu(914400)
tf.margin_top = Emu(457200); tf.margin_bottom = Emu(457200)
tf.vertical_anchor = 1
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = "V3 explicit Arial 18pt"
r.font.size = Pt(18); r.font.name = "Arial"

# V4: title placeholder -> majorFont expected (Georgia)
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "V4 title placeholder"

# V5: body placeholder -> minorFont expected (Times New Roman)
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.placeholders[1].text = "V5 body placeholder"

# V6: size-only 24pt textbox -> minorFont at another size
slide = prs.slides.add_slide(prs.slide_layouts[6])
box = slide.shapes.add_textbox(Emu(914400), Emu(914400), Inches(5), Inches(2))
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Emu(914400); tf.margin_right = Emu(914400)
tf.margin_top = Emu(457200); tf.margin_bottom = Emu(457200)
tf.vertical_anchor = 1
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = "V6 size only 24pt"
r.font.size = Pt(24)

out = os.path.join(OUT_DIR, "theme_default2.pptx")
prs.save(out)

# rewrite theme1.xml: minorFont latin -> Times New Roman, majorFont latin -> Georgia
import zipfile, shutil
import re
tmp = out + ".tmp"
with zipfile.ZipFile(out, "r") as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
    for item in zin.namelist():
        data = zin.read(item)
        if item == "ppt/theme/theme1.xml":
            s = data.decode("utf-8")
            s = s.replace('<a:minorFont><a:latin typeface="Calibri"', '<a:minorFont><a:latin typeface="Times New Roman"')
            s = s.replace('<a:majorFont><a:latin typeface="Calibri"', '<a:majorFont><a:latin typeface="Georgia"')
            data = s.encode("utf-8")
        zout.writestr(item, data)
shutil.move(tmp, out)
print("saved (theme modified)", out)
