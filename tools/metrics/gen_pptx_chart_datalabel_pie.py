# -*- coding: utf-8 -*-
"""Chart data-labels probe for PIE charts: PIE with data labels so Word's
PDF exposes the exact data-label placement/format/font rules for the pie
branch (the bar data-label rule was derived from chart_dlbls S1-S5 in
2026-08-06; the line/pie placement must be measured separately).

  P1: pie, show_value=True          (default = OUTSIDE_END category labels)
  P2: pie, show_value=True + number_format="0.0%"
  P3: pie, show_value=True + position=CENTER

All on blank slides (layout 6), default Office theme, frame 72,72,396,288,
categories [East, West, Midwest], series values (19.2, 21.4, 16.7)."""
import sys, os
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION

base = r"pipeline_data\pptx_probes\chart_datalabel_pie"
os.makedirs(base, exist_ok=True)

prs = Presentation()
blanks = [prs.slide_layouts[6] for _ in range(3)]


def add_pie_slide(layout, num_fmt=None, pos=None):
    slide = prs.slides.add_slide(layout)
    cd = CategoryChartData()
    cd.categories = ["East", "West", "Midwest"]
    cd.add_series("Series 1", (19.2, 21.4, 16.7))
    x, y, cx, cy = Inches(1.0), Inches(1.0), Inches(5.5), Inches(4.0)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, cd)
    plot = gframe.chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_value = True
    if num_fmt is not None:
        dl.number_format_is_linked = False
        dl.number_format = num_fmt
    if pos is not None:
        dl.position = pos
    return slide


add_pie_slide(blanks[0])
add_pie_slide(blanks[1], num_fmt="0.0%")
add_pie_slide(blanks[2], pos=XL_LABEL_POSITION.CENTER)

out = os.path.join(base, "chart_datalabel_pie.pptx")
prs.save(out)
print("saved:", out, os.path.getsize(out))
