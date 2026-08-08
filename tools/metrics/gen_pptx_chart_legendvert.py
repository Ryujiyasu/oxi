# -*- coding: utf-8 -*-
"""Chart spec: the LEGEND BLOCK VERTICAL POSITION when labels wrap.

The doughnut residual probe gave 3 points (L=1,2,4) and two competing models
each explained only 2 of them.  Separate the levers:
  L  = max wrapped line count      n = entry count      which entry wraps.

All arms: doughnut, 1 series (auto-title), legend right,
include_in_layout=False, frame 72,72,396,288 -> wrap cap = 396/3-21.4 = 110.6pt.
Each ~63pt word therefore gets its own line."""
import sys, os

sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

base = r"pipeline_data\pptx_probes\chart_legendvert"
os.makedirs(base, exist_ok=True)

W2 = "Abcdefg Hijklmn"
W3 = "Abcdefg Hijklmn Opqrstu"
W5 = "Abcdefg Hijklmn Opqrstu Vwxyzab Cdefghi"
X2 = "Opqrstu Vwxyzab"
Y2 = "Cdefghi Jklmnop"

ARMS = [
    ("B1 L=2 n=3 first",  [W2, "Ef", "Gh"]),
    ("B2 L=3 n=3 first",  [W3, "Ef", "Gh"]),
    ("B3 L=2 n=3 LAST",   ["Ab", "Ef", W2]),
    ("B4 L=2 n=3 MID",    ["Ab", W2, "Gh"]),
    ("B5 L=2 n=2 first",  [W2, "Ef"]),
    ("B6 L=2 n=4 first",  [W2, "Ef", "Gh", "Ij"]),
    ("B7 L=2 n=3 ALL",    [W2, X2, Y2]),
    ("B8 L=5 n=3 first",  [W5, "Ef", "Gh"]),
]
VALS = (19.2, 21.4, 16.7, 12.3)

prs = Presentation()
prs.slide_width = Pt(720)
prs.slide_height = Pt(540)
for label, cats in ARMS:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("Series 1", VALS[: len(cats)])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, Pt(72), Pt(72), Pt(396), Pt(288), cd
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False

out = os.path.join(base, "chart_legendvert.pptx")
prs.save(out)
for i, (label, cats) in enumerate(ARMS, 1):
    print(f"  S{i}: {label:20s} cats={cats}")
print("\nsaved:", out, os.path.getsize(out))
