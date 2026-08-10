# -*- coding: utf-8 -*-
"""Bubble SIZE-scale sweep (chart_bubble_scale).

The first probe showed the largest bubble's radius is INDEPENDENT of the plot
width (U1 396pt vs U5 500pt gave identical radii) but moves with the vertical
extent.  This deck sweeps the frame HEIGHT with everything else fixed so the
law can be fitted, plus two arms that move the plot top instead (no title /
explicit title) to confirm the same variable drives both.

Run:  python tools/metrics/gen_pptx_chart_bubble_scale.py
"""
import os
import sys

from pptx import Presentation
from pptx.chart.data import BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Emu

sys.stdout.reconfigure(encoding="utf-8")

OUT = r"pipeline_data\pptx_probes\chart_bubble_scale\chart_bubble_scale.pptx"


def pt(v):
    return Emu(int(round(v * 12700)))


def add(prs, h, series, title=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = BubbleChartData()
    for name, pts in series:
        s = cd.add_series(name)
        for (xv, yv, sz) in pts:
            s.add_data_point(xv, yv, sz)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.BUBBLE, pt(72), pt(72),
                                pt(396), pt(h), cd)
    ch = gf.chart
    ch.has_legend = False
    if title is not None:
        ch.has_title = True
        ch.chart_title.text_frame.text = title


def main():
    prs = Presentation()
    prs.slide_width = pt(720)
    prs.slide_height = pt(560)

    S1 = [(1.0, 19.2, 1.0), (2.0, 21.4, 2.0), (3.0, 16.7, 4.0)]
    T2 = [(1.0, 10.5, 2.0), (2.0, 11.2, 3.0), (3.0, 8.5, 1.0)]

    # frame-height sweep, 1 series (auto title, so the top offset is fixed)
    for h in (160, 200, 240, 288, 340, 400, 460):
        add(prs, h, [("Ser1", S1)])
    # same height, no auto title (2 series) -> plot top moves up
    add(prs, 288, [("Ser1", S1), ("Ser2", T2)])
    # same height, explicit title -> plot top in between
    add(prs, 288, [("Ser1", S1)], title="T")
    # tall + 2 series, to separate "frame height" from "plot height"
    add(prs, 460, [("Ser1", S1), ("Ser2", T2)])

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print("wrote", OUT, len(prs.slides.__iter__.__self__._sldIdLst), "slides")


if __name__ == "__main__":
    main()
