# -*- coding: utf-8 -*-
"""Probe: how tall does PowerPoint grow a table row that is declared too short?

Oxi grows a row to `marT + marB + 1.2 x size` per line. Solving d35 s35's two
measurable baselines instead gives **1.3784 x size** for the text part, which
looks like the font's full ascent+descent rather than the 1.2 line advance —
but that is one slide and two equations. This deck asks PowerPoint directly.

Every table has rows declared 1pt (so growth always binds) holding one line at a
known size, in a face whose metrics can be read locally. Consecutive rows are
identical, so the baseline PITCH in PowerPoint's PDF is the grown row height,
read straight off without any model in between.

Arms sweep size x margin x face; a second table per arm holds two-line cells so
the per-line term can be separated from the constant.

Usage:
    python tools/metrics/gen_pptx_rowgrow.py
    python tools/metrics/measure_pptx_word.py pipeline_data/pptx_probes/rowgrow/rowgrow.pptx pipeline_data/pptx_probes/rowgrow
    python tools/metrics/read_pptx_rowgrow.py
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\rowgrow").resolve()
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
PT = 12700

# (face, size pt, margin pt, lines per cell)
ARMS = [
    ("Arial", 8.0, 7.2, 1),
    ("Arial", 12.0, 7.2, 1),
    ("Arial", 18.0, 7.2, 1),
    ("Arial", 24.0, 7.2, 1),
    ("Arial", 12.0, 0.0, 1),
    ("Arial", 12.0, 3.6, 1),
    ("Arial", 12.0, 7.2, 2),
    ("Segoe Script", 12.0, 7.2, 1),
    ("Segoe Script", 18.0, 7.2, 1),
    ("Comic Sans MS", 12.0, 7.2, 1),
    ("Comic Sans MS", 18.0, 7.2, 1),
]
ROWS = 5


def q(tag: str) -> str:
    return f"{{{A}}}{tag}"


def build_table(slide, face: str, size: float, margin: float, lines: int):
    # python-pptx's add_table gives a graphicFrame we then rewrite by hand, so
    # the cell properties are exactly what the probe means to state.
    shape = slide.shapes.add_table(
        ROWS, 2, Emu(457200), Emu(457200), Emu(4572000), Emu(ROWS * PT)
    )
    tbl = shape._element.graphic.graphicData.tbl
    # Drop the table style so no banding or default text style interferes.
    for pr in tbl.findall(q("tblPr")):
        tbl.remove(pr)
    pr = etree.SubElement(tbl, q("tblPr"))
    pr.set("firstRow", "0")
    pr.set("bandRow", "0")
    etree.SubElement(pr, q("noFill"))
    tbl.insert(0, pr)
    for tr in tbl.findall(q("tr")):
        tr.set("h", str(int(1 * PT)))  # 1pt: always too short
        for tc in tr.findall(q("tc")):
            txbody = tc.find(q("txBody"))
            for p in txbody.findall(q("p")):
                txbody.remove(p)
            for i in range(lines):
                p = etree.SubElement(txbody, q("p"))
                ppr = etree.SubElement(p, q("pPr"))
                ln = etree.SubElement(ppr, q("lnSpc"))
                etree.SubElement(ln, q("spcPct")).set("val", "100000")
                etree.SubElement(ppr, q("buNone"))
                r = etree.SubElement(p, q("r"))
                rpr = etree.SubElement(r, q("rPr"))
                rpr.set("lang", "en-US")
                rpr.set("sz", str(int(round(size * 100))))
                rpr.set("kern", "0")
                etree.SubElement(rpr, q("latin")).set("typeface", face)
                etree.SubElement(r, q("t")).text = f"Hxy{i}"
            tcpr = tc.find(q("tcPr"))
            if tcpr is None:
                tcpr = etree.SubElement(tc, q("tcPr"))
            for name in ("marT", "marB", "marL", "marR"):
                tcpr.set(name, str(int(round(margin * PT))))
            tcpr.set("anchor", "ctr")
    return shape


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    manifest = []
    for face, size, margin, lines in ARMS:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cap = slide.shapes.add_textbox(Emu(228600), Emu(114300), Emu(6400800), Emu(300000))
        cap.text_frame.text = f"{face} {size}pt mar={margin} lines={lines}"
        build_table(slide, face, size, margin, lines)
        manifest.append({
            "slide": len(manifest) + 1,
            "face": face,
            "size": size,
            "margin": margin,
            "lines": lines,
            "declared_h": 1.0,
            "rows": ROWS,
        })
        print(f"{face} {size}pt mar={margin} lines={lines}")
    prs.save(OUT / "rowgrow.pptx")
    (OUT / "rowgrow_manifest.json").write_text(
        json.dumps(manifest, indent=1), encoding="utf-8")
    print(f"wrote {OUT / 'rowgrow.pptx'} ({len(manifest)} slides)")


if __name__ == "__main__":
    main()
