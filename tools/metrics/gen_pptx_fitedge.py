# -*- coding: utf-8 -*-
"""Author the FIT-EDGE probe: exactly how full may a line be before it breaks?

Deck 47 s2 raised it. `'You have to be signed in to your Google account.'` in
Caladea 12pt bold measures 258.75pt by PowerPoint's own character steps -- which
is 2070 master units EXACTLY -- inside a 258.76pt text area, and PowerPoint
breaks it anyway. The engine's `fits` is inclusive (`mu/8 <= width + 1e-6`), so
it keeps the line whole and comes out one line short.

Inclusive-vs-strict cannot be told apart anywhere except at the boundary, so
this puts the boundary under a sweep: one string, one face, and a box width
stepped in SIXTEENTHS of a point through the exact fit. Whichever arm first
breaks names the rule.

    text width W (master units, exact)
    box = W + k/16 pt for k = -4 .. +6

`wrap="square"`, no insets, `noAutofit`, one run -- so the text area IS the box
width and nothing else can move the answer.

Usage:
    python tools/metrics/gen_pptx_fitedge.py
    python tools/metrics/read_pptx_fitedge_com.py    # COM, no export needed
"""
from __future__ import annotations

import json
import os
import sys
from pathlib import Path

from fontTools.ttLib import TTFont
from lxml import etree
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

REPO = Path(__file__).resolve().parents[2]
OUT = REPO / "pipeline_data" / "pptx_probes" / "fitedge"
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
EMU_PT = 12700

FACE = "Arial"
SIZE = 12.0
TEXT = "You have to be signed in to your Google account."
# Sixteenths of a point around the exact fit, so a HALF master unit either way
# is visible -- the deck-47 case sat 0.08 master units inside the box.
STEPS = list(range(-4, 7))


def master_width_pt(face_file: str, text: str, size: float) -> float:
    """The engine's own model: each advance on the master unit, then summed."""
    f = TTFont(Path(os.environ["WINDIR"]) / "Fonts" / face_file,
               lazy=True, checkChecksums=0)
    upm = f["head"].unitsPerEm
    cmap = f.getBestCmap()
    total = 0
    for ch in text:
        g = cmap[ord(ch)]
        em = f["hmtx"][g][0] / upm
        total += round(em * size * 8.0)
    return total / 8.0


def add_arm(slide, box_pt: float, y_pt: float) -> None:
    box = slide.shapes.add_textbox(Emu(int(20 * EMU_PT)), Emu(int(y_pt * EMU_PT)),
                                   Emu(round(box_pt * EMU_PT)), Emu(int(60 * EMU_PT)))
    tx = box.text_frame._txBody
    body = tx.find(f"{{{A}}}bodyPr")
    for k, v in (("lIns", "0"), ("rIns", "0"), ("tIns", "0"), ("bIns", "0"),
                 ("wrap", "square"), ("anchor", "t")):
        body.set(k, v)
    for tag in ("spAutoFit", "normAutofit", "noAutofit"):
        for e in body.findall(f"{{{A}}}{tag}"):
            body.remove(e)
    etree.SubElement(body, f"{{{A}}}noAutofit")
    for para in tx.findall(f"{{{A}}}p"):
        tx.remove(para)
    p = etree.SubElement(tx, f"{{{A}}}p")
    ppr = etree.SubElement(p, f"{{{A}}}pPr")
    ppr.set("algn", "l")
    ppr.set("marL", "0")
    ppr.set("indent", "0")
    for tag in ("spcBef", "spcAft"):
        e = etree.SubElement(ppr, f"{{{A}}}{tag}")
        etree.SubElement(e, f"{{{A}}}spcPts").set("val", "0")
    etree.SubElement(ppr, f"{{{A}}}buNone")
    r = etree.SubElement(p, f"{{{A}}}r")
    rpr = etree.SubElement(r, f"{{{A}}}rPr")
    rpr.set("lang", "en-US")
    rpr.set("sz", str(int(round(SIZE * 100))))
    rpr.set("kern", "0")
    etree.SubElement(rpr, f"{{{A}}}latin").set("typeface", FACE)
    etree.SubElement(r, f"{{{A}}}t").text = TEXT


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    w = master_width_pt("arial.ttf", TEXT, SIZE)
    print(f"{FACE} {SIZE}pt: {TEXT!r}")
    print(f"  master-unit width = {w:.4f}pt = {int(round(w * 8))} master units")
    pres = Presentation()
    pres.slide_width = Emu(int(720 * EMU_PT))
    pres.slide_height = Emu(int(405 * EMU_PT))
    blank = pres.slide_layouts[6]
    meta = []
    for i, k in enumerate(STEPS):
        box = w + k / 16.0
        slide = pres.slides.add_slide(blank)
        add_arm(slide, box, 40.0)
        meta.append({"slide": i + 1, "k16": k, "box_pt": round(box, 4),
                     "text_pt": round(w, 4), "face": FACE, "sz": SIZE, "text": TEXT})
        print(f"  arm {i+1:2}: box {box:9.4f}pt  = text {k:+d}/16pt "
              f"({(box - w) * 8:+.3f} master units)")
    path = OUT / "probe_fitedge.pptx"
    pres.save(str(path))
    (OUT / "arms.json").write_text(json.dumps(meta, indent=1), encoding="utf-8")
    print(f"wrote {path} -- {len(meta)} arms")


if __name__ == "__main__":
    main()
