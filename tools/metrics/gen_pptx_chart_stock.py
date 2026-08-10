# -*- coding: utf-8 -*-
"""Stock-chart probe deck (chart_stock).

python-pptx cannot author <c:stockChart> ("XML writer for chart type STOCK_HLC
not yet implemented"), so the deck is built as a LINE chart and each chart part
is then rewritten in the zip:

    <c:lineChart>            -> <c:stockChart>
    drop chart-level <c:grouping>/<c:varyColors>/<c:marker>
    insert <c:hiLowLines/>   (and <c:upDownBars> for the OHLC arms)

CT_StockChart's child order is  ser*, dLbls?, dropLines?, hiLowLines?,
upDownBars?, axId, axId  -- so the inserts go immediately before the first
<c:axId>.  Its <c:ser> is the same CT_LineSer python-pptx already emits, so the
series bodies need no surgery.

The arms separate the levers:

  K1  HLC  4 cats, no legend            baseline geometry + bar/tick shape
  K2  HLC  4 cats, legend               legend band
  K3  HLC  7 cats, no legend            category pitch
  K4  HLC  4 cats, frame 500 wide       width lever
  K5  HLC  4 cats, frame 400 tall       height lever
  K6  OHLC 4 cats, no legend            open tick / up-down bars
  K7  OHLC 4 cats, legend + title       title geometry
  K8  HLC  4 cats, explicit title       plot_top with an explicit title

Word wants the series in order High, Low, Close (HLC) or
Open, High, Low, Close (OHLC).

Run:  python tools/metrics/gen_pptx_chart_stock.py
"""
import os
import re
import shutil
import sys
import zipfile

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Emu

sys.stdout.reconfigure(encoding="utf-8")

OUT = r"pipeline_data\pptx_probes\chart_stock\chart_stock.pptx"

# 1-based chart part index -> True when the arm needs <c:upDownBars>
UPDOWN = {6: True, 7: True}


def pt(v):
    return Emu(int(round(v * 12700)))


def add(prs, w, h, cats, series, legend=False, title=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals)
    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, pt(72), pt(72), pt(w), pt(h), cd)
    ch = gf.chart
    ch.has_legend = legend
    if legend:
        ch.legend.include_in_layout = False
    if title is not None:
        ch.has_title = True
        ch.chart_title.text_frame.text = title
    return ch


def to_stock(xml: str, updown: bool) -> str:
    """Rewrite one chart part from a line chart into a stock chart."""
    assert "<c:lineChart>" in xml, "expected a lineChart body"
    body_lo = xml.index("<c:lineChart>")
    body_hi = xml.index("</c:lineChart>") + len("</c:lineChart>")
    body = xml[body_lo:body_hi]

    # chart-level knobs that CT_StockChart does not allow
    for pat in (r"<c:grouping[^>]*/>", r"<c:varyColors[^>]*/>"):
        body = re.sub(pat, "", body, count=1)
    # the chart-level <c:marker val="1"/> and <c:smooth val="0"/> sit after
    # the last </c:ser>; neither is legal in CT_StockChart
    tail_lo = body.rindex("</c:ser>") + len("</c:ser>")
    head, tail = body[:tail_lo], body[tail_lo:]
    tail = re.sub(r"<c:marker[^>]*/>", "", tail, count=1)
    tail = re.sub(r"<c:smooth[^>]*/>", "", tail, count=1)
    body = head + tail

    # A real stock chart hides the connecting polyline on every series -- the
    # hi-low bars carry the data.  CT_LineSer order is idx, order, tx, spPr,
    # marker, ... so the spPr goes straight after </c:tx>.
    body = body.replace(
        "</c:tx>", "</c:tx><c:spPr><a:ln><a:noFill/></a:ln></c:spPr>")

    extra = "<c:hiLowLines/>"
    if updown:
        extra += '<c:upDownBars><c:gapWidth val="150"/></c:upDownBars>'
    ax = body.index("<c:axId")
    body = body[:ax] + extra + body[ax:]

    body = body.replace("<c:lineChart>", "<c:stockChart>")
    body = body.replace("</c:lineChart>", "</c:stockChart>")
    return xml[:body_lo] + body + xml[body_hi:]


def patch(path: str) -> None:
    tmp = path + ".tmp"
    with zipfile.ZipFile(path) as zin, \
            zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for it in zin.infolist():
            data = zin.read(it.filename)
            m = re.fullmatch(r"ppt/charts/chart(\d+)\.xml", it.filename)
            if m:
                idx = int(m.group(1))
                data = to_stock(
                    data.decode("utf-8"), UPDOWN.get(idx, False)
                ).encode("utf-8")
                print("  chart%d.xml -> stockChart%s"
                      % (idx, " + upDownBars" if UPDOWN.get(idx) else ""))
            zout.writestr(it, data)
    shutil.move(tmp, path)


def main():
    prs = Presentation()
    prs.slide_width = pt(720)
    prs.slide_height = pt(540)

    C4 = ["Q1", "Q2", "Q3", "Q4"]
    C7 = ["Q1", "Q2", "Q3", "Q4", "Q5", "Q6", "Q7"]

    HI4 = (24.0, 26.5, 22.8, 27.2)
    LO4 = (18.2, 19.6, 16.4, 20.1)
    CL4 = (21.5, 25.0, 18.9, 24.4)
    OP4 = (19.0, 20.3, 21.7, 21.0)

    HI7 = HI4 + (25.9, 23.4, 28.0)
    LO7 = LO4 + (19.0, 17.2, 21.3)
    CL7 = CL4 + (23.1, 19.8, 26.6)

    HLC4 = [("High", HI4), ("Low", LO4), ("Close", CL4)]
    HLC7 = [("High", HI7), ("Low", LO7), ("Close", CL7)]
    OHLC4 = [("Open", OP4), ("High", HI4), ("Low", LO4), ("Close", CL4)]

    add(prs, 396, 288, C4, HLC4)
    add(prs, 396, 288, C4, HLC4, legend=True)
    add(prs, 396, 288, C7, HLC7)
    add(prs, 500, 288, C4, HLC4)
    add(prs, 396, 400, C4, HLC4)
    add(prs, 396, 288, C4, OHLC4)
    add(prs, 396, 288, C4, OHLC4, legend=True, title="Quarterly Range")
    add(prs, 396, 288, C4, HLC4, title="Price Band")

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    patch(OUT)
    print("wrote", OUT)


if __name__ == "__main__":
    main()
