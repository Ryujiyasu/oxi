# -*- coding: utf-8 -*-
"""What horizontal scale does PowerPoint set a line at, as a function of size?

A line's drawn width is its /Widths sum times one scale `s` (see
`pptx_line_condense.py`), and the corpus says `s` depends on the FACE and the
SIZE and on nothing else -- the same face at the same size gives the same
number in different decks, for different sentences, to five decimals.

`size * s` then lands within 0.04pt of `round(size)` on every Arial size the
corpus offers, which reads as: PowerPoint lays the line out at the INTEGER
point size and the PDF draws that layout at the exact size. But the corpus
cannot settle it, because its decks are scaled to A4 and their sizes are
mostly NOT integers (15.96, 20.04, 12.96); the few integer ones still carry a
residual of -0.2% to -0.5% that the integer story does not explain.

So this asks directly: one sentence, one face, one line, at sizes that
deliberately mix

  * exact integers                     -- where round(size) == size, so the
                                          dominant term is 1 and only the
                                          residual is left
  * quarter and half points            -- where round(size) moves in known
                                          steps
  * the corpus's own A4-scaled values  -- to check the corpus reading

`wrap="none"` on every box, so no line can ever be squeezed to fit and the
demand-driven mechanism (`pptx-line-squeeze-law`) is out of the picture. The
slide is 720x540pt so the PDF export is 1:1 and `Tf` equals the declared `sz`.

    python tools/metrics/gen_pptx_sizesweep.py
    python tools/metrics/export_pptx_sizesweep.py     # PowerPoint must not
    python tools/metrics/read_pptx_sizesweep.py       # overlap a pptx render
"""
import os
import sys

from pptx import Presentation
from pptx.util import Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

BASE = os.path.join("pipeline_data", "pptx_probes", "sizesweep")
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# ★Long, because a scale is fitted to the line and a short line has no
# leverage: 18 characters left the fit's own residual at 0.2pt, which is the
# size of the effect being measured. The corpus runs that gave a scale
# reproducible to five decimals were 30 to 90 characters. The box is widened
# to match, and `wrap="none"` keeps it one line whatever the size.
TEXT = ("Handgloves quiz for the jumping wizard, box 37 pack my bag with "
        "five dozen liquor jugs")

INTEGERS = [6, 8, 9, 10, 11, 12, 14, 16, 18, 20, 21, 24]
FRACTIONS = [11.5, 12.25, 12.5, 12.75, 13.5, 18.5, 20.5]
CORPUS = [11.04, 12.024, 12.96, 15.96, 18.024, 20.04, 21.024, 23.04]
FACES = ["Arial", "Calibri"]


def main() -> None:
    os.makedirs(BASE, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(540)
    blank = prs.slide_layouts[6]

    arms = [(f, s) for f in FACES for s in INTEGERS + FRACTIONS + CORPUS]
    for face, size in arms:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Pt(20), Pt(200), Pt(680), Pt(120))
        tf = box.text_frame
        body = tf._txBody.find(f"{{{NS_A}}}bodyPr")
        for attr, value in (("lIns", "0"), ("rIns", "0"), ("tIns", "0"),
                            ("bIns", "0"), ("anchor", "t"),
                            # Never wrap: a line that cannot overflow cannot be
                            # squeezed to fit, which leaves only the constant.
                            ("wrap", "none")):
            body.set(attr, value)
        run = tf.paragraphs[0].add_run()
        run.text = TEXT
        # `sz` is hundredths of a point, so a quarter point is exact.
        run.font.size = Pt(size)
        run.font.name = face

    out = os.path.join(BASE, "sizesweep.pptx")
    prs.save(out)
    print(f"wrote {out}: {len(arms)} arms "
          f"({len(FACES)} faces x {len(INTEGERS)+len(FRACTIONS)+len(CORPUS)} sizes)")


if __name__ == "__main__":
    main()
