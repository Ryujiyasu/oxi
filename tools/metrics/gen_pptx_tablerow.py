# -*- coding: utf-8 -*-
"""Author the TABLE ROW-HEIGHT probe.

A table row is normally rendered at exactly its declared `<a:tr h>` -- measured
across the corpus, most tables match to 0.00pt. Two do not, and they are the
whole of d10 s12's defect (MIN 0.8830):

    d04 s35 row0   declared 18.55 -> rendered 24.00   (+5.45)
    d10 s12 row0   declared 81.31 -> rendered 83.84   (+2.53)

So the height is `max(declared, content)` and the open question is `content`.

d04's arm pins one exact fit -- `marT=marB=7.198pt`, one EMPTY 8pt paragraph, no
`lnSpc`:

    2 * 7.198 + 1.2 * 8 = 24.00      exact

d10 fits nothing yet. Its cell is `marT=marB=19.711pt`, one 29.99pt run in the
EMBEDDED face Jua, `<a:lnSpc><a:spcPct val="140013"/>`, borders 2.124pt:

    2*19.711 + 1.2*29.99*1.40013 = 89.81     too tall
    2*19.711 + 1.00*29.99*1.40013 = 81.41    equals the DECLARED height, not the rendered
    rendered 83.84  =>  (83.84 - 39.42) / 29.99 = 1.4811 em, or 1.0578 em once
                        lnSpc is divided out -- neither 1.2 nor any of Jua's
                        OS/2 ratios (typo 1.25, win 1.113)

Rather than keep guessing, this varies ONE thing at a time and reads the answer
off PowerPoint's own PDF, where table rules are vector rectangles and the row
heights are exact to 0.01pt.

Every arm is a 3-row table whose declared heights are deliberately TOO SMALL
(4pt), so the content height always wins and is what gets measured. Row 0 varies;
rows 1 and 2 are a fixed 8pt Arial control, so any arm that moves them is
telling us the effect is not row-local.

Usage:
    python tools/metrics/gen_pptx_tablerow.py
    python tools/metrics/export_pptx_tablerow.py   # PowerPoint COM -> PDF
    python tools/metrics/read_pptx_tablerow.py     # read the PDF back
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

REPO = Path(__file__).resolve().parents[2]
OUT = REPO / "pipeline_data" / "pptx_probes" / "tablerow"
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
EMU_PT = 12700

# (name, size_pt, marT_pt, lnSpc_pct or None, typeface, text)
ARMS = [
    # the size sweep at zero margin isolates the multiplier itself
    ("sz08_mar0",        8.0,  0.0,  None,   "Arial", "Ag"),
    ("sz16_mar0",       16.0,  0.0,  None,   "Arial", "Ag"),
    ("sz30_mar0",       30.0,  0.0,  None,   "Arial", "Ag"),
    # margins: does the row add 2*marT exactly?
    ("sz16_mar7",       16.0,  7.2,  None,   "Arial", "Ag"),
    ("sz16_mar20",      16.0, 19.711, None,  "Arial", "Ag"),
    # lnSpc percentage -- the term d10 needs and d04 does not have
    ("sz16_mar0_ln100", 16.0,  0.0,  100000, "Arial", "Ag"),
    ("sz16_mar0_ln140", 16.0,  0.0,  140013, "Arial", "Ag"),
    ("sz30_mar20_ln140",30.0, 19.711, 140013,"Arial", "Ag"),
    # an EMPTY paragraph, which is what d04's exact-fitting arm actually holds
    ("sz08_mar7_empty",  8.0,  7.198, None,  "Arial", ""),
    # a second face: if the multiplier is the FONT's line height rather than a
    # constant, these two must differ at the same size
    ("sz30_mar0_georgia",30.0, 0.0,  None,   "Georgia", "Ag"),
    ("sz30_mar0_verdana",30.0, 0.0,  None,   "Verdana", "Ag"),
    # two lines, to separate "per line" from "per row"
    ("sz16_mar0_2line", 16.0,  0.0,  None,   "Arial", "Ag\nAg"),
]

TINY_H = int(4 * EMU_PT)          # declared far too small, so content always wins
CTRL = ("Arial", 8.0)             # rows 1-2


def cell(text: str, sz: float, mar_t: float, ln_pct, face: str) -> etree._Element:
    tc = etree.Element(f"{{{A}}}tc")
    tx = etree.SubElement(tc, f"{{{A}}}txBody")
    etree.SubElement(tx, f"{{{A}}}bodyPr")
    etree.SubElement(tx, f"{{{A}}}lstStyle")
    for line in (text.split("\n") if text else [""]):
        p = etree.SubElement(tx, f"{{{A}}}p")
        ppr = etree.SubElement(p, f"{{{A}}}pPr")
        ppr.set("algn", "ctr")
        if ln_pct is not None:
            ls = etree.SubElement(ppr, f"{{{A}}}lnSpc")
            etree.SubElement(ls, f"{{{A}}}spcPct").set("val", str(ln_pct))
        for tag in ("spcBef", "spcAft"):
            e = etree.SubElement(ppr, f"{{{A}}}{tag}")
            etree.SubElement(e, f"{{{A}}}spcPts").set("val", "0")
        etree.SubElement(ppr, f"{{{A}}}buNone")
        if line:
            r = etree.SubElement(p, f"{{{A}}}r")
            rpr = etree.SubElement(r, f"{{{A}}}rPr")
            rpr.set("lang", "en-US")
            rpr.set("sz", str(int(round(sz * 100))))
            etree.SubElement(rpr, f"{{{A}}}latin").set("typeface", face)
            t = etree.SubElement(r, f"{{{A}}}t")
            t.text = line
        else:
            epr = etree.SubElement(p, f"{{{A}}}endParaRPr")
            epr.set("sz", str(int(round(sz * 100))))
            etree.SubElement(epr, f"{{{A}}}latin").set("typeface", face)
    pr = etree.SubElement(tc, f"{{{A}}}tcPr")
    m = str(int(round(mar_t * EMU_PT)))
    pr.set("marT", m)
    pr.set("marB", m)
    pr.set("marL", "0")
    pr.set("marR", "0")
    pr.set("anchor", "ctr")
    return tc


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    manifest = []
    for i, (name, sz, mar, ln, face, text) in enumerate(ARMS):
        slide = prs.slides.add_slide(blank)
        shp = slide.shapes.add_table(3, 2, Emu(457200), Emu(457200), Emu(4572000), Emu(457200))
        tbl = shp._element.graphic.graphicData.tbl
        for ri, tr in enumerate(tbl.findall(f"{{{A}}}tr")):
            tr.set("h", str(TINY_H))
            for tc in tr.findall(f"{{{A}}}tc"):
                tr.remove(tc)
            if ri == 0:
                tr.append(cell(text, sz, mar, ln, face))
                tr.append(cell(text, sz, mar, ln, face))
            else:
                tr.append(cell("Ag", CTRL[1], 0.0, None, CTRL[0]))
                tr.append(cell("Ag", CTRL[1], 0.0, None, CTRL[0]))
        manifest.append(
            {
                "slide": i + 1, "name": name, "sz_pt": sz, "marT_pt": mar,
                "lnSpc_pct": ln, "typeface": face, "text": text,
                "declared_h_pt": TINY_H / EMU_PT,
            }
        )
    path = OUT / "probe_tablerow.pptx"
    prs.save(str(path))
    (OUT / "manifest.json").write_text(json.dumps(manifest, indent=1), encoding="utf-8")
    print(f"wrote {path} ({len(ARMS)} arms, declared h = {TINY_H/EMU_PT}pt each)")
    for m in manifest:
        print(f"  s{m['slide']:>2} {m['name']:<20} sz={m['sz_pt']:<5} marT={m['marT_pt']:<7} "
              f"lnSpc={m['lnSpc_pct']} face={m['typeface']} text={m['text']!r}")


if __name__ == "__main__":
    main()
