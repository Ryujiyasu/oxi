# -*- coding: utf-8 -*-
"""Author the slide-number field probe: what number does <a:fld type="slidenum"> print?

Three arms, each a 6-slide deck whose every slide carries one text box holding
a `<a:fld type="slidenum">` between two literal markers, so the substituted
string can be read straight out of PowerPoint's PDF:

  arm 0  no firstSlideNum on p:presentation  (the schema default)
  arm 1  firstSlideNum="5"
  arm 2  firstSlideNum="100"

Each deck also states a STALE cached value inside the field on slide 3
(`<a:t>777</a:t>` instead of the usual `<a:t>#</a:t>`), which pins whether
PowerPoint prints the cache or recomputes.

Usage:
    python tools/metrics/gen_pptx_slidenum.py
    python tools/metrics/export_pptx_slidenum.py       # PowerPoint COM -> PDF
    python tools/metrics/read_pptx_slidenum.py         # read the PDFs back
"""
from __future__ import annotations

import re
import shutil
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

REPO = Path(__file__).resolve().parents[2]
OUT = REPO / "pipeline_data" / "pptx_probes" / "slidenum"
N_SLIDES = 6
ARMS = [None, 5, 100]

FLD = (
    '<a:fld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
    'id="{{2E6D6EF8-1F41-4B4E-9E2C-7B1F9D0A0001}}" type="slidenum">'
    '<a:rPr lang="en" sz="2400"/><a:t>{cached}</a:t></a:fld>'
)


def build(path: Path) -> None:
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for i in range(N_SLIDES):
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Emu(457200), Emu(1371600), Emu(7772400), Emu(914400))
        tf = box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "["
        r.font.size = Pt(24)
        # the field itself, spliced in as raw DrawingML after the marker run
        from lxml import etree

        cached = "777" if i == 2 else "#"
        fld = etree.fromstring(FLD.format(cached=cached))
        p._p.append(fld)
        r2 = p.add_run()
        r2.text = "]"
        r2.font.size = Pt(24)
    prs.save(str(path))


def set_first_slide_num(path: Path, n: int | None) -> None:
    """Rewrite ppt/presentation.xml with (or without) @firstSlideNum."""
    if n is None:
        return
    tmp = path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(path) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "ppt/presentation.xml":
                xml = data.decode("utf-8")
                xml = re.sub(r'\sfirstSlideNum="\d+"', "", xml)
                xml = xml.replace("<p:presentation ", f'<p:presentation firstSlideNum="{n}" ', 1)
                data = xml.encode("utf-8")
            zout.writestr(item, data)
    shutil.move(str(tmp), str(path))


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    for k, first in enumerate(ARMS):
        path = OUT / f"probe_slidenum_a{k}.pptx"
        build(path)
        set_first_slide_num(path, first)
        print(f"wrote {path}  firstSlideNum={first}")


if __name__ == "__main__":
    main()
