# -*- coding: utf-8 -*-
"""Probe: the geometry of a:highlight, the run-level text highlight.

Oxi does not draw it at all -- d11 slide 38's "and many more..." is white text
on a dk1 highlight, and Oxi drew the text with no box behind it. 65 runs over
19 slides in 8 of the 40 dev decks carry one, so the box needs deriving rather
than guessing: how tall is it against the font, where does its top sit against
the baseline, does line spacing move it, and does it cover the run's trailing
space.

PowerPoint exports the highlight as a filled rectangle, so pymupdf's drawing
list gives its exact bounds, and the text spans in the same page give the
baseline to measure them against.
"""
from __future__ import annotations

import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\highlight").resolve()
A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# (label, size pt, lnSpc pct, plain prefix, highlighted text, plain suffix,
#  typeface)
ARMS = [
    ("mid18", 18, 100, "before ", "HIGH", " after", "Arial"),
    ("mid36", 36, 100, "before ", "HIGH", " after", "Arial"),
    ("mid18_ln150", 18, 150, "before ", "HIGH", " after", "Arial"),
    ("mid18_ln70", 18, 70, "before ", "HIGH", " after", "Arial"),
    ("start18", 18, 100, "", "HIGH", " after", "Arial"),
    ("end18", 18, 100, "before ", "HIGH", "", "Arial"),
    ("space18", 18, 100, "before ", "HIGH ", "after", "Arial"),
    ("desc18", 18, 100, "before ", "gjpqy", " after", "Arial"),
    # Size and face are separated here: a constant ratio across sizes but a
    # different one per face means the box comes from the font's own vertical
    # metrics, and the same ratio everywhere means it is a PowerPoint constant.
    ("arial96", 96, 100, "b ", "HIGH", " a", "Arial"),
    ("times18", 18, 100, "before ", "HIGH", " after", "Times New Roman"),
    ("times96", 96, 100, "b ", "HIGH", " a", "Times New Roman"),
    ("courier18", 18, 100, "before ", "HIGH", " after", "Courier New"),
    ("segoe18", 18, 100, "before ", "HIGH", " after", "Segoe UI"),
    ("georgia96", 96, 100, "b ", "HIGH", " a", "Georgia"),
    # Does a taller neighbour on the same line grow the box? If it does, the
    # highlight is the LINE box; if not, it is the run's own font box, and a
    # mixed-size line needs the run's metrics, not the line's.
    ("tallneighbour", 18, 100, "", "HIGH", " BIG", "Arial"),
    # Two faces that set fsSelection bit 7, where the baseline model reads the
    # typo metrics plus the line gap instead of the win pair. If the box
    # follows the same pair, the height rule is one rule and not two.
    ("bahn96", 96, 100, "b ", "HIGH", " a", "Bahnschrift"),
    ("cascadia96", 96, 100, "b ", "HIGH", " a", "Cascadia Mono"),
]
# The tall neighbour is this many points; a separate list keeps the arm tuple
# uniform.
TALL = {"tallneighbour": 48}


def q(tag: str) -> str:
    return f"{{{A}}}{tag}"


def run(p, text, sz, highlight, face="Arial"):
    r = etree.SubElement(p, q("r"))
    rpr = etree.SubElement(r, q("rPr"))
    rpr.set("lang", "en-US")
    rpr.set("sz", str(sz * 100))
    fill = etree.SubElement(rpr, q("solidFill"))
    etree.SubElement(fill, q("srgbClr")).set("val", "FFFFFF" if highlight else "000000")
    if highlight:
        hl = etree.SubElement(rpr, q("highlight"))
        etree.SubElement(hl, q("srgbClr")).set("val", "FF0000")
    etree.SubElement(rpr, q("latin")).set("typeface", face)
    etree.SubElement(r, q("t")).text = text


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for label, sz, ln, pre, hi, post, face in ARMS:
        slide = prs.slides.add_slide(blank)
        cap = slide.shapes.add_textbox(Emu(228600), Emu(114300), Emu(6400800), Emu(300000))
        cap.text_frame.text = label
        box = slide.shapes.add_textbox(Emu(914400), Emu(1200000), Emu(5486400), Emu(1800000))
        tf = box.text_frame
        tf.word_wrap = False
        body = tf._txBody
        for pel in body.findall(q("p")):
            body.remove(pel)
        p = etree.SubElement(body, q("p"))
        ppr = etree.SubElement(p, q("pPr"))
        etree.SubElement(ppr, q("buNone"))
        lns = etree.SubElement(ppr, q("lnSpc"))
        etree.SubElement(lns, q("spcPct")).set("val", str(ln * 1000))
        if pre:
            run(p, pre, sz, False, face)
        run(p, hi, sz, True, face)
        if post:
            run(p, post, TALL.get(label, sz), False, face)
    prs.save(OUT / "highlight.pptx")
    print(f"wrote {OUT / 'highlight.pptx'}  ({len(ARMS)} arms)")


if __name__ == "__main__":
    main()
