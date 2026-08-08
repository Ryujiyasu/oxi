# -*- coding: utf-8 -*-
"""Chart spec: DOUGHNUT chart probe (<c:doughnutChart>).

The parser only recognises pieChart / lineChart / barChart, so a doughnut
chart's <c:ser> is never collected and the chart renders EMPTY.  This probe
exposes Word's doughnut geometry for measurement:

  S1: DOUGHNUT 1 series, no title/legend/labels  (ring radii + slice angles)
  S2: DOUGHNUT 1 series + legend                 (legend geometry)
  S3: DOUGHNUT 1 series + data labels            (label placement in the ring)
  S4: DOUGHNUT 1 series + explicit title         (ring top under a title)
  S5: DOUGHNUT 1 series, holeSize 25             (does the hole scale?)

Frame 72,72,396,288, categories [East, West, Midwest], default Office theme."""
import sys, os, re, shutil, zipfile

sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

base = r"pipeline_data\pptx_probes\chart_doughnut"
os.makedirs(base, exist_ok=True)

prs = Presentation()


def add(title=None, legend=False, dlbls=False, cats=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = cats or ["East", "West", "Midwest"]
    cd.add_series("Series 1", (19.2, 21.4, 16.7))
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, Pt(72), Pt(72), Pt(396), Pt(288), cd
    )
    chart = gframe.chart
    if title is not None:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    if legend:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
    if dlbls:
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.show_value = True


add()
add(legend=True)
add(dlbls=True)
add(title="Quarterly Revenue")
add()  # slide 5: holeSize patched to 25 below
# The legend band displaces the ring horizontally, so these arms vary it:
add()  # slide 6: <c:legend> deleted below -> no band at all
add(cats=["A", "B", "C"])                                     # 7: short labels
add(cats=["Northeastern Region", "Southwestern Region", "C"])  # 8: long labels

out = os.path.join(base, "chart_doughnut.pptx")
prs.save(out)

# Patch specific chart parts: slide 5 hole size, slide 6 legend removal.
tmp = out + ".tmp"
shutil.copy(out, tmp)
zin = zipfile.ZipFile(tmp, "r")
names = zin.namelist()
charts = sorted((n for n in names if re.match(r"ppt/charts/chart(\d+)\.xml$", n)),
                key=lambda n: int(re.search(r"(\d+)\.xml$", n).group(1)))
with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
    for n in names:
        data = zin.read(n)
        if n == charts[4]:  # slide 5
            xml = data.decode("utf-8")
            xml = re.sub(r'(<c:holeSize\s+val=")\d+(")', r"\g<1>25\g<2>", xml)
            data = xml.encode("utf-8")
            print("  slide5 holeSize ->",
                  re.search(r"<c:holeSize[^>]*/>", xml).group(0))
        elif n == charts[5]:  # slide 6
            xml = data.decode("utf-8")
            xml2 = re.sub(r"<c:legend>.*?</c:legend>", "", xml, flags=re.S)
            print("  slide6 legend removed:", xml != xml2)
            data = xml2.encode("utf-8")
        zout.writestr(n, data)
zin.close()
os.remove(tmp)

with zipfile.ZipFile(out) as z:
    x = z.read(charts[0]).decode("utf-8")
    m = re.search(r"<c:doughnutChart>.*?</c:ser>", x, re.S)
    print("\nslide1 doughnutChart head:")
    print("  " + (m.group(0)[:400] if m else "(not found)"))
    for tag in ("holeSize", "firstSliceAng", "varyColors"):
        mm = re.search(r"<c:%s[^>]*/>" % tag, x)
        print(f"  {tag}: {mm.group(0) if mm else '(absent)'}")

print("\nsaved:", out, os.path.getsize(out))
