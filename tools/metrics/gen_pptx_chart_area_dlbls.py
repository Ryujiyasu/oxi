# -*- coding: utf-8 -*-
"""Chart spec: AREA chart DATA LABELS (<c:dLbls>) + the legend overlay
discriminator.

The area branch (73cf2fa3) renders fills/axes/legend but has no data-label
support, and the legend band rule was carried over from pie/doughnut without
an area-specific overlay arm.  Levers kept orthogonal:

  D1 standard n=1            show_value            (default label position)
  D2 standard n=2            show_value            (per-series placement)
  D3 stacked  n=2            show_value            (segment placement)
  D4 100%     n=2            show_value            (percent formatting)
  D5 standard n=1            show_value + "0.0%"   (number format)
  D6 standard n=2 + legend   show_value            (include_in_layout=False)
  D7 standard n=2 + legend   show_value            (bare legend -> overlay?)

Frame 72,72,396,288 / cats Q1..Q4 so the geometry lines up with chart_area."""
import sys, os

sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

base = r"pipeline_data\pptx_probes\chart_area_dlbls"
os.makedirs(base, exist_ok=True)

CATS = ["Q1", "Q2", "Q3", "Q4"]
S1 = (19.2, 21.4, 16.7, 22.0)
S2 = (10.5, 11.2, 8.5, 12.3)

# (label, chart kind, series list, legend mode, number format)
#   legend mode: None = no legend / "band" = include_in_layout False / "bare"
ARMS = [
    ("D1 area  n=1 dlbls",        XL_CHART_TYPE.AREA,             [S1],     None,   None),
    ("D2 area  n=2 dlbls",        XL_CHART_TYPE.AREA,             [S1, S2], None,   None),
    ("D3 stack n=2 dlbls",        XL_CHART_TYPE.AREA_STACKED,     [S1, S2], None,   None),
    ("D4 100%  n=2 dlbls",        XL_CHART_TYPE.AREA_STACKED_100, [S1, S2], None,   None),
    ("D5 area  n=1 dlbls 0.0%",   XL_CHART_TYPE.AREA,             [S1],     None,   "0.0%"),
    ("D6 area  n=2 dlbls band",   XL_CHART_TYPE.AREA,             [S1, S2], "band", None),
    ("D7 area  n=2 dlbls bare",   XL_CHART_TYPE.AREA,             [S1, S2], "bare", None),
]

prs = Presentation()
prs.slide_width = Pt(720)
prs.slide_height = Pt(540)
for label, kind, series, legend, num_fmt in ARMS:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = CATS
    for i, vals in enumerate(series, 1):
        cd.add_series(f"Ser{i}", vals)
    chart = slide.shapes.add_chart(
        kind, Pt(72), Pt(72), Pt(396), Pt(288), cd
    ).chart
    if legend == "band":
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
    elif legend == "bare":
        chart.has_legend = True
    else:
        chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_value = True
    if num_fmt is not None:
        dl.number_format_is_linked = False
        dl.number_format = num_fmt

out = os.path.join(base, "chart_area_dlbls.pptx")
prs.save(out)
for i, (label, *_rest) in enumerate(ARMS, 1):
    print(f"  S{i}: {label}")
print("\nsaved:", out, os.path.getsize(out))
