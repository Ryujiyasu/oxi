# -*- coding: utf-8 -*-
"""Author the cloud-font resolution probe: which families does PowerPoint draw?

`%LOCALAPPDATA%\\Microsoft\\FontCache\\4\\CloudFonts` holds font files Office
downloads on demand. GDI never sees them, so Oxi substitutes for every run that
names one -- d06 asks for `IBM Plex Sans` in 2370 runs and PowerPoint's own PDF
names that face, while d19 asks for `Nunito`, which is ALSO in the cache, and
PowerPoint's PDF is Calibri throughout. Presence in the cache is therefore not
sufficient, and the difference has to be pinned before the cache is wired in.

One text box per candidate family, same string and size, so the exported PDF
names the face PowerPoint actually used for each. `Arial` and `Calibri` are
controls (installed); `Mali` is the negative control (named by d19, absent from
the cache and from the system).

Usage:
    python tools/metrics/gen_pptx_cloudfont.py
    python tools/metrics/export_pptx_cloudfont.py    # PowerPoint COM -> PDF
    python tools/metrics/read_pptx_cloudfont.py      # read the PDF back
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

REPO = Path(__file__).resolve().parents[2]
OUT = REPO / "pipeline_data" / "pptx_probes" / "cloudfont"

# (label, typeface) -- label is drawn too so the reader can pair span to arm
# even when the face falls back.
CASES = [
    ("arial", "Arial"),
    ("calibri", "Calibri"),
    ("plexsans", "IBM Plex Sans"),
    ("plexcond", "IBM Plex Sans Condensed"),
    ("nunito", "Nunito"),
    ("montserrat", "Montserrat"),
    ("lobster", "Lobster"),
    ("firasans", "Fira Sans"),
    ("mali", "Mali"),
]
SAMPLE = "Hamburgefonstiv 123"
SIZE = 24
BOX_X = Emu(457200)      # 0.50in
BOX_W = Emu(8229600)     # 9.00in
BOX_H = Emu(457200)      # 0.50in
TOP0 = Emu(457200)
PITCH = Emu(548640)      # 0.60in


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for index, (label, face) in enumerate(CASES):
        box = slide.shapes.add_textbox(BOX_X, Emu(int(TOP0) + index * int(PITCH)), BOX_W, BOX_H)
        para = box.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = f"{label}: {SAMPLE}"
        run.font.size = Pt(SIZE)
        run.font.name = face
    path = OUT / "probe_cloudfont.pptx"
    prs.save(str(path))
    print(f"wrote {path}  {len(CASES)} arms")


if __name__ == "__main__":
    main()
