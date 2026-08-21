# -*- coding: utf-8 -*-
"""Probe: what does PowerPoint do with b="1" on a weight-named family, and
where does an EMBEDDED face put a cell's first baseline?

Two open questions, one deck:

1. On the SlidesCarnival "Instructions" template PowerPoint draws every b="1"
   run of an embedded "Barlow Light" in the LIGHT face — although the deck
   embeds a bold part for it. Same on d24's "Fira Sans Light". Is that about
   the family NAME carrying a weight, about the embedded part, or neither?
2. The cell first-baseline model (1.2 x asc/(asc+desc) x size) is confirmed on
   installed faces but still regresses two decks whose cells use an EMBEDDED
   Open Sans. Which OS/2 reading matches for an embedded face?

Arms pair the same text at b=0 and b=1 in a weight-named family (Segoe UI
Light / Semibold / Black) and in a plain one (Segoe UI, Arial), as free text
AND as centred table cells whose rows are declared too short so the row height
is the known 2*margin + 1.2*size.

The deck is saved twice: once plain, once through PowerPoint COM with
EmbedTrueTypeFonts, so the same content can be compared with and without
embedding.

Usage:
    python tools/metrics/gen_pptx_embedbold.py
    python tools/metrics/embed_pptx_fonts.py pipeline_data/pptx_probes/embedbold/embedbold.pptx
    python tools/metrics/measure_pptx_word.py pipeline_data/pptx_probes/embedbold/embedbold_embedded.pptx pipeline_data/pptx_probes/embedbold/emb
    python tools/metrics/read_pptx_embedbold.py
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\embedbold").resolve()
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
PT = 12700
FACES = ["Segoe UI Light", "Segoe UI Semibold", "Segoe UI", "Arial"]
SIZE = 14.0
MARGIN = 7.2


def q(tag: str) -> str:
    return f"{{{A}}}{tag}"


def add_text(slide, face: str, bold: bool, y_pt: float):
    box = slide.shapes.add_textbox(
        Emu(int(36 * PT)), Emu(int(y_pt * PT)), Emu(int(300 * PT)), Emu(int(28 * PT))
    )
    tf = box.text_frame
    body = tf._txBody
    bodypr = body.find(q("bodyPr"))
    for ins in ("lIns", "tIns", "rIns", "bIns"):
        bodypr.set(ins, "0")
    for pel in body.findall(q("p")):
        body.remove(pel)
    p = etree.SubElement(body, q("p"))
    r = etree.SubElement(p, q("r"))
    rpr = etree.SubElement(r, q("rPr"))
    rpr.set("lang", "en-US")
    rpr.set("sz", str(int(SIZE * 100)))
    rpr.set("kern", "0")
    if bold:
        rpr.set("b", "1")
    etree.SubElement(rpr, q("latin")).set("typeface", face)
    etree.SubElement(r, q("t")).text = f"Hamburgefonstiv {'B' if bold else 'R'}"


def add_cell_table(slide, face: str, y_pt: float):
    shape = slide.shapes.add_table(
        3, 1, Emu(int(380 * PT)), Emu(int(y_pt * PT)), Emu(int(280 * PT)), Emu(3 * PT)
    )
    tbl = shape._element.graphic.graphicData.tbl
    for pr in tbl.findall(q("tblPr")):
        tbl.remove(pr)
    pr = etree.SubElement(tbl, q("tblPr"))
    pr.set("firstRow", "0")
    pr.set("bandRow", "0")
    etree.SubElement(pr, q("noFill"))
    tbl.insert(0, pr)
    for tr in tbl.findall(q("tr")):
        tr.set("h", str(1 * PT))
        for tc in tr.findall(q("tc")):
            txbody = tc.find(q("txBody"))
            for p in txbody.findall(q("p")):
                txbody.remove(p)
            p = etree.SubElement(txbody, q("p"))
            ppr = etree.SubElement(p, q("pPr"))
            ln = etree.SubElement(ppr, q("lnSpc"))
            etree.SubElement(ln, q("spcPct")).set("val", "100000")
            etree.SubElement(ppr, q("buNone"))
            r = etree.SubElement(p, q("r"))
            rpr = etree.SubElement(r, q("rPr"))
            rpr.set("lang", "en-US")
            rpr.set("sz", str(int(SIZE * 100)))
            rpr.set("kern", "0")
            etree.SubElement(rpr, q("latin")).set("typeface", face)
            etree.SubElement(r, q("t")).text = "Hxy"
            tcpr = tc.find(q("tcPr"))
            if tcpr is None:
                tcpr = etree.SubElement(tc, q("tcPr"))
            for name in ("marT", "marB", "marL", "marR"):
                tcpr.set(name, str(int(MARGIN * PT)))
            tcpr.set("anchor", "ctr")


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    manifest = []
    for face in FACES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cap = slide.shapes.add_textbox(Emu(228600), Emu(114300), Emu(6400800), Emu(300000))
        cap.text_frame.text = f"{face} @ {SIZE}pt"
        add_text(slide, face, False, 60)
        add_text(slide, face, True, 100)
        add_cell_table(slide, face, 60)
        manifest.append({
            "slide": len(manifest) + 1,
            "face": face,
            "size": SIZE,
            "margin": MARGIN,
            "table_top": 60.0,
            "rows": 3,
        })
        print(f"{face}")
    prs.save(OUT / "embedbold.pptx")
    (OUT / "embedbold_manifest.json").write_text(
        json.dumps(manifest, indent=1), encoding="utf-8")
    print(f"wrote {OUT / 'embedbold.pptx'} ({len(manifest)} slides)")


if __name__ == "__main__":
    main()
