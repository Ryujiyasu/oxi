# -*- coding: utf-8 -*-
"""Find the exact width at which PowerPoint breaks a given line.

The break test is a comparison, so the only way to read PowerPoint's own
measurement of a line is to squeeze the box until it gives: the narrowest width
that still holds the text on ONE line IS that measurement, to within the master
unit PowerPoint measures in (1/8 pt).

Two lines from dev deck d15 sit on that boundary and break the other way from
the engine:

    'our office'          Barlow Light 8pt, box 33.224pt -> PowerPoint 2 lines,
                          engine 1 (its master-unit sum is 265 = 33.125)
    '...template". You will'  Barlow Light 12pt, box 273.47pt -> PowerPoint fits
                          `will`, the engine does not

They fail in OPPOSITE directions, so it is not a scale error but a per-glyph
rounding difference. This sweeps the width in master units around each, so the
answer comes back as a number rather than an inference.

    python tools/metrics/gen_pptx_breakwidth.py
    python tools/metrics/read_pptx_breakwidth_com.py
"""
from __future__ import annotations

import json
import os

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

OUT = os.path.join("tools", "metrics", "breakwidth.pptx")
PLAN = os.path.join("tools", "metrics", "breakwidth.json")
MU = 0.125  # one master unit, in points

# Each arm is a string whose exact design width is known, swept in master units
# so PowerPoint's own measurement of it comes back as a number. The set is
# chosen to separate the candidate rounding rules: short and long, with and
# without an `ffi` ligature, at two sizes.
CASES = [
    ("our_office", "our office", "Barlow Light", 8.0, 33.25, 10),
    ("office8", "office", "Barlow Light", 8.0, 21.0, 10),
    ("our8", "our", "Barlow Light", 8.0, 11.5, 10),
    ("offices12", "offices", "Barlow Light", 12.0, 37.0, 10),
    ("noligature8", "our oxice", "Barlow Light", 8.0, 34.5, 10),
    ("control", "Click on the button under the presentation preview",
     "Barlow Light", 12.0, 261.0, 10),
    ("short12", "presentation preview", "Barlow Light", 12.0, 108.0, 10),
    ("mid12", "You have to be signed in to your Google account.",
     "Barlow Light", 12.0, 248.0, 12),
    ("caps8", "EDIT IN GOOGLE SLIDES", "Barlow Light", 8.0, 78.0, 10),
    ("digits12", "0123456789", "Barlow Light", 12.0, 68.0, 10),
]


def main() -> None:
    pres = Presentation()
    pres.slide_width = Emu(9144000)
    pres.slide_height = Emu(6858000)
    blank = pres.slide_layouts[6]
    plan = []
    for label, text, face, size, centre, span in CASES:
        widths = [round(centre + i * MU, 3) for i in range(-span, span + 1)]
        # One slide per case, one box per width, stacked down the slide.
        per_slide = 18
        for chunk in range(0, len(widths), per_slide):
            slide = pres.slides.add_slide(blank)
            for k, w in enumerate(widths[chunk:chunk + per_slide]):
                top = Emu(int((0.2 + k * 0.36) * 914400))
                box = slide.shapes.add_textbox(Emu(int(0.3 * 914400)), top,
                                               Emu(int(w * 12700)), Emu(int(0.3 * 914400)))
                tf = box.text_frame
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.NONE
                tf.margin_left = tf.margin_right = 0
                tf.margin_top = tf.margin_bottom = 0
                para = tf.paragraphs[0]
                para.alignment = PP_ALIGN.LEFT
                run = para.add_run()
                run.text = text
                run.font.name = face
                run.font.size = Pt(size)
                plan.append({"slide": len(pres.slides), "box": k, "label": label,
                             "width": w, "face": face, "size": size, "text": text})
    pres.save(OUT)
    with open(PLAN, "w", encoding="utf-8") as fh:
        json.dump(plan, fh, ensure_ascii=False, indent=1)
    print("wrote %s (%d slides) and %s (%d arms)"
          % (OUT, len(pres.slides), PLAN, len(plan)))


if __name__ == "__main__":
    main()
