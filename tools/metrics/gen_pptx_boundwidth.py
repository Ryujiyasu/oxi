# -*- coding: utf-8 -*-
"""What does `Lines(j).BoundWidth` actually measure?

The width half of `pptx_line_audit_com.py` finds that 47% of lines agree with
the engine's pen advance to under 0.05pt and 19% are 2pt or more out -- and the
19% carry a per-line CONSTANT of about 3pt, not an error that grows with the
line (deck 9 s2: +3.00pt on a 136pt line and +2.87pt on a 570pt one, same
shape). PowerPoint's own per-character steps sum to the engine's number
(`pptx_char_pos_com.py 9 9 Yellow`), so the 3pt is in the BOX, not the advances.

Neither the face, nor a trailing space, nor whether the paragraph wrapped
separates the two groups. So this sweeps the properties of the BOX with the
text held fixed: one word, one face, one size, one arm per property.

    base        wrap on, no autofit, left
    autofit     spAutoFit
    nowrap      wrap="none"
    centre      algn="ctr"
    right       algn="r"
    inset       non-zero lIns/rIns
    bigger      the same word at 28pt
    bold        b="1"
    italic      i="1"

Whichever arms carry the constant name it.

    python tools/metrics/gen_pptx_boundwidth.py
    python tools/metrics/read_pptx_boundwidth_com.py
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

REPO = Path(__file__).resolve().parents[2]
OUT = REPO / "pipeline_data" / "pptx_probes" / "boundwidth"
WORD = "Yellow"
FACE = "Arial"

ARMS = [
    ("base", {}),
    ("autofit", {"autofit": True}),
    ("nowrap", {"wrap": False}),
    ("centre", {"align": PP_ALIGN.CENTER}),
    ("right", {"align": PP_ALIGN.RIGHT}),
    ("inset", {"inset": True}),
    ("bigger", {"size": 28}),
    ("bold", {"bold": True}),
    ("italic", {"italic": True}),
    # ★The structural difference between this probe and the corpus lines that
    # would not explain themselves: those paragraphs are not alone in their
    # shape. d09 s9's 'Yellow' is the FIRST of two paragraphs, and its box
    # measures 3.4pt wider than the sum of PowerPoint's own character steps --
    # while every arm above measures narrower than the pen, as an ink box
    # should.
    ("twopara_1st", {"tail": "Is the color of gold, butter and ripe lemons."}),
    ("twopara_2nd", {"head": "Yellow"}),
    ("bold_twopara", {"bold": True, "tail": "Is the color of gold, butter."}),
    # Does the mark scale with the text? A space is 3.889pt at 14 and 7.778 at
    # 28, so a step that doubles says "space" and one that does not says
    # something else wearing a space's clothes at one size.
    ("twopara_28", {"size": 28, "tail": "Is the color of gold, butter."}),
    # ★The remaining width offenders are all last lines of BOLD paragraphs in
    # families whose bold PowerPoint has to synthesise (Rubik, Gruppo,
    # Montserrat Medium), so the obvious guess was that a synthesised bold
    # widens the advance and the engine keeps the regular one. These two arms
    # say NO: Bebas Neue has no bold and is in the cloud cache, and at 14pt
    #
    #     synth_base   design 33.63   box 33.76   engine 33.76
    #     synth_bold   design 33.63   box **33.01**   engine 32.96
    #
    # PowerPoint's synthesised bold measures NARROWER than the regular design,
    # not wider, and the engine follows it to 0.05pt. So the offenders' cause is
    # still open -- and it is not this.
    ("synth_base", {"face": "Bebas Neue"}),
    ("synth_bold", {"face": "Bebas Neue", "bold": True}),
]


def build(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i, (label, opt) in enumerate(ARMS):
        col, row = i % 3, i // 3
        left = Emu(457200 + col * 2800000)
        top = Emu(685800 + row * 1600000)
        tag = slide.shapes.add_textbox(left, Emu(top.emu - 300000), Emu(2400000), Emu(250000))
        tag.text_frame.paragraphs[0].add_run().text = label
        tag.text_frame.paragraphs[0].runs[0].font.size = Pt(9)

        box = slide.shapes.add_textbox(left, top, Emu(2400000), Emu(900000))
        tf = box.text_frame
        tf.word_wrap = opt.get("wrap", True)
        if opt.get("autofit"):
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        else:
            tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.TOP
        if opt.get("inset"):
            tf.margin_left = Pt(18)
            tf.margin_right = Pt(18)
        # A paragraph before the measured one, when the arm asks for it.
        if opt.get("head"):
            first = tf.paragraphs[0]
            r0 = first.add_run()
            r0.text = opt["head"]
            r0.font.size = Pt(opt.get("size", 14))
            r0.font.name = FACE
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        if opt.get("align"):
            p.alignment = opt["align"]
        run = p.add_run()
        run.text = WORD
        run.font.size = Pt(opt.get("size", 14))
        run.font.name = opt.get("face", FACE)
        run.font.bold = bool(opt.get("bold"))
        run.font.italic = bool(opt.get("italic"))
        # A paragraph after it, likewise.
        if opt.get("tail"):
            nxt = tf.add_paragraph()
            r1 = nxt.add_run()
            r1.text = opt["tail"]
            r1.font.size = Pt(opt.get("size", 14))
            r1.font.name = FACE
    prs.save(str(path))


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    deck = OUT / "boundwidth.pptx"
    build(deck)
    print(f"wrote {deck}  ({WORD!r} in {FACE})")
    for label, opt in ARMS:
        print(f"  {label:<10} {opt}")


if __name__ == "__main__":
    main()
