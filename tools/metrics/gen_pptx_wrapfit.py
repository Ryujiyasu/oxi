# -*- coding: utf-8 -*-
"""Probe: what width sum does PowerPoint's line-break test actually use?

Four dev-corpus lines break although their float design-advance sum fits with
0.01-0.10pt to spare, while others hold at 0.002pt slack -- and the trailing
GDI-ABC-ink term that would explain the breaks regresses the corpus (S-TRAILINK).
This deck asks PowerPoint directly: one two-word string per arm, box width
swept around the string's design advance sum, insets zeroed, kerning off. The
exported PDF's band count per slide locates the exact threshold W*, and
W* - design_sum is the term PowerPoint adds for that ending glyph.

The Arial pair is the scalpel: 'f' and '.' have the SAME advance (0.2778 em)
but opposite ink (f reaches +0.0347 em past its advance, '.' stops 0.087 em
short), so any threshold difference between "west warf" and "west war." is a
pure trailing-ink term. Segoe Script 't' (+0.2197 em ink) is the loud version:
if PowerPoint counted full trailing ink its threshold would sit ~8.8pt above
the design sum at 40pt.

Usage:
    python tools/metrics/gen_pptx_wrapfit.py
    python tools/metrics/measure_pptx_word.py pipeline_data/pptx_probes/wrapfit/wrapfit.pptx pipeline_data/pptx_probes/wrapfit
    python tools/metrics/read_pptx_wrapfit.py
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from fontTools.ttLib import TTFont
from lxml import etree
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\wrapfit").resolve()
A = "http://schemas.openxmlformats.org/drawingml/2006/main"

FONT_FILES = {
    "Arial": r"C:\Windows\Fonts\arial.ttf",
    "Segoe Script": r"C:\Windows\Fonts\segoesc.ttf",
}

# (label, font, size pt, text, delta_lo, delta_hi, step) -- deltas in pt
# around the design sum. Windows sized so a full-trailing-ink threshold
# still lands inside the swept range.
ARMS = [
    ("arial_f",   "Arial",        40.0, "west warf", -1.5,  3.0, 0.15),
    ("arial_dot", "Arial",        40.0, "west war.", -1.5,  1.5, 0.15),
    ("segsc_t",   "Segoe Script", 40.0, "meno mint", -1.5,  9.9, 0.30),
    ("segsc_dot", "Segoe Script", 40.0, "meno min.", -1.5,  1.5, 0.15),
]


def q(tag: str) -> str:
    return f"{{{A}}}{tag}"


def design_sum_pt(font_path: str, text: str, fs: float) -> float:
    f = TTFont(font_path, lazy=True)
    upm = f["head"].unitsPerEm
    cmap = f.getBestCmap()
    hmtx = f["hmtx"]
    units = sum(hmtx[cmap[ord(c)]][0] for c in text)
    return units / upm * fs


def add_arm_slide(prs, label: str, font: str, fs: float, text: str, width_pt: float):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(
        Emu(457200), Emu(914400), Emu(int(round(width_pt * 12700))),
        Emu(int(round(fs * 1.2 * 3 * 12700))),
    )
    tf = box.text_frame
    tf.word_wrap = True
    body = tf._txBody
    bodypr = body.find(q("bodyPr"))
    for ins in ("lIns", "tIns", "rIns", "bIns"):
        bodypr.set(ins, "0")
    for pel in body.findall(q("p")):
        body.remove(pel)
    p = etree.SubElement(body, q("p"))
    r = etree.SubElement(p, q("r"))
    rpr = etree.SubElement(r, q("rPr"))
    rpr.set("lang", "en-US")
    rpr.set("sz", str(int(round(fs * 100))))
    rpr.set("kern", "0")
    etree.SubElement(rpr, q("latin")).set("typeface", font)
    etree.SubElement(r, q("t")).text = text


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    manifest = []
    for label, font, fs, text, lo, hi, step in ARMS:
        s = design_sum_pt(FONT_FILES[font], text, fs)
        n = int(round((hi - lo) / step)) + 1
        for k in range(n):
            delta = lo + k * step
            w = s + delta
            add_arm_slide(prs, label, font, fs, text, w)
            manifest.append({
                "slide": len(manifest) + 1,
                "arm": label,
                "font": font,
                "fs": fs,
                "text": text,
                "design_sum_pt": round(s, 4),
                "delta_pt": round(delta, 4),
                "width_pt": round(w, 4),
            })
        print(f"{label}: design sum {s:.3f}pt, {n} slides")
    prs.save(OUT / "wrapfit.pptx")
    (OUT / "wrapfit_manifest.json").write_text(
        json.dumps(manifest, indent=1), encoding="utf-8")
    print(f"wrote {OUT / 'wrapfit.pptx'} ({len(manifest)} slides)")


if __name__ == "__main__":
    main()
