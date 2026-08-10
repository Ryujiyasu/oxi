# -*- coding: utf-8 -*-
"""Chart data-labels probe: 5 slides of a clustered/stacked column chart with
different data-label configurations, so Word's PDF exposes the exact data
label placement/format/font rules for measurement.

  S1: clustered column, show_value=True        (default OUTSIDE_END position)
  S2: clustered column, show_value=True + number_format="0.0%"  (percent)
  S3: clustered column, show_value=True + position=CENTER
  S4: clustered column, show_value=True + position=INSIDE_END
  S5: stacked column,    show_value=True       (stacked segment labels)

All on blank slides (layout 6), default Office theme, frame 72,72,396,288,
categories [East, West, Midwest], one series (19.2, 21.4, 16.7)."""
import sys, os
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION

base = r"pipeline_data\pptx_probes\chart_datalabel"
os.makedirs(base, exist_ok=True)

prs = Presentation()
blanks = [prs.slide_layouts[6] for _ in range(5)]


def add_chart_slide(layout, ctype, show_val, num_fmt=None, pos=None):
    slide = prs.slides.add_slide(layout)
    cd = CategoryChartData()
    cd.categories = ["East", "West", "Midwest"]
    cd.add_series("Series 1", (19.2, 21.4, 16.7))
    x, y, cx, cy = Inches(1.0), Inches(1.0), Inches(5.5), Inches(4.0)
    gframe = slide.shapes.add_chart(ctype, x, y, cx, cy, cd)
    plot = gframe.chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_value = True
    if num_fmt is not None:
        dl.number_format_is_linked = False
        dl.number_format = num_fmt
    if pos is not None:
        dl.position = pos
    return slide


add_chart_slide(blanks[0], XL_CHART_TYPE.COLUMN_CLUSTERED, True)
add_chart_slide(blanks[1], XL_CHART_TYPE.COLUMN_CLUSTERED, True, num_fmt="0.0%")
add_chart_slide(blanks[2], XL_CHART_TYPE.COLUMN_CLUSTERED, True, pos=XL_LABEL_POSITION.CENTER)
add_chart_slide(blanks[3], XL_CHART_TYPE.COLUMN_CLUSTERED, True, pos=XL_LABEL_POSITION.INSIDE_END)
add_chart_slide(blanks[4], XL_CHART_TYPE.COLUMN_STACKED, True)

out = os.path.join(base, "chart_datalabel.pptx")
prs.save(out)
print("saved:", out, os.path.getsize(out))
