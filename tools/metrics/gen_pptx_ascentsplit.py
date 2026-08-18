# -*- coding: utf-8 -*-
"""Probe: WHICH font metric splits the 1.2 line height into ascent and descent?

`mixedpitch` established that the step between two lines of different size is
`d * s1 + a * s2` with `a + d = 1.2`, and that Arial / Georgia / Calibri /
Verdana each land on their own `1.2 * winAscent / (winAscent + winDescent)`.
Those four agree with almost every other candidate metric, so they cannot say
WHICH one it is. d28's embedded Calistoga can: its win metrics are 1300/400 of
a 1000 em (1.7em tall) while its typo and hhea metrics are 1000/-300, and
PowerPoint's own PDF puts a at 0.9400 -- neither the win ratio (0.9176) nor the
typo ratio (0.9231).

These faces are the installed fonts whose win and typo ratios disagree most, so
each formula predicts a visibly different `a` for them. Arial is the control.
"""
from __future__ import annotations

import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\ascentsplit").resolve()
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
FONTS = [
    "Arial",              # control: win 0.8103 / typo 0.7758
    "Goudy Stout",        # win 0.7448 / typo 0.9757, win box 1.368 em
    "Castellar",          # win 0.7688 / typo 0.9946
    "Stencil",            # win 0.7664 / typo 1.0000
    "Maiandra GD",        # win 0.8033 / typo 1.4624
    "Lucida Handwriting", # win 0.7468 / typo 1.6133
    "Haettenschweiler",   # win 0.8845 / typo 1.2071
    # The seven above all have OS/2 fsSelection bit 7 (USE_TYPO_METRICS) CLEAR.
    # These have it SET, and their win and typo ratios disagree, so they say
    # whether PowerPoint honours the bit -- which is what d28's embedded
    # Calistoga (bit set, win 0.9176, measured 0.9377) is asking.
    "Cambria Math",       # bit set: win 0.6702 / typo 0.9334
    "Noto Serif",         # bit set: win 0.8798 / typo 0.9419
    "Noto Sans",          # bit set: win 0.8880 / typo 0.9419
    "Reem Kufi",          # bit set: win 0.9333 / typo 0.8800 (the other way round)
    "Liberation Sans Narrow",  # bit set: win 0.9770 / typo 0.9309
]
# A big size gap makes `a` precise: the PDF quantises to 0.03pt, so a 40pt gap
# resolves a to better than 0.001.
PAIRS = [(20, 60), (60, 20), (20, 20)]


def q(tag: str) -> str:
    return f"{{{A}}}{tag}"


def add_para(body, text: str, pt: float, font: str) -> None:
    p = etree.SubElement(body, q("p"))
    ppr = etree.SubElement(p, q("pPr"))
    ln = etree.SubElement(ppr, q("lnSpc"))
    etree.SubElement(ln, q("spcPct")).set("val", "100000")
    for name, val in (("spcBef", "0"), ("spcAft", "0")):
        el = etree.SubElement(ppr, q(name))
        etree.SubElement(el, q("spcPts")).set("val", val)
    r = etree.SubElement(p, q("r"))
    rpr = etree.SubElement(r, q("rPr"))
    rpr.set("lang", "en-US")
    rpr.set("sz", str(int(pt * 100)))
    etree.SubElement(rpr, q("latin")).set("typeface", font)
    etree.SubElement(r, q("t")).text = text


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for font in FONTS:
        for s1, s2 in PAIRS:
            slide = prs.slides.add_slide(blank)
            cap = slide.shapes.add_textbox(Emu(228600), Emu(114300), Emu(6400800), Emu(300000))
            cap.text_frame.text = f"{font} {s1}->{s2}"
            box = slide.shapes.add_textbox(Emu(457200), Emu(1400000), Emu(7500000), Emu(3000000))
            body = box.text_frame._txBody
            for pel in body.findall(q("p")):
                body.remove(pel)
            add_para(body, "AAA", s1, font)
            add_para(body, "BBB", s2, font)
    prs.save(OUT / "ascentsplit.pptx")
    print(f"wrote {OUT / 'ascentsplit.pptx'}  ({len(FONTS)}x{len(PAIRS)} arms)")


if __name__ == "__main__":
    main()
