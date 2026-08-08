# -*- coding: utf-8 -*-
"""Chart spec: the LEGEND LABEL WRAP WIDTH.

A single very long word cannot wrap at a space, so Word force-breaks it at
the last character that fits -> the widest resulting line IS the wrap cap.
Sweep the frame WIDTH (and one HEIGHT arm) to see what the cap scales with.

All arms: doughnut, legend right, include_in_layout=False, frame at 72,72."""
import sys, os

sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

base = r"pipeline_data\pptx_probes\chart_legendwrap"
os.makedirs(base, exist_ok=True)

LONG = "W" * 30  # ~ 30 * 12.6pt Calibri18 = 380pt, far past any cap
ARMS = [
    (200.0, 288.0),
    (250.0, 288.0),
    (320.0, 288.0),
    (396.0, 288.0),
    (500.0, 288.0),
    (600.0, 288.0),
    (396.0, 180.0),  # height arm: does the cap follow the frame height?
    (396.0, 400.0),
]

prs = Presentation()
prs.slide_width = Pt(720)
prs.slide_height = Pt(540)
for sw, shh in ARMS:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = [LONG, "Ef", "Gh"]
    cd.add_series("Series 1", (19.2, 21.4, 16.7))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, Pt(72), Pt(72), Pt(sw), Pt(shh), cd
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False

out = os.path.join(base, "chart_legendwrap.pptx")
prs.save(out)
for i, (sw, shh) in enumerate(ARMS, 1):
    print(f"  S{i}: frame {sw:6.1f} x {shh:6.1f}   right edge = {72+sw:.1f}")
print("\nsaved:", out, os.path.getsize(out))
