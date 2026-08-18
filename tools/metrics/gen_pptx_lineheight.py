# -*- coding: utf-8 -*-
"""Probe: what line advance does PowerPoint use for single (100%) spacing?

Oxi models it as `fs * 1.2`, derived on Calibri -- whose typo metrics happen to
give 1.2207. Arial does not agree: d24's 60pt title steps 64.82pt, a ratio of
1.0800, and its typo metrics give 1.0884 while its hhea metrics give 1.1499.
So the advance is font-dependent and the constant is wrong for everything that
is not Calibri-shaped.

Each arm is one installed font, three lines of plain text at 40pt with an
explicit `lnSpc 100%`, in a box wide enough to avoid wrapping surprises. The
reader compares the measured baseline pitch with the font's own metrics.
"""
from __future__ import annotations

import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\lineheight").resolve()
FONTS = ["Arial", "Calibri", "Times New Roman", "Verdana", "Georgia",
         "Tahoma", "Segoe UI", "Trebuchet MS"]
SIZE = 40


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    tmp, dst = OUT / "_stage.pptx", OUT / "lineheight.pptx"
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for fam in FONTS:
        s = prs.slides.add_slide(blank)
        cap = s.shapes.add_textbox(Emu(228600), Emu(114300), Emu(6400800), Emu(300000))
        cap.text_frame.text = fam
        box = s.shapes.add_textbox(Emu(457200), Emu(914400), Emu(7772400), Emu(4114800))
        tf = box.text_frame
        tf.word_wrap = False
        for i, line in enumerate(["Handgloves one", "Handgloves two", "Handgloves three"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run()
            r.text = line
            r.font.size = Pt(SIZE)
            r.font.name = fam
    prs.save(tmp)

    with zipfile.ZipFile(tmp) as zin:
        names = zin.namelist()
        data = {n: zin.read(n) for n in names}
    # Pin every paragraph to single spacing and zero space-before/after so the
    # only thing under test is the font's own line advance.
    for i in range(1, len(FONTS) + 1):
        part = f"ppt/slides/slide{i}.xml"
        xml = data[part].decode("utf-8")
        xml = xml.replace(
            "<a:pPr>",
            '<a:pPr><a:lnSpc><a:spcPct val="100000"/></a:lnSpc>'
            '<a:spcBef><a:spcPts val="0"/></a:spcBef>'
            '<a:spcAft><a:spcPts val="0"/></a:spcAft>',
        )
        xml = xml.replace(
            "<a:p>",
            '<a:p><a:pPr><a:lnSpc><a:spcPct val="100000"/></a:lnSpc>'
            '<a:spcBef><a:spcPts val="0"/></a:spcBef>'
            '<a:spcAft><a:spcPts val="0"/></a:spcAft></a:pPr>',
        )
        data[part] = xml.encode("utf-8")
    with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:
        for n in names:
            zout.writestr(n, data[n])
    tmp.unlink()
    print(f"wrote {dst}  ({len(FONTS)} arms)")


if __name__ == "__main__":
    main()
