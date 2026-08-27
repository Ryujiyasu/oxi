# -*- coding: utf-8 -*-
"""Author the a:srcRect / geometry PICTURE-CROP probe.

d44 s23 is the corpus's worst real defect after the MIN re-rank (0.8200, heavy
10.22%). Its portrait declares

    <a:srcRect b="5384" l="790" r="16864" t="1782"/>  <a:stretch/>
    prstGeom chord, square shape, 800x800 source

Oxi applies all four edges and stretches the surviving 659x742 region into the
square. PowerPoint shows a WIDER view. A search over which subset of the edges
reproduces PowerPoint favoured "top and bottom only" (mean|err| 21.91 against
52.63 for all four) -- but the best residual was far too high to call a law, so
this asks PowerPoint directly instead.

The source is a COORDINATE GRID, not a photograph: 10x10 cells, each a distinct
colour, with the (row, col) recoverable from the colour alone. Reading back which
colours land where gives the visible source rectangle exactly, with no
registration or masking guesswork -- which is what made the pixel search
inconclusive.

Arms isolate one thing at a time: each edge alone, all four together, the same
crop under rect / ellipse / chord, and a NON-SQUARE shape (where "preserve the
aspect" and "stretch to fill" must disagree).

Usage:
    python tools/metrics/gen_pptx_srcrect.py
    python tools/metrics/export_pptx_srcrect.py   # PowerPoint COM -> PDF
    python tools/metrics/read_pptx_srcrect.py     # read the PDF back
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from lxml import etree
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

REPO = Path(__file__).resolve().parents[2]
OUT = REPO / "pipeline_data" / "pptx_probes" / "srcrect"
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
EMU_IN = 914400

N = 10          # grid cells per side
CELL = 80       # px per cell -> 800x800 source

# (name, prst, srcRect dict or None, shape w_in, shape h_in)
CROP = {"l": 10000, "t": 10000, "r": 20000, "b": 5000}   # deliberately asymmetric
ARMS = [
    ("rect_none",      "rect",    None,                     3.0, 3.0),
    ("rect_l10",       "rect",    {"l": 10000},             3.0, 3.0),
    ("rect_r10",       "rect",    {"r": 10000},             3.0, 3.0),
    ("rect_t10",       "rect",    {"t": 10000},             3.0, 3.0),
    ("rect_b10",       "rect",    {"b": 10000},             3.0, 3.0),
    ("rect_all",       "rect",    CROP,                     3.0, 3.0),
    ("ellipse_all",    "ellipse", CROP,                     3.0, 3.0),
    ("chord_all",      "chord",   CROP,                     3.0, 3.0),
    ("ellipse_none",   "ellipse", None,                     3.0, 3.0),
    ("rect_all_wide",  "rect",    CROP,                     5.0, 2.5),
    ("rect_all_tall",  "rect",    CROP,                     2.5, 5.0),
]


def cell_colour(r: int, c: int) -> tuple[int, int, int]:
    """A colour that encodes (r, c) unambiguously and survives JPEG-ish noise."""
    return (20 + r * 23, 20 + c * 23, 128)


def make_source(path: Path) -> None:
    img = Image.new("RGB", (N * CELL, N * CELL), (255, 255, 255))
    d = ImageDraw.Draw(img)
    for r in range(N):
        for c in range(N):
            d.rectangle(
                [c * CELL, r * CELL, (c + 1) * CELL - 1, (r + 1) * CELL - 1],
                fill=cell_colour(r, c),
                outline=(255, 255, 255),
                width=2,
            )
    img.save(path)


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    grid = OUT / "grid.png"
    make_source(grid)

    prs = Presentation()
    blank = prs.slide_layouts[6]
    manifest = []
    for i, (name, prst, crop, w_in, h_in) in enumerate(ARMS):
        slide = prs.slides.add_slide(blank)
        pic = slide.shapes.add_picture(
            str(grid), Emu(int(0.5 * EMU_IN)), Emu(int(0.5 * EMU_IN)),
            Emu(int(w_in * EMU_IN)), Emu(int(h_in * EMU_IN)),
        )
        el = pic._element
        blip_fill = el.find(f"{{{A}}}blipFill") or el.find(
            "{http://schemas.openxmlformats.org/presentationml/2006/main}blipFill"
        )
        if crop:
            sr = etree.SubElement(blip_fill, f"{{{A}}}srcRect")
            for k, v in crop.items():
                sr.set(k, str(v))
            # srcRect must precede a:stretch
            st = blip_fill.find(f"{{{A}}}stretch")
            if st is not None:
                blip_fill.remove(st)
                blip_fill.append(st)
        sp_pr = el.find("{http://schemas.openxmlformats.org/presentationml/2006/main}spPr")
        geom = sp_pr.find(f"{{{A}}}prstGeom")
        if geom is not None:
            geom.set("prst", prst)
        manifest.append(
            {
                "slide": i + 1, "name": name, "prst": prst, "srcRect": crop,
                "x_in": 0.5, "y_in": 0.5, "w_in": w_in, "h_in": h_in,
                "grid": N, "cell_px": CELL,
            }
        )
    path = OUT / "probe_srcrect.pptx"
    prs.save(str(path))
    (OUT / "manifest.json").write_text(json.dumps(manifest, indent=1), encoding="utf-8")
    print(f"wrote {path} ({len(ARMS)} arms), source grid {N}x{N} at {grid}")
    for m in manifest:
        print(f"  s{m['slide']:>2} {m['name']:<16} prst={m['prst']:<8} "
              f"srcRect={m['srcRect']} shape={m['w_in']}x{m['h_in']}in")


if __name__ == "__main__":
    main()
