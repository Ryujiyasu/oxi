# -*- coding: utf-8 -*-
"""HORIZONTAL bar residual probe (chart_bar ship-note items 1 and 2).

chart_bar shipped with two values that were NOT measured:
  (1) plot_top when an explicit <c:title> is present -- the vertical chart's
      45.69 minus 5.0 was adopted by analogy.
  (2) data-label placement for a STACKED horizontal bar -- only the clustered
      OUTSIDE_END rule was measured.

  S1: BAR_CLUSTERED 1 series + explicit title            (plot_top, auto-title case)
  S2: BAR_CLUSTERED 2 series + explicit title            (plot_top, no auto title)
  S3: BAR_STACKED   2 series + data labels               (stacked label placement)
  S4: BAR_STACKED   2 series + data labels + legend      (stacked labels + legend)
  S5: BAR_CLUSTERED 1 series + explicit title + labels   (title & labels together)

Frame 72,72,396,288, categories [East, West, Midwest], default Office theme."""
import sys, os

sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

base = r"pipeline_data\pptx_probes\chart_bar_resid"
os.makedirs(base, exist_ok=True)

prs = Presentation()


def add(ctype, series, title=None, legend=False, dlbls=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = ["East", "West", "Midwest"]
    for name, vals in series:
        cd.add_series(name, vals)
    gframe = slide.shapes.add_chart(
        ctype, Pt(72), Pt(72), Pt(396), Pt(288), cd
    )
    chart = gframe.chart
    if title is not None:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    if legend:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
    if dlbls:
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.show_value = True


ONE = [("Series 1", (19.2, 21.4, 16.7))]
TWO = [("Revenue", (19.2, 21.4, 16.7)), ("Cost", (10.5, 15.0, 12.3))]

add(XL_CHART_TYPE.BAR_CLUSTERED, ONE, title="Quarterly Revenue")
add(XL_CHART_TYPE.BAR_CLUSTERED, TWO, title="Quarterly Revenue")
add(XL_CHART_TYPE.BAR_STACKED, TWO, dlbls=True)
add(XL_CHART_TYPE.BAR_STACKED, TWO, legend=True, dlbls=True)
add(XL_CHART_TYPE.BAR_CLUSTERED, ONE, title="Quarterly Revenue", dlbls=True)

out = os.path.join(base, "chart_bar_resid.pptx")
prs.save(out)
print("saved:", out, os.path.getsize(out))
