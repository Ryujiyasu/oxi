# gen_pptx_spec6_baseline.py — probe: first-line baseline offset from text-area top vs font size
# Word baseline origin sweep across font sizes (Arial 10..28pt), single 1-line paragraph, Left.
# Derives the em constant Word uses for the first baseline (hhea asc+gap? win? typo? something else).
# Usage: python gen_pptx_spec6_baseline.py
import os
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.enum.text import PP_ALIGN

OUT_DIR = r"C:\Users\ryuji\oxi-main\pipeline_data\pptx_probes\spec6_baseline"
os.makedirs(OUT_DIR, exist_ok=True)

# One slide per font size, so the PDF baseline of each "FS<n>" label is unambiguous.
SIZES = [10, 12, 14, 18, 24, 28]
prs = Presentation()
prs.slide_width = Emu(9144000)   # 10in
prs.slide_height = Emu(6858000)  # 7.5in

for i, fs in enumerate(SIZES):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Emu(914400)   # 0.1in = 7.2pt
    top = Emu(914400)
    width = Inches(5)
    height = Inches(2)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(914400)    # 0.1in = 7.2pt
    tf.margin_right = Emu(914400)
    tf.margin_top = Emu(457200)     # 0.05in = 3.6pt
    tf.margin_bottom = Emu(457200)
    tf.vertical_anchor = 1          # MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "FS{0} baseline".format(fs)
    r.font.size = Pt(fs)
    r.font.name = "Arial"

out = os.path.join(OUT_DIR, "spec6_baseline.pptx")
prs.save(out)
print("saved", out)
