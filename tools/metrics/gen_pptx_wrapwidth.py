# -*- coding: utf-8 -*-
"""Probe: what width does each line of an indented paragraph wrap against?

Oxi wraps every line of a paragraph against one width -- the shape's inner
width -- and only afterwards shifts line 0 by the bullet/indent geometry. So a
hanging-indent line 0, drawn 18pt further left, is allowed to run 18pt past the
right inset. The `bulletph` probe caught it once S-ADVWIDTH stopped
over-measuring: PowerPoint breaks its line 0 after "to" and Oxi now fits
"second" as well.

The arms vary marL / indent / bullet with text long enough to wrap three times.
Reading each line's first and last pen position out of PowerPoint's PDF says
whether the lines share a right EDGE (width = edge - start, different per line)
or a WIDTH (same for all, which is what Oxi assumes).
"""
from __future__ import annotations

import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\wrapwidth").resolve()
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
# Short equal-ish words make the break granularity fine, so the right edge is
# pinned to within one word's width.
TEXT = ("alpha beta gamma delta epsilon zeta eta theta iota kappa lambda mu "
        "nu xi omicron pi rho sigma tau upsilon phi chi psi omega alpha beta "
        "gamma delta epsilon zeta eta theta iota kappa lambda mu nu xi")
# (label, marL pt, indent pt, bullet)
ARMS = [
    ("plain", 0, 0, None),
    ("hang18", 18, -18, "•"),
    ("hang36", 36, -36, "•"),
    ("hang18_nobu", 18, -18, None),
    ("firstind18", 18, 18, None),
    ("marL36_ind0", 36, 0, None),
    ("marL36_hang18", 36, -18, "•"),
]


def q(tag: str) -> str:
    return f"{{{A}}}{tag}"


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for label, mar_l, indent, bullet in ARMS:
        slide = prs.slides.add_slide(blank)
        cap = slide.shapes.add_textbox(Emu(228600), Emu(114300), Emu(6400800), Emu(300000))
        cap.text_frame.text = f"{label} marL={mar_l} ind={indent} bu={'yes' if bullet else 'no'}"
        # a deliberately narrow box so the indent is a big fraction of the width
        box = slide.shapes.add_textbox(Emu(914400), Emu(1200000), Emu(3200400), Emu(3000000))
        tf = box.text_frame
        tf.word_wrap = True
        body = tf._txBody
        for pel in body.findall(q("p")):
            body.remove(pel)
        p = etree.SubElement(body, q("p"))
        ppr = etree.SubElement(p, q("pPr"))
        ppr.set("marL", str(int(mar_l * 12700)))
        ppr.set("indent", str(int(indent * 12700)))
        ln = etree.SubElement(ppr, q("lnSpc"))
        etree.SubElement(ln, q("spcPct")).set("val", "100000")
        if bullet:
            etree.SubElement(ppr, q("buFont")).set("typeface", "Arial")
            etree.SubElement(ppr, q("buChar")).set("char", bullet)
        else:
            etree.SubElement(ppr, q("buNone"))
        r = etree.SubElement(p, q("r"))
        rpr = etree.SubElement(r, q("rPr"))
        rpr.set("lang", "en-US")
        rpr.set("sz", "1400")
        etree.SubElement(rpr, q("latin")).set("typeface", "Arial")
        etree.SubElement(r, q("t")).text = TEXT
    prs.save(OUT / "wrapwidth.pptx")
    print(f"wrote {OUT / 'wrapwidth.pptx'}  ({len(ARMS)} arms)")


if __name__ == "__main__":
    main()
