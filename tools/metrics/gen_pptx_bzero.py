# -*- coding: utf-8 -*-
"""Does `b="0"` on a run turn OFF a bold that the LEVEL turned on?

`SlideRun.bold` is a plain `bool`, so "explicitly not bold" and "said nothing"
reach the renderer as the same value and it resolves them with
`run.bold || default_bold`. The 2026-08-31 audit counted 3675 candidate runs
and found no case that was actually drawn wrong, so it declined to change the
IR on speculation and asked for an example.

d15 slide 11 and d11 slide 11 are that example -- one shape, one level, two
paragraphs, and PowerPoint draws `<a:rPr b="0"/>` in Barlow and the run beside
it that says nothing in Barlow,Bold. This isolates the same question so the
answer does not depend on those two decks sharing a template.

    arm  level bold  run says      prediction
     1        yes    b="0"         upright   <- the question
     1        yes    (nothing)     bold      <- its control, same shape
     2        yes    b="1"         bold      <- the attribute is read at all
     3        no     b="0"         upright   <- the level is what supplies bold
     4        yes    two runs in ONE paragraph, b="0" then nothing

Arial is used because this machine has a real Arial Bold, so the truth PDF
names the two faces apart (`ArialMT` against `Arial-BoldMT`) instead of leaving
the answer to a synthesised weight.

    python tools/metrics/gen_pptx_bzero.py
    python tools/metrics/export_pptx_bzero.py
    python tools/metrics/read_pptx_bzero.py
"""
import os
import sys

from lxml import etree
from pptx import Presentation
from pptx.util import Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

BASE = os.path.join("pipeline_data", "pptx_probes", "bzero")
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
SIZE_PT = 28.0
BOX = (72.0, 72.0, 500.0, 320.0)

# (arm name, level is bold, [(text, bold attribute or None), ...] per paragraph)
ARMS = [
    ("level-bold, run b=0 beside a silent run", True,
     [[("ZERO ZERO ZERO", False)], [("SILENT SILENT", None)]]),
    ("level-bold, run b=1 beside a silent run", True,
     [[("ONE ONE ONE", True)], [("SILENT SILENT", None)]]),
    ("level not bold, run b=0 beside a silent run", False,
     [[("ZERO ZERO ZERO", False)], [("SILENT SILENT", None)]]),
    ("level-bold, both runs in ONE paragraph", True,
     [[("ZERO ZERO ", False), ("SILENT SILENT", None)]]),
]


def main() -> None:
    os.makedirs(BASE, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(540)
    blank = prs.slide_layouts[6]

    for _, level_bold, paras in ARMS:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Pt(BOX[0]), Pt(BOX[1]), Pt(BOX[2]), Pt(BOX[3]))
        tf = box.text_frame
        tf.word_wrap = True
        body = tf._txBody.find(f"{{{NS_A}}}bodyPr")
        for attr, value in (("lIns", "0"), ("rIns", "0"), ("tIns", "0"),
                            ("bIns", "0"), ("anchor", "t"), ("wrap", "square")):
            body.set(attr, value)

        # The shape's OWN list style is the level here: no master, no layout,
        # nothing outside the shape can be the reason an arm differs.
        lst = tf._txBody.find(f"{{{NS_A}}}lstStyle")
        if lst is None:
            # It has to sit between bodyPr and the paragraphs, so it is
            # inserted, never appended.
            lst = etree.Element(f"{{{NS_A}}}lstStyle")
            body.addnext(lst)
        lvl = etree.SubElement(lst, f"{{{NS_A}}}lvl1pPr")
        dpr = etree.SubElement(lvl, f"{{{NS_A}}}defRPr")
        dpr.set("sz", str(int(SIZE_PT * 100)))
        dpr.set("b", "1" if level_bold else "0")
        latin = etree.SubElement(dpr, f"{{{NS_A}}}latin")
        latin.set("typeface", "Arial")

        for i, runs in enumerate(paras):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            for text, bold in runs:
                run = para.add_run()
                run.text = text
                run.font.size = Pt(SIZE_PT)
                run.font.name = "Arial"
                # None leaves the attribute out; False writes b="0".
                run.font.bold = bold

    out = os.path.join(BASE, "bzero.pptx")
    prs.save(out)
    print(f"wrote {out} with {len(ARMS)} arms")
    for i, (name, _, _) in enumerate(ARMS, 1):
        print(f"   slide {i}: {name}")


if __name__ == "__main__":
    main()
