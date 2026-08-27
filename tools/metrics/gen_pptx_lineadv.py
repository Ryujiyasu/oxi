# -*- coding: utf-8 -*-
"""Author the SUB-100% LINE-ADVANCE probe.

Oxi advances every line by `fs * 1.2 * n` (`n` = `a:lnSpc/a:spcPct`). The
`firstline` probe (2026-08-23) derived where the FIRST baseline sits for
`n` = 0.70 .. 1.20 across six faces and matched within 0.072pt -- but it only
ever measured ONE baseline, so the line-to-line advance below 100% was assumed
linear rather than measured.

d40 s1 is the specimen and the corpus's MIN slide for that deck: 216.15pt
**Grand Hotel** (embedded), `<a:lnSpc><a:spcPct val="60000"/>`, two lines. Its
second line sits 8.16pt HIGH, i.e. PowerPoint's advance is 9.60pt LONGER than
`1.2 * fs * n`:

    Oxi  1.2   * 216.15 * 0.60 = 155.63
    PPT  1.2739 * 216.15 * 0.60 = 165.23   (measured, ink-top to ink-top)

1.2739 em is close to a script face's OWN natural line height, which suggests
the multiplied quantity is `max(1.2, face_natural)` rather than a flat 1.2 --
consistent with `firstline`, whose six faces all sit BELOW 1.2 (Arial 1.1172,
Georgia 1.1362, Times 1.1074) and so could not tell the two apart.

So the faces here deliberately STRADDLE 1.2:

    Arial          1.1172      Verdana        1.2158
    Times New Rom  1.1074      Comic Sans MS  1.3936
    Georgia        1.1362      Segoe Script   (deepest installed)

If the advance is a flat 1.2, every face at a given `n` advances identically.
If it is the face's own natural height with a 1.2 floor, the right-hand column
advances further and the left-hand one does not move at all.

Three lines per arm, so the 2nd->3rd advance is measured clear of the first
baseline rule. `noAutofit` + `wrap="none"` keep the box from resizing.

Usage:
    python tools/metrics/gen_pptx_lineadv.py
    python tools/metrics/export_pptx_lineadv.py   # PowerPoint COM -> PDF
    python tools/metrics/read_pptx_lineadv.py     # read the PDF back
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

REPO = Path(__file__).resolve().parents[2]
OUT = REPO / "pipeline_data" / "pptx_probes" / "lineadv"
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
P = "http://schemas.openxmlformats.org/presentationml/2006/main"
EMU_PT = 12700

FACES = ["Arial", "Times New Roman", "Georgia", "Verdana", "Comic Sans MS", "Segoe Script"]
PCTS = [60000, 80000, 100000]
SZ = 40.0
LINES = ["Hxg", "Hxg", "Hxg"]

ARMS = [(f"{f.replace(' ', '')}_{p // 1000}", f, SZ, p) for f in FACES for p in PCTS]
# one large-size arm: the rule must not depend on the point size
ARMS.append(("ComicSansMS_60_sz216", "Comic Sans MS", 216.15, 60000))
ARMS.append(("Arial_60_sz216", "Arial", 216.15, 60000))

# The same question inside a TABLE CELL. Oxi answers it differently there --
# `max(1.2, 1.0625 * pct)`, carried over from the d10 s12 ROW-HEIGHT derivation
# -- so a cell at lnSpc 140% advances 29.75pt where d03 s15 and d33 s9 both show
# PowerPoint at exactly 33.60 = 1.2 * 1.40 * 20.
TABLE_ARMS = [
    (f"tbl_{f.replace(' ', '')}_{p // 1000}", f, 20.0, p)
    for f in ("Arial", "Comic Sans MS")
    for p in (60000, 140000)
]


def cell_arm(slide, name: str, face: str, sz: float, pct: int, idx: int) -> dict:
    """One 2x1 table whose single cell holds the same three paragraphs."""
    shp = slide.shapes.add_table(1, 2, Emu(int(36 * EMU_PT)), Emu(int(36 * EMU_PT)),
                                 Emu(int(500 * EMU_PT)), Emu(int(20 * EMU_PT)))
    tbl = shp._element.graphic.graphicData.tbl
    tr = tbl.find(f"{{{A}}}tr")
    tr.set("h", str(int(4 * EMU_PT)))       # too small, so content wins
    for tc in tr.findall(f"{{{A}}}tc"):
        tx = tc.find(f"{{{A}}}txBody")
        for para in tx.findall(f"{{{A}}}p"):
            tx.remove(para)
        for line in LINES:
            p = etree.SubElement(tx, f"{{{A}}}p")
            ppr = etree.SubElement(p, f"{{{A}}}pPr")
            ppr.set("algn", "ctr")
            ls = etree.SubElement(ppr, f"{{{A}}}lnSpc")
            etree.SubElement(ls, f"{{{A}}}spcPct").set("val", str(pct))
            for tag in ("spcBef", "spcAft"):
                e = etree.SubElement(ppr, f"{{{A}}}{tag}")
                etree.SubElement(e, f"{{{A}}}spcPts").set("val", "0")
            etree.SubElement(ppr, f"{{{A}}}buNone")
            r = etree.SubElement(p, f"{{{A}}}r")
            rpr = etree.SubElement(r, f"{{{A}}}rPr")
            rpr.set("lang", "en-US")
            rpr.set("sz", str(int(round(sz * 100))))
            etree.SubElement(rpr, f"{{{A}}}latin").set("typeface", face)
            etree.SubElement(r, f"{{{A}}}t").text = line
        pr = tc.find(f"{{{A}}}tcPr")
        if pr is None:
            pr = etree.SubElement(tc, f"{{{A}}}tcPr")
        for k, v in (("marT", "57150"), ("marB", "57150"),
                     ("marL", "0"), ("marR", "0"), ("anchor", "ctr")):
            pr.set(k, v)
    return {"slide": idx, "name": name, "typeface": face, "sz_pt": sz,
            "lnSpc_pct": pct, "n_lines": len(LINES), "in_table": True}


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    manifest = []
    for i, (name, face, sz, pct) in enumerate(ARMS + TABLE_ARMS):
        slide = prs.slides.add_slide(blank)
        if name.startswith("tbl_"):
            manifest.append(cell_arm(slide, name, face, sz, pct, i + 1))
            continue
        box = slide.shapes.add_textbox(
            Emu(int(36 * EMU_PT)), Emu(int(36 * EMU_PT)),
            Emu(int(600 * EMU_PT)), Emu(int(460 * EMU_PT)),
        )
        tx = box._element.find(f"{{{P}}}txBody")
        body = tx.find(f"{{{A}}}bodyPr")
        for k, v in (("lIns", "0"), ("rIns", "0"), ("tIns", "0"), ("bIns", "0"),
                     ("anchor", "t"), ("wrap", "none")):
            body.set(k, v)
        for fit in ("spAutoFit", "normAutofit", "noAutofit"):
            for e in body.findall(f"{{{A}}}{fit}"):
                body.remove(e)
        etree.SubElement(body, f"{{{A}}}noAutofit")
        for para in tx.findall(f"{{{A}}}p"):
            tx.remove(para)
        for line in LINES:
            p = etree.SubElement(tx, f"{{{A}}}p")
            ppr = etree.SubElement(p, f"{{{A}}}pPr")
            ppr.set("algn", "l")
            ppr.set("marL", "0")
            ppr.set("indent", "0")
            ls = etree.SubElement(ppr, f"{{{A}}}lnSpc")
            etree.SubElement(ls, f"{{{A}}}spcPct").set("val", str(pct))
            for tag in ("spcBef", "spcAft"):
                e = etree.SubElement(ppr, f"{{{A}}}{tag}")
                etree.SubElement(e, f"{{{A}}}spcPts").set("val", "0")
            etree.SubElement(ppr, f"{{{A}}}buNone")
            r = etree.SubElement(p, f"{{{A}}}r")
            rpr = etree.SubElement(r, f"{{{A}}}rPr")
            rpr.set("lang", "en-US")
            rpr.set("sz", str(int(round(sz * 100))))
            etree.SubElement(rpr, f"{{{A}}}latin").set("typeface", face)
            t = etree.SubElement(r, f"{{{A}}}t")
            t.text = line
        manifest.append({"slide": i + 1, "name": name, "typeface": face,
                         "sz_pt": sz, "lnSpc_pct": pct, "n_lines": len(LINES)})
    path = OUT / "probe_lineadv.pptx"
    prs.save(str(path))
    (OUT / "manifest.json").write_text(json.dumps(manifest, indent=1), encoding="utf-8")
    print(f"wrote {path} ({len(ARMS)} arms, {len(LINES)} lines each)")
    for m in manifest:
        print(f"  s{m['slide']:>2} {m['name']:<24} {m['typeface']:<16} "
              f"sz={m['sz_pt']:<7} lnSpc={m['lnSpc_pct'] / 1000:.0f}%")


if __name__ == "__main__":
    main()
