"""buAutoNum fourth measurement round — 16 supported schemes + startAt transitions.

Slide E1/E2: 16 latin/arabic autonumber schemes that PowerPoint accepts
             (the 4 `*Plain` schemes crash PowerPoint's file open and are
             excluded), 2 paragraphs each (formats + continuation).
Slide F: startAt value change [5][3][none].
Slide G: none -> startAt -> none  [none][5][none].
"""
import sys

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from lxml import etree

A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def add_textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(Inches(left / 72.0), Inches(top / 72.0),
                                  Inches(width / 72.0), Inches(height / 72.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(91440)    # 0.1 in (correct)
    tf.margin_right = Emu(91440)
    tf.margin_top = Emu(45720)
    tf.margin_bottom = Emu(45720)
    return tb, tf


def set_autonum(paragraph, kind, mar_l_pt, indent_pt, lvl=0, start_at=None,
                bufont="Arial"):
    p = paragraph._p
    pPr = p.get_or_add_pPr()
    pPr.set("lvl", str(lvl))
    pPr.set("marL", str(int(mar_l_pt * 12700)))
    pPr.set("indent", str(int(indent_pt * 12700)))
    for tag in ("buNone", "buChar", "buAutoNum", "buFont", "buClr", "buSzPct"):
        for el in pPr.findall(f"{{{A}}}{tag}"):
            pPr.remove(el)
    bf = etree.SubElement(pPr, f"{{{A}}}buFont")
    bf.set("typeface", bufont)
    bn = etree.SubElement(pPr, f"{{{A}}}buAutoNum")
    bn.set("type", kind)
    if start_at is not None:
        bn.set("startAt", str(start_at))


def add_para(tf, text, first, font=None, size=18):
    if first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    if font is not None:
        r.font.name = font
    return p


def fill(prs, paras, box_h=470, top=40):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb, tf = add_textbox(s, 72, top, 500, box_h)
    first = True
    for spec in paras:
        text = spec["text"]
        p = add_para(tf, text, first, spec.get("font"), spec.get("size", 18))
        first = False
        if "kind" in spec:
            set_autonum(p, spec["kind"], spec.get("mar", 72),
                        spec.get("ind", -36), spec.get("lvl", 0),
                        spec.get("start_at"))
    return s


# PowerPoint crashes opening a file that uses any `*Plain` scheme; only these
# 16 are supported (verified empirically).
SCHEMES16 = [
    "arabicPeriod", "arabicParenR", "arabicParenBoth", "arabicPlain",
    "romanUcPeriod", "romanLcPeriod", "romanUcParenR", "romanLcParenR",
    "romanUcParenBoth", "romanLcParenBoth",
    "alphaUcPeriod", "alphaLcPeriod", "alphaUcParenR", "alphaLcParenR",
    "alphaUcParenBoth", "alphaLcParenBoth",
]

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

paras_e1 = []
paras_e2 = []
for i, kind in enumerate(SCHEMES16):
    for n in (1, 2):
        spec = {"text": f"{kind}-{n}", "kind": kind, "mar": 72, "ind": -36}
        (paras_e1 if i < 8 else paras_e2).append(spec)
fill(prs, paras_e1)
fill(prs, paras_e2)

# Slide F: startAt value change
fill(prs, [
    {"text": "F1", "kind": "arabicPeriod", "mar": 72, "ind": -36,
     "start_at": 5},
    {"text": "F2", "kind": "arabicPeriod", "mar": 72, "ind": -36,
     "start_at": 3},
    {"text": "F3", "kind": "arabicPeriod", "mar": 72, "ind": -36},
])

# Slide G: none -> startAt -> none
fill(prs, [
    {"text": "G1", "kind": "arabicPeriod", "mar": 72, "ind": -36},
    {"text": "G2", "kind": "arabicPeriod", "mar": 72, "ind": -36,
     "start_at": 5},
    {"text": "G3", "kind": "arabicPeriod", "mar": 72, "ind": -36},
])

prs.save(sys.argv[1] if len(sys.argv) > 1 else "autonum4.pptx")
print("saved")
