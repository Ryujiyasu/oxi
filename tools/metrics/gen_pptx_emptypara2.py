# -*- coding: utf-8 -*-
"""Probe: the height of an empty paragraph, measured as a difference.

The first emptypara probe pinned the case that matters most -- an empty
paragraph carrying `<a:endParaRPr sz="..."/>` advances by sz x 1.2 x lnSpc,
exactly, and beats a run rPr on the same (textless) run. What it could not
answer was the case with NO endParaRPr, because each arm's first baseline sits
at a different place (a 24pt first line pushes it down 13pt) and that offset
contaminated the subtraction.

So each question here is asked twice: once with the empty paragraph and once
without it, with everything else identical. The AAA->BBB baseline gap differs
by exactly the empty paragraph's advance, and the first-line placement cancels.
"""
from __future__ import annotations

import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\emptypara2").resolve()
A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# (label, first pt, last pt, endParaRPr sz or None, lnSpc pct on the empty para)
QUESTIONS = [
    ("ctl_10_10_epr1000", 10.0, 10.0, 1000, 100),   # control: known answer 12.0
    ("prev24_next10_none", 24.0, 10.0, None, 100),
    ("prev10_next24_none", 10.0, 24.0, None, 100),
    ("prev32_next10_none", 32.0, 10.0, None, 100),
    ("prev10_next10_none", 10.0, 10.0, None, 100),
    ("prev24_next10_epr1000", 24.0, 10.0, 1000, 100),
]


def q(tag: str) -> str:
    return f"{{{A}}}{tag}"


def set_lnspc(p, pct: int) -> None:
    ppr = p.find(q("pPr"))
    if ppr is None:
        ppr = etree.Element(q("pPr"))
        p.insert(0, ppr)
    ln = etree.SubElement(ppr, q("lnSpc"))
    etree.SubElement(ln, q("spcPct")).set("val", str(pct * 1000))
    ppr.insert(0, ln)


def build(slide, label: str, first_pt: float, last_pt: float,
          epr: int | None, ls: int, with_empty: bool) -> None:
    cap = slide.shapes.add_textbox(Emu(228600), Emu(114300), Emu(6400800), Emu(300000))
    cap.text_frame.text = f"{label} {'WITH' if with_empty else 'WITHOUT'}"
    box = slide.shapes.add_textbox(Emu(457200), Emu(914400), Emu(6000000), Emu(3000000))
    tf = box.text_frame
    tf.word_wrap = True
    body = tf._txBody
    for pel in body.findall(q("p")):
        body.remove(pel)

    def text_para(text: str, pt: float) -> None:
        p = etree.SubElement(body, q("p"))
        r = etree.SubElement(p, q("r"))
        rpr = etree.SubElement(r, q("rPr"))
        rpr.set("lang", "en-US")
        rpr.set("sz", str(int(pt * 100)))
        etree.SubElement(r, q("t")).text = text
        set_lnspc(p, 100)

    text_para("AAA", first_pt)
    if with_empty:
        p = etree.SubElement(body, q("p"))
        if epr is not None:
            e = etree.SubElement(p, q("endParaRPr"))
            e.set("lang", "en-US")
            e.set("sz", str(epr))
        set_lnspc(p, ls)
    text_para("BBB", last_pt)


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for label, fp, lp, epr, ls in QUESTIONS:
        for with_empty in (True, False):
            build(prs.slides.add_slide(blank), label, fp, lp, epr, ls, with_empty)
    prs.save(OUT / "emptypara2.pptx")
    print(f"wrote {OUT / 'emptypara2.pptx'}  ({len(QUESTIONS)} questions x 2)")


if __name__ == "__main__":
    main()
