# -*- coding: utf-8 -*-
"""Build tests/fixtures/highlight_test.pptx.

Three runs in one paragraph: no highlight, a highlight whose colour element is
self-closing, and one whose colour element has a child. quick-xml routes the
last two to different events, and the colour a highlight carries is the same
element `a:solidFill` uses -- so a parser that does not know it is inside
`a:highlight` reads it as the run's TEXT colour, which is what Oxi did to d11
slide 38's white-on-dk1 "and many more...".
"""
from __future__ import annotations

import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path("tests/fixtures/highlight_test.pptx").resolve()
A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def q(tag: str) -> str:
    return f"{{{A}}}{tag}"


def main() -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Emu(914400), Emu(914400), Emu(5486400), Emu(914400))
    body = box.text_frame._txBody
    for pel in body.findall(q("p")):
        body.remove(pel)
    p = etree.SubElement(body, q("p"))

    def add(text, color, highlight, highlight_child):
        r = etree.SubElement(p, q("r"))
        rpr = etree.SubElement(r, q("rPr"))
        rpr.set("lang", "en-US")
        rpr.set("sz", "1800")
        fill = etree.SubElement(rpr, q("solidFill"))
        etree.SubElement(fill, q("srgbClr")).set("val", color)
        if highlight:
            hl = etree.SubElement(rpr, q("highlight"))
            clr = etree.SubElement(hl, q("srgbClr"))
            clr.set("val", highlight)
            if highlight_child:
                # A child makes the colour element Event::Start, not Empty.
                etree.SubElement(clr, q("lumMod")).set("val", "100000")
        etree.SubElement(r, q("t")).text = text

    add("plain ", "112233", None, False)
    add("empty ", "FFFFFF", "FF0000", False)
    add("start", "000000", "00FF00", True)
    # `<a:br>` is a soft line break INSIDE the paragraph. Every one of the 70 in
    # the dev corpus carries an a:rPr, so it arrives as Event::Start.
    br = etree.SubElement(p, q("br"))
    rpr = etree.SubElement(br, q("rPr"))
    rpr.set("lang", "en-US")
    rpr.set("sz", "1800")
    add("after the break", "445566", None, False)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT}")


if __name__ == "__main__":
    main()
