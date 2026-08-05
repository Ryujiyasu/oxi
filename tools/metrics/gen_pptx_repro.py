# -*- coding: utf-8 -*-
"""Generate controlled .pptx repros for the OxiSlides Ra loop (the pptx analog
of the docx _pb_* probe generators).

The FIRST wave targets the most foundational specs:
  1. slide size (EMU -> pt conversion)
  2. shape position/size (xfrm off/ext -> pt, rotation)
  3. placeholder (title/body) presence
  4. table dimensions
  5. text runs (font size / bold / italic / color / family)
"""
import os
import sys

sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


def gen_geometry(out_path):
    """Slides covering slide-size + shape geometry."""
    prs = Presentation()  # default 4:3 (10 x 7.5 in)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]  # blank

    # Slide 1: positioned autoshapes at known points.
    s = prs.slides.add_slide(blank)
    shapes = [
        (Inches(1.0), Inches(1.0), Inches(3.0), Inches(1.0), MSO_SHAPE.RECTANGLE),
        (Inches(5.0), Inches(2.0), Inches(2.0), Inches(2.0), MSO_SHAPE.OVAL),
        (Inches(9.0), Inches(0.5), Inches(3.0), Inches(0.75), MSO_SHAPE.ROUNDED_RECTANGLE),
        (Inches(2.0), Inches(4.0), Inches(5.0), Inches(2.5), MSO_SHAPE.CHEVRON),
    ]
    for i, (x, y, w, h, kind) in enumerate(shapes):
        sh = s.shapes.add_shape(kind, x, y, w, h)
        sh.name = "shape%d" % i

    # Rotated rectangle (45 deg).
    rot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(5.5), Inches(2.0), Inches(1.0))
    rot.name = "rotated45"
    rot.rotation = 45.0

    # Slide 2: title + body placeholders (default layout) + a text box.
    s2 = prs.slides.add_slide(prs.slide_layouts[0])
    s2.shapes.title.text = "Title Text"
    s2.placeholders[1].text = "Body line one\nBody line two"

    tb = s2.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(6.0), Inches(2.0))
    tf = tb.text_frame
    tf.text = "First paragraph"
    p = tf.add_paragraph()
    p.text = "Second paragraph"
    p.alignment = PP_ALIGN.CENTER

    # Slide 3: table.
    s3 = prs.slides.add_slide(blank)
    rows, cols = 3, 4
    tbl_shape = s3.shapes.add_table(rows, cols, Inches(1.0), Inches(1.0), Inches(8.0), Inches(2.0))
    tbl = tbl_shape.table
    for r in range(rows):
        for c in range(cols):
            tbl.cell(r, c).text = "R%dC%d" % (r, c)

    # Slide 4: formatted text runs.
    s4 = prs.slides.add_slide(blank)
    tb4 = s4.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(8.0), Inches(3.0))
    tf4 = tb4.text_frame
    tf4.word_wrap = True
    tf4.text = "Plain 18pt "
    r1 = tf4.paragraphs[0].runs[0]
    r1.font.size = Pt(18)
    p2 = tf4.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Bold red 24pt"
    r2.font.size = Pt(24)
    r2.font.bold = True
    r2.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    p3 = tf4.add_paragraph()
    r3 = p3.add_run()
    r3.text = "Italic 12pt"
    r3.font.size = Pt(12)
    r3.font.italic = True

    prs.save(out_path)


def gen_slide_sizes(out_dir):
    """One FILE per common slide size, to pin EMU->pt conversion.

    PowerPoint's PageSetup.SlideWidth/Height are single per-presentation
    values, so a multi-size test must be split across files (setting
    prs.slide_width mid-file overwrites the whole presentation).
    """
    for w_in, h_in in [(10.0, 7.5), (13.333, 7.5), (12.0, 6.75)]:
        prs = Presentation()
        prs.slide_width = Inches(w_in)
        prs.slide_height = Inches(h_in)
        # The layout must come from THIS presentation (a layout from another
        # Presentation object duplicates zip parts -> corrupted file).
        layout = prs.slide_layouts[6]
        s = prs.slides.add_slide(layout)
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        sh.name = "fullsize"
        out_path = os.path.join(out_dir, "slide_size_%dx%d.pptx" % (w_in * 72, h_in * 72))
        prs.save(out_path)
        print("wrote", out_path)


if __name__ == "__main__":
    out_dir = sys.argv[1] if len(sys.argv) > 1 else r"c:\Users\ryuji\oxi-main\pipeline_data\pptx_probes"
    os.makedirs(out_dir, exist_ok=True)
    gen_geometry(os.path.join(out_dir, "geometry.pptx"))
    gen_slide_sizes(out_dir)
