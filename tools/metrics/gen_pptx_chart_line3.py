# -*- coding: utf-8 -*-
"""3-series line chart legend repro (see gen_pptx_chart_line2.py)."""
import sys, os
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

base = r"pipeline_data\pptx_probes\chart_line3"
os.makedirs(base, exist_ok=True)

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

chart_data = CategoryChartData()
chart_data.categories = ["East", "West", "Midwest", "North", "South"]
chart_data.add_series("Series 1", (19.2, 21.4, 16.7, 22.0, 18.5))
chart_data.add_series("Series 2", (15.3, 17.5, 20.1, 13.9, 16.2))
chart_data.add_series("Series 3", (10.8, 12.4, 9.7, 14.1, 11.3))

x, y, cx, cy = Inches(1.0), Inches(1.0), Inches(5.5), Inches(4.0)
gframe = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data)

out = os.path.join(base, "chart_line3.pptx")
prs.save(out)
print("saved:", out, os.path.getsize(out))
