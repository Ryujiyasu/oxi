# -*- coding: utf-8 -*-
"""Probe: how does PowerPoint rotate a picture / a picture-filled shape?

The renderer's image path says in a comment "rotation=0 only for now", and the
dev corpus carries 1395 shapes with `a:xfrm/@rot != 0` on 362 slides. d28
slide 3 (the floor deck's worst slide, 0.4027) is exactly this: a 116.6-degree
rotated torn-paper image that Oxi paints axis-aligned, so it covers the wrong
part of the portrait underneath.

The rotation rule for a shape OUTLINE is already derived (rotate about the box
centre, verified to 0.01pt on connectors and presets). What is NOT measured is
what happens to the raster inside, and whether `rotWithShape` changes it.

The source is a 2x2 colour grid (TL red, TR green, BL blue, BR yellow), so the
colour at a sampled point names which source corner landed there and the
orientation can be read off rather than inferred.

Arms (4in box, centred; all sampled at the same points):
  E1 p:pic, rot=0                     baseline
  E2 p:pic, rot=90                    a right angle: unambiguous corner map
  E3 p:pic, rot=30
  E4 shape blipFill, rot=90, rotWithShape=1
  E5 shape blipFill, rot=90, rotWithShape=0    does the raster stay upright?
  E6 p:pic, rot=90 + flipH            order of flip vs rotation
  E7 p:pic, flipH only                the mirror on its own (141 corpus shapes)
  E8 shape blipFill, rot=30 + flipH   the same composition off a right angle
"""
from __future__ import annotations

import io
import sys
import zipfile
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\img_rotation").resolve()
SPACE = 1000
SHAPE_EMU = 3657600
SHAPE_XY = (2743200, 1600200)


def grid_png() -> bytes:
    img = Image.new("RGB", (400, 400), "white")
    d = ImageDraw.Draw(img)
    d.rectangle([0, 0, 199, 199], fill=(255, 0, 0))
    d.rectangle([200, 0, 399, 199], fill=(0, 200, 0))
    d.rectangle([0, 200, 199, 399], fill=(0, 0, 255))
    d.rectangle([200, 200, 399, 399], fill=(255, 220, 0))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


RECT_GEOM = (
    "<a:custGeom><a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>"
    '<a:rect b="b" l="l" r="r" t="t"/><a:pathLst>'
    f'<a:path extrusionOk="0" h="{SPACE}" w="{SPACE}">'
    '<a:moveTo><a:pt x="0" y="0"/></a:moveTo>'
    f'<a:lnTo><a:pt x="{SPACE}" y="0"/></a:lnTo>'
    f'<a:lnTo><a:pt x="{SPACE}" y="{SPACE}"/></a:lnTo>'
    f'<a:lnTo><a:pt x="0" y="{SPACE}"/></a:lnTo>'
    "<a:close/></a:path></a:pathLst></a:custGeom>"
)

# (label, kind, rot_deg, flip_h, rot_with_shape)
ARMS = [
    ("E1 pic rot=0", "pic", 0, False, None),
    ("E2 pic rot=90", "pic", 90, False, None),
    ("E3 pic rot=30", "pic", 30, False, None),
    ("E4 shape fill rot=90 rotWithShape=1", "shape", 90, False, 1),
    ("E5 shape fill rot=90 rotWithShape=0", "shape", 90, False, 0),
    ("E6 pic rot=90 flipH", "pic", 90, True, None),
    ("E7 pic flipH only", "pic", 0, True, None),
    ("E8 shape fill rot=30 flipH", "shape", 30, True, 1),
]


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    tmp, dst = OUT / "_stage.pptx", OUT / "img_rotation.pptx"
    png = grid_png()
    (OUT / "source.png").write_bytes(png)

    prs = Presentation()
    blank = prs.slide_layouts[6]
    for label, kind, *_ in ARMS:
        s = prs.slides.add_slide(blank)
        box = s.shapes.add_textbox(Emu(228600), Emu(228600), Emu(6400800), Emu(400050))
        box.text_frame.text = label
        s.shapes.add_picture(io.BytesIO(png), Emu(SHAPE_XY[0]), Emu(SHAPE_XY[1]),
                             Emu(SHAPE_EMU), Emu(SHAPE_EMU))
        if kind == "shape":
            sp = s.shapes.add_shape(1, Emu(SHAPE_XY[0]), Emu(SHAPE_XY[1]),
                                    Emu(SHAPE_EMU), Emu(SHAPE_EMU))
            sp.line.fill.background()
            sp.shadow.inherit = False
    prs.save(tmp)

    with zipfile.ZipFile(tmp) as zin:
        names = zin.namelist()
        data = {n: zin.read(n) for n in names}

    for i, (label, kind, rot, flip, rws) in enumerate(ARMS, start=1):
        part = f"ppt/slides/slide{i}.xml"
        xml = data[part].decode("utf-8")
        rels = data[f"ppt/slides/_rels/slide{i}.xml.rels"].decode("utf-8")
        rid = next(t.split('"')[0] for t in rels.split('Id="')[1:] if "/image" in t)
        attrs = ""
        if rot:
            attrs += f' rot="{int(round(rot * 60000))}"'
        if flip:
            attrs += ' flipH="1"'
        if kind == "pic":
            # The picture's own xfrm carries the rotation.
            head = xml.index("<p:pic>")
            seg = xml[head:]
            seg = seg.replace("<a:xfrm>", f"<a:xfrm{attrs}>", 1)
            xml = xml[:head] + seg
        else:
            # Delete the helper picture; the SHAPE carries geometry + blipFill.
            ps = xml.index("<p:pic>")
            pe = xml.index("</p:pic>") + len("</p:pic>")
            xml = xml[:ps] + xml[pe:]
            gs = xml.rindex("<a:prstGeom")
            ge = xml.rindex("</a:prstGeom>") + len("</a:prstGeom>")
            fill = (
                f'<a:blipFill rotWithShape="{rws}"><a:blip r:embed="{rid}">'
                "<a:alphaModFix/></a:blip><a:stretch><a:fillRect/></a:stretch></a:blipFill>"
            )
            xml = xml[:gs] + RECT_GEOM + fill + xml[ge:]
            tail = xml.rindex("<a:xfrm")
            xml = xml[:tail] + xml[tail:].replace("<a:xfrm>", f"<a:xfrm{attrs}>", 1)
        data[part] = xml.encode("utf-8")

    with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:
        for n in names:
            zout.writestr(n, data[n])
    tmp.unlink()
    print(f"wrote {dst}  ({len(ARMS)} arms)")


if __name__ == "__main__":
    main()
