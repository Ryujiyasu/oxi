# -*- coding: utf-8 -*-
"""Line chart spec repro: minimal LINE_MARKERS chart on a blank slide,
default theme. Single series over 5 categories so the line/point geometry
is unambiguous for Word-PDF measurement (plot area, axes, line points,
markers, gridlines)."""
import sys, os
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

base = r"pipeline_data\pptx_probes\chart_line"
os.makedirs(base, exist_ok=True)

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

chart_data = CategoryChartData()
chart_data.categories = ["East", "West", "Midwest", "North", "South"]
chart_data.add_series("Series 1", (19.2, 21.4, 16.7, 22.0, 18.5))

x, y, cx, cy = Inches(1.0), Inches(1.0), Inches(5.5), Inches(4.0)
gframe = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data)

out = os.path.join(base, "chart_line.pptx")
prs.save(out)
print("saved:", out, os.path.getsize(out))
