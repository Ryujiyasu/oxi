# -*- coding: utf-8 -*-
"""Build the SCATTER (XY) chart probe.

Scatter is structurally different from every chart type shipped so far:
its series carry <c:xVal>/<c:yVal> (numeric on BOTH axes) instead of
<c:cat>/<c:val>, so the category-band geometry that bar/line/area share
does not apply.

The arms keep every other lever fixed (frame, data, theme) and vary one
thing each so the rules can be read off a single Word export:

  S1  markers only, 1 series           (baseline geometry + auto title)
  S2  markers only, 2 series           (series colours, no auto title)
  S3  straight lines + markers, 1 ser  (does a line get drawn?)
  S4  smooth lines + markers, 1 ser    (smooth vs straight)
  S5  lines only (no markers), 2 ser   (marker suppression)
  S6  markers, 2 series + legend       (legend band / entry order)
  S7  markers, 1 series + data labels  (label placement)
  S8  markers, 1 series, no title      (plot_top without a title)
"""
import os
from pptx import Presentation
from pptx.chart.data import XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt

OUT_DIR = r"pipeline_data\pptx_probes\chart_scatter"
OUT = os.path.join(OUT_DIR, "chart_scatter.pptx")

X = [1.0, 2.0, 3.0, 4.0]
S1 = [19.2, 21.4, 16.7, 22.0]
S2 = [10.5, 11.2, 8.5, 12.3]

ARMS = [
    ("S1 markers 1ser", XL_CHART_TYPE.XY_SCATTER, 1, False, False, False),
    ("S2 markers 2ser", XL_CHART_TYPE.XY_SCATTER, 2, False, False, False),
    ("S3 lines+mk 1ser", XL_CHART_TYPE.XY_SCATTER_LINES, 1, False, False, False),
    ("S4 smooth+mk 1ser", XL_CHART_TYPE.XY_SCATTER_SMOOTH, 1, False, False, False),
    ("S5 lines only 2ser", XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS, 2, False, False, False),
    ("S6 markers 2ser legend", XL_CHART_TYPE.XY_SCATTER, 2, True, False, False),
    ("S7 markers 1ser dlbls", XL_CHART_TYPE.XY_SCATTER, 1, False, True, False),
    ("S8 markers 1ser notitle", XL_CHART_TYPE.XY_SCATTER, 1, False, False, True),
]


def main() -> None:
    os.makedirs(OUT_DIR, exist_ok=True)
    prs = Presentation()
    blank = prs.slide_layouts[6]

    for name, ctype, nser, legend, dlbls, notitle in ARMS:
        slide = prs.slides.add_slide(blank)
        cd = XyChartData()
        s = cd.add_series("Ser1")
        for x, y in zip(X, S1):
            s.add_data_point(x, y)
        if nser >= 2:
            s2 = cd.add_series("Ser2")
            for x, y in zip(X, S2):
                s2.add_data_point(x, y)

        gf = slide.shapes.add_chart(
            ctype, Pt(72), Pt(72), Pt(396), Pt(288), cd
        )
        chart = gf.chart
        chart.has_legend = legend
        if legend:
            chart.legend.include_in_layout = False
        # python-pptx cannot write <c:dLbls> on a scatterChart
        # (CT_ScatterChart has no dLbls property), so S7 gets it by
        # rewriting the chart part below -- the same technique the
        # doughnut probe uses for holeSize / bare legend.
        if notitle:
            chart.has_title = False
        print(f"  {name}: type={ctype} nser={nser} legend={legend} "
              f"dlbls={dlbls} notitle={notitle}")

    prs.save(OUT)
    print("saved", OUT)
    _inject_dlbls(OUT, "ppt/charts/chart7.xml")


DLBLS = (
    '<c:dLbls><c:showLegendKey val="0"/><c:showVal val="1"/>'
    '<c:showCatName val="0"/><c:showSerName val="0"/>'
    '<c:showPercent val="0"/><c:showBubbleSize val="0"/></c:dLbls>'
)


def _inject_dlbls(path: str, part: str) -> None:
    """Insert <c:dLbls> after the last </c:ser> of a scatterChart.

    CT_ScatterChart's child order is scatterStyle, varyColors, ser*,
    dLbls, axId* -- so the label block has to land after every series
    and before the first axId."""
    import shutil, zipfile

    tmp = path + ".tmp"
    with zipfile.ZipFile(path) as zin, zipfile.ZipFile(
        tmp, "w", zipfile.ZIP_DEFLATED
    ) as zout:
        for item in zin.infolist():
            blob = zin.read(item.filename)
            if item.filename == part:
                xml = blob.decode("utf-8")
                if "<c:dLbls>" in xml:
                    print(f"  {part}: already has dLbls")
                else:
                    i = xml.rindex("</c:ser>") + len("</c:ser>")
                    xml = xml[:i] + DLBLS + xml[i:]
                    print(f"  {part}: injected <c:dLbls> showVal=1")
                blob = xml.encode("utf-8")
            zout.writestr(item, blob)
    shutil.move(tmp, path)


if __name__ == "__main__":
    main()
