# -*- coding: utf-8 -*-
"""Generate the Spec #4 (text frame layout) controlled repros.

Wave 1 targets the text-frame foundations the current renderer lacks:
  - multi-paragraph line advance (currently 1 paragraph = 1 line, fs*1.2)
  - line spacing (lnSpc: 100/150/200%)
  - paragraph spacing (spcBef / spcAft)
  - word wrap (word_wrap on/off)

The actual line pitch / wrap positions are NOT readable via COM (PowerPoint
has no Information(6) analog); they are measured from PowerPoint's own PDF
export (deck.pdf) with fitz baselines. COM supplies structure: paragraph
count, per-paragraph text, shape geometry, wrap setting.
"""
import os
import sys

sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt


def gen(out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # blank

    # ---- Slide 1: multi-paragraph, default settings, no wrap ----
    # A wide box so none of the lines wrap. Pins the DEFAULT line advance
    # (multiple paragraphs in one text frame).
    s1 = prs.slides.add_slide(blank)
    tb1 = s1.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.0), Inches(4.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    tf1.text = "Line one"
    for t in ("Line two", "Line three", "Line four"):
        p = tf1.add_paragraph()
        p.text = t
    for p in tf1.paragraphs:
        for r in p.runs:
            r.font.size = Pt(18)

    # ---- Slide 2: word wrap ON, a long paragraph in a narrow box ----
    s2 = prs.slides.add_slide(blank)
    tb2 = s2.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(4.0), Inches(3.0))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    tf2.text = ("The quick brown fox jumps over the lazy dog. "
                "Pack my box with five dozen liquor jugs.")
    for r in tf2.paragraphs[0].runs:
        r.font.size = Pt(18)

    # ---- Slide 3: line spacing 100 / 150 / 200% ----
    s3 = prs.slides.add_slide(blank)
    tb3 = s3.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.0), Inches(4.0))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    tf3.text = "Single spacing"
    for t, ls in (("One and a half", 1.5), ("Double", 2.0)):
        p = tf3.add_paragraph()
        p.text = t
        p.line_spacing = ls
    for p in tf3.paragraphs:
        for r in p.runs:
            r.font.size = Pt(18)

    # ---- Slide 4: paragraph space before / after ----
    s4 = prs.slides.add_slide(blank)
    tb4 = s4.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.0), Inches(4.0))
    tf4 = tb4.text_frame
    tf4.word_wrap = True
    tf4.text = "First para"
    p2 = tf4.add_paragraph()
    p2.text = "Second para"
    p2.space_before = Pt(12)
    p2.space_after = Pt(24)
    p3 = tf4.add_paragraph()
    p3.text = "Third para"
    p3.space_before = Pt(6)
    for p in tf4.paragraphs:
        for r in p.runs:
            r.font.size = Pt(18)

    prs.save(out_path)
    print("wrote", out_path)


if __name__ == "__main__":
    out_dir = sys.argv[1] if len(sys.argv) > 1 else r"c:\Users\ryuji\oxi-main\pipeline_data\pptx_probes"
    os.makedirs(out_dir, exist_ok=True)
    gen(os.path.join(out_dir, "spec4_textframe.pptx"))
