# -*- coding: utf-8 -*-
"""Generate the Spec #4 wave-3 repro: FINE line-spacing sweep (f(n)).

Wave-2 (spec4b) gave 8 coarse points per font and falsified H1 (1.2*fs*n+c),
H2 (single + (n-1)*k) and H3 (1.2*fs*n). Key structural facts:
  - single spacing (n=1.0, spcPct 100%) is FONT-DEPENDENT (Calibri r=1.26,
    Arial r=1.22, TNR r=1.227; none match hhea)
  - n>=1.1 is font-INDEPENDENT and fs-linear but NON-linear in n
To pin the exact function we need a FINE sweep (snap detection) and the
sub-1.0 region (0.9) to see how the 1.0 font-dependent value connects to the
n>=1.1 font-independent curve.

Design: 4 slides x 9 paragraphs, Calibri fs=18 only (n>=1.1 font-independent
+ fs-linear already proven). Each slide's first paragraph is a duplicate of
the previous slide's last n (so it carries no measured advance), the rest
are the sweep. Baseline deltas between consecutive paragraphs pin the
advance of the paragraph BEFORE the delta, so per slide we get 8 advances:

  Slide 1: [1.0, 0.9, 1.0, 1.05, 1.1, 1.15, 1.2, 1.25, 1.3]   -> 0.9..1.3
  Slide 2: [1.3, 1.35, 1.4, 1.45, 1.5, 1.55, 1.6, 1.65, 1.7] -> 1.35..1.7
  Slide 3: [1.7, 1.75, 1.8, 1.85, 1.9, 1.95, 2.0, 2.1, 2.2]   -> 1.75..2.2
  Slide 4: [2.2, 2.3, 2.4, 2.5, 2.6, 2.7, 2.8, 2.9, 3.0]      -> 2.3..3.0

Total 32 measured advances: 0.9, 1.0, 1.05..2.0 (step 0.05), 2.0..3.0
(step 0.1). Same text box geometry as wave-2 (off 72,72 ext 792x432).
"""
import os
import sys

sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
_REPO = Path(__file__).resolve().parents[2]

SLIDES = [
    [1.0, 0.9, 1.0, 1.05, 1.1, 1.15, 1.2, 1.25, 1.3],
    [1.3, 1.35, 1.4, 1.45, 1.5, 1.55, 1.6, 1.65, 1.7],
    [1.7, 1.75, 1.8, 1.85, 1.9, 1.95, 2.0, 2.1, 2.2],
    [2.2, 2.3, 2.4, 2.5, 2.6, 2.7, 2.8, 2.9, 3.0],
]


def add_slide(prs, spacings):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tb = s.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.0), Inches(6.0))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, spc in enumerate(spacings):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "s%05.2f" % spc
        p.line_spacing = spc
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(18)
            r.font.name = "Calibri"
    return s


def gen(out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for spcs in SLIDES:
        add_slide(prs, spcs)

    prs.save(out_path)
    print("wrote", out_path)


if __name__ == "__main__":
    out_dir = sys.argv[1] if len(sys.argv) > 1 else str(_REPO / r"pipeline_data\pptx_probes")
    os.makedirs(out_dir, exist_ok=True)
    gen(os.path.join(out_dir, "spec4c_lspacing.pptx"))
