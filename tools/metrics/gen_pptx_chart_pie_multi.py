#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generate a multi-series pie chart pptx for Ra-loop measurement.

Two series (the same categories East/West/Midwest) to probe how Word/PPT
renders a MULTI-series pie (python-pptx emits 2 <c:ser>).  Word PDF ->
pipeline_data/pptx_probes/chart_pie_multi/chart_pie_multi.pdf (via COM export).
"""
import os
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

OUT = r"pipeline_data\pptx_probes\chart_pie_multi\chart_pie_multi.pptx"


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    chart_data = CategoryChartData()
    chart_data.categories = ["East", "West", "Midwest"]
    chart_data.add_series("Rev", (19.2, 21.4, 16.7))
    chart_data.add_series("Cost", (10.5, 11.2, 8.5))

    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(1), Inches(1), Inches(5.5), Inches(4.0), chart_data
    )
    chart = gframe.chart
    chart.has_legend = False

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print(f"Saved {OUT}")


if __name__ == "__main__":
    main()
