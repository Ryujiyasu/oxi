# -*- coding: utf-8 -*-
"""PIE OUTSIDE_END data-label placement probe: sweep label text WIDTH at a
FIXED slice mid-angle so Word's placement law (outer-edge-constant vs
radial-constant vs something else) can be read from the PDF.

Rationale: chart_datalabel_pie P1 (short labels) vs P2 (long '1920.0%')
showed label CENTER distances of 75.4/80.4/75.4 vs 60.8/67.3/60.9pt from
the circle centre — long labels sit ~14pt INWARD. To separate the models we
vary ONLY the label width at a fixed slice angle:

  - slice 0 angle is fixed by the value ratio (1:1 -> 180deg, mid 90deg;
    2:1 -> 240deg, mid 120deg; 3:1 -> 270deg, mid 135deg from 12 o'clock).
  - label width is swept via number_format decimal places ("0","0.0",
    "0.00","0.000","0.0000") and via literal-text prefixes ('"A: "0',
    '"AAAAAAAAAAAAAAAA: "0.00') which change the rendered width WITHOUT
    changing the value/angle.

All OUTSIDE_END (default position), default Office theme, frame 72,72,396,288,
series "Series 1" (auto-title), so the circle shrink rule (15.78pt both sides)
and circle centre (270, 233.68) match the shipped chart_datalabel_pie probe."""
import sys, os
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

base = r"pipeline_data\pptx_probes\pie_label_probe"
os.makedirs(base, exist_ok=True)

# (values, num_fmt)  -- slice0 angle / mid-angle noted per slide
SLIDES = [
    # width ramp at mid 90deg (1:1)
    ((1.0, 1.0), "0"),            #  0 '1'
    ((1.0, 1.0), "0.0"),          #  1 '1.0'
    ((1.0, 1.0), "0.00"),         #  2 '1.00'
    ((1.0, 1.0), "0.000"),        #  3 '1.000'
    ((1.0, 1.0), '"A: "0'),       #  4 'A: 1' (literal prefix)
    ((1.0, 1.0), '"AAAAAAAAAAAAAAAA: "0.00'),  #  5 long literal
    # width ramp at mid 120deg (2:1)
    ((2.0, 1.0), "0.0"),          #  6 '2.0'
    ((2.0, 1.0), "0.0000"),       #  7 '2.0000'
    # width ramp at mid 135deg (3:1)
    ((3.0, 1.0), "0.0"),          #  8 '3.0'
    ((3.0, 1.0), "0.0000"),       #  9 '3.0000'
]

prs = Presentation()
blank = prs.slide_layouts[6]

for i, (values, fmt) in enumerate(SLIDES):
    slide = prs.slides.add_slide(blank)
    cd = CategoryChartData()
    cd.categories = ["A", "B"]
    cd.add_series("Series 1", values)
    x, y, cx, cy = Inches(1.0), Inches(1.0), Inches(5.5), Inches(4.0)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, cd)
    plot = gframe.chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_value = True
    dl.number_format_is_linked = False
    dl.number_format = fmt
    # position left default = OUTSIDE_END

out = os.path.join(base, "pie_label_probe.pptx")
prs.save(out)
print("saved:", out, os.path.getsize(out), "slides:", len(prs.slides._sldIdLst))
