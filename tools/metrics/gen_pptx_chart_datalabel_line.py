# -*- coding: utf-8 -*-
"""Chart data-labels probe for LINE charts: LINE_MARKERS with data labels so
Word's PDF exposes the exact data-label placement/format/font rules for the
line branch (the bar data-label rule was derived from chart_dlbls S1-S5 in
2026-08-06; the line/pie placement must be measured separately).

  L1: single-series line, show_value=True        (default label position)
  L2: single-series line, show_value=True + number_format="0.0%"
  L3: multi-series line (2), show_value=True

All on blank slides (layout 6), default Office theme, frame 72,72,396,288,
categories [East, West, Midwest], series values (19.2, 21.4, 16.7)."""
import sys, os
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

base = r"pipeline_data\pptx_probes\chart_datalabel_line"
os.makedirs(base, exist_ok=True)

prs = Presentation()
blanks = [prs.slide_layouts[6] for _ in range(3)]


def add_line_slide(layout, n_series, num_fmt=None):
    slide = prs.slides.add_slide(layout)
    cd = CategoryChartData()
    cd.categories = ["East", "West", "Midwest"]
    cd.add_series("Series 1", (19.2, 21.4, 16.7))
    if n_series == 2:
        cd.add_series("Series 2", (15.3, 17.5, 20.1))
    x, y, cx, cy = Inches(1.0), Inches(1.0), Inches(5.5), Inches(4.0)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, cd)
    plot = gframe.chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_value = True
    if num_fmt is not None:
        dl.number_format_is_linked = False
        dl.number_format = num_fmt
    return slide


add_line_slide(blanks[0], 1)
add_line_slide(blanks[1], 1, num_fmt="0.0%")
add_line_slide(blanks[2], 2)

out = os.path.join(base, "chart_datalabel_line.pptx")
prs.save(out)
print("saved:", out, os.path.getsize(out))
