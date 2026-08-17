# -*- coding: utf-8 -*-
"""Probe: which fill rule does PowerPoint use for a multi-subpath a:custGeom?

901 of the dev corpus's 11470 custGeom shapes have more than one `a:moveTo`
AND a fill, so the rule decides real ink -- even-odd leaves a hole where
non-zero winding fills it. GDI defaults to ALTERNATE (even-odd) and the two
rules disagree exactly on the arms below, so one export answers it.

Arms (one shape per slide, 4in square, sampled at its centre):
  C1 outer + inner square, SAME winding      even-odd: hole   nonzero: filled
  C2 outer + inner square, OPPOSITE winding  both: hole
  C3 pentagram (one self-intersecting path)  even-odd: hollow nonzero: filled
  C4 two disjoint squares                    both: both filled (sanity)
  C5 three nested squares, same winding      even-odd: fill/hole/fill
python-pptx cannot write custGeom, so each shape's prstGeom is swapped for one
by zip surgery -- the same technique the gradient / doughnut probes use.
"""
from __future__ import annotations

import math
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\custgeom_fillrule").resolve()
SPACE = 1000  # a:path/@w and @h


def poly(points: list[tuple[float, float]]) -> str:
    """One closed subpath through `points`."""
    head = points[0]
    body = "".join(
        f'<a:lnTo><a:pt x="{int(round(x))}" y="{int(round(y))}"/></a:lnTo>'
        for x, y in points[1:]
    )
    return (
        f'<a:moveTo><a:pt x="{int(round(head[0]))}" y="{int(round(head[1]))}"/></a:moveTo>'
        f"{body}<a:close/>"
    )


def geom(subpaths: list[str]) -> str:
    return (
        "<a:custGeom><a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>"
        '<a:rect b="b" l="l" r="r" t="t"/><a:pathLst>'
        f'<a:path extrusionOk="0" h="{SPACE}" w="{SPACE}">{"".join(subpaths)}</a:path>'
        "</a:pathLst></a:custGeom>"
    )


def square(lo: float, hi: float, clockwise: bool) -> list[tuple[float, float]]:
    cw = [(lo, lo), (hi, lo), (hi, hi), (lo, hi)]
    return cw if clockwise else list(reversed(cw))


def star(points: int = 5) -> list[tuple[float, float]]:
    """A pentagram: every other vertex of a regular pentagon, one closed path."""
    r, c = SPACE / 2, SPACE / 2
    verts = [
        (c + r * math.cos(-math.pi / 2 + 2 * math.pi * i / points),
         c + r * math.sin(-math.pi / 2 + 2 * math.pi * i / points))
        for i in range(points)
    ]
    return [verts[(i * 2) % points] for i in range(points)]


ARMS: list[tuple[str, str]] = [
    ("C1 nested squares, same winding", geom([poly(square(0, SPACE, True)),
                                              poly(square(250, 750, True))])),
    ("C2 nested squares, opposite winding", geom([poly(square(0, SPACE, True)),
                                                  poly(square(250, 750, False))])),
    ("C3 pentagram, single path", geom([poly(star())])),
    ("C4 two disjoint squares", geom([poly(square(0, 400, True)),
                                      poly(square(600, SPACE, True))])),
    ("C5 three nested squares, same winding", geom([poly(square(0, SPACE, True)),
                                                    poly(square(200, 800, True)),
                                                    poly(square(400, 600, True))])),
]

SHAPE_EMU = 3657600  # 4in
SHAPE_XY = (2743200, 1600200)  # centred on a 10x7.5in slide


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    tmp, dst = OUT / "_stage.pptx", OUT / "custgeom_fillrule.pptx"

    prs = Presentation()
    blank = prs.slide_layouts[6]
    for label, _ in ARMS:
        s = prs.slides.add_slide(blank)
        box = s.shapes.add_textbox(Emu(228600), Emu(228600), Emu(6400800), Emu(400050))
        box.text_frame.text = label
        shape = s.shapes.add_shape(1, Emu(SHAPE_XY[0]), Emu(SHAPE_XY[1]),
                                   Emu(SHAPE_EMU), Emu(SHAPE_EMU))
        # A flat red fill with no outline and no inherited theme shadow, so the
        # centre pixel answers "filled or not" with no other ink in the way.
        shape.fill.solid()
        shape.fill.fore_color.rgb = __import__("pptx.dml.color", fromlist=["RGBColor"]).RGBColor(0xFF, 0x00, 0x00)
        shape.line.fill.background()
        shape.shadow.inherit = False
    prs.save(tmp)

    with zipfile.ZipFile(tmp) as zin:
        names = zin.namelist()
        data = {n: zin.read(n) for n in names}

    for i, (label, g) in enumerate(ARMS, start=1):
        part = f"ppt/slides/slide{i}.xml"
        xml = data[part].decode("utf-8")
        # rindex, not index: python-pptx gives the caption TEXTBOX a prstGeom
        # too, and it comes first. Injecting there leaves the probe shape a
        # plain rectangle -- which is exactly what the first export showed.
        start = xml.rindex("<a:prstGeom")
        end = xml.rindex("</a:prstGeom>") + len("</a:prstGeom>")
        data[part] = (xml[:start] + g + xml[end:]).encode("utf-8")

    with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:
        for n in names:
            zout.writestr(n, data[n])
    tmp.unlink()
    print(f"wrote {dst}  ({len(ARMS)} arms)")


if __name__ == "__main__":
    main()
