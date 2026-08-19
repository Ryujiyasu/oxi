# -*- coding: utf-8 -*-
"""Probe: where does PowerPoint break a "word" that is wider than its line?

Oxi's wrap splits on spaces only, so a run with no space in it stays one line
and runs off the shape. d11 and d24 slide 38 are the visible case -- 53 emoji
in a 490pt box that PowerPoint lays out in four rows -- and 45 paragraphs
across nine dev decks carry a space-free run of 30 characters or more, long
URLs being the other kind.

Splitting at "the last character that fits" is the obvious guess and it is not
obviously right: a break could prefer punctuation (the way a browser breaks a
URL after `/`), could refuse to break a Latin word at all and let it overflow,
or could keep at least one character on the line even when that one character
does not fit. Each arm here is one unbroken run in a narrow box; reading the
line starts out of PowerPoint's PDF says which rule it follows.
"""
from __future__ import annotations

import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\charwrap").resolve()
A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# (label, text, size pt, latin typeface)
ARMS = [
    # plain Latin letters: is an overlong word broken at all?
    ("latin", "abcdefghijklmnopqrstuvwxyzabcdefghijklmnopqrstuvwxyz", 18, "Arial"),
    # same but the box could break after a hyphen instead
    ("hyphen", "alpha-beta-gamma-delta-epsilon-zeta-eta-theta-iota", 18, "Arial"),
    # a URL: slashes and dots are the candidate break points
    ("url", "https://www.example.com/some/rather/long/path/index.html", 18, "Arial"),
    # digits only, no break opportunity of any kind
    ("digits", "12345678901234567890123456789012345678901234567890", 18, "Arial"),
    # one word short enough to fit, then an overlong one after a space
    ("mixed", "ok abcdefghijklmnopqrstuvwxyzabcdefghijklmnopqrstuvwxyz", 18, "Arial"),
    # non-BMP: every character is a surrogate pair
    ("emoji", "\U0001F44B\U0001F446\U0001F449\U0001F44D\U0001F464\U0001F466"
              "\U0001F467\U0001F468\U0001F469\U0001F46A\U0001F483\U0001F3C3"
              "\U0001F491\U0001F602\U0001F609\U0001F60B\U0001F612\U0001F476", 28, "Arial"),
    # CJK, which every layout engine breaks per character
    ("cjk", "本日は晴天なり本日は晴天なり本日は晴天なり本日は晴天なり", 18, "Arial"),
    # a single character far wider than the whole line
    ("giant", "MMMMMMMMMMMMMMMMMMMMMMMM", 54, "Arial"),
]


def q(tag: str) -> str:
    return f"{{{A}}}{tag}"


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for label, text, sz, face in ARMS:
        slide = prs.slides.add_slide(blank)
        cap = slide.shapes.add_textbox(Emu(228600), Emu(114300), Emu(6400800), Emu(300000))
        cap.text_frame.text = label
        box = slide.shapes.add_textbox(Emu(914400), Emu(1200000), Emu(2286000), Emu(3600000))
        tf = box.text_frame
        tf.word_wrap = True
        body = tf._txBody
        for pel in body.findall(q("p")):
            body.remove(pel)
        p = etree.SubElement(body, q("p"))
        ppr = etree.SubElement(p, q("pPr"))
        etree.SubElement(ppr, q("buNone"))
        ln = etree.SubElement(ppr, q("lnSpc"))
        etree.SubElement(ln, q("spcPct")).set("val", "100000")
        r = etree.SubElement(p, q("r"))
        rpr = etree.SubElement(r, q("rPr"))
        rpr.set("lang", "en-US")
        rpr.set("sz", str(sz * 100))
        etree.SubElement(rpr, q("latin")).set("typeface", face)
        etree.SubElement(rpr, q("ea")).set("typeface", "Yu Gothic")
        etree.SubElement(r, q("t")).text = text
    prs.save(OUT / "charwrap.pptx")
    print(f"wrote {OUT / 'charwrap.pptx'}  ({len(ARMS)} arms)")


if __name__ == "__main__":
    main()
