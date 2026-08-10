"""100%-STACKED column chart on a blank slide, default theme.

Same 3 categories x 2 series data as chart_stacked (clustered twin chart2) so
the 100%-stacking rule can be measured against Word's PDF render (fitz
get_drawings + rawdict): how the per-series segments stack to 100%, the value
axis scale (0..100?), the auto-title / legend conditions.
"""
import os
import sys

from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

sys.stdout.reconfigure(encoding="utf-8")


def main():
    base = r"pipeline_data\pptx_probes\chart_stacked100"
    os.makedirs(base, exist_ok=True)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3"]
    chart_data.add_series("Revenue", (19.2, 21.4, 16.7))
    chart_data.add_series("Cost", (10.5, 15.0, 12.3))

    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED_100,
        Inches(1.0), Inches(1.0), Inches(5.5), Inches(4.0),
        chart_data,
    )

    out = os.path.join(base, "chart_stacked100.pptx")
    prs.save(out)
    print("saved:", out, os.path.getsize(out))

    # Verify the chart XML declares grouping="percentStacked"
    import zipfile
    with zipfile.ZipFile(out) as z:
        names = [n for n in z.namelist() if n.startswith("ppt/charts/") and n.endswith(".xml")]
        print("chart parts:", names)
        for n in names:
            xml = z.read(n).decode("utf-8")
            for key in ("grouping", "percentStacked", "stacked"):
                if key in xml:
                    i = xml.find(key)
                    print(f"  [{key}] ...{xml[max(0, i-60):i+80]!r}...")


if __name__ == "__main__":
    main()
