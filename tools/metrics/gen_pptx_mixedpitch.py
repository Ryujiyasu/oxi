# -*- coding: utf-8 -*-
"""Probe: the baseline step between two lines of DIFFERENT font size.

Oxi advances a paragraph by `fs * 1.2 * lnSpc` where fs is that paragraph's own
size, so the step from a 55pt line to a 66pt line is 55*1.2 = 66pt. d28's title
slide says PowerPoint uses 76.3pt there, while agreeing exactly on 66->66
(79.2pt = 66*1.2). A line-box model explains both: the step is the DESCENT part
of the line being left plus the ASCENT part of the line being entered, and when
the sizes are equal that collapses to (a+d)*s = 1.2*s.

Each size arm is two one-line paragraphs whose sizes are what varies, so the
step is read straight off the two baselines. The four faces sit far apart in
ascent:descent ratio (Verdana 0.827:0.173, Arial 0.810:0.190, Georgia
0.807:0.193, Calibri 0.780:0.220), so they cannot share one constant if the
split really is the font's own. The trailing arms ask whose `lnSpc` scales the
step -- the line being left, the one being entered, or both.
"""
from __future__ import annotations

import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\mixedpitch").resolve()
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
PAIRS = [(10, 10), (10, 20), (20, 10), (10, 40), (40, 10), (24, 66), (55, 66), (66, 55)]
FONTS = ["Arial", "Georgia", "Calibri", "Verdana"]
LNSPC = [("ls150_both", 150, 150), ("ls150_next", 100, 150), ("ls150_prev", 150, 100)]


def q(tag: str) -> str:
    return f"{{{A}}}{tag}"


def add_para(body, text: str, pt: float, font: str, ls: int) -> None:
    p = etree.SubElement(body, q("p"))
    ppr = etree.SubElement(p, q("pPr"))
    ln = etree.SubElement(ppr, q("lnSpc"))
    etree.SubElement(ln, q("spcPct")).set("val", str(ls * 1000))
    r = etree.SubElement(p, q("r"))
    rpr = etree.SubElement(r, q("rPr"))
    rpr.set("lang", "en-US")
    rpr.set("sz", str(int(pt * 100)))
    etree.SubElement(rpr, q("latin")).set("typeface", font)
    etree.SubElement(r, q("t")).text = text


def arm(prs, blank, caption: str):
    slide = prs.slides.add_slide(blank)
    cap = slide.shapes.add_textbox(Emu(228600), Emu(114300), Emu(6400800), Emu(300000))
    cap.text_frame.text = caption
    box = slide.shapes.add_textbox(Emu(457200), Emu(1200000), Emu(7000000), Emu(3000000))
    body = box.text_frame._txBody
    for pel in body.findall(q("p")):
        body.remove(pel)
    return body


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for font in FONTS:
        for s1, s2 in PAIRS:
            body = arm(prs, blank, f"{font} {s1}->{s2}")
            add_para(body, "AAA", s1, font, 100)
            add_para(body, "BBB", s2, font, 100)
    for label, l1, l2 in LNSPC:
        body = arm(prs, blank, f"Arial 10->40 {label}")
        add_para(body, "AAA", 10, "Arial", l1)
        add_para(body, "BBB", 40, "Arial", l2)
    print(f"wrote {OUT / 'mixedpitch.pptx'}  "
          f"({len(FONTS)}x{len(PAIRS)} size arms + {len(LNSPC)} lnSpc arms)")
    prs.save(OUT / "mixedpitch.pptx")


if __name__ == "__main__":
    main()
