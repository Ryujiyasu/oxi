# -*- coding: utf-8 -*-
"""Generate the Spec #4 wave-2 repros: line-spacing function sweep.

Wave-1 pinned single spacing = fs*1.2 (21.60pt @ 18pt) but the multiple
spacing was NON-linear (1.5x -> 29.16, 2.0x -> 40.47; neither fs*1.2*n nor
a line in n). Wave-2 sweeps line_spacing finely and varies font + size to
derive the actual function:

  - Slide A: theme-default font, fs=18, line_spacing 1.0..3.0
  - Slide B-D: explicit Calibri / Arial / Times New Roman, fs=18, same sweep
  - Slide E-F: explicit Calibri, fs=10 / fs=36, same sweep

Each slide has ONE text box with 9 paragraphs ("sp <n>"), one line each
(short, wide box -> no wrap). Baseline deltas between consecutive paragraphs
pin the per-paragraph line advance.

Also lets us confirm the theme-default font (via PDF span fonts) and the
first-baseline offset (top + ascent relationship).
"""
import os
import sys

sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt

SPACINGS = [1.0, 1.1, 1.2, 1.3, 1.5, 1.8, 2.0, 2.5, 3.0]


def add_slide(prs, font_name, fs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tb = s.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.0), Inches(6.0))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, spc in enumerate(SPACINGS):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "sp %s" % spc
        p.line_spacing = spc
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(fs)
            if font_name:
                r.font.name = font_name
    return s


def gen(out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_slide(prs, None, 18)                       # A: theme default, fs 18
    for fn in ("Calibri", "Arial", "Times New Roman"):
        add_slide(prs, fn, 18)                     # B-D: explicit font, fs 18
    add_slide(prs, "Calibri", 10)                  # E: Calibri fs 10
    add_slide(prs, "Calibri", 36)                  # F: Calibri fs 36

    prs.save(out_path)
    print("wrote", out_path)


if __name__ == "__main__":
    out_dir = sys.argv[1] if len(sys.argv) > 1 else r"c:\Users\ryuji\oxi-main\pipeline_data\pptx_probes"
    os.makedirs(out_dir, exist_ok=True)
    gen(os.path.join(out_dir, "spec4b_lspacing.pptx"))
