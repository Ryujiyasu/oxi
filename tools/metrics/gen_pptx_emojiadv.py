# -*- coding: utf-8 -*-
"""Probe: how wide is an emoji, in each of its two presentations?

`emojipres` showed PowerPoint paints U+2764 monochrome in the run colour and
U+2764 U+FE0F in colour, and that the two arrive in the PDF as images of
different width (86.8pt against 102.8pt at 40pt). Image rectangles carry bleed,
so they are not the advance. Bracketing the character between two letters makes
the advance readable exactly: the trailing letter IS a text span, and its x
minus the leading letter's gives letter + emoji.
"""
from __future__ import annotations

import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\emojiadv").resolve()
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
VS16 = "️"

ARMS = [
    ("none", "AA"),
    ("heart_text", "A❤A"),
    ("heart_color", "A❤" + VS16 + "A"),
    ("hand", "A✋A"),
    ("thermo_text", "A🌡A"),
    ("thermo_color", "A🌡" + VS16 + "A"),
    ("grin", "A😀A"),
    ("watch", "A⌚A"),
]


def q(tag: str) -> str:
    return f"{{{A}}}{tag}"


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for label, text in ARMS:
        slide = prs.slides.add_slide(blank)
        cap = slide.shapes.add_textbox(Emu(228600), Emu(114300), Emu(6400800), Emu(300000))
        cap.text_frame.text = label
        box = slide.shapes.add_textbox(Emu(914400), Emu(1200000), Emu(5486400), Emu(900000))
        tf = box.text_frame
        tf.word_wrap = False
        body = tf._txBody
        for pel in body.findall(q("p")):
            body.remove(pel)
        p = etree.SubElement(body, q("p"))
        ppr = etree.SubElement(p, q("pPr"))
        etree.SubElement(ppr, q("buNone"))
        r = etree.SubElement(p, q("r"))
        rpr = etree.SubElement(r, q("rPr"))
        rpr.set("lang", "en-US")
        rpr.set("sz", "4000")
        etree.SubElement(rpr, q("latin")).set("typeface", "Arial")
        etree.SubElement(r, q("t")).text = text
    prs.save(OUT / "emojiadv.pptx")
    print(f"wrote {OUT / 'emojiadv.pptx'}  ({len(ARMS)} arms)")


if __name__ == "__main__":
    main()
