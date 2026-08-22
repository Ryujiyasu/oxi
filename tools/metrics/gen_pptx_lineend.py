# -*- coding: utf-8 -*-
"""Author a minimal repro for DrawingML line-end (arrowhead) sizing.

The dev corpus shows PowerPoint drawing filled circles at connector ends that
Oxi omits entirely, and the two corpus data points do not fit a single
"size = k x line width" law: an `oval` at `w/len=med` on a 0.75pt line came out
5.76pt across, while `lg` on a 3.00pt line came out 14.88pt.  Either the
factors differ per size token in a way those two points cannot separate, or the
scaling clamps for thin lines.  This deck separates the variables: one slide per
end TYPE, and on each slide a 5x5 grid of line WIDTH against the (w, len) token
pair, so the law can be read off directly.

Every connector is horizontal and black on white, well clear of its neighbours,
so the measuring script can flood-fill one line's ink without leaking.

    python tools/metrics/gen_pptx_lineend.py [out.pptx]
"""
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

OUT = Path(sys.argv[1] if len(sys.argv) > 1
           else r'pipeline_data\pptx_probe\probe_lineend.pptx')
TYPES = ['oval', 'triangle', 'stealth', 'arrow', 'diamond']
WIDTHS_PT = [0.75, 1.5, 3.0, 4.5, 6.0]
TOKENS = [('sm', 'sm'), ('med', 'med'), ('lg', 'lg'), ('sm', 'lg'), ('lg', 'sm')]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

manifest = []
for si, kind in enumerate(TYPES, start=1):
    slide = prs.slides.add_slide(blank)
    for row, (w_tok, len_tok) in enumerate(TOKENS):
        for col, w_pt in enumerate(WIDTHS_PT):
            x_in = 0.30 + col * 2.60
            y_in = 0.80 + row * 1.30
            x = Inches(x_in)
            y = Inches(y_in)
            manifest.append({
                'slide': si, 'type': kind, 'w_tok': w_tok, 'len_tok': len_tok,
                'line_pt': w_pt,
                'x0_pt': x_in * 72, 'x1_pt': (x_in + 2.00) * 72, 'y_pt': y_in * 72,
            })
            cxn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, x, y, x + Inches(2.00), y)
            # python-pptx writes prst="line"; every connector in the dev corpus
            # is a straightConnector1, and the renderer routes on that name, so
            # the repro has to state the form the corpus actually uses.
            geom = cxn._element.spPr.find(qn('a:prstGeom'))
            geom.set('prst', 'straightConnector1')
            line = cxn.line
            line.width = Pt(w_pt)
            line.color.rgb = RGBColor(0x00, 0x00, 0x00)
            ln = cxn.line._get_or_add_ln()
            # python-pptx has no API for the ends; the schema orders them
            # head then tail, after prstDash/round, so append and let Word's
            # own round-trip normalise if it must.
            for tag in ('headEnd', 'tailEnd'):
                el = ln.makeelement(qn('a:' + tag), {})
                el.set('type', kind)
                el.set('w', w_tok)
                el.set('len', len_tok)
                ln.append(el)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
OUT.with_suffix('.json').write_text(json.dumps(manifest, indent=1), encoding='utf-8')
print('wrote', OUT, 'and', OUT.with_suffix('.json'))
print(f'{len(TYPES)} slides x {len(TOKENS)} token pairs x {len(WIDTHS_PT)} widths'
      f' = {len(TYPES) * len(TOKENS) * len(WIDTHS_PT)} connectors')
