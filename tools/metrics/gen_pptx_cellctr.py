# -*- coding: utf-8 -*-
"""Probe: does a centred cell's baseline depend on how much room the row has?

The cell first-baseline model (`1.2 x asc/(asc+desc) x size` below the line
box top) is confirmed on five faces — but d25 s7 wants it and d35 s35 wants
GDI's TextOut top instead, and the visible difference between those two is how
much taller the row is than the text block (111.8pt around a 44.5pt block
versus 30.17 around 9.65).

So: one table, one face, one size, rows declared at increasing heights, every
cell `anchor="ctr"` with a single line. All rows exceed the minimum
(2*margin + 1.2*size, confirmed by the rowgrow probe), so no row grows and the
row tops accumulate exactly — which makes the baseline offset readable per row.

If the offset is constant, centring is not the variable and d35's difference is
about its face. If it drifts with the row's spare room, the drift IS the rule.

Usage:
    python tools/metrics/gen_pptx_cellctr.py
    python tools/metrics/measure_pptx_word.py pipeline_data/pptx_probes/cellctr/cellctr.pptx pipeline_data/pptx_probes/cellctr
    python tools/metrics/read_pptx_cellctr.py
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\cellctr").resolve()
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
PT = 12700
SIZE = 12.0
MARGIN = 7.2
TABLE_TOP = 60.0
# every height is >= 2*MARGIN + 1.2*SIZE = 28.8, so nothing grows
HEIGHTS = [28.8, 33.0, 40.0, 52.0, 72.0, 100.0]
FACES = ["Arial", "Segoe Script", "Comic Sans MS"]


def q(tag: str) -> str:
    return f"{{{A}}}{tag}"


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    manifest = []
    for face in FACES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cap = slide.shapes.add_textbox(Emu(228600), Emu(114300), Emu(6400800), Emu(300000))
        cap.text_frame.text = f"{face} {SIZE}pt anchor=ctr, rows {HEIGHTS}"
        shape = slide.shapes.add_table(
            len(HEIGHTS), 1, Emu(int(60 * PT)), Emu(int(TABLE_TOP * PT)),
            Emu(int(240 * PT)), Emu(int(sum(HEIGHTS) * PT)),
        )
        tbl = shape._element.graphic.graphicData.tbl
        for pr in tbl.findall(q("tblPr")):
            tbl.remove(pr)
        pr = etree.SubElement(tbl, q("tblPr"))
        pr.set("firstRow", "0")
        pr.set("bandRow", "0")
        etree.SubElement(pr, q("noFill"))
        tbl.insert(0, pr)
        for tr, h in zip(tbl.findall(q("tr")), HEIGHTS):
            tr.set("h", str(int(round(h * PT))))
            for tc in tr.findall(q("tc")):
                txbody = tc.find(q("txBody"))
                for p in txbody.findall(q("p")):
                    txbody.remove(p)
                p = etree.SubElement(txbody, q("p"))
                ppr = etree.SubElement(p, q("pPr"))
                ln = etree.SubElement(ppr, q("lnSpc"))
                etree.SubElement(ln, q("spcPct")).set("val", "100000")
                etree.SubElement(ppr, q("buNone"))
                r = etree.SubElement(p, q("r"))
                rpr = etree.SubElement(r, q("rPr"))
                rpr.set("lang", "en-US")
                rpr.set("sz", str(int(SIZE * 100)))
                rpr.set("kern", "0")
                etree.SubElement(rpr, q("latin")).set("typeface", face)
                etree.SubElement(r, q("t")).text = "Hxy"
                tcpr = tc.find(q("tcPr"))
                if tcpr is None:
                    tcpr = etree.SubElement(tc, q("tcPr"))
                for name in ("marT", "marB", "marL", "marR"):
                    tcpr.set(name, str(int(MARGIN * PT)))
                tcpr.set("anchor", "ctr")
        manifest.append({
            "slide": len(manifest) + 1,
            "face": face,
            "size": SIZE,
            "margin": MARGIN,
            "table_top": TABLE_TOP,
            "heights": HEIGHTS,
        })
        print(f"{face}: rows {HEIGHTS}")
    prs.save(OUT / "cellctr.pptx")
    (OUT / "cellctr_manifest.json").write_text(
        json.dumps(manifest, indent=1), encoding="utf-8")
    print(f"wrote {OUT / 'cellctr.pptx'}")


if __name__ == "__main__":
    main()
