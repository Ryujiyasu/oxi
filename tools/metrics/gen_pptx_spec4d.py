# -*- coding: utf-8 -*-
"""Spec #4 wave-4: measure line spacing (a:lnSpc spcPct) INSIDE a single
multi-line paragraph.

The 1-para-1-line probes (spec4b/spec4c) could NOT isolate the line advance:
each paragraph carries its own line spacing, so the baseline delta between
adjacent paragraphs is a MIX of the two paragraphs' line heights (leading is
split around each line). wave-4 fixes this by making ONE paragraph wrap over
several lines (word_wrap=on, long text) — the baseline deltas WITHIN a single
paragraph are that paragraph's own line advance, cleanly.

Each slide = one TextBox (11in wide, 6in tall) with a SINGLE paragraph whose
line_spacing = n. The paragraph text is long enough to wrap ~6 lines. We then
measure per-line baselines in the PDF; consecutive deltas within the slide are
the line advance for that n. The first line is skipped (frame-top baseline
offset mixes in).

n sweep: 0.5,0.6,0.7,0.8,0.9,1.0,1.1,1.2,1.3,1.4,1.5,1.6,1.8,2.0,2.5,3.0
Output: pipeline_data\\pptx_probes\\spec4d_multiline.pptx
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUT = os.path.join(
    os.path.dirname(os.path.abspath(__file__)), "..", "..",
    "pipeline_data", "pptx_probes", "spec4d_multiline.pptx"
)
OUT = os.path.abspath(OUT)

NSWEEP = [0.5, 0.6, 0.7, 0.8, 0.9, 1.0, 1.1, 1.2, 1.3, 1.4, 1.5, 1.6, 1.8, 2.0, 2.5, 3.0]

# Long text: ~7 words per repeat; ~8 repeats => ~460 chars => ~5-6 lines at
# Calibri 18pt in an 11in-wide box (effective width ~777.6pt).
WORDS = "lorem ipsum dolor sit amet consectetur adipiscing elit "
TEXT = (WORDS * 8).strip()


def add_slide(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.0), Inches(6.0))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.line_spacing = n
    r = p.add_run()
    r.text = TEXT
    r.font.size = Pt(18)
    r.font.name = "Calibri"
    r.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(12.5)
    prs.slide_height = Inches(7.5)
    for n in NSWEEP:
        add_slide(prs, n)
    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print("wrote", OUT)


if __name__ == "__main__":
    main()
