# -*- coding: utf-8 -*-
"""Spec #5 (paragraph alignment) wave-1 repro: one TextBox with 4 single-line
paragraphs (Left/Center/Right/Justify), same font/size, no wrap. 12.5x7.5in."""
import sys, os
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

base = r"pipeline_data\pptx_probes\spec5_align"
os.makedirs(base, exist_ok=True)

prs = Presentation()  # default 12.5x7.5in
slide = prs.slides.add_slide(prs.slide_layouts[6])

# x=72 y=72 w=288 h=108 (pt)
tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(4.0), Inches(1.5))
tf = tb.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.1)
tf.margin_right = Inches(0.1)
tf.margin_top = Inches(0.05)
tf.margin_bottom = Inches(0.05)

items = [
    ("Left-Align", PP_ALIGN.LEFT),
    ("Center-Align", PP_ALIGN.CENTER),
    ("Right-Align", PP_ALIGN.RIGHT),
    ("Justify-Align", PP_ALIGN.JUSTIFY),
]
first = True
for text, align in items:
    if first:
        p = tf.paragraphs[0]
        first = False
    else:
        p = tf.add_paragraph()
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(18)
    r.font.name = "Arial"

out = os.path.join(base, "spec5_align.pptx")
prs.save(out)
print("saved:", out, os.path.getsize(out))
