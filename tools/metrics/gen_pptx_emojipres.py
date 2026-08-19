# -*- coding: utf-8 -*-
"""Probe: which characters does PowerPoint paint in COLOUR?

Oxi now draws COLR layers for any character the requested face lacks, and that
is one character too many: d11 slide 38 shows PowerPoint drawing U+2764 as a
BLACK heart while Oxi drew the colour one. Unicode gives every emoji an
Emoji_Presentation property -- U+2764 defaults to text presentation and only
becomes the colour glyph after U+FE0F -- so the question is whether PowerPoint
honours it or simply paints whatever the fallback font offers.

The read side uses PowerPoint's own PDF export: a colour emoji comes out as an
image with no text span (the charwrap probe's emoji arm returned no spans at
all), a monochrome one as ordinary text. Each arm is one character alone in a
box so the two cases cannot be confused.
"""
from __future__ import annotations

import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\emojipres").resolve()
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
VS16 = "️"

# (label, text) -- the Emoji_Presentation=No cases with and without VS16, the
# =Yes cases that should need no selector, and one plain letter as a control.
ARMS = [
    ("heart_plain", "❤"),
    ("heart_vs16", "❤" + VS16),
    ("hand_yes", "✋"),
    ("watch_yes", "⌚"),
    ("smile_no", "☺"),
    ("smile_vs16", "☺" + VS16),
    ("thermo_no", "🌡"),
    ("thermo_vs16", "🌡" + VS16),
    ("grin_yes", "😀"),
    ("eye_no", "👁"),
    ("copyright_no", "©"),
    ("letter_ctl", "A"),
]


def q(tag: str) -> str:
    return f"{{{A}}}{tag}"


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for label, text in ARMS:
        slide = prs.slides.add_slide(blank)
        cap = slide.shapes.add_textbox(Emu(228600), Emu(114300), Emu(6400800), Emu(300000))
        cap.text_frame.text = label
        box = slide.shapes.add_textbox(Emu(914400), Emu(1200000), Emu(2286000), Emu(900000))
        tf = box.text_frame
        body = tf._txBody
        for pel in body.findall(q("p")):
            body.remove(pel)
        p = etree.SubElement(body, q("p"))
        ppr = etree.SubElement(p, q("pPr"))
        etree.SubElement(ppr, q("buNone"))
        r = etree.SubElement(p, q("r"))
        rpr = etree.SubElement(r, q("rPr"))
        rpr.set("lang", "en-US")
        rpr.set("sz", "4000")
        etree.SubElement(rpr, q("latin")).set("typeface", "Arial")
        etree.SubElement(r, q("t")).text = text
    prs.save(OUT / "emojipres.pptx")
    print(f"wrote {OUT / 'emojipres.pptx'}  ({len(ARMS)} arms)")


if __name__ == "__main__":
    main()
