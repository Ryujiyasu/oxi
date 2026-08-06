#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Multi-slide pie probe to isolate title/legend influence on circle geometry.

slide1: pie, NO title, NO legend   (control vs chart_pie which had auto-title)
slide2: pie, NO title, HAS legend  (does legend shift the circle left?)
slide3: pie, HAS title, HAS legend (title + legend together)
slide4: 3-series pie, no title     (per-series colors at 3 series)

Values (19.2,21.4,16.7) + extra series (8.7,6.4,4.4) keep axis proportions.
"""
import os
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

OUT = r"pipeline_data\pptx_probes\chart_pie2\chart_pie2.pptx"


def add_pie(prs, series_names, has_title, has_legend):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    cd = CategoryChartData()
    cd.categories = ["East", "West", "Midwest"]
    for i, name in enumerate(series_names):
        cd.add_series(name, (19.2, 21.4, 16.7) if i == 0 else (8.7, 6.4, 4.4))
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(1), Inches(1), Inches(5.5), Inches(4.0), cd
    )
    ch = gframe.chart
    ch.has_title = has_title
    ch.has_legend = has_legend
    return ch


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    add_pie(prs, ["Series 1"], False, False)          # slide1
    add_pie(prs, ["Series 1"], False, True)           # slide2
    add_pie(prs, ["Series 1"], True, True)            # slide3
    add_pie(prs, ["Series 1", "Series 2", "Series 3"], False, False)  # slide4
    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print(f"Saved {OUT}")


if __name__ == "__main__":
    main()
