# -*- coding: utf-8 -*-
"""Spec #4e: first-line baseline offset vs (font, n) - control probe.

wave-4 (spec4d, Calibri) showed:
  - n != 1.0: first_baseline - text_area_top ~= 0.75 * (fs*1.2*n)   (ratio ~0.75)
  - n == 1.0: offset is +0.62pt deeper (ratio 0.7787)
Open question: is the 0.75 ratio font-independent (PowerPoint internal
"ascent=0.75" model) or Calibri-specific?  And is n=1.0 special for all
fonts?

Probe: Arial + Times New Roman x n in {0.5, 1.0, 2.0}, one multi-line
paragraph per slide (wave-4 methodology - single paragraph, intra-para
baseline deltas).
"""
import os
import sys

sys.stdout.reconfigure(encoding="utf-8")

sys.path.insert(0, r"c:\Users\ryuji\oxi-main\tools\metrics")

OUT_DIR = r"c:\Users\ryuji\oxi-main\pipeline_data\pptx_probes\spec4e_multifont"
OUT_PPTX = os.path.join(OUT_DIR, "spec4e_multifont.pptx")

WORDS = "lorem ipsum dolor sit amet consectetur adipiscing elit "
TEXT = (WORDS * 8).strip()

FONTS = ["Arial", "Times New Roman"]
NSWEEP = [0.5, 1.0, 2.0]

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except Exception as e:  # noqa
    print("python-pptx import failed: %s" % e)
    HAS_PPTX = False


def build():
    os.makedirs(OUT_DIR, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(12.5)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for font in FONTS:
        for n in NSWEEP:
            slide = prs.slides.add_slide(blank)
            tb = slide.shapes.add_textbox(
                Inches(1.0), Inches(1.0), Inches(11.0), Inches(6.0))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)
            p = tf.paragraphs[0]
            p.line_spacing = n
            run = p.add_run()
            run.text = TEXT
            run.font.size = Pt(18)
            run.font.name = font
            run.font.color.rgb = RGBColor(0, 0, 0)
    prs.save(OUT_PPTX)
    print("saved %s" % OUT_PPTX)


if __name__ == "__main__":
    build()
