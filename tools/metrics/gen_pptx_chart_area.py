# -*- coding: utf-8 -*-
"""Chart spec: AREA charts (<c:areaChart>).

Levers: grouping (standard / stacked / percentStacked), series count,
legend presence, explicit vs auto title.  Frame 72,72,396,288 throughout
so the geometry can be compared against the line-chart derivation."""
import sys, os

sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

base = r"pipeline_data\pptx_probes\chart_area"
os.makedirs(base, exist_ok=True)

CATS = ["Q1", "Q2", "Q3", "Q4"]
S1 = (19.2, 21.4, 16.7, 22.0)
S2 = (10.5, 11.2, 8.5, 12.3)
S3 = (5.1, 6.4, 7.2, 4.8)

ARMS = [
    ("A1 area  n=1 title",      XL_CHART_TYPE.AREA,             [S1],         False, True),
    ("A2 area  n=2",            XL_CHART_TYPE.AREA,             [S1, S2],     False, True),
    ("A3 area  n=2 legend",     XL_CHART_TYPE.AREA,             [S1, S2],     True,  True),
    ("A4 stack n=2",            XL_CHART_TYPE.AREA_STACKED,     [S1, S2],     False, True),
    ("A5 stack n=2 legend",     XL_CHART_TYPE.AREA_STACKED,     [S1, S2],     True,  True),
    ("A6 100%  n=2 legend",     XL_CHART_TYPE.AREA_STACKED_100, [S1, S2],     True,  True),
    ("A7 area  n=1 NOtitle",    XL_CHART_TYPE.AREA,             [S1],         False, False),
    ("A8 area  n=3 legend",     XL_CHART_TYPE.AREA,             [S1, S2, S3], True,  True),
]

prs = Presentation()
prs.slide_width = Pt(720)
prs.slide_height = Pt(540)
for label, kind, series, legend, title in ARMS:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = CATS
    for i, vals in enumerate(series, 1):
        cd.add_series(f"Ser{i}", vals)
    chart = slide.shapes.add_chart(
        kind, Pt(72), Pt(72), Pt(396), Pt(288), cd
    ).chart
    if legend:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
    if not title:
        chart.has_title = False

out = os.path.join(base, "chart_area.pptx")
prs.save(out)
for i, (label, *_rest) in enumerate(ARMS, 1):
    print(f"  S{i}: {label}")
print("\nsaved:", out, os.path.getsize(out))
