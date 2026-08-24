# -*- coding: utf-8 -*-
"""Author the EMBEDDED-ITALIC-PART SELECTION probe.

S-ITALADV is parked because PowerPoint does not always use the italic part a
deck embeds. Both d15 and d16 embed regular / bold / italic / boldItalic for the
family their body text asks for, and yet:

    d16  level-italic, regular run -> SourceSansPro-Italic       (embedded part)
    d16  level-italic, bold run    -> SourceSansPro-BlackItalic  (embedded part)
    d15  level-italic, bold        -> Barlow-Bold, SKEWED        (upright!)

The question this probe answers:

> Given a family with all four parts embedded, which combinations of (run bold,
> run italic, level bold, level italic) make PowerPoint use the embedded ITALIC
> part, and which make it skew the upright one instead?

The arms are authored INSIDE the corpus decks themselves rather than in a fresh
presentation, because the `.fntdata` parts are EOT (MicroType Express) and
cannot be written by hand -- PowerPoint only embeds fonts it has installed, and
neither Source Sans Pro nor Barlow is. Appending slides to the deck keeps its
`p:embeddedFontLst` and every part byte-for-byte, so the arms are asking the
same font resolver the corpus slides ask.

The reader takes each line's POSTSCRIPT name out of PowerPoint's PDF. The subset
family names are useless -- `43,Italic` is an INDEX, not a family
(`font-audit-three-sources`) -- but name id 6 is exact.

Usage:
    python tools/metrics/gen_pptx_italface.py
    python tools/metrics/export_pptx_italface.py    # PowerPoint COM -> PDF
    python tools/metrics/read_pptx_italface.py      # read the PDF back
"""
from __future__ import annotations

import json
import shutil
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

# These live under `pptx_derive`, NOT `pptx_probes`: each arm deck is a
# 50-slide copy of a corpus deck, and the probe corpus is a byte-identity
# gate that renders every deck twice. Five of these would add ~500 renders
# per gate run for no regression value the corpus does not already give.
REPO = Path(__file__).resolve().parents[2]
DEV = REPO / "pipeline_data" / "pptx_benchmark" / "dev" / "pptx"
OUT = REPO / "pipeline_data" / "pptx_derive" / "italface"

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
P = "http://schemas.openxmlformats.org/presentationml/2006/main"
EMU_IN = 914400
SIZE = 32
TEXT = "Hamburgefonstiv"

# (deck glob, the family whose four parts that deck embeds)
SUBJECTS = [("d16*", "Source Sans Pro"), ("d15*", "Barlow")]

# (name, run_bold, run_italic, lvl_bold, lvl_italic)
ARMS = [
    ("run_b0i0", 0, 0, None, None),
    ("run_b1i0", 1, 0, None, None),
    ("run_b0i1", 0, 1, None, None),
    ("run_b1i1", 1, 1, None, None),
    ("lvl_b0i0", None, None, 0, 0),
    ("lvl_b1i0", None, None, 1, 0),
    ("lvl_b0i1", None, None, 0, 1),
    ("lvl_b1i1", None, None, 1, 1),
    # Does it matter WHERE the flag comes from? d16's real case is the second.
    ("mix_lvlI_runB", 1, None, None, 1),
    ("mix_lvlB_runI", None, 1, 1, None),
]


def add_arm(prs, family, name, rb, ri, lb, li):
    slide = prs.slides.add_slide(prs.slide_layouts[-1])
    box = slide.shapes.add_textbox(Emu(EMU_IN // 2), Emu(EMU_IN), Emu(9 * EMU_IN), Emu(EMU_IN))
    tf = box.text_frame
    body = tf._txBody.find(f"{{{A}}}bodyPr")
    for child in list(body):
        body.remove(child)
    etree.SubElement(body, f"{{{A}}}noAutofit")
    # The shape's OWN level style carries the "level" flags.
    if lb is not None or li is not None:
        ls = tf._txBody.find(f"{{{A}}}lstStyle")
        if ls is None:
            ls = etree.SubElement(tf._txBody, f"{{{A}}}lstStyle")
            body.addnext(ls)
        lvl = etree.SubElement(ls, f"{{{A}}}lvl1pPr")
        d = etree.SubElement(lvl, f"{{{A}}}defRPr")
        d.set("sz", str(SIZE * 100))
        if lb is not None:
            d.set("b", str(lb))
        if li is not None:
            d.set("i", str(li))
        etree.SubElement(d, f"{{{A}}}latin").set("typeface", family)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"{name} {TEXT}"
    r.font.size = Pt(SIZE)
    r.font.name = family
    r.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    rpr = r._r.get_or_add_rPr()
    if rb is not None:
        rpr.set("b", str(rb))
    if ri is not None:
        rpr.set("i", str(ri))
    return {"arm": name, "family": family, "run_b": rb, "run_i": ri,
            "lvl_b": lb, "lvl_i": li, "text": f"{name} {TEXT}"}


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    index = []
    for glob, family in SUBJECTS:
        src = sorted(DEV.glob(f"{glob}.pptx"))[0]
        dst = OUT / f"italface_{family.replace(' ', '')}.pptx"
        shutil.copyfile(src, dst)
        prs = Presentation(str(dst))
        base = len(prs.slides.__iter__.__self__._sldIdLst)
        arms = []
        for i, (name, rb, ri, lb, li) in enumerate(ARMS, 1):
            rec = add_arm(prs, family, name, rb, ri, lb, li)
            rec["slide"] = base + i
            arms.append(rec)
        prs.save(str(dst))
        index.append({"pptx": dst.name, "family": family, "base": base, "arms": arms})
        print(f"wrote {dst}  ({base} original slides + {len(arms)} arms)")
    (OUT / "arms.json").write_text(json.dumps(index, indent=1), encoding="utf-8")
    print(f"wrote {OUT / 'arms.json'}")


if __name__ == "__main__":
    main()
