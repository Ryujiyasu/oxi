# -*- coding: utf-8 -*-
"""Probe: how does Word choose the value axis when the data contains NEGATIVE values?

Current Oxi behaviour (all chart types): the value axis starts at 0 and
nice_axis_max() picks the top.  A negative datum has no representation.

8 arms, one chart per slide, frame (72,72,396,288):
  N1 column clustered   mixed  (19.2, -8.5, 16.7)
  N2 column clustered   all negative (-19.2, -8.5, -16.7)
  N3 line + markers     mixed
  N4 scatter            Y mixed, X positive
  N5 scatter            X and Y both mixed (negative X)
  N6 area               mixed
  N7 bar (horizontal)   mixed
  N8 column stacked     mixed (two series)
"""
import os
from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt

OUT_DIR = os.path.join("pipeline_data", "pptx_probes", "chart_negative")
OUT = os.path.join(OUT_DIR, "chart_negative.pptx")

CATS = ["Q1", "Q2", "Q3"]
MIXED = (19.2, -8.5, 16.7)
ALLNEG = (-19.2, -8.5, -16.7)
MIXED2 = (10.5, -11.2, 8.5)

FRAME = (Pt(72), Pt(72), Pt(396), Pt(288))


def add_cat(slide, ctype, series, cats=CATS):
    cd = CategoryChartData()
    cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals)
    gf = slide.shapes.add_chart(ctype, *FRAME, cd)
    return gf.chart


def add_xy(slide, ctype, series):
    cd = XyChartData()
    for name, pts in series:
        s = cd.add_series(name)
        for x, y in pts:
            s.add_data_point(x, y)
    gf = slide.shapes.add_chart(ctype, *FRAME, cd)
    return gf.chart


def main() -> None:
    os.makedirs(OUT_DIR, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Pt(540)
    prs.slide_height = Pt(405)
    blank = prs.slide_layouts[6]

    # N1 column clustered, mixed
    add_cat(prs.slides.add_slide(blank), XL_CHART_TYPE.COLUMN_CLUSTERED,
            [("Ser1", MIXED)])
    # N2 column clustered, all negative
    add_cat(prs.slides.add_slide(blank), XL_CHART_TYPE.COLUMN_CLUSTERED,
            [("Ser1", ALLNEG)])
    # N3 line + markers, mixed
    add_cat(prs.slides.add_slide(blank), XL_CHART_TYPE.LINE_MARKERS,
            [("Ser1", MIXED)])
    # N4 scatter, Y mixed
    add_xy(prs.slides.add_slide(blank), XL_CHART_TYPE.XY_SCATTER,
           [("Ser1", [(1, 19.2), (2, -8.5), (3, 16.7)])])
    # N5 scatter, X and Y mixed
    add_xy(prs.slides.add_slide(blank), XL_CHART_TYPE.XY_SCATTER,
           [("Ser1", [(-2, 19.2), (-1, -8.5), (1, 16.7), (2, -4.0)])])
    # N6 area, mixed
    add_cat(prs.slides.add_slide(blank), XL_CHART_TYPE.AREA,
            [("Ser1", MIXED)])
    # N7 bar (horizontal), mixed
    add_cat(prs.slides.add_slide(blank), XL_CHART_TYPE.BAR_CLUSTERED,
            [("Ser1", MIXED)])
    # N8 column stacked, two series mixed
    add_cat(prs.slides.add_slide(blank), XL_CHART_TYPE.COLUMN_STACKED,
            [("Ser1", MIXED), ("Ser2", MIXED2)])

    prs.save(OUT)
    print("saved", OUT, len(prs.slides.__iter__.__self__._sldIdLst), "slides")


if __name__ == "__main__":
    main()
