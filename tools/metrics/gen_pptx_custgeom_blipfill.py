# -*- coding: utf-8 -*-
"""Probe: how does PowerPoint place an `a:blipFill` that sits on a shape?

The dev corpus has 2141 shape-level blipFills on 337 slides in 27 decks, and
EVERY one of them is on a custGeom shape. 191 declare a negative
`a:stretch/a:fillRect`, i.e. a destination that reaches outside the shape box
-- d28 (corpus floor 0.5217) is built out of them, and Oxi paints those images
2.45 x 6.75 times oversized, across the whole page.

Two questions, one export:
  (1) is the fill CLIPPED to the shape's outline, or drawn over its box?
  (2) does `a:fillRect` inset/expand the DESTINATION the way the p:pic path
      already models (l/t/r/b as fractions, negative = expand)?

The source image is a 2x2 colour grid (TL red, TR green, BL blue, TR yellow)
with a black frame, so a sampled point names the source quadrant it came from
and the mapping can be read off directly rather than inferred.

Arms (shape box 4in square, centred):
  D1 rect path,     fillRect 0            baseline: quadrants fill the box
  D2 triangle path, fillRect 0            the box corners answer question (1)
  D3 rect path,     fillRect r=-100%      destination twice as wide
  D4 triangle path, fillRect r=-100% b=-100%
  D5 rect path,     srcRect l=25% r=25%   source crop, for comparison
  D6 rect path,     fillRect r=-145.344% b=-574.764%   the literal d28 values
"""
from __future__ import annotations

import io
import sys
import zipfile
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\custgeom_blipfill").resolve()
SPACE = 1000
SHAPE_EMU = 3657600  # 4in
SHAPE_XY = (2743200, 1600200)


def grid_png() -> bytes:
    img = Image.new("RGB", (400, 400), "white")
    d = ImageDraw.Draw(img)
    d.rectangle([0, 0, 199, 199], fill=(255, 0, 0))
    d.rectangle([200, 0, 399, 199], fill=(0, 200, 0))
    d.rectangle([0, 200, 199, 399], fill=(0, 0, 255))
    d.rectangle([200, 200, 399, 399], fill=(255, 220, 0))
    d.rectangle([0, 0, 399, 399], outline=(0, 0, 0), width=8)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def path_xml(points: list[tuple[int, int]]) -> str:
    body = "".join(f'<a:lnTo><a:pt x="{x}" y="{y}"/></a:lnTo>' for x, y in points[1:])
    return (
        "<a:custGeom><a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>"
        '<a:rect b="b" l="l" r="r" t="t"/><a:pathLst>'
        f'<a:path extrusionOk="0" h="{SPACE}" w="{SPACE}">'
        f'<a:moveTo><a:pt x="{points[0][0]}" y="{points[0][1]}"/></a:moveTo>'
        f"{body}<a:close/></a:path></a:pathLst></a:custGeom>"
    )


RECT = [(0, 0), (SPACE, 0), (SPACE, SPACE), (0, SPACE)]
TRI = [(SPACE // 2, 0), (SPACE, SPACE), (0, SPACE)]


def blip_fill(rid: str, src: str = "", fill: str = "") -> str:
    src_el = f"<a:srcRect {src}/>" if src else ""
    fill_el = f"<a:fillRect {fill}/>" if fill else "<a:fillRect/>"
    return (
        f'<a:blipFill rotWithShape="1"><a:blip r:embed="{rid}"><a:alphaModFix/></a:blip>'
        f"{src_el}<a:stretch>{fill_el}</a:stretch></a:blipFill>"
    )


# (label, geometry, blipFill body)
ARMS = [
    ("D1 rect, fillRect 0", path_xml(RECT), blip_fill("RID")),
    ("D2 triangle, fillRect 0", path_xml(TRI), blip_fill("RID")),
    ("D3 rect, fillRect r=-100%", path_xml(RECT), blip_fill("RID", fill='b="0" l="0" r="-100000" t="0"')),
    ("D4 triangle, fillRect r=-100% b=-100%", path_xml(TRI),
     blip_fill("RID", fill='b="-100000" l="0" r="-100000" t="0"')),
    ("D5 rect, srcRect l=25% r=25%", path_xml(RECT), blip_fill("RID", src='l="25000" r="25000"')),
    ("D6 rect, the d28 fillRect", path_xml(RECT),
     blip_fill("RID", fill='b="-574764" l="0" r="-145344" t="0"')),
]


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    tmp, dst = OUT / "_stage.pptx", OUT / "custgeom_blipfill.pptx"
    png = grid_png()
    (OUT / "source.png").write_bytes(png)

    prs = Presentation()
    blank = prs.slide_layouts[6]
    for label, _, _ in ARMS:
        s = prs.slides.add_slide(blank)
        box = s.shapes.add_textbox(Emu(228600), Emu(228600), Emu(6400800), Emu(400050))
        box.text_frame.text = label
        # The picture is added only so python-pptx registers the media part and
        # a relationship on this slide; it is deleted below and its rId reused
        # by the shape fill.
        pic = s.shapes.add_picture(io.BytesIO(png), Emu(0), Emu(0), Emu(91440), Emu(91440))
        shape = s.shapes.add_shape(1, Emu(SHAPE_XY[0]), Emu(SHAPE_XY[1]),
                                   Emu(SHAPE_EMU), Emu(SHAPE_EMU))
        shape.line.fill.background()
        shape.shadow.inherit = False
        del pic  # keep the reference alive only until save
    prs.save(tmp)

    with zipfile.ZipFile(tmp) as zin:
        names = zin.namelist()
        data = {n: zin.read(n) for n in names}

    for i, (label, geom, fill) in enumerate(ARMS, start=1):
        part = f"ppt/slides/slide{i}.xml"
        xml = data[part].decode("utf-8")
        rels = data[f"ppt/slides/_rels/slide{i}.xml.rels"].decode("utf-8")
        rid = rels.split('Id="')[1].split('"')[0] if 'Id="' in rels else "rId2"
        for token in rels.split('Id="')[1:]:
            if "/image" in token:
                rid = token.split('"')[0]
                break
        # Drop the helper <p:pic> entirely: only the shape must carry ink.
        pic_start = xml.index("<p:pic>")
        pic_end = xml.index("</p:pic>") + len("</p:pic>")
        xml = xml[:pic_start] + xml[pic_end:]
        # rindex: the caption textbox has a prstGeom of its own and comes first.
        gs = xml.rindex("<a:prstGeom")
        ge = xml.rindex("</a:prstGeom>") + len("</a:prstGeom>")
        # The blipFill goes straight after the geometry, which is where spPr's
        # sequence puts a fill. python-pptx's add_shape writes NO explicit fill
        # (the shape inherits one through <p:style>), so there is nothing to
        # replace -- an spPr fill overrides the style reference.
        xml = xml[:gs] + geom + fill.replace("RID", rid) + xml[ge:]
        data[part] = xml.encode("utf-8")

    with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:
        for n in names:
            zout.writestr(n, data[n])
    tmp.unlink()
    print(f"wrote {dst}  ({len(ARMS)} arms)")


if __name__ == "__main__":
    main()
