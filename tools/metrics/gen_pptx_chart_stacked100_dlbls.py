# -*- coding: utf-8 -*-
"""100%-STACKED column chart data-labels probe.

Extends chart_stacked100 (COLUMN_STACKED_100, 3 categories x 2 series, same
data Revenue (19.2, 21.4, 16.7) / Cost (10.5, 15.0, 12.3)) with data labels
so Word's PDF exposes the exact label placement/format/font rules for the
percentStacked branch.

  S1: labels on, show_value=True, default position (no dLblPos)  + no numFmt
  S2: labels on, show_value=True, default position               + numFmt "0%"

python-pptx writes <c:dLbls> (START tag with children) when has_data_labels
is True; leaving dl.position unset means no <c:dLblPos> is written, so Word
applies its default (expected OUTSIDE_END / CENTER like stacked)."""
import sys, os
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

base = r"pipeline_data\pptx_probes\chart_stacked100_dlbls"
os.makedirs(base, exist_ok=True)

prs = Presentation()
blanks = [prs.slide_layouts[6] for _ in range(2)]


def add_stack100_slide(layout, num_fmt=None):
    slide = prs.slides.add_slide(layout)
    cd = CategoryChartData()
    cd.categories = ["Q1", "Q2", "Q3"]
    cd.add_series("Revenue", (19.2, 21.4, 16.7))
    cd.add_series("Cost", (10.5, 15.0, 12.3))
    x, y, cx, cy = Inches(1.0), Inches(1.0), Inches(5.5), Inches(4.0)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED_100, x, y, cx, cy, cd)
    plot = gframe.chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_value = True
    if num_fmt is not None:
        dl.number_format_is_linked = False
        dl.number_format = num_fmt
    return slide


add_stack100_slide(blanks[0])
add_stack100_slide(blanks[1], num_fmt="0%")

out = os.path.join(base, "chart_stacked100_dlbls.pptx")
prs.save(out)
print("saved:", out, os.path.getsize(out))

# Dump the c:dLbls XML of the first chart part to confirm what python-pptx wrote
import zipfile
with zipfile.ZipFile(out) as z:
    for n in z.namelist():
        if n.startswith("ppt/charts/") and n.endswith(".xml"):
            xml = z.read(n).decode("utf-8")
            i = xml.find("dLbls")
            if i >= 0:
                print(f"[{n}] ...{xml[max(0, i-120):i+320]!r}...")
