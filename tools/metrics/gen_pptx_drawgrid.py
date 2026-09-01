# -*- coding: utf-8 -*-
"""Author the DRAW-GRID probe: does the pen move by the exact advance or by 1/8pt?

`pptx-master-unit-break-law` established the master unit (1/8pt, 576 to the
inch) as the unit PowerPoint BREAKS on. The corpus says it is also the unit
PowerPoint DRAWS on -- d35 s25's 114-character line ends 2.461pt from exact
accumulation and 0.230pt from the grid -- but a corpus line carries wrapping,
alignment and possibly the line squeeze along with the question.

This asks it alone. Each arm is ONE run of ONE repeated glyph in a box far
wider than the text, `wrap="none"`, no autofit, no insets, left-aligned, so the
Nth character's x is N steps of a single advance and nothing else:

    exact   x[N] = N * em * sz
    grid    x[N] = N * round(em * sz * 8) / 8

Over 40 glyphs the two differ by up to 2.5pt, which no reading error can hide.
The faces are INSTALLED ones (no embedding, no cloud copy: see
`pptx-local-copy-beats-embedded`), and the sizes deliberately include
non-integer ones -- d30 s2 sets 15.99pt and its advances do NOT land on the
grid, so the size may be part of the rule rather than a detail.

A second paragraph per arm alternates two glyphs, so a kern pair would show as
a step that the single-glyph paragraph does not have.

Usage:
    python tools/metrics/gen_pptx_drawgrid.py
    python tools/metrics/export_pptx_drawgrid.py   # PowerPoint COM -> PDF
    python tools/metrics/read_pptx_drawgrid.py     # read the PDF back
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

REPO = Path(__file__).resolve().parents[2]
OUT = REPO / "pipeline_data" / "pptx_probes" / "drawgrid"
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
EMU_PT = 12700

# Installed faces the metrics tables also carry, so the probe's answer can be
# compared against the same advances the engine would use.
FACES = ["Arial", "Georgia", "Verdana"]

# Integer and non-integer, small and large. 15.99 is d30 s2's own size.
SIZES = [8.0, 10.0, 12.0, 12.5, 15.99, 18.0, 24.0, 32.0]

REPEAT = "n" * 40
ALTERNATE = "VA" * 20        # a pair most faces kern

ARMS = [(f, s) for f in FACES for s in SIZES]


def add_box(slide, y_pt: float, face: str, sz: float, text: str) -> None:
    """One left-aligned, non-wrapping, inset-free line of `text`."""
    box = slide.shapes.add_textbox(
        Emu(int(20 * EMU_PT)), Emu(int(y_pt * EMU_PT)),
        Emu(int(1200 * EMU_PT)), Emu(int(50 * EMU_PT)))
    tx = box.text_frame._txBody
    body = tx.find(f"{{{A}}}bodyPr")
    for k, v in (("lIns", "0"), ("rIns", "0"), ("tIns", "0"), ("bIns", "0"),
                 ("wrap", "none"), ("anchor", "t")):
        body.set(k, v)
    # ★`add_textbox` already puts `<a:spAutoFit/>` in the bodyPr, and the three
    # autofit elements are a CHOICE: appending `noAutofit` beside it makes a
    # file PowerPoint refuses to open at all -- with a COM error that names
    # nothing ("an exception occurred"), so it reads as a broken environment
    # rather than a broken document.
    for tag in ("spAutoFit", "normAutofit", "noAutofit"):
        for e in body.findall(f"{{{A}}}{tag}"):
            body.remove(e)
    etree.SubElement(body, f"{{{A}}}noAutofit")
    for para in tx.findall(f"{{{A}}}p"):
        tx.remove(para)
    p = etree.SubElement(tx, f"{{{A}}}p")
    ppr = etree.SubElement(p, f"{{{A}}}pPr")
    ppr.set("algn", "l")
    ppr.set("marL", "0")
    ppr.set("indent", "0")
    for tag in ("spcBef", "spcAft"):
        e = etree.SubElement(ppr, f"{{{A}}}{tag}")
        etree.SubElement(e, f"{{{A}}}spcPts").set("val", "0")
    etree.SubElement(ppr, f"{{{A}}}buNone")
    r = etree.SubElement(p, f"{{{A}}}r")
    rpr = etree.SubElement(r, f"{{{A}}}rPr")
    rpr.set("lang", "en-US")
    rpr.set("sz", str(int(round(sz * 100))))
    # ★Kerning off by declaration, not by hope: `kern` names the smallest size
    # kerning applies from, and 0 turns it off. The ALTERNATE paragraph is here
    # to show what that declaration is worth.
    rpr.set("kern", "0")
    etree.SubElement(rpr, f"{{{A}}}latin").set("typeface", face)
    etree.SubElement(r, f"{{{A}}}t").text = text


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    pres = Presentation()
    pres.slide_width = Emu(int(1280 * EMU_PT))
    pres.slide_height = Emu(int(400 * EMU_PT))
    blank = pres.slide_layouts[6]
    meta = []
    for i, (face, sz) in enumerate(ARMS):
        slide = pres.slides.add_slide(blank)
        add_box(slide, 40.0, face, sz, REPEAT)
        add_box(slide, 160.0, face, sz, ALTERNATE)
        meta.append({"slide": i + 1, "typeface": face, "sz_pt": sz,
                     "repeat": REPEAT, "alternate": ALTERNATE})
    path = OUT / "probe_drawgrid.pptx"
    pres.save(str(path))
    (OUT / "arms.json").write_text(json.dumps(meta, indent=1), encoding="utf-8")
    print(f"wrote {path} -- {len(meta)} arms")


if __name__ == "__main__":
    main()
