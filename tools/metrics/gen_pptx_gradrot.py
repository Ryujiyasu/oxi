# -*- coding: utf-8 -*-
"""Author the shape-gradient ORIENTATION probe.

`paint_shape_gradient` runs the ramp at `a:lin@ang` and ignores the shape's own
`rot` / `flipH` / `flipV`.  d06's layout band (`ang="5400012"`, `rot="10800000"`,
`flipH="1"`) comes out mirrored top-for-bottom against PowerPoint's own PDF --
same amplitude, exactly reversed -- which is what put this probe on the queue.

The question is which of the shape's own transforms the ramp rides:

  A  the ang -> screen-direction mapping with no transform at all (8 arms)
  B  ang x rot x rotWithShape (18 arms) -- `rotWithShape` is the attribute that
     is SUPPOSED to govern this, and its default is not the same in every
     reading of ECMA-376, so absent / "0" / "1" are all measured
  C  ang x flipH/flipV (6 arms)
  D  the d06 geometry itself: a wide band at ang 90, rot 180, flipH 1 (2 arms,
     with and without the flip)

Every arm is one shape on its own slide, filled BLACK at pos 0 and WHITE at
pos 100000 with no outline, so the reader can fit a plane to the ink and read
the ramp's direction straight off it.  Blocks A-C use a SQUARE shape so that
`rot` leaves the footprint alone and `scaled` cannot matter; block D uses the
real proportions.

Usage:
    python tools/metrics/gen_pptx_gradrot.py
    python tools/metrics/export_pptx_gradrot.py    # PowerPoint COM -> PDF
    python tools/metrics/read_pptx_gradrot.py      # read the PDF back
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

REPO = Path(__file__).resolve().parents[2]
OUT = REPO / "pipeline_data" / "pptx_probes" / "gradrot"

A = "http://schemas.openxmlformats.org/drawingml/2006/main"

EMU_IN = 914400
# Square block for A-C: 4.00in on a side, centred left-of-middle so the reader
# can find it without depending on the slide size.
SQ_X, SQ_Y, SQ_W, SQ_H = 2 * EMU_IN, EMU_IN, 4 * EMU_IN, 4 * EMU_IN
# Block D: d06's layout band proportions, scaled to fit the default 10x7.5in.
BD_X, BD_Y, BD_W, BD_H = 0, 3 * EMU_IN, 10 * EMU_IN, 2 * EMU_IN

DEG = 60000  # a:lin@ang units per degree


def grad_fill(ang_deg: float, rot_with_shape: str | None) -> etree._Element:
    """A black->white two-stop linear ramp at `ang_deg`."""
    gf = etree.Element(f"{{{A}}}gradFill")
    if rot_with_shape is not None:
        gf.set("rotWithShape", rot_with_shape)
    gs_lst = etree.SubElement(gf, f"{{{A}}}gsLst")
    for pos, rgb in ((0, "000000"), (100000, "FFFFFF")):
        gs = etree.SubElement(gs_lst, f"{{{A}}}gs")
        gs.set("pos", str(pos))
        etree.SubElement(gs, f"{{{A}}}srgbClr").set("val", rgb)
    lin = etree.SubElement(gf, f"{{{A}}}lin")
    lin.set("ang", str(int(round(ang_deg * DEG))))
    lin.set("scaled", "0")
    return gf


def add_arm(prs, name, geom, ang, rot=0, flip_h=False, flip_v=False, rws=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    x, y, w, h = geom
    sp = slide.shapes.add_shape(1, Emu(x), Emu(y), Emu(w), Emu(h))  # 1 = rect
    sp.line.fill.background()
    sp_pr = sp._element.spPr
    # Drop python-pptx's own solid fill / theme effects before the ramp goes in.
    for tag in ("solidFill", "gradFill", "noFill", "blipFill", "pattFill", "grpFill"):
        for el in sp_pr.findall(f"{{{A}}}{tag}"):
            sp_pr.remove(el)
    prst = sp_pr.find(f"{{{A}}}prstGeom")
    prst.addnext(grad_fill(ang, rws))
    xfrm = sp_pr.find(f"{{{A}}}xfrm")
    if rot:
        xfrm.set("rot", str(int(round(rot * DEG))))
    if flip_h:
        xfrm.set("flipH", "1")
    if flip_v:
        xfrm.set("flipV", "1")
    # A style reference would repaint the shape; strip it.
    style = sp._element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}style")
    if style is not None:
        sp._element.remove(style)
    return {
        "arm": name,
        "slide": len(prs.slides.__iter__.__self__._sldIdLst),
        "ang": ang,
        "rot": rot,
        "flipH": int(flip_h),
        "flipV": int(flip_v),
        "rotWithShape": rws,
        "box": [x, y, w, h],
    }


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    arms = []
    n = 0

    def emit(name, geom, **kw):
        nonlocal n
        n += 1
        rec = add_arm(prs, name, geom, **kw)
        rec["slide"] = n
        arms.append(rec)

    # A: the bare ang -> screen mapping.
    for ang in (0, 45, 90, 135, 180, 225, 270, 315):
        emit(f"A_ang{ang}", (SQ_X, SQ_Y, SQ_W, SQ_H), ang=ang)

    # B: rotation, with every reading of rotWithShape.
    for ang in (0, 90):
        for rot in (90, 180, 270):
            for rws in (None, "0", "1"):
                tag = "abs" if rws is None else rws
                emit(f"B_a{ang}_r{rot}_rws{tag}",
                     (SQ_X, SQ_Y, SQ_W, SQ_H), ang=ang, rot=rot, rws=rws)

    # C: the flips.
    for ang in (0, 90):
        for fh, fv in ((True, False), (False, True), (True, True)):
            emit(f"C_a{ang}_h{int(fh)}v{int(fv)}",
                 (SQ_X, SQ_Y, SQ_W, SQ_H), ang=ang, flip_h=fh, flip_v=fv)

    # D: d06's own case.
    emit("D_band_r180", (BD_X, BD_Y, BD_W, BD_H), ang=90, rot=180)
    emit("D_band_r180_h", (BD_X, BD_Y, BD_W, BD_H), ang=90, rot=180, flip_h=True)

    # E: the ORDER of the two transforms. Every corpus arm is rot=180, where
    # flip-then-rotate and rotate-then-flip agree, so blocks B-D cannot tell
    # them apart. These four can, and the prediction is written down first:
    #   flip-then-rotate (the DrawingML xfrm order)  a' = rot + f(ang)
    #   rotate-then-flip                             a' = f(ang + rot)
    # with f_H(x) = 180 - x and f_V(x) = -x.
    #   E1 ang45 rot90  flipH ->  225  vs   45
    #   E2 ang45 rot90  flipV ->   45  vs  225
    #   E3 ang45 rot270 flipH ->   45  vs  225
    #   E4 ang135 rot90 flipV ->  315  vs  135
    emit("E_a45_r90_h", (SQ_X, SQ_Y, SQ_W, SQ_H), ang=45, rot=90, flip_h=True)
    emit("E_a45_r90_v", (SQ_X, SQ_Y, SQ_W, SQ_H), ang=45, rot=90, flip_v=True)
    emit("E_a45_r270_h", (SQ_X, SQ_Y, SQ_W, SQ_H), ang=45, rot=270, flip_h=True)
    emit("E_a135_r90_v", (SQ_X, SQ_Y, SQ_W, SQ_H), ang=135, rot=90, flip_v=True)

    path = OUT / "probe_gradrot.pptx"
    prs.save(str(path))
    (OUT / "arms.json").write_text(json.dumps(arms, indent=1), encoding="utf-8")
    print(f"wrote {path}  {len(arms)} slides")
    print(f"wrote {OUT / 'arms.json'}")


if __name__ == "__main__":
    main()
