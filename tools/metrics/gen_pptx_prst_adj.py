# -*- coding: utf-8 -*-
"""Round 3 of the preset-geometry probe: is `adj` linear?

Rounds 1-2 measured DEFAULT adjustments only.  A census of the dev corpus shows
which values are actually overridden:

  ellipse    567   explicit adj:  0        -> defaults cover 100%
  roundRect   45   explicit adj:  0        -> defaults cover 100%
  homePlate  100   explicit adj: 14        -> adj=50000 x1, adj=32030 x13
  teardrop    48   explicit adj:  6        -> adj=100000 x6 (= the default)

The lone explicit homePlate adj=50000 renders like the default (measured
d = ss/2), which pins the default at 50000 and suggests d = ss*adj/100000.  That
reading is worth one deck before it is trusted for the 13 shapes at 32030:

  J1 homePlate adj=0        -> predicted d = 0      (a plain rectangle)
  J2 homePlate adj=32030    -> predicted d = 92.25  (the corpus value)
  J3 homePlate adj=50000    -> predicted d = 144    (= the default, a control)
  J4 homePlate adj=100000   -> predicted d = 288    (a triangle)
  J5 roundRect adj=32030    -> predicted r = 92.25  (same law, second preset)
  J6 teardrop  adj=100000   -> should equal the default arm exactly

All on a 396x288 box, so ss = 288 throughout.
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

OUT = os.path.abspath(r"pipeline_data\pptx_probes\prst_adj")

ARMS = [
    ("J1_homeplate_adj0", MSO_SHAPE.PENTAGON, 0.0),
    ("J2_homeplate_adj32030", MSO_SHAPE.PENTAGON, 0.32030),
    ("J3_homeplate_adj50000", MSO_SHAPE.PENTAGON, 0.50000),
    ("J4_homeplate_adj100000", MSO_SHAPE.PENTAGON, 1.00000),
    ("J5_roundrect_adj32030", MSO_SHAPE.ROUNDED_RECTANGLE, 0.32030),
    ("J6_teardrop_adj100000", MSO_SHAPE.TEAR, 1.00000),
]


def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(540)
    blank = prs.slide_layouts[6]
    for name, kind, adj in ARMS:
        sl = prs.slides.add_slide(blank)
        sh = sl.shapes.add_shape(kind, Pt(72), Pt(72), Pt(396), Pt(288))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(0x4F, 0x81, 0xBD)
        sh.line.fill.background()
        sh.adjustments[0] = adj
        sh.text_frame.text = ""
        tb = sl.shapes.add_textbox(Pt(500), Pt(20), Pt(210), Pt(30))
        tb.text_frame.text = name
    p = os.path.join(OUT, "prst_adj.pptx")
    prs.save(p)
    print("wrote %s  (%d slides)" % (p, len(ARMS)))


if __name__ == "__main__":
    main()
