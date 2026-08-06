# -*- coding: utf-8 -*-
"""Chart spec (item 8) wave-3 repro: LEGEND-DECLARING multi-series
clustered-column chart on a blank slide, default theme. Same data as chart2
(2 series x 3 categories) but with <c:legend> enabled via python-pptx's
has_legend, so Word's legend drawing (swatch + series names) can be measured
against the PDF render (fitz get_drawings + rawdict)."""
import sys, os
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

base = r"pipeline_data\pptx_probes\chart_legend"
os.makedirs(base, exist_ok=True)

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

chart_data = CategoryChartData()
chart_data.categories = ["Q1", "Q2", "Q3"]
chart_data.add_series("Revenue", (19.2, 21.4, 16.7))
chart_data.add_series("Cost", (10.5, 15.0, 12.3))

x, y, cx, cy = Inches(1.0), Inches(1.0), Inches(5.5), Inches(4.0)
gframe = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
gframe.chart.has_legend = True

out = os.path.join(base, "chart_legend.pptx")
prs.save(out)
print("saved:", out, os.path.getsize(out))
