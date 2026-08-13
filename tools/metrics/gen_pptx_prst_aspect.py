# -*- coding: utf-8 -*-
"""Round 2 of the preset-geometry probe: separate `ss` from `w` and `h`.

Round 1 measured a single 396x288 box, where the derived quantities collide:
  roundRect  corner r   = 48   = min(w,h)*0.16667  = h*0.16667  = w*0.1212
  homePlate  point inset= 144  = min(w,h)/2        = h/2        = w*0.3636
  teardrop   tip                reaches the corner exactly

ECMA-376 writes these against `ss` (the SHORTER side), so a portrait arm (h>w)
and a flat arm both have to be measured before any formula can be claimed.

Arms, all at 72,72:
  A1 roundRect  396x288  landscape (round-1 repeat, as a control)
  A2 roundRect  288x396  portrait   -- ss is now w
  A3 roundRect  396x108  flat       -- ss much smaller than w
  A4 homePlate  396x288  landscape (control)
  A5 homePlate  288x396  portrait
  A6 homePlate  396x108  flat
  A7 teardrop   396x288  landscape (control)
  A8 teardrop   288x396  portrait
  A9 ellipse    288x396  portrait   -- confirm the round-1 ellipse law holds
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

OUT = os.path.abspath(r"pipeline_data\pptx_probes\prst_aspect")

ARMS = [
    ("A1_roundrect_land", MSO_SHAPE.ROUNDED_RECTANGLE, 396, 288),
    ("A2_roundrect_port", MSO_SHAPE.ROUNDED_RECTANGLE, 288, 396),
    ("A3_roundrect_flat", MSO_SHAPE.ROUNDED_RECTANGLE, 396, 108),
    ("A4_homeplate_land", MSO_SHAPE.PENTAGON, 396, 288),
    ("A5_homeplate_port", MSO_SHAPE.PENTAGON, 288, 396),
    ("A6_homeplate_flat", MSO_SHAPE.PENTAGON, 396, 108),
    ("A7_teardrop_land", MSO_SHAPE.TEAR, 396, 288),
    ("A8_teardrop_port", MSO_SHAPE.TEAR, 288, 396),
    ("A9_ellipse_port", MSO_SHAPE.OVAL, 288, 396),
]


def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(540)
    blank = prs.slide_layouts[6]
    for name, kind, w, h in ARMS:
        sl = prs.slides.add_slide(blank)
        sh = sl.shapes.add_shape(kind, Pt(72), Pt(72), Pt(w), Pt(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(0x4F, 0x81, 0xBD)
        sh.line.fill.background()
        sh.text_frame.text = ""
        tb = sl.shapes.add_textbox(Pt(500), Pt(20), Pt(210), Pt(30))
        tb.text_frame.text = name
    p = os.path.join(OUT, "prst_aspect.pptx")
    prs.save(p)
    print("wrote %s  (%d slides)" % (p, len(ARMS)))


if __name__ == "__main__":
    main()
