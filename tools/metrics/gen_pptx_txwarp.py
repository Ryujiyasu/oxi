# -*- coding: utf-8 -*-
"""Probe: how does PowerPoint size text that carries a:prstTxWarp?

The corpus says a `textPlain` warp scales the text to the shape — d35 s4's ink
width matches its box width to a pixel — but the deck's font is an embedded
Oswald whose digits PowerPoint outlined, so its glyph metrics cannot be read
back. This probe uses INSTALLED faces so every metric is available locally.

Arms vary the box aspect and the text length at a fixed face, plus two more
faces at one box, all with no `sz` anywhere (as the corpus shapes have).

Usage:
    python tools/metrics/gen_pptx_txwarp.py
    python tools/metrics/measure_pptx_word.py pipeline_data/pptx_probes/txwarp/txwarp.pptx pipeline_data/pptx_probes/txwarp
    python tools/metrics/read_pptx_txwarp_probe.py
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\txwarp").resolve()
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
PT = 12700

# (label, face, text, x, y, w, h) in points
ARMS = [
    ("tall1", "Arial", "1", 60, 60, 60, 300),
    ("wide1", "Arial", "1", 200, 60, 150, 300),
    ("short1", "Arial", "1", 420, 60, 60, 120),
    ("wordH", "Arial", "HELLO", 60, 400, 300, 90),
    ("faceA", "Segoe UI", "1", 420, 250, 60, 200),
    ("faceB", "Comic Sans MS", "1", 540, 250, 90, 200),
]


def q(tag: str) -> str:
    return f"{{{A}}}{tag}"


def add_warp(slide, face: str, text: str, x, y, w, h, autoshape: bool = False):
    if autoshape:
        from pptx.enum.shapes import MSO_SHAPE
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(int(x * PT)), Emu(int(y * PT)),
            Emu(int(w * PT)), Emu(int(h * PT))
        )
        # strip the theme fill/line so only the text shows
        sp_pr = box._element.spPr
        for tag in ("solidFill", "ln"):
            for el in sp_pr.findall(q(tag)):
                sp_pr.remove(el)
        etree.SubElement(sp_pr, q("noFill"))
        ln = etree.SubElement(sp_pr, q("ln"))
        etree.SubElement(ln, q("noFill"))
    else:
        box = slide.shapes.add_textbox(
            Emu(int(x * PT)), Emu(int(y * PT)), Emu(int(w * PT)), Emu(int(h * PT))
        )
    body = box.text_frame._txBody
    bodypr = body.find(q("bodyPr"))
    for ins in ("lIns", "tIns", "rIns", "bIns"):
        bodypr.set(ins, "0")
    etree.SubElement(bodypr, q("prstTxWarp")).set("prst", "textPlain")
    for pel in body.findall(q("p")):
        body.remove(pel)
    p = etree.SubElement(body, q("p"))
    ppr = etree.SubElement(p, q("pPr"))
    ppr.set("algn", "ctr")
    r = etree.SubElement(p, q("r"))
    rpr = etree.SubElement(r, q("rPr"))
    rpr.set("lang", "en-US")
    rpr.set("kern", "0")
    # an autoshape's text inherits the theme's light colour, which is invisible
    # on the probe's white page, so the colour is stated (the SIZE never is)
    fill = etree.SubElement(rpr, q("solidFill"))
    etree.SubElement(fill, q("srgbClr")).set("val", "000000")
    etree.SubElement(rpr, q("latin")).set("typeface", face)
    etree.SubElement(r, q("t")).text = text


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    manifest = []
    # slide 1: text boxes; slide 2: the same boxes as AUTOSHAPES, which is what
    # the corpus shapes are (`<p:cNvSpPr/>` with a prstGeom, no txBox="1")
    for slide_no, autoshape in ((1, False), (2, True)):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for label, face, text, x, y, w, h in ARMS:
            add_warp(slide, face, text, x, y, w, h, autoshape)
            manifest.append({
                "slide": slide_no,
                "label": label + ("_shape" if autoshape else "_box"),
                "face": face, "text": text,
                "x": x, "y": y, "w": w, "h": h,
            })
            print(f"s{slide_no} {label}{'(shape)' if autoshape else '(box)'}: "
                  f"{face} {text!r} box {w}x{h}")
    prs.save(OUT / "txwarp.pptx")
    (OUT / "txwarp_manifest.json").write_text(
        json.dumps(manifest, indent=1), encoding="utf-8")
    print(f"wrote {OUT / 'txwarp.pptx'}")


if __name__ == "__main__":
    main()
