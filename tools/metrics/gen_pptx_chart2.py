# -*- coding: utf-8 -*-
"""Chart spec (item 8) wave-2 repro: MULTI-SERIES clustered-column chart on a
blank slide, default theme. 3 categories x 2 series (with a THIRD empty-ish
series kept small so the cluster layout and per-series accent colours can be
measured against Word's PDF render (fitz get_drawings + rawdict)."""
import sys, os
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

base = r"pipeline_data\pptx_probes\chart2"
os.makedirs(base, exist_ok=True)

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

chart_data = CategoryChartData()
chart_data.categories = ["Q1", "Q2", "Q3"]
chart_data.add_series("Revenue", (19.2, 21.4, 16.7))
chart_data.add_series("Cost", (10.5, 15.0, 12.3))

x, y, cx, cy = Inches(1.0), Inches(1.0), Inches(5.5), Inches(4.0)
gframe = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

out = os.path.join(base, "chart2.pptx")
prs.save(out)
print("saved:", out, os.path.getsize(out))
