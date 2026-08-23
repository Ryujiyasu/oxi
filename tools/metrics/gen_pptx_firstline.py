# -*- coding: utf-8 -*-
"""Probe: where does PowerPoint put the FIRST baseline when lnSpc is not 100%?

Oxi places it at `0.75 * fs * 1.2 * n` for any `n != 1` and at the face's own
`1.2 * asc / (asc + desc) * fs` when n == 1 -- two unrelated models either side
of a single value of n. d10 slide 11 says the first is wrong: a 40pt Mali title
with `lnSpc 90%` in a box whose top is 159.250pt has PowerPoint's first
baseline at 193.270 (offset **34.020**), where Oxi computes 0.75 x 40 x 1.2 x
0.9 = 32.400. The PITCH is right (both step 43.20 = 40 x 1.2 x 0.9), so only
the first line is misplaced.

Each arm is one (font, lnSpc, size) with THREE lines in a fixed box: `anchor=t`,
every inset 0 and `<a:noAutofit/>`, so the box top IS the text-area top and the
first baseline's offset from it can be read straight out of the PDF. The fonts
are installed ones with known metrics, so the reader can test a model against
`asc`, `desc` and the line gap rather than curve-fitting.

Usage:
    python tools/metrics/gen_pptx_firstline.py
    python tools/metrics/export_pptx_firstline.py    # PowerPoint COM -> PDF
    python tools/metrics/read_pptx_firstline.py      # read the PDF back
"""
from __future__ import annotations

import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

REPO = Path(__file__).resolve().parents[2]
OUT = REPO / "pipeline_data" / "pptx_probes" / "firstline"

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
BOX_X = Emu(457200)          # 0.50in
BOX_Y = Emu(1828800)         # 2.00in  = 144pt, the reader's origin
BOX_W = Emu(8229600)         # 9.00in
BOX_H = Emu(5486400)         # 6.00in, tall enough for three lines at 144pt
LINES = 3

# (font, lnSpc percent, size pt). The percents bracket 100 on both sides; the
# extra Arial sizes test whether the offset is linear in the size.
# The last four bracket the disputed band. d39's 58 -> 144pt Bebas Neue titles
# (face 0.8769, descent 0.2678 of the box) take the QUARTER at lnSpc 93%, while
# Segoe Script (0.8249, descent 0.3126) keeps its own -- so the flip, if it is a
# threshold at all, sits between them: Courier New 0.2651, Segoe Print 0.2841,
# Lucida Sans Unicode 0.2863, MV Boli 0.2930.
# Segoe Script is the DISCRIMINATOR: its 1.2*asc/(asc+desc) is 0.83, below the
# 0.9 that `0.75 * 1.2 * n` takes at n == 1, so it is the only arm that can say
# whether that term is a floor at 100% or only away from it. d04 slide 1 needs
# the answer -- its 58pt Satisfy title reads 0.7877 and PowerPoint puts the
# baseline there, 6.5pt above where a floor would.
ARMS = [
    (font, pct, 40)
    for font in ("Arial", "Calibri", "Verdana", "Georgia", "Segoe Script", "Comic Sans MS",
                 "Courier New", "Segoe Print", "MV Boli", "Lucida Sans Unicode")
    for pct in (100, 90, 80, 70)
] + [("Arial", 90, 20), ("Arial", 90, 60), ("Arial", 120, 40), ("Calibri", 120, 40),
     ("Segoe Script", 120, 40), ("Segoe Script", 110, 40), ("Segoe Script", 95, 40),
     # d39's title is 143.64pt Bebas Neue at 93.01%, and its measured offset is
     # the QUARTER where Courier New at 40pt and almost the same descent share
     # is not -- so SIZE and the odd percent are the two variables left.
     ("Segoe Script", 93, 40), ("Courier New", 93, 40),
     ("Segoe Script", 93, 144), ("Courier New", 93, 144),
     ("Segoe Script", 90, 144), ("Courier New", 90, 144),
     ("Arial", 90, 144)]


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for font, pct, size in ARMS:
        slide = prs.slides.add_slide(blank)
        # A caption OUTSIDE the measured box so the reader can name the arm.
        cap = slide.shapes.add_textbox(BOX_X, Emu(228600), BOX_W, Emu(300000))
        cap.text_frame.text = f"arm {font}|{pct}|{size}"
        box = slide.shapes.add_textbox(BOX_X, BOX_Y, BOX_W, BOX_H)
        tf = box.text_frame
        body = tf._txBody.find(f"{{{A}}}bodyPr")
        for child in list(body):
            body.remove(child)
        body.set("anchor", "t")
        body.set("anchorCtr", "0")
        body.set("wrap", "square")
        for key in ("lIns", "rIns", "tIns", "bIns"):
            body.set(key, "0")
        body.append(etree.SubElement(body, f"{{{A}}}noAutofit"))
        for index in range(LINES):
            para = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
            ppr = para._pPr if para._pPr is not None else para._p.get_or_add_pPr()
            spc = etree.SubElement(ppr, f"{{{A}}}lnSpc")
            pct_el = etree.SubElement(spc, f"{{{A}}}spcPct")
            pct_el.set("val", str(pct * 1000))
            ppr.insert(0, spc)
            run = para.add_run()
            run.text = f"Hxg{index + 1}"
            run.font.size = Pt(size)
            run.font.name = font
    path = OUT / "probe_firstline.pptx"
    prs.save(str(path))
    print(f"wrote {path}  {len(ARMS)} arms, box top {int(BOX_Y) / 12700:.1f}pt")


if __name__ == "__main__":
    main()
