# -*- coding: utf-8 -*-
"""Probe: does Oxi put on exactly the pen PowerPoint states for a faked bold?

`pptx_fauxbold_stroke.py` reads the amount off PowerPoint's own PDF -- text
rendering mode 2 with a pen of `size/35`. This probe is the other half: it puts
the same string on the page twice, once PLAIN and once BOLD, in a family that
carries no bold, at four sizes, so the two questions can be told apart:

  the PLAIN arm is the CONTROL. It measures how much ink Oxi's rasteriser lays
  down for glyphs PowerPoint drew with no synthesis at all, i.e. the constant
  bias between the two rasterisers. Without it a 4% ink shortfall on the bold
  arm cannot be read -- it might be the pen, or it might be antialiasing, and
  they call for opposite fixes ([[measurement_instrument_traps]]: do not
  interpret the difference until the control arm agrees with the truth).

  the BOLD arm is the TEST. Its ink minus the control's, against PowerPoint's
  own bold-minus-plain, is the pen and nothing else.

    python tools/metrics/gen_pptx_fauxpen.py
    python tools/metrics/export_pptx_fauxpen.py    # PowerPoint COM -> PDF
    python tools/metrics/read_pptx_fauxpen.py

★Never export while the renderer is producing PNGs
([[pptx_com_render_must_not_overlap]]), and `embedTrueTypeFonts="true"` is
required or the whole `embeddedFontLst` is ignored.
"""
from __future__ import annotations

import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
from gen_pptx_boldslot import FAMILY, source_parts  # noqa: E402

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\fauxpen").resolve()
TEXT = "Handgloves"
SIZES = (12, 24, 48, 96)


def build(dst: Path, bold: bool, regular: bytes) -> None:
    import re
    import zipfile

    stage = dst.with_suffix(".stage.pptx")
    prs = Presentation()
    for size in SIZES:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        box = s.shapes.add_textbox(Emu(457200), Emu(1371600), Emu(8229600), Emu(2000000))
        tf = box.text_frame
        tf.word_wrap = False
        r = tf.paragraphs[0].add_run()
        r.text = TEXT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = FAMILY
    prs.save(stage)

    with zipfile.ZipFile(stage) as zin:
        data = {n: zin.read(n) for n in zin.namelist()}
    stage.unlink()

    data["ppt/fonts/font1.fntdata"] = regular
    rels = data["ppt/_rels/presentation.xml.rels"].decode("utf-8")
    data["ppt/_rels/presentation.xml.rels"] = rels.replace(
        "</Relationships>",
        '<Relationship Id="rIdFont1" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/font" '
        'Target="fonts/font1.fntdata"/></Relationships>',
    ).encode("utf-8")

    pres = data["ppt/presentation.xml"].decode("utf-8")
    pres = re.sub(r"<p:presentation ", '<p:presentation embedTrueTypeFonts="true" ', pres, count=1)
    lst = (f'<p:embeddedFontLst><p:embeddedFont><p:font typeface="{FAMILY}"/>'
           f'<p:regular r:id="rIdFont1"/></p:embeddedFont></p:embeddedFontLst>')
    if "<p:defaultTextStyle>" in pres:
        pres = pres.replace("<p:defaultTextStyle>", lst + "<p:defaultTextStyle>", 1)
    else:
        pres = pres.replace("</p:presentation>", lst + "</p:presentation>", 1)
    data["ppt/presentation.xml"] = pres.encode("utf-8")

    ct = data["[Content_Types].xml"].decode("utf-8")
    if 'Extension="fntdata"' not in ct:
        ct = ct.replace("<Default", '<Default Extension="fntdata" '
                        'ContentType="application/x-fontdata"/><Default', 1)
    data["[Content_Types].xml"] = ct.encode("utf-8")

    with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:
        for n in list(data):
            zout.writestr(n, data[n])


def main() -> None:
    shutil.rmtree(OUT, ignore_errors=True)
    OUT.mkdir(parents=True, exist_ok=True)
    regular, _ = source_parts()
    build(OUT / "plain.pptx", False, regular)
    build(OUT / "bold.pptx", True, regular)
    print(f"wrote {OUT} -- plain.pptx and bold.pptx "
          f"({FAMILY}, {TEXT!r}, sizes {', '.join(str(s) for s in SIZES)}pt)")


if __name__ == "__main__":
    main()
