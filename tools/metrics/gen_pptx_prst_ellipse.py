# -*- coding: utf-8 -*-
"""Probe deck for preset geometry: ellipse (and its neighbours).

`a:prstGeom` is drawn as a plain rectangle by the renderer today, and `ellipse`
is by far the most common non-rect preset in the dev corpus (497 shapes).  No
corpus slide carries an ellipse ALONE, so the PDF vectors there can't be
attributed -- hence this purpose-built deck: one preset per slide, isolated, so
PowerPoint's exported paths belong to it unambiguously.

Levers, one per slide (frame 72,72 unless noted):
  E1  ellipse    396x288  fill only, no line     -- the base geometry
  E2  ellipse    396x288  fill + 3pt line        -- where the stroke sits
  E3  ellipse    288x288  (a circle)             -- aspect independence
  E4  ellipse    396x108  wide and flat          -- aspect independence
  E5  ellipse    396x288  rot 30                 -- rotation about the centre
  E6  ellipse    396x288  flipH                  -- flip is a no-op on an ellipse
  E7  ellipse    396x288  + text                 -- text rect vs the box
  E8  roundRect  396x288  fill + line            -- adj (corner radius) default
  E9  homePlate  396x288  fill                   -- second most common
  E10 teardrop   396x288  fill                   -- third

Run:  python tools/metrics/gen_pptx_prst_ellipse.py
Then: PowerPoint COM -> PDF, and read the vector paths with fitz.
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

OUT = os.path.abspath(r"pipeline_data\pptx_probes\prst_ellipse")

# (name, MSO_SHAPE, left, top, w, h, rot, flip_h, line_pt, text)
ARMS = [
    ("E1_ellipse_fill", MSO_SHAPE.OVAL, 72, 72, 396, 288, 0, False, None, None),
    ("E2_ellipse_line", MSO_SHAPE.OVAL, 72, 72, 396, 288, 0, False, 3.0, None),
    ("E3_circle", MSO_SHAPE.OVAL, 72, 72, 288, 288, 0, False, None, None),
    ("E4_ellipse_flat", MSO_SHAPE.OVAL, 72, 72, 396, 108, 0, False, None, None),
    ("E5_ellipse_rot30", MSO_SHAPE.OVAL, 72, 72, 396, 288, 30, False, None, None),
    ("E6_ellipse_fliph", MSO_SHAPE.OVAL, 72, 72, 396, 288, 0, True, None, None),
    ("E7_ellipse_text", MSO_SHAPE.OVAL, 72, 72, 396, 288, 0, False, None, "Mg"),
    ("E8_roundrect", MSO_SHAPE.ROUNDED_RECTANGLE, 72, 72, 396, 288, 0, False, 3.0, None),
    ("E9_homeplate", MSO_SHAPE.PENTAGON, 72, 72, 396, 288, 0, False, None, None),
    ("E10_teardrop", MSO_SHAPE.TEAR, 72, 72, 396, 288, 0, False, None, None),
]


def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(540)
    blank = prs.slide_layouts[6]
    for name, kind, x, y, w, h, rot, flip, lw, txt in ARMS:
        sl = prs.slides.add_slide(blank)
        sh = sl.shapes.add_shape(kind, Pt(x), Pt(y), Pt(w), Pt(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(0x4F, 0x81, 0xBD)
        if lw is None:
            sh.line.fill.background()
        else:
            sh.line.color.rgb = RGBColor(0xC0, 0x50, 0x4D)
            sh.line.width = Pt(lw)
        if rot:
            sh.rotation = rot
        if flip:
            # python-pptx has no flip setter; write the attribute directly.
            sh._element.spPr.xfrm.set("flipH", "1")
        tf = sh.text_frame
        tf.text = txt if txt else ""
        # A label OUTSIDE the shape identifies the arm in the PDF.
        tb = sl.shapes.add_textbox(Pt(500), Pt(20), Pt(200), Pt(30))
        tb.text_frame.text = name
    p = os.path.join(OUT, "prst_ellipse.pptx")
    prs.save(p)
    print("wrote %s  (%d slides)" % (p, len(ARMS)))


if __name__ == "__main__":
    main()
