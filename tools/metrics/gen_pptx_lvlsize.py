# -*- coding: utf-8 -*-
"""Minimal repro: which size does a run that declares none resolve to?

Each arm is one paragraph of three runs -- silent, `sz`-declared, silent -- in a
placeholder whose LEVEL declares its own size. If the silent runs resolve to
the level, PowerPoint reports the level's number for runs 1 and 3 and the
declared one for run 2; if `paragraph_font_size`'s model were right, all three
would come back at the declared size (the paragraph's largest).

Arms vary what could plausibly matter: the placeholder kind, whether the
sibling's size is smaller or larger than the level's, which run declares, and
two controls (all silent, all declared).

    python tools/metrics/gen_pptx_lvlsize.py
    python tools/metrics/read_pptx_lvlsize_com.py
"""
from __future__ import annotations

import copy
import json
import os

from pptx import Presentation
from pptx.util import Emu, Pt

OUT = os.path.join("tools", "metrics", "lvlsize.pptx")
PLAN = os.path.join("tools", "metrics", "lvlsize.json")

# (label, layout index, placeholder index, [run sizes in pt, None = silent])
# Layout 0 of the default template is "Title Slide" (ctrTitle 44pt + subtitle
# 32pt); layout 1 is "Title and Content" (title 44pt + body 28/24/20pt).
CASES = [
    ("title_small_mid", 0, 0, [None, 24, None]),
    ("title_large_mid", 0, 0, [None, 60, None]),
    ("title_first_declares", 0, 0, [24, None, None]),
    ("title_last_declares", 0, 0, [None, None, 24]),
    ("body_small_mid", 1, 1, [None, 12, None]),
    ("body_large_mid", 1, 1, [None, 40, None]),
    ("subtitle_small_mid", 0, 1, [None, 12, None]),
    ("all_silent", 0, 0, [None, None, None]),
    ("all_declared", 0, 0, [20, 24, 28]),
]

WORDS = ["alpha", "beta", "gamma"]


def main() -> None:
    pres = Presentation()
    plan = []
    for label, layout, ph_idx, sizes in CASES:
        slide = pres.slides.add_slide(pres.slide_layouts[layout])
        ph = slide.placeholders[ph_idx]
        tf = ph.text_frame
        para = tf.paragraphs[0]
        for word, size in zip(WORDS, sizes):
            run = para.add_run()
            run.text = word + " "
            if size is not None:
                run.font.size = Pt(size)
        plan.append({"slide": len(pres.slides), "label": label,
                     "ph_idx": ph_idx, "sizes": sizes})
    pres.save(OUT)
    with open(PLAN, "w", encoding="utf-8") as fh:
        json.dump(plan, fh, indent=1)
    print("wrote %s (%d slides) and %s" % (OUT, len(pres.slides), PLAN))


if __name__ == "__main__":
    main()
