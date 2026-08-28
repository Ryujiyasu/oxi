# -*- coding: utf-8 -*-
"""Probe: what does PowerPoint do with a fractional fixed line spacing?

`<a:lnSpc><a:spcPts val="3220"/>` is 32.20pt and Oxi sets exactly that. The blind
corpus says PowerPoint does not: 31 s24 renders **32.000**, 31 s12's 33.59 renders
34.00, 36/37's 37.79 renders 37.99 and their 29.40 renders 29.03 -- while every
integer value (105.00, 28.00) is honoured to 0.005pt. Rounding to the nearest
whole point explains all of them; flooring and half-point rounding do not.

Each arm is one slide with a single 6-line frame at a different fraction, so the
tie at .5 (half-up against half-even) is decided too:

    20.00 control   20.10 20.40 20.49   20.50 <- the tie   20.51 20.60 20.90
    21.50           33.59 (the corpus value)              120.50 (large tie)

    python tools/metrics/gen_pptx_lnspcround.py
    python tools/metrics/export_pptx_lnspcround.py     # PowerPoint COM -> PDF
    python tools/metrics/read_pptx_lnspcround.py

Arial at 12pt, `wrap="none"`, no space-before/after, so the pitch is the only
thing the arm can move.
"""
from __future__ import annotations

import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\lnspcround").resolve()
ARMS = [2000, 2010, 2040, 2049, 2050, 2051, 2060, 2090, 2150, 3359, 12050]
LINES = 6
SIZE = 12


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    tmp, dst = OUT / "_stage.pptx", OUT / "lnspcround.pptx"
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for val in ARMS:
        s = prs.slides.add_slide(blank)
        cap = s.shapes.add_textbox(Emu(228600), Emu(114300), Emu(6400800), Emu(300000))
        cap.text_frame.text = f"spcPts {val}"
        box = s.shapes.add_textbox(Emu(457200), Emu(914400), Emu(7772400), Emu(4114800))
        tf = box.text_frame
        tf.word_wrap = False
        for i in range(LINES):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run()
            r.text = f"Handgloves line {i + 1}"
            r.font.size = Pt(SIZE)
            r.font.name = "Arial"
    prs.save(tmp)

    with zipfile.ZipFile(tmp) as zin:
        names = zin.namelist()
        data = {n: zin.read(n) for n in names}
    for i, val in enumerate(ARMS, start=1):
        part = f"ppt/slides/slide{i}.xml"
        xml = data[part].decode("utf-8")
        ppr = (f'<a:pPr><a:lnSpc><a:spcPts val="{val}"/></a:lnSpc>'
               '<a:spcBef><a:spcPts val="0"/></a:spcBef>'
               '<a:spcAft><a:spcPts val="0"/></a:spcAft>')
        # The caption box has no run properties to disturb; only the probe frame
        # carries <a:p> elements with runs, and every one of them gets the value.
        xml = xml.replace("<a:pPr>", ppr).replace("<a:p><a:r>", f"<a:p>{ppr}</a:pPr><a:r>")
        data[part] = xml.encode("utf-8")
    with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:
        for n in names:
            zout.writestr(n, data[n])
    tmp.unlink()
    print(f"wrote {dst}  ({len(ARMS)} arms)")


if __name__ == "__main__":
    main()
