# -*- coding: utf-8 -*-
"""Author the vertical-anchor overflow probe.

`compute_shape_anchor_off` clamps the BOTTOM anchor's offset at zero, so a text
block taller than its box pins to the box TOP. The CENTRE anchor was measured
in 2026-08 and does NOT clamp -- a tall block still centres on the box and
overflows equally. This probe asks the same question of `t`, `ctr` and `b`
across blocks that fit and blocks that do not.

One shape per slide so nothing overlaps: a fixed 1.00in box at a known offset,
`<a:noAutofit/>` so PowerPoint cannot resize it, and N lines of 32pt text
(line pitch 38.4pt, so 2 lines already overflow a 72pt box). Each line carries
its own marker text so the PDF reader can tell them apart.

Usage:
    python tools/metrics/gen_pptx_anchorb.py
    python tools/metrics/export_pptx_anchorb.py     # PowerPoint COM -> PDF
    python tools/metrics/read_pptx_anchorb.py       # read the PDF back
"""
from __future__ import annotations

import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

REPO = Path(__file__).resolve().parents[2]
OUT = REPO / "pipeline_data" / "pptx_probes" / "anchorb"

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
BOX_X = Emu(914400)          # 1.00in
BOX_Y = Emu(1828800)         # 2.00in
BOX_W = Emu(6400800)         # 7.00in
BOX_H = Emu(914400)          # 1.00in  = 72pt
INS = Emu(91440)             # 0.10in  = 7.2pt on every side
SIZE = 32                    # pt -> 1.2 x 32 = 38.4pt a line
ANCHORS = ["t", "ctr", "b"]
LINE_COUNTS = [1, 2, 3, 4]


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for anchor in ANCHORS:
        for n in LINE_COUNTS:
            slide = prs.slides.add_slide(blank)
            box = slide.shapes.add_textbox(BOX_X, BOX_Y, BOX_W, BOX_H)
            tf = box.text_frame
            body = tf._txBody.find(f"{{{A}}}bodyPr")
            # Fixed geometry: no autofit, explicit anchor, known insets.
            for child in list(body):
                body.remove(child)
            body.set("anchor", anchor)
            body.set("anchorCtr", "0")
            body.set("wrap", "square")
            for k in ("lIns", "rIns", "tIns", "bIns"):
                body.set(k, str(int(INS)))
            body.append(etree.SubElement(body, f"{{{A}}}noAutofit"))
            for i in range(n):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                r = p.add_run()
                r.text = f"{anchor}{n}L{i + 1}"
                r.font.size = Pt(SIZE)
                r.font.name = "Arial"
            # A hairline rectangle on the same geometry marks the box in the PDF.
            frame = slide.shapes.add_shape(1, BOX_X, BOX_Y, BOX_W, BOX_H)
            frame.fill.background()
            frame.line.color.rgb = RGBColor(0xFF, 0x00, 0x00)
            frame.line.width = Pt(0.75)
            frame.text_frame.text = ""
    path = OUT / "probe_anchorb.pptx"
    prs.save(str(path))
    print(f"wrote {path}  {len(ANCHORS) * len(LINE_COUNTS)} slides, "
          f"box {int(BOX_H) / 12700:.1f}pt, line pitch {1.2 * SIZE:.1f}pt")


if __name__ == "__main__":
    main()
