# -*- coding: utf-8 -*-
"""Probe: what size does an EMPTY paragraph use for its line height?

d28 slide 13 puts an empty paragraph between two 10pt body paragraphs. Its
`<a:endParaRPr sz="1000"/>` matches the body, but the deck's inherited default
(presentation defaultTextStyle lvl1 and the master otherStyle) is 14pt. The
rendered gap is 35px at 150dpi = 16.8pt = 10 x 1.2 x 1.4, i.e. PowerPoint sized
that line from `endParaRPr`, not from the inherited default -- Oxi took the
default and drew 49px, pushing everything below by 15px. 8858 empty paragraphs
across the 40-deck dev corpus carry an `endParaRPr` sz.

Each arm is `AAA / <empty> / BBB` in one box; the empty paragraph's declared
size is what varies. The last two arms disambiguate: E has no endParaRPr at all
(so the reader can see the fallback), and F puts a *run* rPr sz on the empty
paragraph that disagrees with its endParaRPr.
"""
from __future__ import annotations

import copy
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\emptypara").resolve()
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS = {"a": A}
BODY_PT = 10.0
# (label, endParaRPr sz in hundredths or None, run rPr sz or None, lnSpc pct)
ARMS = [
    ("A_epr700", 700, None, 100),
    ("B_epr1000", 1000, None, 100),
    ("C_epr2400", 2400, None, 100),
    ("D_epr4000", 4000, None, 100),
    ("E_none", None, None, 100),
    ("F_run1000_epr4000", 4000, 1000, 100),
    ("G_epr1000_ls140", 1000, None, 140),
    # E showed the fallback is not the inherited 18pt default. H/I ask whether
    # an endParaRPr-less empty paragraph takes the size of the paragraph before
    # it or the one after it; J removes the "before" so only "after" is left.
    ("H_prev24_none", None, None, 100),
    ("I_next24_none", None, None, 100),
    ("J_first_none", None, None, 100),
]
# per-arm body sizes: (first paragraph pt, last paragraph pt, drop first para)
BODY = {
    "H_prev24_none": (24.0, 10.0, False),
    "I_next24_none": (10.0, 24.0, False),
    "J_first_none": (10.0, 24.0, True),
}


def q(tag: str) -> str:
    return f"{{{A}}}{tag}"


def set_lnspc(p, pct: int) -> None:
    ppr = p.find(q("pPr"))
    if ppr is None:
        ppr = etree.SubElement(p, q("pPr"))
        p.insert(0, ppr)
    ln = etree.SubElement(ppr, q("lnSpc"))
    etree.SubElement(ln, q("spcPct")).set("val", str(pct * 1000))
    ppr.insert(0, ln)


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for label, epr_sz, run_sz, ls in ARMS:
        s = prs.slides.add_slide(blank)
        cap = s.shapes.add_textbox(Emu(228600), Emu(114300), Emu(6400800), Emu(300000))
        cap.text_frame.text = label
        box = s.shapes.add_textbox(Emu(457200), Emu(914400), Emu(6000000), Emu(3000000))
        tf = box.text_frame
        tf.word_wrap = True
        txbody = tf._txBody
        for pel in txbody.findall(q("p")):
            txbody.remove(pel)

        first_pt, last_pt, drop_first = BODY.get(label, (BODY_PT, BODY_PT, False))

        def add_text_para(text: str, pt: float):
            p = etree.SubElement(txbody, q("p"))
            r = etree.SubElement(p, q("r"))
            rpr = etree.SubElement(r, q("rPr"))
            rpr.set("lang", "en-US")
            rpr.set("sz", str(int(pt * 100)))
            t = etree.SubElement(r, q("t"))
            t.text = text
            set_lnspc(p, ls)
            return p

        if not drop_first:
            add_text_para("AAA first paragraph", first_pt)
        # the empty paragraph under test
        p = etree.SubElement(txbody, q("p"))
        if run_sz is not None:
            r = etree.SubElement(p, q("r"))
            rpr = etree.SubElement(r, q("rPr"))
            rpr.set("lang", "en-US")
            rpr.set("sz", str(run_sz))
            etree.SubElement(r, q("t")).text = ""
        if epr_sz is not None:
            e = etree.SubElement(p, q("endParaRPr"))
            e.set("lang", "en-US")
            e.set("sz", str(epr_sz))
        set_lnspc(p, ls)
        add_text_para("BBB third paragraph", last_pt)

    prs.save(OUT / "emptypara.pptx")
    print(f"wrote {OUT / 'emptypara.pptx'}  ({len(ARMS)} arms)")


if __name__ == "__main__":
    main()
