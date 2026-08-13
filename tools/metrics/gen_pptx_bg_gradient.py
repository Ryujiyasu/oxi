# -*- coding: utf-8 -*-
"""Build a deck whose ONLY content is a gradient slide background.

The dev corpus has 75 gradient-background slides, but every one of them also
carries shapes, so sampling their PDFs measures the shapes, not the gradient
(d15 slide1's bottom-right corner reads #FF4092 -- that is artwork, not the
fill).  So author the gradients ourselves, one per slide, with nothing on top.

python-pptx cannot write <p:bg>, so each slide's XML is patched in the zip --
the same surgery already used for doughnut holeSize / bubbleScale / stock.

Arms sweep the axes that OOXML exposes:
  linear  : angle 0 / 90 / 45 / 270, two stops and three stops, scaled 0/1
  path    : circle focused bottom-right (the d15 shape) and top-left, plus rect
Colours are deliberately saturated primaries so the ramp direction is readable
from a handful of pixels.
"""
from __future__ import annotations

import re
import shutil
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path(r"pipeline_data\pptx_probes\bg_gradient").resolve()

# (label, <p:bg> inner XML)
def lin(ang_deg: float, stops: list[tuple[int, str]], scaled: int = 0) -> str:
    gs = "".join(
        f'<a:gs pos="{pos}"><a:srgbClr val="{col}"/></a:gs>' for pos, col in stops
    )
    return (f'<a:gradFill><a:gsLst>{gs}</a:gsLst>'
            f'<a:lin ang="{int(round(ang_deg * 60000))}" scaled="{scaled}"/></a:gradFill>')


def path(kind: str, l: int, t: int, r: int, b: int, stops: list[tuple[int, str]]) -> str:
    gs = "".join(
        f'<a:gs pos="{pos}"><a:srgbClr val="{col}"/></a:gs>' for pos, col in stops
    )
    return (f'<a:gradFill><a:gsLst>{gs}</a:gsLst>'
            f'<a:path path="{kind}"><a:fillToRect l="{l}" t="{t}" r="{r}" b="{b}"/></a:path>'
            f'<a:tileRect/></a:gradFill>')


RED, BLU, GRN = "FF0000", "0000FF", "00FF00"
TWO = [(0, RED), (100000, BLU)]
THREE = [(0, RED), (50000, GRN), (100000, BLU)]

ARMS: list[tuple[str, str]] = [
    ("B1 lin ang=0 (2 stops)",      lin(0, TWO)),
    ("B2 lin ang=90 (2 stops)",     lin(90, TWO)),
    ("B3 lin ang=45 (2 stops)",     lin(45, TWO)),
    ("B4 lin ang=270 (2 stops)",    lin(270, TWO)),
    ("B5 lin ang=0 (3 stops)",      lin(0, THREE)),
    ("B6 lin ang=45 scaled=1",      lin(45, TWO, scaled=1)),
    ("B7 path circle -> BR",        path("circle", 100000, 100000, 0, 0, TWO)),
    ("B8 path circle -> TL",        path("circle", 0, 0, 100000, 100000, TWO)),
    ("B9 path circle centred",      path("circle", 50000, 50000, 50000, 50000, TWO)),
    ("B10 path rect -> BR",         path("rect", 100000, 100000, 0, 0, TWO)),
]


def main():
    OUT.mkdir(parents=True, exist_ok=True)
    tmp = OUT / "_stage.pptx"
    dst = OUT / "bg_gradient.pptx"

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(9144000), Emu(6858000)
    blank = prs.slide_layouts[6]
    for label, _ in ARMS:
        s = prs.slides.add_slide(blank)
        # A tiny caption well away from the corners we sample.
        box = s.shapes.add_textbox(Emu(228600), Emu(228600), Emu(3200400), Emu(400050))
        box.text_frame.text = label
    prs.save(tmp)

    with zipfile.ZipFile(tmp) as zin:
        names = zin.namelist()
        data = {n: zin.read(n) for n in names}

    for i, (label, bg) in enumerate(ARMS, start=1):
        part = f"ppt/slides/slide{i}.xml"
        x = data[part].decode("utf-8")
        # <p:bg> is the FIRST child of <p:cSld>, before <p:spTree>.
        assert "<p:bg>" not in x, part
        x2 = x.replace("<p:cSld>", f"<p:cSld><p:bg><p:bgPr>{bg}<a:effectLst/></p:bgPr></p:bg>", 1)
        assert x2 != x, f"could not inject background into {part}"
        data[part] = x2.encode("utf-8")

    with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:
        for n in names:
            zout.writestr(n, data[n])
    tmp.unlink()

    print(f"wrote {dst}  ({len(ARMS)} arms)")
    for i, (label, bg) in enumerate(ARMS, start=1):
        kind = re.search(r"<a:(lin|path)[^>]*", bg).group(0)
        print(f"  slide{i:<3d} {label:28s} {kind}")


if __name__ == "__main__":
    main()
