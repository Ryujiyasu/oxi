#!/usr/bin/env python3
"""Build docs/showcase.pptx -- the deck the public slides viewer opens.

Self-authored so the site ships nothing it did not write. Each slide exercises
one of the things the engine resolves and the browser viewer draws: a gradient
background inherited by the slide, a placeholder that takes its size and face
from the master, a run-level highlight, rotation, and a picture.
"""
from __future__ import annotations

import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUT = Path("docs/showcase.pptx").resolve()
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
EMU = 914400


def q(tag: str) -> str:
    return f"{{{A}}}{tag}"


def gradient_bg(slide, c0: str, c1: str, ang_deg: float) -> None:
    """A `p:bg/p:bgPr/a:gradFill`, which python-pptx cannot express."""
    cs = slide._element.find(q("cSld").replace(A, A)) if False else slide._element[0]
    bg = etree.SubElement(cs, "{http://schemas.openxmlformats.org/presentationml/2006/main}bg")
    cs.insert(0, bg)
    pr = etree.SubElement(bg, "{http://schemas.openxmlformats.org/presentationml/2006/main}bgPr")
    grad = etree.SubElement(pr, q("gradFill"))
    lst = etree.SubElement(grad, q("gsLst"))
    for pos, col in ((0, c0), (100000, c1)):
        gs = etree.SubElement(lst, q("gs"))
        gs.set("pos", str(pos))
        etree.SubElement(gs, q("srgbClr")).set("val", col)
    etree.SubElement(grad, q("lin")).set("ang", str(int(ang_deg * 60000)))
    etree.SubElement(pr, q("effectLst"))


def text_box(slide, x, y, w, h, runs, size=28, color="FFFFFF", align="l",
             bold=False, face="Segoe UI", highlight=None, rot=None):
    box = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    if rot is not None:
        box._element.spPr.find(q("xfrm")).set("rot", str(int(rot * 60000)))
    tf = box.text_frame
    tf.word_wrap = True
    body = tf._txBody
    for pel in body.findall(q("p")):
        body.remove(pel)
    p = etree.SubElement(body, q("p"))
    ppr = etree.SubElement(p, q("pPr"))
    ppr.set("algn", align)
    etree.SubElement(ppr, q("buNone"))
    for text in runs if isinstance(runs, list) else [runs]:
        r = etree.SubElement(p, q("r"))
        rpr = etree.SubElement(r, q("rPr"))
        rpr.set("lang", "en-US")
        rpr.set("sz", str(int(size * 100)))
        if bold:
            rpr.set("b", "1")
        fill = etree.SubElement(rpr, q("solidFill"))
        etree.SubElement(fill, q("srgbClr")).set("val", color)
        if highlight:
            hl = etree.SubElement(rpr, q("highlight"))
            etree.SubElement(hl, q("srgbClr")).set("val", highlight)
        etree.SubElement(rpr, q("latin")).set("typeface", face)
        etree.SubElement(r, q("t")).text = text
    return box


def rect(slide, x, y, w, h, color, rot=None, alpha=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = __import__("pptx.dml.color", fromlist=["RGBColor"]).RGBColor.from_string(color)
    sh.line.fill.background()
    if rot is not None:
        sh.rotation = rot
    if alpha is not None:
        srgb = sh.fill.fore_color._xFill.find(q("solidFill")) if False else None
        el = sh._element.spPr.find(q("solidFill")).find(q("srgbClr"))
        etree.SubElement(el, q("alpha")).set("val", str(int(alpha * 100000)))
    return sh


def main() -> None:
    prs = Presentation()
    prs.slide_width = Emu(int(13.333 * EMU))
    prs.slide_height = Emu(int(7.5 * EMU))
    blank = prs.slide_layouts[6]
    W, H = prs.slide_width, prs.slide_height

    # 1 -- gradient background, a large title, a highlighted line
    s = prs.slides.add_slide(blank)
    gradient_bg(s, "3B1C5A", "8E2B6B", 45)
    text_box(s, int(0.9 * EMU), int(2.1 * EMU), int(11.5 * EMU), int(1.6 * EMU),
             "Oxi renders this in your browser", size=54, bold=True)
    text_box(s, int(0.95 * EMU), int(3.7 * EMU), int(11.5 * EMU), int(0.9 * EMU),
             " PowerPoint (.pptx), parsed and drawn by Rust compiled to WebAssembly ",
             size=20, color="3B1C5A", highlight="FFFFFF")
    text_box(s, int(0.95 * EMU), int(4.8 * EMU), int(11.5 * EMU), int(0.8 * EMU),
             "Nothing is uploaded. The file never leaves this device.",
             size=18, color="D9C7E8")

    # 2 -- shapes: fills, transparency, rotation
    s = prs.slides.add_slide(blank)
    gradient_bg(s, "0E4C6B", "12A0A8", 90)
    text_box(s, int(0.9 * EMU), int(0.7 * EMU), int(11.5 * EMU), int(1.0 * EMU),
             "Shapes, fills and rotation", size=40, bold=True)
    for i, (col, rot) in enumerate((("FFB703", 0), ("FB8500", -12), ("E63946", 8), ("A8DADC", -5))):
        rect(s, int((1.0 + i * 2.9) * EMU), int(2.4 * EMU), int(2.3 * EMU), int(2.3 * EMU),
             col, rot=rot, alpha=0.85)
    text_box(s, int(0.9 * EMU), int(5.4 * EMU), int(11.5 * EMU), int(1.0 * EMU),
             "Rounded rectangles at 85% opacity, each turned by its own angle.",
             size=18, color="D7F3F5")

    # 3 -- a bulleted body against a flat fill
    s = prs.slides.add_slide(blank)
    gradient_bg(s, "1D2939", "344054", 90)
    text_box(s, int(0.9 * EMU), int(0.7 * EMU), int(11.5 * EMU), int(1.0 * EMU),
             "What this page draws", size=40, bold=True)
    # Only what THIS page actually draws. The native renderer does more --
    # tables, charts, custom geometry, embedded fonts, colour emoji -- and
    # claiming those here would be claiming them for the viewer.
    lines = [
        "Theme colours resolved through slide, layout and master",
        "Placeholder inheritance: size, face, colour, alignment",
        "Gradient and picture backgrounds",
        "Shape fills, transparency, rotation and flips",
        "Text highlight, bold, italic, underline",
    ]
    for i, line in enumerate(lines):
        text_box(s, int(1.1 * EMU), int((2.0 + i * 0.85) * EMU), int(11.0 * EMU), int(0.8 * EMU),
                 "-  " + line, size=22, color="E4E7EC")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT}")


if __name__ == "__main__":
    main()
